# -*- coding: utf-8 -*-
"""
editor.py  --  Box / Step 05 helper: read a slide's editable text, and write
edits back into a .pptx while preserving formatting.

We let the user edit the two text fields that matter most per slide — the
TITLE and the SUBTITLE/headline — because those are single-line and edit
cleanly. Body text is shown for context but not edited here (editing richly
formatted multi-run body text safely is a later step).
"""

from pptx import Presentation

from build_library import read_id


def _text_shapes(slide):
    """(index, shape) for every shape on the slide that has non-empty text."""
    out = []
    for i, sh in enumerate(slide.shapes):
        if sh.has_text_frame and sh.text_frame.text.strip():
            out.append((i, sh))
    return out


def editable_fields(slide):
    """Return [(shape_index, label, current_text)] for Title + Subtitle."""
    shapes = _text_shapes(slide)
    fields = []
    title_idx = None

    # Prefer the real title placeholder; else the first text shape.
    try:
        title_shape = slide.shapes.title
    except (ValueError, AttributeError):
        title_shape = None
    if title_shape is not None and title_shape.text.strip():
        for i, sh in enumerate(slide.shapes):
            if sh is title_shape:
                title_idx = i
                fields.append((i, "Title", title_shape.text_frame.text.strip()))
                break
    if title_idx is None and shapes:
        i, sh = shapes[0]
        title_idx = i
        fields.append((i, "Title", sh.text_frame.text.strip()))

    # Subtitle = first text shape that isn't the title.
    for i, sh in shapes:
        if i != title_idx:
            fields.append((i, "Subtitle", sh.text_frame.text.strip()))
            break

    return fields


def set_text(shape, text):
    """Replace a shape's text with `text`, keeping the FIRST run's formatting
    (font, size, colour). Extra runs/paragraphs are removed."""
    tf = shape.text_frame
    paras = tf.paragraphs
    first = paras[0]
    if first.runs:
        first.runs[0].text = text
        for r in first.runs[1:]:
            r._r.getparent().remove(r._r)
    else:
        first.text = text
    for p in list(paras[1:]):
        p._p.getparent().remove(p._p)


def full_text(slide):
    """All of a slide's text, for read-only context in the review screen."""
    return [sh.text_frame.text.strip() for _, sh in _text_shapes(slide)]


def apply_edits(path, edits):
    """edits = {slide_id: {shape_index(int): new_text}} -> write into the deck."""
    prs = Presentation(path)
    for slide in prs.slides:
        sid = read_id(slide)
        if sid in edits:
            shapes = list(slide.shapes)
            for idx, text in edits[sid].items():
                if 0 <= idx < len(shapes) and shapes[idx].has_text_frame:
                    set_text(shapes[idx], text)
    prs.save(path)


def replace_tokens(path, tokens):
    """Replace literal tokens (e.g. {'[CLIENT]': 'Acme Bank'}) in every run.
    Run-level replace keeps formatting intact."""
    prs = Presentation(path)
    for slide in prs.slides:
        for sh in slide.shapes:
            if not sh.has_text_frame:
                continue
            for p in sh.text_frame.paragraphs:
                for r in p.runs:
                    for old, new in tokens.items():
                        if old in r.text:
                            r.text = r.text.replace(old, new)
    prs.save(path)
