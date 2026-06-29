# -*- coding: utf-8 -*-
"""
skills.py  --  Three data-driven slides for Workforce-only + RFI decks.

Source data = J2W_Delivery_Footprint_Organized_Latest.xlsx
  Sheet: "Clean - Co x Func x Skill"
  Columns: Industry | Company | Function Delivered | Normalized Skill | Count

Gate: candidates() returns [] unless BOTH hold:
  1. Work type = pure Workforce (no other type selected)
  2. Transcript/notes mentions "RFI" or "request for information"

Three slide types (templates live in templates.pptx):
  industry_strength   -- overview of J2W's presence in the deck's industry
  skill_deepdive      -- one combined slide for all skills matched in the notes
  company_footprint   -- existing relationship if client fuzzy-matches a company
"""

import re
from collections import defaultdict

import openpyxl
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION
from pptx.dml.color import RGBColor

EXCEL            = "J2W_Delivery_Footprint_Organized_Latest.xlsx"
SHEET            = "Clean - Co x Func x Skill"
SKILLS_TEMPLATES = "skills_templates.pptx"   # branded template slides (from import_skills_templates.py)

BRAND = [
    RGBColor(0x2C, 0x6E, 0x66),
    RGBColor(0x7F, 0xB2, 0xA9),
    RGBColor(0xCF, 0xE7, 0xE2),
    RGBColor(0x11, 0x11, 0x10),
    RGBColor(0x4A, 0x9E, 0x94),
    RGBColor(0xA5, 0xCC, 0xC6),
    RGBColor(0x1E, 0x4D, 0x47),
]

# Distinct multi-color palette for the industry pie chart (no teal/green)
CHART_COLORS_MULTI = [
    RGBColor(0xE0, 0x6C, 0x1F),  # burnt orange
    RGBColor(0x7B, 0x4E, 0xA5),  # purple
    RGBColor(0x2A, 0x7E, 0xBC),  # steel blue
    RGBColor(0xB8, 0x2E, 0x2E),  # crimson red
    RGBColor(0xF5, 0xA6, 0x23),  # amber yellow
    RGBColor(0x8B, 0x5C, 0x2A),  # warm brown
    RGBColor(0xD4, 0x5E, 0x9A),  # raspberry pink
]

# Words to strip when normalizing company names for fuzzy matching
_CO_STRIP = re.compile(
    r'\b(pvt|ltd|llp|llc|inc|corp|co|pty|pte|sdn|bhd|gmbh|ag|nv|sa|ab|'
    r'consulting|advisory|payroll|ops|coe|technologies|technology|'
    r'software|solutions|services|group|india|global|international|'
    r'c2h|routing|sgp|benz|mercedes)\b',
    re.IGNORECASE
)

# Form industry codes -> Excel Industry column values
_IND_MAP = {
    "BANKING":          "Banking & Financial Services",
    "BANKING_FINANCE":  "Banking & Financial Services",
    "FINANCE":          "Banking & Financial Services",
    "TECH_IT":          "IT Services & Consulting",
    "TECH":             "IT Services & Consulting",
    "IT":               "IT Services & Consulting",
    "HEALTHCARE":       "Healthcare & Life Sciences",
    "HEALTH":           "Healthcare & Life Sciences",
    "INSURANCE":        "Insurance",
    "ENERGY":           "Energy & Utilities",
    "UTILITIES":        "Energy & Utilities",
    "RETAIL":           "E-commerce & Retail",
    "ECOMMERCE":        "E-commerce & Retail",
    "E_COMMERCE":       "E-commerce & Retail",
    "MANUFACTURING":    "Manufacturing & Materials",
    "TELECOM":          "Telecommunications",
    "AUTOMOTIVE":       "Automotive & Industrial",
    "INDUSTRIAL":       "Automotive & Industrial",
    "AEROSPACE":        "Aerospace & Defense",
    "DEFENSE":          "Aerospace & Defense",
    "REALESTATE":       "Real Estate & Professional Services",
    "REAL_ESTATE":      "Real Estate & Professional Services",
    "CONSULTING":       "Consulting & Professional Services",
    "PROFESSIONAL":     "Consulting & Professional Services",
    "SEMICONDUCTORS":   "Semiconductors",
    "SEMICONDUCTOR":    "Semiconductors",
    "SOFTWARE":         "Software & Cloud",
    "CLOUD":            "Software & Cloud",
    "CYBERSECURITY":    "Cybersecurity",
    "CYBER":            "Cybersecurity",
    "HARDWARE":         "Technology Hardware",
    "TECHNOLOGY_HARDWARE": "Technology Hardware",
}

# Generic words that match too many skills; skip them in keyword matching
_SKILL_STOPWORDS = {
    'with', 'that', 'this', 'from', 'have', 'will', 'been', 'test', 'data',
    'back', 'your', 'tech', 'work', 'code', 'team', 'tion', 'ment', 'able',
    'ware', 'base', 'ness', 'ding', 'over', 'into', 'ting', 'ring', 'port',
}


# ------------------------------------------------------------------ #
# Data loading (module-level cache)
# ------------------------------------------------------------------ #
_cache = None


def _load():
    global _cache
    if _cache is not None:
        return _cache
    wb = openpyxl.load_workbook(EXCEL, data_only=True)
    ws = wb[SHEET]
    rows = []
    for r in ws.iter_rows(min_row=2, values_only=True):
        industry, company, function, skill, count = (r + (None,) * 5)[:5]
        if not company:
            continue
        rows.append({
            "industry": (industry or "").strip(),
            "company":  (company  or "").strip(),
            "function": (function or "").strip(),
            "skill":    (skill    or "").strip(),
            "count":    int(count or 0),
        })
    _cache = rows
    return rows


# ------------------------------------------------------------------ #
# Helpers
# ------------------------------------------------------------------ #
def _is_rfi(transcript):
    t = (transcript or "").lower()
    return bool(re.search(r'\brfi\b|request\s+for\s+information', t))


def _industry_label(deck_industry):
    code = (deck_industry or "").upper().strip()
    code = re.sub(r'[\s\-&]+', '_', code)
    if code in _IND_MAP:
        return _IND_MAP[code]
    # Partial match fallback
    for k, v in _IND_MAP.items():
        if code.startswith(k) or k.startswith(code):
            return v
    return deck_industry  # pass through as-is for an exact Excel match


def _normalize_co(name):
    n = (name or "").lower()
    n = _CO_STRIP.sub("", n)
    n = re.sub(r'[^a-z0-9]', '', n)
    return n.strip()


def _match_companies(client_name, all_companies):
    cn = _normalize_co(client_name)
    if len(cn) < 3:
        return []
    matched = []
    for co in all_companies:
        co_n = _normalize_co(co)
        if not co_n:
            continue
        if cn in co_n or co_n in cn:
            matched.append(co)
    return matched


def _match_skills(transcript, all_skills):
    """Return skills whose name (or significant words) appear in the transcript."""
    notes = (transcript or "").lower()
    matched = []
    for skill in all_skills:
        sk_lower = skill.lower()
        if sk_lower in notes:
            matched.append(skill)
            continue
        words = [w for w in re.findall(r'[a-z]{4,}', sk_lower)
                 if w not in _SKILL_STOPWORDS]
        if len(words) >= 2:
            if all(re.search(r'\b' + re.escape(w) + r'\b', notes) for w in words[:2]):
                matched.append(skill)
        elif len(words) == 1:
            if re.search(r'\b' + re.escape(words[0]) + r'\b', notes):
                matched.append(skill)
    seen = set()
    out = []
    for s in matched:
        if s not in seen:
            seen.add(s)
            out.append(s)
    return out


# ------------------------------------------------------------------ #
# Slide data builders
# ------------------------------------------------------------------ #
def _industry_slide_data(rows, industry_label):
    ind_rows = [r for r in rows
                if r["industry"].lower() == (industry_label or "").lower()]
    if not ind_rows:
        return None

    companies = {r["company"] for r in ind_rows}
    functions = {r["function"] for r in ind_rows}
    skills    = {r["skill"]   for r in ind_rows}
    total     = sum(r["count"] for r in ind_rows)

    # Top 3 frequently hired: sort by # distinct companies, then total headcount
    sk_cos = defaultdict(set)
    sk_cnt = defaultdict(int)
    for r in ind_rows:
        sk_cos[r["skill"]].add(r["company"])
        sk_cnt[r["skill"]] += r["count"]
    top3 = sorted(sk_cos.keys(), key=lambda s: (-len(sk_cos[s]), -sk_cnt[s]))[:3]

    fn_totals = defaultdict(int)
    for r in ind_rows:
        fn_totals[r["function"]] += r["count"]
    fn_sorted = sorted(fn_totals.items(), key=lambda x: -x[1])

    return {
        "industry":      industry_label,
        "total":         total,
        "num_companies": len(companies),
        "num_functions": len(functions),
        "num_skills":    len(skills),
        "top3":          top3,
        "top3_cos":      [len(sk_cos[s]) for s in top3],
        "fn_chart":      fn_sorted,
    }


def _skill_slide_data(rows, matched_skills):
    result = []
    for skill in matched_skills:
        sk_rows = [r for r in rows if r["skill"].lower() == skill.lower()]
        if not sk_rows:
            continue
        total = sum(r["count"] for r in sk_rows)
        cos   = defaultdict(int)
        inds  = set()
        for r in sk_rows:
            cos[r["company"]] += r["count"]
            inds.add(r["industry"])
        result.append({
            "skill":      skill,
            "total":      total,
            "companies":  sorted(cos.items(), key=lambda x: -x[1]),
            "industries": sorted(inds),
        })
    return result


def _company_slide_data(rows, client_name):
    all_cos = list({r["company"] for r in rows})
    matched = _match_companies(client_name, all_cos)
    if not matched:
        return None

    co_rows  = [r for r in rows if r["company"] in set(matched)]
    total    = sum(r["count"] for r in co_rows)
    functions = {r["function"] for r in co_rows}
    skills    = {r["skill"]   for r in co_rows}

    fn_totals = defaultdict(int)
    for r in co_rows:
        fn_totals[r["function"]] += r["count"]
    fn_sorted = sorted(fn_totals.items(), key=lambda x: -x[1])

    # Use the shortest matched name as the display name (most likely the clean variant)
    display_name = min(matched, key=len)

    return {
        "company":       display_name,
        "total":         total,
        "num_functions": len(functions),
        "num_skills":    len(skills),
        "fn_breakdown":  fn_sorted,
    }


# ------------------------------------------------------------------ #
# Public API
# ------------------------------------------------------------------ #
def candidates(context):
    """Return slide candidates. Returns [] unless Workforce + RFI gate passes."""
    wts = {str(w).upper() for w in (context.get("work_types") or [])}
    if wts != {"WORKFORCE"}:
        return []
    transcript = context.get("transcript", "")
    if not _is_rfi(transcript):
        return []

    rows          = _load()
    industry_lbl  = _industry_label(context.get("industry", ""))
    client        = (context.get("client_name", "") or "").strip()
    out           = []

    # Slide 1: Industry Strength
    ind_data = _industry_slide_data(rows, industry_lbl)
    if ind_data:
        out.append({
            "id":       "SK:industry",
            "kind":     "industry_strength",
            "label":    f"Industry strength — {industry_lbl}",
            "template": "industry_strength",
            "data":     ind_data,
            "stale":    False,
        })

    # Slide 2: Combined Skill Deep-dive
    all_skills = list({r["skill"] for r in rows})
    matched    = _match_skills(transcript, all_skills)
    if matched:
        sk_data = _skill_slide_data(rows, matched)
        if sk_data:
            names = [s["skill"] for s in sk_data[:3]]
            label = "Skills deployed — " + "  ·  ".join(names)
            if len(sk_data) > 3:
                label += f"  +{len(sk_data)-3} more"
            out.append({
                "id":       "SK:skills",
                "kind":     "skill_deepdive",
                "label":    label,
                "template": "skill_deepdive",
                "data":     sk_data,
                "stale":    False,
            })

    # Slide 3: Company Relationship
    if client:
        co_data = _company_slide_data(rows, client)
        if co_data:
            out.append({
                "id":       f"FP:{co_data['company']}",
                "kind":     "company_footprint",
                "label":    f"Client relationship — {co_data['company']}",
                "template": "company_footprint",
                "data":     co_data,
                "stale":    False,
            })

    return out


def by_id(context, sid):
    return next((c for c in candidates(context) if c["id"] == sid), None)


# ------------------------------------------------------------------ #
# Marker filling
# ------------------------------------------------------------------ #
def fill_markers(slide, mapping):
    for sh in slide.shapes:
        if not sh.has_text_frame:
            continue
        for p in sh.text_frame.paragraphs:
            full = "".join(run.text for run in p.runs)
            if "{{" not in full:
                continue
            new = full
            for m, v in mapping.items():
                new = new.replace("{{" + m + "}}", v)
            if new != full and p.runs:
                p.runs[0].text = new
                for run in p.runs[1:]:
                    run.text = ""


def _mapping_industry(data):
    top3     = data["top3"]
    top3_cos = data["top3_cos"]
    def _sk(i):
        if i < len(top3):
            return f"{top3[i]}  ({top3_cos[i]} {'company' if top3_cos[i]==1 else 'companies'})"
        return "—"
    return {
        "INDUSTRY_NAME":     data["industry"],
        "TOTAL_CONSULTANTS": f"{data['total']:,}",
        "NUM_COMPANIES":     str(data["num_companies"]),
        "NUM_FUNCTIONS":     str(data["num_functions"]),
        "NUM_SKILLS":        str(data["num_skills"]),
        "TOP_SKILL_1":       _sk(0),
        "TOP_SKILL_2":       _sk(1),
        "TOP_SKILL_3":       _sk(2),
    }


def _mapping_skills(sk_list):
    lines = []
    for s in sk_list:
        top_cos = ", ".join(f"{co} ({cnt})" for co, cnt in s["companies"][:4])
        if len(s["companies"]) > 4:
            top_cos += f"  +{len(s['companies'])-4} more"
        lines.append(
            f"▸ {s['skill']}  —  {s['total']:,} consultants  ·  "
            f"{len(s['companies'])} {'company' if len(s['companies'])==1 else 'companies'}\n"
            f"   {top_cos}"
        )
    header = "  ·  ".join(s["skill"] for s in sk_list)
    return {
        "SKILLS_HEADER": header,
        "SKILL_SUMMARY": "\n\n".join(lines),
    }


def _mapping_company(data):
    fn_lines = "\n".join(
        f"▸ {fn}:  {cnt:,}" for fn, cnt in data["fn_breakdown"]
    )
    return {
        "COMPANY_NAME":      data["company"],
        "TOTAL_DEPLOYED":    f"{data['total']:,}",
        "NUM_FUNCTIONS_CO":  str(data["num_functions"]),
        "NUM_SKILLS_CO":     str(data["num_skills"]),
        "ENGAGEMENT_TYPE":   "Existing client",
        "FUNCTION_BREAKDOWN": fn_lines,
    }


# ------------------------------------------------------------------ #
# Charts
# ------------------------------------------------------------------ #
def _find_chart_placeholder(slide):
    for sh in slide.shapes:
        if sh.has_text_frame and "{{CHART}}" in sh.text_frame.text:
            return sh
    return None


def _add_pie_chart(slide, categories, values):
    """Pie chart for the industry_strength slide.
    Uses a distinct multi-color palette (no teal/green).
    Legend on the RIGHT shows category names; percentage labels sit INSIDE
    each slice — this avoids label overlap when there are many categories."""
    from pptx.enum.chart import XL_LABEL_POSITION

    holder = _find_chart_placeholder(slide)
    if holder is None:
        return False
    left, top, w, h = holder.left, holder.top, holder.width, holder.height
    holder._element.getparent().remove(holder._element)

    pairs = [(c, v) for c, v in zip(categories, values) if v]
    if not pairs:
        return False
    cats, vals = zip(*pairs)

    cd = CategoryChartData()
    cd.categories = list(cats)
    cd.add_series("Consultants", list(vals))

    gf    = slide.shapes.add_chart(XL_CHART_TYPE.PIE, left, top, w, h, cd)
    chart = gf.chart
    chart.has_title  = False
    # Legend on the RIGHT — shows function names with color squares, no overlap
    chart.has_legend = True
    chart.legend.position = XL_LEGEND_POSITION.RIGHT
    chart.legend.include_in_layout = False

    # Percentage labels inside each slice; at this position PowerPoint
    # never generates overlapping labels regardless of how many categories exist
    plot = chart.plots[0]
    plot.has_data_labels = True
    dl = plot.data_labels
    dl.show_category_name = False   # category shown in legend; no duplication
    dl.show_percentage    = True    # e.g. "27%"
    dl.show_value         = False
    dl.show_legend_key    = False
    dl.position           = XL_LABEL_POSITION.INSIDE_END

    # Apply distinct multi-color fill per slice
    for i, pt in enumerate(plot.series[0].points):
        pt.format.fill.solid()
        pt.format.fill.fore_color.rgb = CHART_COLORS_MULTI[i % len(CHART_COLORS_MULTI)]
    return True


def _add_bar(slide, categories, values):
    """Horizontal bar chart — categories on Y-axis, counts on X-axis."""
    holder = _find_chart_placeholder(slide)
    if holder is None:
        return False
    left, top, w, h = holder.left, holder.top, holder.width, holder.height
    holder._element.getparent().remove(holder._element)

    # Cap at 10 bars for readability
    cats = list(categories[:10])
    vals = list(values[:10])
    if not vals:
        return False

    cd = CategoryChartData()
    cd.categories = cats
    cd.add_series("Consultants", vals)
    gf = slide.shapes.add_chart(XL_CHART_TYPE.BAR_CLUSTERED, left, top, w, h, cd)
    chart = gf.chart
    chart.has_legend = False
    series = chart.plots[0].series[0]
    series.format.fill.solid()
    series.format.fill.fore_color.rgb = BRAND[0]
    return True


# ------------------------------------------------------------------ #
# Build slides into the assembled deck
# ------------------------------------------------------------------ #
def find_template(prs, name):
    tag = "J2W_TEMPLATE: " + name
    for s in prs.slides:
        if s.has_notes_slide:
            for line in (s.notes_slide.notes_text_frame.text or "").splitlines():
                if line.strip() == tag:
                    return s
    return None


def build_into(deck_path, order, cands, master_path="WORKING_COPY_Master_Deck.pptx"):
    """Copy filled skills slides into the assembled deck and reorder to match `order`."""
    from pptx import Presentation
    import slide_generator
    import assembler
    from build_library import read_id

    cand_by_id = {c["id"]: c for c in cands if c["id"] in order}
    if not cand_by_id:
        return 0

    prs        = Presentation(deck_path)
    sld_id_lst = prs.slides._sldIdLst
    tfile      = Presentation(SKILLS_TEMPLATES)

    skill_elem = {}
    for sid, c in cand_by_id.items():
        t = find_template(tfile, c["template"])
        if t is None:
            print(f"  WARNING: template '{c['template']}' not found — run create_skills_templates.py")
            continue

        new  = slide_generator._copy_slide(prs, t)
        kind = c["kind"]
        data = c["data"]

        if kind == "industry_strength":
            fill_markers(new, _mapping_industry(data))
            fn_names = [fn for fn, _ in data["fn_chart"]]
            fn_vals  = [cnt for _, cnt in data["fn_chart"]]
            _add_pie_chart(new, fn_names, fn_vals)

        elif kind == "skill_deepdive":
            fill_markers(new, _mapping_skills(data))
            sk_names = [s["skill"] for s in data]
            sk_vals  = [s["total"] for s in data]
            _add_bar(new, sk_names, sk_vals)

        elif kind == "company_footprint":
            fill_markers(new, _mapping_company(data))
            fn_names = [fn for fn, _ in data["fn_breakdown"]]
            fn_vals  = [cnt for _, cnt in data["fn_breakdown"]]
            _add_bar(new, fn_names, fn_vals)

        skill_elem[sid] = list(sld_id_lst)[-1]

    # Reorder the whole deck to match `order`; use explicit None checks (lxml falsy-when-empty).
    cs_elem = {read_id(s): e for s, e in zip(prs.slides, list(sld_id_lst)) if read_id(s)}
    for sid in order:
        e = cs_elem.get(sid)
        if e is None:
            e = skill_elem.get(sid)
        if e is not None:
            sld_id_lst.append(e)

    assembler._atomic_save(prs, deck_path)
    return len(skill_elem)
