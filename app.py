# -*- coding: utf-8 -*-
"""
app.py  --  Box 3: the web form the sales/tech team uses.

Run it:   py app.py
Then open: http://127.0.0.1:5000

Flow: fill in the client context -> see the slide picks -> DRAG them into
any order (or use the up/down buttons) -> download the tailored .pptx in
exactly that order. Reuses matcher.py (brain) + assembler.py (.pptx builder).
No transcript / no AI yet.
"""

import os
import re
import json
import glob
import uuid
from collections import Counter
from zipfile import BadZipFile

from flask import Flask, request, render_template_string, send_file, abort, redirect
from pptx import Presentation

import tagger
import matcher
import assembler
import editor
import slide_generator
import staging
import meeting_log
import skills
from build_library import read_id


def _maybe_generate(path):
    """Fill flagged gaps with slides. REUSE a previously generated slide for the
    same slot if one exists (the learning loop); otherwise draft a fresh one with
    AI and stage it for review. Then append all of them to the deck.

    Returns the list of UNVERIFIED slides added (still need expert verification),
    so the caller can warn the salesperson. Each: {title, work_type, keywords}."""
    gen_wts = request.form.getlist("gen")
    if not gen_wts:
        return []
    industry = request.form.get("industry", "")
    transcript = request.form.get("transcript", "")
    client = request.form.get("client_name", "")
    items, unverified = [], []
    for wt in gen_wts:
        existing = staging.find(wt, industry)        # reuse before regenerate
        if existing:
            verified = existing.get("status") == "approved"   # approved == client-ready
            item = {"title": existing["title"], "keywords": existing["keywords"],
                    "bullets": existing["bullets"], "template": existing["template"],
                    "verified": verified}
        else:
            content = slide_generator.draft({"type": "needs_case_study", "work_type": wt},
                                            {"industry": industry, "transcript": transcript})
            staging.add(content, wt, industry, client)   # save as 'pending' for review
            item = dict(content, verified=False)         # freshly written == unverified
        items.append(item)
        if not item["verified"]:
            unverified.append({"title": item["title"], "work_type": wt,
                               "keywords": item.get("keywords", "")})
    slide_generator.append_generated(path, items)
    return unverified

app = Flask(__name__)

OUTPUT_DIR = "output"
INDUSTRIES = list(tagger.INDUSTRY.keys())
FUNCTIONS = list(tagger.FUNCTION.keys())
WORK_TYPES = ["WORKFORCE", "AI_POD", "MS"]
DEAL_STAGES = ["Intro", "Proposal", "Negotiation", "Follow-up"]
MEETING_TYPES = ["Discovery", "Demo", "Pitch", "Follow-up", "Other"]
OWNER = "Athithia"

# Friendly labels for the work-type codes (used in dropdowns / tables).
WT_LABELS = {"WORKFORCE": "Workforce", "AI_POD": "AI Pods", "MS": "Managed services"}

# Deck phase — fixed list, pick exactly one, in this order.
PHASES = [
    "Pre-read",        # sent to the client before a meeting, as preparation
    "First Meeting",   # introductory — industry-specific, tuned to the stakeholder
    "Second Meeting",  # focused on the specific thing the client showed interest in
    "Proposal",        # formal proposal stage
]


def current_salesperson():
    """Who generated the deck. No login exists yet, so return a clearly-marked
    placeholder. At deploy time, wire this to the logged-in user."""
    return "[NOT LOGGED IN - wire to login at deploy]"

SHELL = """
<!doctype html><html lang="en"><head><meta charset="utf-8">
<meta name="viewport" content="width=device-width, initial-scale=1">
<title>{{ title }}</title>
<link rel="stylesheet" href="/static/app.css">
<link rel="stylesheet" href="https://cdn.jsdelivr.net/npm/@tabler/icons-webfont@3.24.0/dist/tabler-icons.min.css">
</head><body>
<div class="app">
  <aside class="sidebar">
    <div class="logo"><i class="ti ti-stack-2"></i></div>
    <a class="nav-btn {{ 'active' if active=='new' else '' }}" href="/"><i class="ti ti-sparkles"></i><span class="nav-tip">New deck</span></a>
    <a class="nav-btn {{ 'active' if active=='library' else '' }}" href="/library"><i class="ti ti-books"></i><span class="nav-tip">Slide library</span></a>
    <a class="nav-btn {{ 'active' if active=='templates' else '' }}" href="/templates"><i class="ti ti-template"></i><span class="nav-tip">Templates</span></a>
    <a class="nav-btn {{ 'active' if active=='staging' else '' }}" href="/staging"><i class="ti ti-checkup-list"></i><span class="nav-tip">AI history</span></a>
    <a class="nav-btn {{ 'active' if active=='meetings' else '' }}" href="/meetings"><i class="ti ti-history"></i><span class="nav-tip">Deck repository</span></a>
    <div class="nav-sp"></div>
    <a class="nav-btn {{ 'active' if active=='home' else '' }}" href="/dashboard"><i class="ti ti-layout-dashboard"></i><span class="nav-tip">Dashboard</span></a>
    <a class="nav-btn" href="#"><i class="ti ti-settings"></i><span class="nav-tip">Settings</span></a>
  </aside>
  <div class="main">
    <header class="topbar">
      <div class="crumb"><i class="ti ti-folder"></i> {{ crumb|safe }}</div>
      {% if tabs %}<div class="tabs">{% for t in tabs %}<span class="tab {{ 'active' if loop.first else '' }}">{{ t }}</span>{% endfor %}</div>{% endif %}
      <div class="top-sp"></div>
      <button class="icon-btn" aria-label="Search"><i class="ti ti-search"></i></button>
      <button class="icon-btn" aria-label="Adjust"><i class="ti ti-adjustments-horizontal"></i></button>
      <div class="avatars"><span class="av">AT</span><span class="av" style="background:#dfe7e5">J2W</span></div>
    </header>
    <main class="page">{{ body|safe }}</main>
  </div>
</div>
</body></html>
"""


def shell(body, active="home", crumb="<b>Home</b> / Overview", title="J2W Pre-sales Engine", tabs=None):
    return render_template_string(SHELL, body=body, active=active, crumb=crumb, title=title, tabs=tabs)


def _dash_stats():
    try:
        recs = json.load(open("tagged_library.json", encoding="utf-8"))
    except Exception:
        recs = []
    total = len(recs)
    cases = sum(1 for r in recs if r.get("tags", {}).get("kind", {}).get("value") == "CASE_STUDY")
    wt, ind = Counter(), Counter()
    for r in recs:
        t = r.get("tags", {})
        if t.get("work_type", {}).get("value"):
            wt[t["work_type"]["value"]] += 1
        if t.get("industry", {}).get("value"):
            ind[t["industry"]["value"]] += 1
    try:
        templates = list(slide_generator.list_templates().keys())
    except Exception:
        templates = []
    decks = []
    for p in sorted(glob.glob("output/*.pptx"), key=os.path.getmtime, reverse=True)[:6]:
        nm = os.path.basename(p)
        decks.append({"name": nm.replace("Tailored_Deck_", "").replace(".pptx", "").replace("_", " "),
                      "size": "%.1f MB" % (os.path.getsize(p) / 1048576)})
    wt_items = [("Workforce", wt.get("WORKFORCE", 0)), ("AI Pods", wt.get("AI_POD", 0)),
                ("Managed services", wt.get("MS", 0))]
    wt_max = max([c for _, c in wt_items] + [1])
    ind_items = ind.most_common(5)
    ind_max = max([c for _, c in ind_items] + [1])
    return dict(total=total, cases=cases, templates=templates, decks=decks,
                wt_items=wt_items, wt_max=wt_max, ind_items=ind_items, ind_max=ind_max)


DASHBOARD_BODY = """
<div style="display:flex;align-items:flex-end;justify-content:space-between;gap:20px;margin:14px 0 26px">
  <div>
    <div class="eyebrow">Pre-sales engine</div>
    <h1 class="display">Welcome back, Athithia</h1>
    <p class="lede">Turn client context into a tailored, on-brand deck in seconds — pick the right slides, fill the gaps, ship.</p>
  </div>
  <a class="btn btn-primary btn-lg" href="/new"><i class="ti ti-sparkles"></i> New deck</a>
</div>

<div class="grid g-4">
  <div class="card stat"><div class="iwrap"><i class="ti ti-stack-2"></i></div>
    <div class="k">Slides in library</div><div class="v">{{ total }}</div>
    <span class="chip up"><i class="ti ti-arrow-up-right"></i> live</span></div>
  <div class="card stat"><div class="iwrap"><i class="ti ti-briefcase"></i></div>
    <div class="k">Case studies</div><div class="v">{{ cases }}</div>
    <div class="foot">ready to match</div></div>
  <div class="card stat"><div class="iwrap"><i class="ti ti-template"></i></div>
    <div class="k">Templates</div><div class="v">{{ templates|length }}</div>
    <div class="foot">for generated slides</div></div>
  <div class="card card-dark stat"><div class="iwrap"><i class="ti ti-file-download"></i></div>
    <div class="k" style="color:#b9b9b4">Decks built</div><div class="v">{{ decks|length }}</div>
    <span class="chip up"><i class="ti ti-clock"></i> recent</span></div>
</div>

<div class="section grid g-2">
  <div class="card">
    <div class="section-head"><h2 class="sec-title">Library composition</h2><span class="tag">by work type</span></div>
    <div class="bars">
      {% for name,c in wt_items %}
      <div class="bar-row"><span class="lbl">{{ name }}</span>
        <div class="bar-track"><div class="bar-fill" style="width:{{ (c*100//wt_max) if wt_max else 0 }}%"></div></div>
        <span class="num">{{ c }}</span></div>
      {% endfor %}
    </div>
    <div class="section-head" style="margin:24px 0 12px"><h2 class="sec-title" style="font-size:16px">Top industries</h2></div>
    <div class="bars">
      {% for name,c in ind_items %}
      <div class="bar-row"><span class="lbl">{{ name|title }}</span>
        <div class="bar-track"><div class="bar-fill" style="width:{{ (c*100//ind_max) if ind_max else 0 }}%;background:var(--teal-300)"></div></div>
        <span class="num">{{ c }}</span></div>
      {% endfor %}
    </div>
  </div>

  <div class="card">
    <div class="section-head"><h2 class="sec-title">Recent decks</h2><a class="hint" href="/new" style="font-size:13px">New →</a></div>
    {% if decks %}
    <div class="list">
      {% for d in decks %}
      <div class="row"><div class="ic"><i class="ti ti-presentation"></i></div>
        <div class="main"><div class="t">{{ d.name }}</div><div class="s">{{ d.size }}</div></div>
        <span class="meta"><i class="ti ti-download"></i></span></div>
      {% endfor %}
    </div>
    {% else %}
    <p class="muted" style="font-size:14px">No decks yet. Hit <b>New deck</b> to build your first.</p>
    {% endif %}
  </div>
</div>

<div class="section card card-teal" style="display:flex;align-items:center;justify-content:space-between;gap:20px">
  <div><div class="eyebrow" style="color:#cfe7e2">Quick start</div>
    <div class="display" style="font-size:24px;margin-top:4px">Paste a transcript, get matched slides</div>
    <p style="color:#dceeea;margin:6px 0 0;max-width:520px">Drop in your meeting notes — the engine ranks the right case studies, flags gaps, and you ship a tailored deck.</p></div>
  <a class="btn btn-lg" href="/new" style="background:#fff;color:var(--teal-700)"><i class="ti ti-arrow-right"></i> Start</a>
</div>
"""


NEW_FORM_BODY = """
<div style="margin:14px 0 24px">
  <div class="eyebrow">New deck</div>
  <h1 class="display" style="font-size:38px">Tell us about the client</h1>
  <p class="lede">Fill in the context. The engine picks the right slides, flags gaps, and you refine before download.</p>
</div>
{% if error %}<div class="card" style="border-left:4px solid #c0392b;margin-bottom:18px;color:#c0392b"><b>{{ error }}</b></div>{% endif %}
<form id="deckForm" method="post" action="/build" class="form-layout">
  <div>
    <div class="card" style="margin-bottom:18px">
      <h2 class="sec-title" style="margin-bottom:18px">Client &amp; focus</h2>
      <div class="fgrid">
        <div class="field"><label>Client name</label>
          <input class="input" name="client_name" placeholder="Acme Bank" required></div>
        <div class="field"><label>Industry</label>
          <select name="industry" required><option value="">Select…</option>{% for i in industries %}<option>{{ i }}</option>{% endfor %}</select></div>
        <div class="field"><label>Deck phase</label>
          <select name="phase" required><option value="">Select…</option>{% for p in phases %}<option>{{ p }}</option>{% endfor %}</select></div>
        <div class="field"><label>Recipient / stakeholder <span class="hint">— optional</span></label>
          <input class="input" name="recipient" placeholder="Head of Engineering"></div>
      </div>
    </div>

    <div class="card" style="margin-bottom:18px">
      <h2 class="sec-title" style="margin-bottom:6px">Work type</h2>
      <p class="hint" style="font-size:13px;margin:0 0 14px">Pick one or more — each pulls in its standard slides plus matching case studies.</p>
      <div class="chip-row">
        <label class="chip-toggle"><input type="checkbox" name="work_types" value="WORKFORCE"><i class="ti ti-users"></i> Workforce</label>
        <label class="chip-toggle"><input type="checkbox" name="work_types" value="AI_POD"><i class="ti ti-cpu"></i> AI Pods</label>
        <label class="chip-toggle"><input type="checkbox" name="work_types" value="MS"><i class="ti ti-settings-cog"></i> Managed services</label>
      </div>
      <p id="wt-warn" class="hint" style="display:none;color:#c0392b;margin-top:10px">Please select at least one work type.</p>
    </div>

    <div class="card" style="margin-bottom:18px">
      <h2 class="sec-title" style="margin-bottom:6px">Function</h2>
      <p class="hint" style="font-size:13px;margin:0 0 14px">Pick one or more — or choose <b>Any function</b> if it's not function-specific.</p>
      <div class="chip-row" id="fn-chips">
        <button type="button" id="fn-any" class="chip-toggle" onclick="toggleAnyFn()"><i class="ti ti-asterisk"></i> Any function</button>
        {% for f in functions %}<label class="chip-toggle fn-chip"><input type="checkbox" name="functions" value="{{ f }}" onchange="fnChanged()"> {{ f.replace('_',' ')|title }}</label>{% endfor %}
      </div>
    </div>

    <div class="card">
      <h2 class="sec-title" style="margin-bottom:6px">Give me more information</h2>
      <p class="hint" style="font-size:13px;margin:0 0 12px">Paste anything useful — a meeting transcript, minutes (MOM), client or company research, notes. The AI reads it to pick the most relevant slides and write any that are missing.</p>
      <textarea name="transcript" placeholder="Paste transcript, MOM, research, or notes about the client and what they need…"></textarea>
    </div>
  </div>

  <aside class="side-card">
    <div class="card">
      <h2 class="sec-title" style="margin-bottom:16px">How it works</h2>
      <div class="steps">
        <div class="step"><div class="n">1</div><div class="st"><b>Match</b>Right slides picked from {{ library_count }} in the library</div></div>
        <div class="step"><div class="n">2</div><div class="st"><b>Fill gaps</b>Any missing slide is written by AI on the fly</div></div>
        <div class="step"><div class="n">3</div><div class="st"><b>Ship</b>Preview, refine the text, download the .pptx</div></div>
      </div>
      <button class="btn btn-primary btn-lg btn-block" type="submit" style="margin-top:22px"><i class="ti ti-wand"></i> Generate deck</button>
      <a href="/library" class="hint" style="display:block;text-align:center;margin-top:14px;font-size:13px">Browse the library →</a>
    </div>
  </aside>
</form>

<div id="loader" style="display:none;position:fixed;inset:0;background:rgba(17,17,16,.55);z-index:999;align-items:center;justify-content:center">
  <div style="background:#fff;border-radius:16px;padding:30px 38px;text-align:center;box-shadow:0 20px 60px rgba(0,0,0,.3)">
    <div style="width:40px;height:40px;border:4px solid #e6e6e3;border-top-color:#2C6E66;border-radius:50%;margin:0 auto 16px;animation:j2wspin .8s linear infinite"></div>
    <div style="font-family:'Space Grotesk',Inter,sans-serif;font-weight:700;font-size:18px;color:#111110">Generating your deck…</div>
    <div style="font-size:13px;color:#6b7280;margin-top:5px">The AI is matching the right slides — a few seconds.</div>
  </div>
</div>
<style>@keyframes j2wspin{to{transform:rotate(360deg)}}</style>

<script>
(function(){
  var f = document.getElementById('deckForm');
  if(!f) return;
  f.addEventListener('submit', function(e){
    var checked = f.querySelectorAll('input[name="work_types"]:checked').length;
    var warn = document.getElementById('wt-warn');
    if(checked === 0){
      e.preventDefault();
      if(warn) warn.style.display = 'block';
    } else {
      if(warn) warn.style.display = 'none';
      var ld = document.getElementById('loader');
      if(ld) ld.style.display = 'flex';
    }
  });
})();
// "Any function" = no specific function (optional). Highlighting it clears the
// individual function chips; ticking any chip clears the highlight.
function setAnyActive(on){
  var b=document.getElementById('fn-any'); if(!b) return;
  if(on){b.style.background='#2C6E66';b.style.color='#fff';b.style.borderColor='#2C6E66';b.dataset.on='1';}
  else{b.style.background='';b.style.color='';b.style.borderColor='';b.dataset.on='';}
}
function toggleAnyFn(){
  var on=document.getElementById('fn-any').dataset.on==='1';
  if(!on){document.querySelectorAll('#fn-chips input[name="functions"]').forEach(function(c){c.checked=false;});setAnyActive(true);}
  else{setAnyActive(false);}
}
function fnChanged(){
  var any=[...document.querySelectorAll('#fn-chips input[name="functions"]')].some(function(c){return c.checked;});
  if(any) setAnyActive(false);
}
window.addEventListener('pageshow',function(){var l=document.getElementById('loader');if(l)l.style.display='none';});
</script>
"""


FORM_HTML = """
<!doctype html><html><head><meta charset="utf-8"><title>J2W Pre-sales Engine</title>
<style>
 body{font-family:Segoe UI,Arial,sans-serif;max-width:680px;margin:40px auto;color:#222}
 h1{font-size:22px} label{display:block;margin:14px 0 4px;font-weight:600}
 input[type=text],select,textarea{width:100%;padding:8px;font-size:15px;box-sizing:border-box;font-family:inherit}
 .chk{display:inline-block;margin-right:18px;font-weight:400}
 button{margin-top:22px;padding:10px 22px;font-size:16px;background:#1a5;color:#fff;border:0;border-radius:5px;cursor:pointer}
 .hint{color:#777;font-size:13px;font-weight:400}
</style></head><body>
<h1>J2W Pre-sales Engine</h1>
<p class="hint">Fill in the client context. The tool picks the right slides; you can reorder them before download.</p>
<form method="post" action="/build">
  <label>Client name</label>
  <input type="text" name="client_name" placeholder="Acme Bank" required>

  <label>Industry</label>
  <select name="industry">
    <option value="">— select —</option>
    {% for i in industries %}<option value="{{i}}">{{i}}</option>{% endfor %}
  </select>

  <label>Work type(s) <span class="hint">— pick one or more</span></label>
  {% for w in work_types %}
    <label class="chk"><input type="checkbox" name="work_types" value="{{w}}"> {{w}}</label>
  {% endfor %}

  <label>Function <span class="hint">— optional</span></label>
  <select name="function">
    <option value="">— any —</option>
    {% for f in functions %}<option value="{{f}}">{{f}}</option>{% endfor %}
  </select>

  <label>Deal stage <span class="hint">— optional</span></label>
  <select name="deal_stage">
    <option value="">— select —</option>
    {% for d in deal_stages %}<option value="{{d}}">{{d}}</option>{% endfor %}
  </select>

  <label>Meeting type <span class="hint">— optional</span></label>
  <select name="meeting_type">
    <option value="">— select —</option>
    {% for m in meeting_types %}<option value="{{m}}">{{m}}</option>{% endfor %}
  </select>

  <label>Recipient <span class="hint">— optional</span></label>
  <input type="text" name="recipient" placeholder="CTO">

  <label>Meeting transcript / MOM <span class="hint">— optional; paste notes here</span></label>
  <textarea name="transcript" rows="6"
    placeholder="Paste the meeting transcript or minutes here to drive smarter case-study picks."></textarea>

  <label class="chk" style="margin-top:12px"><input type="checkbox" name="use_ai" value="1">
    Use AI to refine picks from the transcript <span class="hint">(uses your API key)</span></label>

  <button type="submit">Build deck</button>
</form>
</body></html>
"""

RESULT_HTML = """
<!doctype html><html><head><meta charset="utf-8"><title>Deck preview</title>
<style>
 body{font-family:Segoe UI,Arial,sans-serif;max-width:780px;margin:40px auto;color:#222}
 h1{font-size:22px}
 .hint{color:#777;font-size:13px}
 ul{list-style:none;padding:0;margin:16px 0}
 li{display:flex;align-items:center;gap:10px;border:1px solid #ddd;border-radius:6px;
    padding:8px 10px;margin:6px 0;background:#fafafa;cursor:grab}
 li.drag{opacity:.4}
 .pos{width:26px;color:#888;font-variant-numeric:tabular-nums}
 .grip{color:#bbb;cursor:grab}
 .id{font-weight:700;width:48px}
 .title{flex:1}
 .reason{color:#888;font-size:12px}
 .reason.weak{color:#b40;font-weight:600}
 .mv{border:1px solid #ccc;background:#fff;border-radius:4px;cursor:pointer;width:26px}
 .rm{border:1px solid #e0b4b4;background:#fff;color:#b40;border-radius:4px;cursor:pointer;width:26px}
 .addbar{display:flex;gap:8px;align-items:center;margin:12px 0 4px}
 .addbar select{flex:1;padding:7px;font-size:14px}
 .addbar button{background:#345;color:#fff;border:0;border-radius:4px;padding:8px 14px;cursor:pointer}
 .dl{margin-top:8px;padding:11px 24px;font-size:16px;background:#1a5;color:#fff;border:0;border-radius:5px;cursor:pointer}
 .back{display:inline-block;margin-top:18px;color:#555}
</style></head><body>
<h1>Deck for {{ ctx.client_name }}</h1>
<p>Context: industry <b>{{ ctx.industry or '—' }}</b>,
   work types <b>{{ ', '.join(ctx.work_types) or '—' }}</b>,
   function <b>{{ ctx.function or 'any' }}</b>,
   deal stage <b>{{ ctx.deal_stage or '—' }}</b>,
   meeting <b>{{ ctx.meeting_type or '—' }}</b>.</p>
{% if ctx.transcript %}<p class="hint">📝 Transcript received ({{ ctx.transcript|length }} chars).
 {% if ai_used %}<span style="color:#1a5">✓ AI refined the case-study picks from it.</span>{% else %}Keyword-matched — tick &ldquo;Use AI&rdquo; on the form for smarter picks.{% endif %}</p>{% endif %}
<p class="hint">{{ picks|length }} slides. Drag a row (or use ↑ ↓) to reorder, then download.</p>

<ul id="list">
 {% for p in picks %}
 <li draggable="true" data-id="{{ p.slide_id }}">
   <span class="pos"></span>
   <span class="grip">⠿</span>
   <span class="id">{{ p.slide_id }}</span>
   <span class="title">{{ titles.get(p.slide_id, '') }}</span>
   <span class="reason {{ 'weak' if 'WEAK' in p.reason else '' }}">{{ p.reason }}</span>
   <button type="button" class="mv" onclick="move(this,-1)">↑</button>
   <button type="button" class="mv" onclick="move(this,1)">↓</button>
   <button type="button" class="rm" onclick="removeRow(this)" title="remove">✕</button>
 </li>
 {% endfor %}
</ul>

<div class="addbar">
  <select id="addsel">
    {% for sid, t in all_slides %}<option value="{{ sid }}">{{ sid }} — {{ t }}</option>{% endfor %}
  </select>
  <button type="button" onclick="addSlide()">+ Add slide</button>
</div>

<form method="post" action="/download" onsubmit="syncOrder()">
  <input type="hidden" name="client_name" value="{{ ctx.client_name }}">
  <input type="hidden" name="order" id="order">
  <input type="hidden" name="industry" value="{{ ctx.industry }}">
  <input type="hidden" name="transcript" value="{{ ctx.transcript }}">

  {% if gaps %}
  <h3 style="font-size:16px;margin-top:22px">⚠ Needs to be created ({{ gaps|length }})</h3>
  <ul style="list-style:none;padding:0">
   {% for g in gaps %}
   <li style="border:1px solid #d9a441;background:#faeeda;border-radius:6px;padding:9px 11px;margin:6px 0;color:#7a4a06">
     {{ g.detail }}
     {% if g.type == 'needs_case_study' %}
       <label style="margin-left:8px;font-size:13px;color:#345"><input type="checkbox" name="gen" value="{{ g.work_type }}"> Generate with AI</label>
     {% endif %}
   </li>
   {% endfor %}
  </ul>
  <p class="hint">Tick &ldquo;Generate with AI&rdquo; to draft the missing slide from the template (uses your API key); it's appended at the end of the deck.</p>
  {% endif %}

  {% if suggestions %}
  <ul style="list-style:none;padding:0;margin-top:10px">
   {% for s in suggestions %}
   <li style="border:1px solid #b9a9d6;background:#f3eefa;border-radius:6px;padding:8px 10px;margin:6px 0;color:#534b6b">💡 {{ s }}</li>
   {% endfor %}
  </ul>
  {% endif %}

  <button class="dl" type="submit" formaction="/review" style="background:#345">Review &amp; edit text →</button>
  <button class="dl" type="submit">⬇ Download now</button>
</form>
<a class="back" href="/">← build another</a>

<script>
 const list = document.getElementById('list');
 const SLIDE_TITLES = {{ titles|tojson }};
 function renumber(){ [...list.children].forEach((li,i)=> li.querySelector('.pos').textContent = (i+1)+'.'); }
 function removeRow(btn){ btn.closest('li').remove(); renumber(); }
 function addSlide(){
   const id = document.getElementById('addsel').value;
   if(!id) return;
   if([...list.children].some(li => li.dataset.id === id)){ alert(id + ' is already in the deck'); return; }
   list.appendChild(makeRow(id, SLIDE_TITLES[id] || '', 'added manually'));
   renumber();
 }
 function makeRow(id, title, reason){
   const li = document.createElement('li');
   li.draggable = true; li.dataset.id = id;
   const mk = (cls, txt) => { const s = document.createElement('span'); s.className = cls; s.textContent = txt; return s; };
   li.append(mk('pos',''), mk('grip','⠿'), mk('id', id), mk('title', title), mk('reason', reason));
   const mv = (t, dir) => { const b = document.createElement('button'); b.type='button'; b.className='mv'; b.textContent=t; b.onclick=()=>move(b,dir); return b; };
   const rm = document.createElement('button'); rm.type='button'; rm.className='rm'; rm.textContent='✕'; rm.onclick=()=>removeRow(rm);
   li.append(mv('↑',-1), mv('↓',1), rm);
   return li;
 }
 function move(btn, dir){
   const li = btn.closest('li');
   if(dir<0 && li.previousElementSibling) list.insertBefore(li, li.previousElementSibling);
   if(dir>0 && li.nextElementSibling) list.insertBefore(li.nextElementSibling, li);
   renumber();
 }
 let dragEl = null;
 list.addEventListener('dragstart', e => { dragEl = e.target.closest('li'); dragEl.classList.add('drag'); });
 list.addEventListener('dragend',   e => { if(dragEl) dragEl.classList.remove('drag'); renumber(); });
 list.addEventListener('dragover',  e => {
   e.preventDefault();
   const li = e.target.closest('li');
   if(!li || li===dragEl) return;
   const r = li.getBoundingClientRect();
   const after = (e.clientY - r.top) / r.height > 0.5;
   list.insertBefore(dragEl, after ? li.nextElementSibling : li);
 });
 function syncOrder(){ document.getElementById('order').value = [...list.children].map(li => li.dataset.id).join(','); }
 renumber();
</script>
</body></html>
"""


REVIEW_HTML = """
<!doctype html><html><head><meta charset="utf-8"><title>Review &amp; edit</title>
<style>
 body{font-family:Segoe UI,Arial,sans-serif;max-width:820px;margin:40px auto;color:#222}
 h1{font-size:22px} .hint{color:#777;font-size:13px}
 .card{border:1px solid #ddd;border-radius:8px;padding:12px 14px;margin:12px 0;background:#fafafa}
 .cid{font-weight:700;color:#345;font-size:13px;margin-bottom:6px}
 label{display:block;font-size:12px;color:#555;margin:8px 0 2px}
 input{width:100%;padding:7px;font-size:14px;box-sizing:border-box}
 .ctx{color:#999;font-size:12px;margin-top:8px;white-space:pre-wrap}
 .dl{margin:8px 14px 0 0;padding:11px 24px;font-size:16px;background:#1a5;color:#fff;border:0;border-radius:5px;cursor:pointer}
 .back{display:inline-block;margin-top:18px;color:#555}
</style></head><body>
<h1>Review &amp; edit — {{ client }}</h1>
<p class="hint">Edit each slide's title / subtitle below. <b>[CLIENT]</b> has been filled in with the client name.
   Body text is shown for context (not edited here). Then download.</p>
<form method="post" action="/finalize">
  <input type="hidden" name="client_name" value="{{ client }}">
  <input type="hidden" name="order" value="{{ order }}">
  <input type="hidden" name="industry" value="{{ industry }}">
  <input type="hidden" name="transcript" value="{{ transcript }}">
  {% for w in gen %}<input type="hidden" name="gen" value="{{ w }}">{% endfor %}
  {% for c in cards %}
  <div class="card">
    <div class="cid">{{ c.id }}</div>
    {% for idx, label, text in c.fields %}
      <label>{{ label }}</label>
      <input type="text" name="edit__{{ c.id }}__{{ idx }}" value="{{ text }}">
    {% endfor %}
    {% if c.context %}<div class="ctx">{{ c.context }}</div>{% endif %}
  </div>
  {% endfor %}
  <button class="dl" type="submit">⬇ Download PowerPoint</button>
</form>
<a class="back" href="/">← build another</a>
</body></html>
"""


BUILD_BODY = """
<div style="margin:14px 0 18px">
  <div class="eyebrow">Suggested deck</div>
  <h1 class="display" style="font-size:38px">Deck for <span id="deck-client">{{ ctx.client_name }}</span></h1>
  <div class="ctx-chips">
    {% if ctx.industry %}<span class="ctx-chip">Industry <b>{{ ctx.industry }}</b></span>{% endif %}
    {% if ctx.work_types %}<span class="ctx-chip">Work <b>{{ ', '.join(ctx.work_types) }}</b></span>{% endif %}
    {% if ctx.functions %}<span class="ctx-chip">Function <b>{{ ', '.join(ctx.functions) }}</b></span>{% endif %}
    {% if persona_labels %}<span class="ctx-chip">Persona <b>{{ ', '.join(persona_labels) }}</b></span>{% endif %}
    {% if ai_used %}<span class="ai-badge"><i class="ti ti-sparkles"></i> AI-refined from your info</span>{% endif %}
  </div>
</div>

<form method="post" action="/review" onsubmit="syncOrder();var l=document.getElementById('loader');if(l)l.style.display='flex'" class="form-layout">
  <div>
    <div class="card">
      <div class="section-head"><h2 class="sec-title">Slides <span class="hint" style="font-size:14px;font-weight:400">— drag to reorder</span></h2>
        <span class="tag" id="count-tag">{{ picks|length }} slides</span></div>
      <ul id="list" class="slide-list">
        {% for p in picks %}
        <li class="slide-item" draggable="true" data-id="{{ p.slide_id }}">
          <span class="grip"><i class="ti ti-grip-vertical"></i></span>
          <span class="pos"></span>
          <span class="sid" {% if p.skill %}style="background:#2C6E66;color:#fff"{% endif %}>{{ p.tag if p.skill else p.slide_id }}</span>
          <div class="s-main">
            <div class="s-title">{{ p.label if p.skill else titles.get(p.slide_id, '') }}</div>
            <div class="s-reason {{ 'is-weak' if 'WEAK' in p.reason else '' }}" {% if p.skill and 'stale' in p.reason %}style="color:#c0392b;font-weight:600"{% endif %}>{{ p.reason }}</div>
          </div>
          <div class="s-actions">
            <button type="button" class="mini" onclick="move(this,-1)" aria-label="up"><i class="ti ti-chevron-up"></i></button>
            <button type="button" class="mini" onclick="move(this,1)" aria-label="down"><i class="ti ti-chevron-down"></i></button>
            <button type="button" class="mini rm" onclick="removeRow(this)" aria-label="remove"><i class="ti ti-x"></i></button>
          </div>
        </li>
        {% endfor %}
      </ul>
      <div id="resume-empty" style="display:none;color:#6b7280;padding:6px 2px">No deck in progress yet. <a href="/new">Start a new deck</a>, or browse the <a href="/library">library</a> and add slides.</div>
      <div class="addbar">
        <select id="addsel">{% for sid, t in all_slides %}<option value="{{ sid }}">{{ sid }} — {{ t }}</option>{% endfor %}</select>
        <button type="button" class="btn" onclick="addSlide()"><i class="ti ti-plus"></i> Add slide</button>
      </div>
    </div>

    {% if gaps %}
    <div class="card" style="margin-top:18px;border-left:4px solid #c47d27">
      <h2 class="sec-title">Heads-up — not in this deck <span class="hint" style="font-size:14px;font-weight:400">— {{ gaps|length }}</span></h2>
      <p class="hint" style="font-size:13px;margin:0 0 14px">Flagged so you're aware. If you want a slide for any of these, use <b>Create a slide with AI</b> below, or add one from the panel.</p>
      {% for g in gaps %}
      {% if g.type == 'missing_capability' %}
      <div class="gap-row" style="display:flex;align-items:flex-start;gap:8px;background:#fff8f0;border:1px solid #f0d9bf;border-radius:8px;padding:10px 12px;margin-bottom:8px">
        <i class="ti ti-alert-triangle" style="color:#c47d27;margin-top:2px"></i>
        <span class="gap-text"><b>Asked but not in the deck:</b> {{ g.detail }}</span>
      </div>
      {% else %}
      <div class="gap-row" style="display:flex;align-items:flex-start;gap:8px;padding:6px 2px;margin-bottom:4px">
        <i class="ti ti-info-circle" style="color:#6E6E69;margin-top:2px"></i>
        <span class="gap-text">{{ g.detail }}</span>
      </div>
      {% endif %}
      {% endfor %}
    </div>
    {% endif %}

    <div class="card" style="margin-top:18px;border-left:4px solid #2C6E66">
      <h2 class="sec-title" style="margin-bottom:4px"><i class="ti ti-sparkles"></i> Create a slide with AI</h2>
      <p class="hint" style="font-size:13px;margin:0 0 14px">Fill in what you know — the AI writes it in J2W's strict format (6 capabilities · 3 results), self-checks it, and shows it here before you add it.</p>
      <div class="fgrid" style="gap:12px;margin-bottom:10px">
        <div>
          <label style="font-size:12px;font-weight:500;color:#6E6E69;margin-bottom:5px;display:block">Topic / Use case <span style="color:#C2503C">*</span></label>
          <input class="input" id="ca-topic" placeholder="e.g. Fraud detection for a retail bank" style="font-size:13px">
        </div>
        <div>
          <label style="font-size:12px;font-weight:500;color:#6E6E69;margin-bottom:5px;display:block">Industry <span style="color:#9C9C97;font-weight:400">(pre-filled)</span></label>
          <input class="input" id="ca-industry" placeholder="e.g. Banking" style="font-size:13px">
        </div>
      </div>
      <div style="margin-bottom:10px">
        <label style="font-size:12px;font-weight:500;color:#6E6E69;margin-bottom:5px;display:block">Problem / Challenge <span style="color:#C2503C">*</span></label>
        <textarea id="ca-problem" style="width:100%;min-height:54px;font-size:13px" placeholder="What was the client's core challenge?"></textarea>
      </div>
      <div class="fgrid" style="gap:12px;margin-bottom:12px">
        <div>
          <label style="font-size:12px;font-weight:500;color:#6E6E69;margin-bottom:5px;display:block">Solution <span style="color:#9C9C97;font-weight:400">(optional)</span></label>
          <textarea id="ca-solution" style="width:100%;min-height:54px;font-size:13px" placeholder="What did J2W deploy or deliver?"></textarea>
        </div>
        <div>
          <label style="font-size:12px;font-weight:500;color:#6E6E69;margin-bottom:5px;display:block">Results <span style="color:#9C9C97;font-weight:400">(optional)</span></label>
          <textarea id="ca-results" style="width:100%;min-height:54px;font-size:13px" placeholder="Key outcomes — numbers if you have them."></textarea>
        </div>
      </div>
      <button type="button" id="ca-genbtn" class="btn btn-primary" onclick="createAI(this)"><i class="ti ti-wand"></i> Generate</button>
      <div id="ca-loader" style="display:none;margin-top:12px">
        <div style="height:4px;background:#e0ece9;border-radius:99px;overflow:hidden">
          <div id="ca-bar" style="height:100%;width:0%;background:#2C6E66;border-radius:99px;transition:width 0.4s ease"></div>
        </div>
        <p style="font-size:12px;color:#6E6E69;margin:6px 0 0"><i class="ti ti-sparkles" style="font-size:13px"></i> Writing your case study…</p>
      </div>
      <div id="ca-preview" style="display:none;margin-top:14px"></div>
    </div>

    {% if suggested %}
    <div class="card" style="margin-top:18px">
      <h2 class="sec-title">You might also include <span class="hint" style="font-size:14px;font-weight:400">— related, lower priority</span></h2>
      <div class="sug-list">
        {% for s in suggested %}
        <div class="sug-item">
          <span class="sid">{{ s.slide_id }}</span>
          <div class="s-main"><div class="s-title">{{ titles.get(s.slide_id, '') }}</div>
            <div class="s-reason">{{ s.reason }}</div></div>
          <button type="button" class="btn" onclick="addSlide('{{ s.slide_id }}', this)"><i class="ti ti-plus"></i> Add</button>
        </div>
        {% endfor %}
      </div>
    </div>
    {% endif %}

    {% if suggestions %}
    <div class="card" style="margin-top:18px">
      {% for s in suggestions %}<div class="sug-row"><i class="ti ti-bulb"></i><span>{{ s }}</span></div>{% endfor %}
    </div>
    {% endif %}
  </div>

  <aside class="side-card">
    <div class="card">
      <input type="hidden" name="client_name" id="h-client" value="{{ ctx.client_name }}">
      <input type="hidden" name="order" id="order">
      <input type="hidden" name="industry" id="h-industry" value="{{ ctx.industry }}">
      <input type="hidden" name="transcript" id="h-transcript" value="{{ ctx.transcript }}">
      <input type="hidden" name="phase" id="h-phase" value="{{ ctx.phase }}">
      <input type="hidden" name="recipient" id="h-recipient" value="{{ ctx.recipient }}">
      <span id="h-multi">{% for f in ctx.functions %}<input type="hidden" name="functions" value="{{ f }}">{% endfor %}{% for w in ctx.work_types %}<input type="hidden" name="work_types" value="{{ w }}">{% endfor %}</span>
      <div class="summary">
        <div class="sum-row"><span>Slides</span><b id="sum-count">{{ picks|length }}</b></div>
        <div class="sum-row"><span>Gaps flagged</span><b>{{ gaps|length }}</b></div>
        <div class="sum-row"><span>Matching</span><b style="{{ 'color:var(--teal)' if ai_used else '' }}">{{ 'AI' if ai_used else 'Rules' }}</b></div>
      </div>
      <button class="btn btn-primary btn-lg btn-block" type="submit" style="margin-top:18px"><i class="ti ti-arrow-right"></i> Next: Review &amp; edit</button>
      <p class="hint" style="font-size:12px;text-align:center;margin:10px 0 0">Edit text, accept AI slides, then preview &amp; download.</p>
      <a href="/new" class="hint" style="display:block;text-align:center;margin-top:10px;font-size:13px">← Start over</a>
    </div>
  </aside>
</form>

<div id="loader" style="display:none;position:fixed;inset:0;background:rgba(17,17,16,.55);z-index:999;align-items:center;justify-content:center">
  <div style="background:#fff;border-radius:16px;padding:30px 38px;text-align:center;box-shadow:0 20px 60px rgba(0,0,0,.3)">
    <div style="width:40px;height:40px;border:4px solid #e6e6e3;border-top-color:#2C6E66;border-radius:50%;margin:0 auto 16px;animation:j2wspin .8s linear infinite"></div>
    <div style="font-family:'Space Grotesk',Inter,sans-serif;font-weight:700;font-size:18px;color:#111110">Preparing your review…</div>
    <div style="font-size:13px;color:#6b7280;margin-top:5px">Writing any AI slides for the gaps — a few seconds.</div>
  </div>
</div>
<style>@keyframes j2wspin{to{transform:rotate(360deg)}}</style>

<script>
 const list=document.getElementById('list');
 const SLIDE_TITLES={{ titles|tojson }};
 const RESUME={{ 'true' if resume else 'false' }};
 const SERVER_CTX={{ ctx|tojson }};
 const BUILD_ID={{ build_id|tojson }};
 const DECK_KEY='j2w_deck';
 function loadDeck(){try{return JSON.parse(localStorage.getItem(DECK_KEY));}catch(e){return null;}}
 function saveDeck(d){localStorage.setItem(DECK_KEY,JSON.stringify(d));}
 function currentOrder(){return [...list.children].map(li=>li.dataset.id);}
 function persist(){var d=loadDeck()||{}; d.active=true; d.order=currentOrder();
   if(!RESUME){d.ctx=SERVER_CTX;d.buildId=BUILD_ID;} saveDeck(d);}
 function setH(id,v){var e=document.getElementById(id);if(e)e.value=v||'';}
 function fillHiddenCtx(c){c=c||{};setH('h-client',c.client_name);setH('h-industry',c.industry);
   setH('h-transcript',c.transcript);setH('h-phase',c.phase);setH('h-recipient',c.recipient);
   var box=document.getElementById('h-multi');if(box){box.innerHTML='';
     (c.functions||[]).forEach(f=>box.insertAdjacentHTML('beforeend','<input type="hidden" name="functions" value="'+f+'">'));
     (c.work_types||[]).forEach(w=>box.insertAdjacentHTML('beforeend','<input type="hidden" name="work_types" value="'+w+'">'));}}
 function rebuild(order){list.innerHTML='';(order||[]).forEach(id=>list.appendChild(makeRow(id,SLIDE_TITLES[id]||'','in your deck')));}
 function renumber(){[...list.children].forEach((li,i)=>li.querySelector('.pos').textContent=(i+1));
   const n=list.children.length;document.getElementById('sum-count').textContent=n;
   document.getElementById('count-tag').textContent=n+' slides';persist();}
 function move(b,d){const li=b.closest('li');
   if(d<0&&li.previousElementSibling)list.insertBefore(li,li.previousElementSibling);
   if(d>0&&li.nextElementSibling)list.insertBefore(li.nextElementSibling,li);renumber();}
 function removeRow(b){b.closest('li').remove();renumber();}
 function addSlide(id,btn){id=id||document.getElementById('addsel').value;if(!id)return;
   if([...list.children].some(li=>li.dataset.id===id)){if(!btn)alert(id+' is already in the deck');return;}
   list.appendChild(makeRow(id,SLIDE_TITLES[id]||'','added manually'));renumber();
   if(btn){btn.disabled=true;btn.innerHTML='<i class="ti ti-check"></i> Added';}}
 function makeRow(id,title,reason){const li=document.createElement('li');
   li.className='slide-item';li.draggable=true;li.dataset.id=id;
   var tg=id; if(id.indexOf('NEW:')===0)tg='AI'; else if(id==='SK:industry')tg='IND'; else if(id==='SK:skills')tg='SKL'; else if(id.indexOf('SK:')===0)tg='CAP'; else if(id.indexOf('FP:')===0)tg='FOOT';
   li.innerHTML='<span class="grip"><i class="ti ti-grip-vertical"></i></span><span class="pos"></span>'+
    '<span class="sid">'+tg+'</span><div class="s-main"><div class="s-title"></div>'+
    '<div class="s-reason">'+reason+'</div></div><div class="s-actions">'+
    '<button type="button" class="mini" onclick="move(this,-1)"><i class="ti ti-chevron-up"></i></button>'+
    '<button type="button" class="mini" onclick="move(this,1)"><i class="ti ti-chevron-down"></i></button>'+
    '<button type="button" class="mini rm" onclick="removeRow(this)"><i class="ti ti-x"></i></button></div>';
   li.querySelector('.s-title').textContent=title;return li;}
 let dragEl=null;
 list.addEventListener('dragstart',e=>{dragEl=e.target.closest('li');dragEl.classList.add('drag');});
 list.addEventListener('dragend',e=>{if(dragEl)dragEl.classList.remove('drag');renumber();});
 list.addEventListener('dragover',e=>{e.preventDefault();const li=e.target.closest('li');if(!li||li===dragEl)return;
   const r=li.getBoundingClientRect();const after=(e.clientY-r.top)/r.height>0.5;
   list.insertBefore(dragEl,after?li.nextElementSibling:li);});
 function syncOrder(){document.getElementById('order').value=[...list.children].map(li=>li.dataset.id).join(',');}
 function caEsc(s){var e=document.createElement('div');e.textContent=(s==null?'':s);return e.innerHTML;}
 function createAI(btn){
   var topic=(document.getElementById('ca-topic')||{}).value||'';
   var problem=(document.getElementById('ca-problem')||{}).value||'';
   topic=topic.trim(); problem=problem.trim();
   if(!topic||!problem){alert('Please fill in Topic / use case and Problem — both are required.');return;}
   var solution=((document.getElementById('ca-solution')||{}).value||'').trim();
   var results=((document.getElementById('ca-results')||{}).value||'').trim();
   var ind=((document.getElementById('ca-industry')||{}).value||'').trim()||(SERVER_CTX&&SERVER_CTX.industry)||'';
   var brief=topic+'. Problem: '+problem;
   if(solution) brief+=' Solution: '+solution;
   if(results) brief+=' Results: '+results;
   btn.disabled=true; var old=btn.innerHTML; btn.innerHTML='<i class="ti ti-loader"></i> Writing…';
   var loader=document.getElementById('ca-loader');
   var bar=document.getElementById('ca-bar');
   if(loader) loader.style.display='block';
   var pct=0; var ticker=setInterval(function(){pct=Math.min(pct+Math.random()*7+3,88);if(bar)bar.style.width=pct+'%';},500);
   var fd=new FormData(); fd.append('brief',brief);
   fd.append('industry',ind);
   fd.append('client_name',(SERVER_CTX&&SERVER_CTX.client_name)||'');
   fetch('/create_ai',{method:'POST',body:fd}).then(r=>r.json()).then(d=>{
     clearInterval(ticker); if(bar)bar.style.width='100%';
     setTimeout(function(){
       btn.disabled=false; btn.innerHTML=old;
       if(loader)loader.style.display='none'; if(bar)bar.style.width='0%';
       if(!d.ok){alert(d.error||'Could not generate');return;}
       window._ca=d; caRender(d);
     },350);
   }).catch(function(){clearInterval(ticker);btn.disabled=false;btn.innerHTML=old;if(loader)loader.style.display='none';alert('Could not generate');});
 }
 function caRegen(){createAI(document.getElementById('ca-genbtn'));}
 function caRender(d){
   var caps=(d.capabilities||[]).map(c=>'<li>'+caEsc(c)+'</li>').join('');
   var res=(d.results||[]).map(c=>'<li>'+caEsc(c)+'</li>').join('');
   var rv=d.review||{}; var strong=((rv.quality||'')+'').toLowerCase().indexOf('strong')>-1;
   document.getElementById('ca-preview').innerHTML=
    '<div class="card" style="background:#f3f8f7;border:1px solid #bcdfd8">'+
    '<div style="font-weight:700;font-size:16px">'+caEsc(d.title)+'</div>'+
    '<div class="hint" style="font-size:12px;margin-bottom:8px">'+caEsc(d.subhead)+'</div>'+
    '<div style="font-size:13px"><b>Challenge:</b> '+caEsc(d.challenge)+'</div>'+
    '<div style="font-size:13px;margin-top:4px"><b>Solution:</b> '+caEsc(d.solution)+'</div>'+
    '<div style="display:flex;gap:18px;margin-top:8px;flex-wrap:wrap">'+
      '<div style="flex:1;min-width:220px"><b style="font-size:12px">Capabilities</b><ul style="margin:4px 0 0;padding-left:18px;font-size:12px">'+caps+'</ul></div>'+
      '<div style="flex:1;min-width:220px"><b style="font-size:12px">Results</b><ul style="margin:4px 0 0;padding-left:18px;font-size:12px">'+res+'</ul></div></div>'+
    '<div style="margin-top:10px;font-size:12px;padding:8px;border-radius:6px;background:'+(strong?'#e8f3f1':'#fdecea')+';color:'+(strong?'#1f5a52':'#8a2a1e')+'">'+
      '<b>Self-review &mdash; '+caEsc(rv.quality||'')+'.</b> Weakest: '+caEsc(rv.weakest||'None')+'. Fix: '+caEsc(rv.fix||'None')+'</div>'+
    '<div style="display:flex;gap:8px;margin-top:10px">'+
      '<button type="button" class="btn btn-primary" onclick="addCreated()"><i class="ti ti-plus"></i> Add to deck</button>'+
      '<button type="button" class="btn" onclick="caRegen()"><i class="ti ti-refresh"></i> Regenerate</button>'+
      '<button type="button" class="btn" onclick="caDiscard()"><i class="ti ti-x"></i> Discard</button></div></div>';
   document.getElementById('ca-preview').style.display='block';
 }
 function addCreated(){
   var d=window._ca; if(!d) return;
   if([...list.children].some(li=>li.dataset.id===d.id)){alert('Already added');return;}
   list.appendChild(makeRow(d.id, d.title, 'created with AI')); renumber();
   caDiscard(); ['ca-topic','ca-problem','ca-solution','ca-results'].forEach(function(id){var el=document.getElementById(id);if(el)el.value='';});
 }
 function caDiscard(){var p=document.getElementById('ca-preview');if(p)p.style.display='none'; window._ca=null;}
 (function init(){
   var d=loadDeck();
   if(RESUME){
     if(d&&d.order&&d.order.length){rebuild(d.order);fillHiddenCtx(d.ctx);
       var dc=document.getElementById('deck-client');
       if(dc&&d.ctx&&d.ctx.client_name) dc.textContent=d.ctx.client_name;}
     else{var emp=document.getElementById('resume-empty');if(emp)emp.style.display='';}
   } else if(d&&d.buildId===BUILD_ID&&d.order&&d.order.length){
     // Same build re-opened (e.g. browser Back) — restore it as left, keeping any library adds.
     rebuild(d.order);
   } // else: brand-new build — keep the server-rendered picks (persist() below seeds the deck).
   renumber();
   var caInd=document.getElementById('ca-industry');
   if(caInd&&SERVER_CTX&&SERVER_CTX.industry) caInd.value=SERVER_CTX.industry;
 })();
 window.addEventListener('pageshow',function(){var l=document.getElementById('loader');if(l)l.style.display='none';});
</script>
"""


REVIEW_BODY = """
<div style="margin:14px 0 22px">
  <div class="eyebrow">Review &amp; edit</div>
  <h1 class="display" style="font-size:38px">Final touches for {{ client }}</h1>
  <p class="lede">Edit each slide's title and subtitle. <b>[CLIENT]</b> is filled with the client name. {% if ai_cards %}<b>AI-written slides are below — read each one and Accept or Reject it.</b>{% endif %}</p>
</div>
<form method="post" action="/finalize" class="form-layout" onsubmit="var l=document.getElementById('loader');if(l)l.style.display='flex'">
  <div>
    {% if ai_cards %}
    <div class="card" style="border-left:5px solid #2C6E66;background:#f3f8f7;margin-bottom:18px">
      <h2 class="sec-title" style="margin-bottom:6px"><i class="ti ti-sparkles"></i> AI-written slides — your call</h2>
      <p class="hint" style="font-size:13px;margin:0">These were written by AI to fill gaps. Edit if you like, then <b>Accept</b> (adds it to the deck &amp; your library, client-ready) or <b>Reject</b> (drops it). Nothing leaves this screen unchecked.</p>
    </div>
    {% for a in ai_cards %}
    <div class="card rev-card" style="border:1px solid #bcdfd8">
      <div class="rev-head"><span class="sid" style="background:#2C6E66;color:#fff">AI · {{ a.work_type }}</span>
        <div style="margin-left:auto;display:flex;gap:14px;font-size:14px">
          <label style="display:flex;align-items:center;gap:5px;cursor:pointer"><input type="radio" name="ai_decision__{{ a.id }}" value="accept" checked> Accept</label>
          <label style="display:flex;align-items:center;gap:5px;cursor:pointer;color:#c0392b"><input type="radio" name="ai_decision__{{ a.id }}" value="reject"> Reject</label>
        </div>
      </div>
      <div class="field"><label>Title</label>
        <input class="input" name="ai_title__{{ a.id }}" value="{{ a.title }}"></div>
      <div class="field"><label>Keywords</label>
        <input class="input" name="ai_keywords__{{ a.id }}" value="{{ a.keywords }}"></div>
      <div class="field" style="margin-bottom:0"><label>Bullets (one per line)</label>
        <textarea name="ai_bullets__{{ a.id }}" style="min-height:96px">{{ a.bullets }}</textarea></div>
    </div>
    {% endfor %}
    {% endif %}

    {% for c in cards %}
    <div class="card rev-card">
      <div class="rev-head"><span class="sid">{{ c.id }}</span></div>
      {% for idx, label, text in c.fields %}
      <div class="field" {% if loop.last %}style="margin-bottom:0"{% endif %}>
        <label>{{ label }}</label>
        <input class="input" name="edit__{{ c.id }}__{{ idx }}" value="{{ text }}">
      </div>
      {% endfor %}
      {% if c.context %}<div class="rev-ctx">{{ c.context }}</div>{% endif %}
    </div>
    {% endfor %}
  </div>

  <aside class="side-card">
    <div class="card">
      <input type="hidden" name="client_name" value="{{ client }}">
      <input type="hidden" name="order" value="{{ order }}">
      <input type="hidden" name="ai_ids" value="{{ ai_ids }}">
      <input type="hidden" name="industry" value="{{ industry }}">
      <input type="hidden" name="transcript" value="{{ transcript }}">
      <input type="hidden" name="phase" value="{{ phase }}">
      <input type="hidden" name="recipient" value="{{ recipient }}">
      {% for f in functions %}<input type="hidden" name="functions" value="{{ f }}">{% endfor %}
      {% for w in work_types %}<input type="hidden" name="work_types" value="{{ w }}">{% endfor %}
      <div class="summary">
        <div class="sum-row"><span>Library slides</span><b>{{ cards|length }}</b></div>
        <div class="sum-row"><span>AI slides</span><b style="{{ 'color:var(--teal)' if ai_cards else '' }}">{{ ai_cards|length }}</b></div>
      </div>
      <button class="btn btn-primary btn-lg btn-block" type="submit" style="margin-top:18px"><i class="ti ti-download"></i> Build &amp; download</button>
      <a href="javascript:history.back()" class="hint" style="display:block;text-align:center;margin-top:14px;font-size:13px">← Back to slides</a>
    </div>
  </aside>
</form>

<div id="loader" style="display:none;position:fixed;inset:0;background:rgba(17,17,16,.55);z-index:999;align-items:center;justify-content:center">
  <div style="background:#fff;border-radius:16px;padding:30px 38px;text-align:center;box-shadow:0 20px 60px rgba(0,0,0,.3)">
    <div style="width:40px;height:40px;border:4px solid #e6e6e3;border-top-color:#2C6E66;border-radius:50%;margin:0 auto 16px;animation:j2wspin .8s linear infinite"></div>
    <div style="font-family:'Space Grotesk',Inter,sans-serif;font-weight:700;font-size:18px;color:#111110">Building your deck…</div>
    <div style="font-size:13px;color:#6b7280;margin-top:5px">Assembling slides and your download — a few seconds.</div>
  </div>
</div>
<style>@keyframes j2wspin{to{transform:rotate(360deg)}}</style>
<script>window.addEventListener('pageshow',function(){var l=document.getElementById('loader');if(l)l.style.display='none';});</script>
"""


LIBRARY_BODY = """
<div style="margin:14px 0 6px">
  <div class="eyebrow">Slide library</div>
  <h1 class="display" style="font-size:38px">All slides <span class="hint" style="font-size:24px">· <span id="lib-count">{{ total }}</span></span></h1>
  <p class="lede">Browse, search and filter the whole warehouse. Keywords are what the matcher runs on.</p>
</div>

<div class="toolbar">
  <div class="search-box"><i class="ti ti-search"></i>
    <input class="input" id="lib-search" placeholder="Search title, keywords, ID…"></div>
  <div class="fchips">
    <button class="fchip active" data-wtf="">All</button>
    <button class="fchip" data-wtf="WORKFORCE">Workforce</button>
    <button class="fchip" data-wtf="AI_POD">AI Pods</button>
    <button class="fchip" data-wtf="MS">Managed</button>
  </div>
  <div class="fchips">
    <button class="fchip active" data-kindf="">Any kind</button>
    <button class="fchip" data-kindf="CASE_STUDY">Case study</button>
    <button class="fchip" data-kindf="STANDARD">Standard</button>
  </div>
  <select id="ind-sel" style="width:auto">
    <option value="">All industries</option>
    {% for i in industries %}<option value="{{ i }}">{{ i|title }}</option>{% endfor %}
  </select>
</div>

<div class="lib-grid" id="lib-grid">
  {% for s in slides %}
  <div class="lib-card" data-wt="{{ s.wt }}" data-kind="{{ s.kind }}" data-ind="{{ s.ind }}" data-text="{{ s.search }}">
    <div class="lib-head"><span class="sid">{{ s.id }}</span>
      <span class="lib-tags">{{ s.wt.replace('_',' ')|title if s.wt else '—' }} · {{ s.kind|title }}</span></div>
    <div class="lib-title">{{ s.title }}</div>
    {% if s.kw %}<div class="kw">{% for k in s.kw %}<span>{{ k }}</span>{% endfor %}</div>{% endif %}
    <div class="lib-foot" style="display:flex;align-items:center;justify-content:space-between;gap:8px">
      <span style="overflow:hidden;text-overflow:ellipsis;white-space:nowrap">{{ s.ind|title if s.ind else 'No industry' }}{% if s.fn %} · {{ s.fn.replace('_',' ')|title }}{% endif %}</span>
      <span style="display:flex;gap:6px;flex:none">
        <button type="button" class="mini deck-add" data-id="{{ s.id }}" onclick="libAdd('{{ s.id }}',this)" title="Add to current deck" aria-label="Add {{ s.id }} to deck"><i class="ti ti-plus"></i></button>
        <a class="mini" href="/slide/{{ s.id }}/download" title="Download this slide" aria-label="Download {{ s.id }}"><i class="ti ti-download"></i></a>
      </span>
    </div>
  </div>
  {% endfor %}
  <div class="lib-empty" id="lib-empty" style="display:none">No slides match your filters.</div>
</div>

<a id="deck-pill" href="/deck" style="position:fixed;right:22px;bottom:22px;display:none;align-items:center;gap:8px;background:#111110;color:#fff;padding:11px 16px;border-radius:999px;box-shadow:0 6px 20px rgba(0,0,0,.25);font-size:14px;font-weight:600;text-decoration:none;z-index:50">
  <i class="ti ti-stack-2"></i> Open deck <b id="deck-pill-n" style="background:#2C6E66;border-radius:999px;padding:1px 9px;font-size:13px">0</b>
</a>

<script>
 const DECK_KEY='j2w_deck';
 function loadDeck(){try{return JSON.parse(localStorage.getItem(DECK_KEY));}catch(e){return null;}}
 function saveDeck(d){localStorage.setItem(DECK_KEY,JSON.stringify(d));}
 function updatePill(){var d=loadDeck();var n=(d&&d.order)?d.order.length:0;
   document.getElementById('deck-pill-n').textContent=n;
   document.getElementById('deck-pill').style.display=n?'inline-flex':'none';}
 function libAdd(id,btn){var d=loadDeck()||{active:true,ctx:{client_name:'',industry:'',transcript:'',phase:'',recipient:'',functions:[],work_types:[]},order:[]};
   d.active=true;if(!d.order)d.order=[];
   if(d.order.indexOf(id)===-1){d.order.push(id);saveDeck(d);}
   if(btn){btn.innerHTML='<i class="ti ti-check"></i>';btn.style.color='#2C6E66';}
   updatePill();}
 (function(){var d=loadDeck();var inDeck=(d&&d.order)?d.order:[];
   document.querySelectorAll('.deck-add').forEach(b=>{if(inDeck.indexOf(b.dataset.id)>-1){b.innerHTML='<i class="ti ti-check"></i>';b.style.color='#2C6E66';}});
   updatePill();})();
 const cards=[...document.querySelectorAll('.lib-card')];
 const f={wt:'',kind:'',ind:'',q:''};
 function apply(){let n=0;
   cards.forEach(c=>{const ok=(!f.wt||c.dataset.wt===f.wt)&&(!f.kind||c.dataset.kind===f.kind)
     &&(!f.ind||c.dataset.ind===f.ind)&&(!f.q||c.dataset.text.indexOf(f.q)>-1);
     c.style.display=ok?'':'none';if(ok)n++;});
   document.getElementById('lib-count').textContent=n;
   document.getElementById('lib-empty').style.display=n?'none':'';}
 document.getElementById('lib-search').addEventListener('input',e=>{f.q=e.target.value.toLowerCase();apply();});
 document.getElementById('ind-sel').addEventListener('change',e=>{f.ind=e.target.value;apply();});
 function wire(attr,key){document.querySelectorAll('[data-'+attr+']').forEach(b=>b.addEventListener('click',()=>{
   f[key]=b.getAttribute('data-'+attr);
   b.parentNode.querySelectorAll('.fchip').forEach(x=>x.classList.remove('active'));
   b.classList.add('active');apply();}));}
 wire('wtf','wt');wire('kindf','kind');
</script>
"""


TEMPLATES_BODY = """
<div style="margin:14px 0 22px">
  <div class="eyebrow">Templates</div>
  <h1 class="display" style="font-size:38px">Slide templates</h1>
  <p class="lede">When a slide is missing, AI writes the content and the engine pours it into one of these templates by replacing markers. Pluggable — add or swap anytime, no code change.</p>
</div>

<div class="grid" style="grid-template-columns:repeat(auto-fill,minmax(300px,1fr))">
  {% for t in items %}
  <div class="card">
    <div class="lib-head"><span class="sid">{{ t.name }}</span>
      <span class="chip flat">{{ t.status }}</span></div>
    <div class="lib-title" style="margin-bottom:10px">For {{ t.name.replace('_',' ') }} slides</div>
    <p class="hint" style="font-size:12px;margin:0 0 9px">Marker tokens it fills</p>
    <div class="kw">{% for m in t.markers %}<span>{{ m }}</span>{% endfor %}{% if not t.markers %}<span>none found</span>{% endif %}</div>
  </div>
  {% endfor %}
  {% if not items %}<div class="lib-empty">No templates found.</div>{% endif %}
</div>

<div class="section card card-dark">
  <h2 class="sec-title" style="color:#fff;margin-bottom:18px">Add or swap a template</h2>
  <div class="steps">
    <div class="step"><div class="n">1</div><div class="st" style="color:#bdbdb8"><b style="color:#fff">Design the slide</b>Build it in PowerPoint with the real J2W look.</div></div>
    <div class="step"><div class="n">2</div><div class="st" style="color:#bdbdb8"><b style="color:#fff">Add markers</b>Place {{ '{{TITLE}}' }}, {{ '{{KEYWORDS}}' }} and {{ '{{BULLETS}}' }} where text should flow in.</div></div>
    <div class="step"><div class="n">3</div><div class="st" style="color:#bdbdb8"><b style="color:#fff">Tag it</b>In the slide's notes add a line: J2W_TEMPLATE: case_study</div></div>
    <div class="step"><div class="n">4</div><div class="st" style="color:#bdbdb8"><b style="color:#fff">Drop into templates.pptx</b>The engine picks it up automatically — no code change.</div></div>
  </div>
</div>
"""


STAGING_BODY = """
<div style="margin:14px 0 22px">
  <div class="eyebrow">AI history</div>
  <h1 class="display" style="font-size:38px">AI-written slides</h1>
  <p class="lede">A read-only record of every slide AI has written to fill a gap. You accept or reject each one on the <b>Review &amp; edit</b> step while building a deck — <b>Accepted</b> slides join the master library (client-ready); <b>Rejected</b> ones are dropped. This page is just the log.</p>
</div>

{% if items %}
<div class="card"><div class="list">
  {% for s in items %}
  <div class="row">
    <div class="ic">
      {% if s.status == 'approved' %}<i class="ti ti-check" style="color:#2C6E66"></i>
      {% elif s.status == 'discarded' %}<i class="ti ti-x" style="color:#c0392b"></i>
      {% else %}<i class="ti ti-clock" style="color:#9aa0a6"></i>{% endif %}
    </div>
    <div class="main">
      <div class="t">{{ s.title }}</div>
      <div class="s">
        {% if s.status == 'approved' %}<b style="color:#2C6E66">Accepted</b> · now {{ s.promoted_id or 'in library' }}
        {% elif s.status == 'discarded' %}<b style="color:#c0392b">Rejected</b>
        {% else %}<b style="color:#9aa0a6">Not yet decided</b>{% endif %}
        · {{ s.work_type }} · {{ s.industry or 'any industry' }}
        {% if s.for_client %} · for {{ s.for_client }}{% endif %}
        {% if s.created_at %} · {{ s.created_at }}{% endif %}
      </div>
    </div>
  </div>
  {% endfor %}
</div></div>
{% else %}
<div class="card"><p class="muted" style="margin:0">Nothing yet. When AI writes a slide to fill a gap, it's logged here after you accept or reject it on the Review step.</p></div>
{% endif %}
"""


MEETINGS_BODY = """
<div style="margin:14px 0 22px">
  <div class="eyebrow">Deck repository</div>
  <h1 class="display" style="font-size:38px">All created decks</h1>
  <p class="lede">Find a previously generated deck. Filter by industry, work type and phase — newest first. Everyone on the team sees all decks.</p>
</div>

<form method="get" action="/meetings" class="card" style="margin-bottom:20px">
  <div class="fgrid">
    <div class="field"><label>Industry</label>
      <select name="industry"><option value="">Any</option>
        {% for i in industries %}<option {{ 'selected' if f_ind==i else '' }}>{{ i }}</option>{% endfor %}
      </select></div>
    <div class="field"><label>Work type</label>
      <select name="work_type"><option value="">Any</option>
        {% for w in work_types %}<option value="{{ w }}" {{ 'selected' if f_wt==w else '' }}>{{ wt_labels.get(w, w) }}</option>{% endfor %}
      </select></div>
    <div class="field"><label>Phase</label>
      <select name="phase"><option value="">Any</option>
        {% for p in phases %}<option {{ 'selected' if f_phase==p else '' }}>{{ p }}</option>{% endfor %}
      </select></div>
    <div class="field" style="display:flex;align-items:flex-end;gap:10px">
      <button class="btn btn-primary" type="submit"><i class="ti ti-search"></i> Search</button>
      <a class="btn" href="/meetings">Clear</a>
    </div>
  </div>
</form>

<div class="card">
  <div class="section-head" style="margin-bottom:14px">
    <h2 class="sec-title">Results</h2><span class="tag">{{ total }} meeting{{ '' if total==1 else 's' }}</span>
  </div>
  {% if rows %}
  <table style="width:100%;border-collapse:collapse;font-size:14px">
    <thead><tr style="text-align:left;color:#6b7280;border-bottom:1px solid #e6e6e3">
      <th style="padding:8px 10px">Client</th><th style="padding:8px 10px">Salesperson</th>
      <th style="padding:8px 10px">Recipient</th><th style="padding:8px 10px">Date</th>
      <th style="padding:8px 10px">Phase</th><th style="padding:8px 10px">Deck</th>
    </tr></thead>
    <tbody>
      {% for r in rows %}
      <tr style="border-bottom:1px solid #f0f0ee">
        <td style="padding:9px 10px"><b>{{ r.client }}</b><div class="hint" style="font-size:12px">{{ r.industry or '—' }}</div></td>
        <td style="padding:9px 10px">{{ r.salesperson }}</td>
        <td style="padding:9px 10px">{{ r.recipient or '—' }}</td>
        <td style="padding:9px 10px;white-space:nowrap">{{ r.generated_at }}</td>
        <td style="padding:9px 10px"><span class="chip flat">{{ r.phase or '—' }}</span></td>
        <td style="padding:9px 10px"><a href="/output/{{ r.deck_file }}">{{ r.deck_file }}</a></td>
      </tr>
      {% endfor %}
    </tbody>
  </table>
  {% else %}
  <p class="muted" style="margin:0">No matching decks yet. Generate a deck and it will appear here automatically.</p>
  {% endif %}
</div>
"""


PREVIEW_BODY = """
<div style="max-width:560px;margin:40px auto;text-align:center">
  <div style="width:74px;height:74px;border-radius:50%;background:#e8f3f1;display:flex;align-items:center;justify-content:center;margin:0 auto 18px">
    <i class="ti ti-circle-check" style="font-size:42px;color:#2C6E66"></i>
  </div>
  <div class="eyebrow">Done</div>
  <h1 class="display" style="font-size:34px">Your deck is ready{% if client %}, for {{ client }}{% endif %}</h1>
  <p class="lede" style="margin:8px 0 22px">{{ count }} slides assembled. Click the button below to download your deck.</p>
  <a class="btn btn-primary btn-lg" href="/output/{{ filename }}?dl=1"><i class="ti ti-download"></i> Download .pptx</a>
  <a class="btn btn-lg" href="/new" style="margin-left:8px"><i class="ti ti-plus"></i> Build another</a>
  <p class="hint" style="font-size:12px;margin-top:18px">If a file of the same name is open in PowerPoint, close it first, then click Download.</p>
</div>
"""


def _safe(name):
    return re.sub(r"[^A-Za-z0-9_-]+", "_", name).strip("_") or "Client"


@app.route("/")
@app.route("/new")
def home():
    try:
        lib_count = len(json.load(open("tagged_library.json", encoding="utf-8")))
    except Exception:
        lib_count = 0
    body = render_template_string(NEW_FORM_BODY, industries=INDUSTRIES, functions=FUNCTIONS,
                                  phases=PHASES, library_count=lib_count, error="")
    return shell(body, active="new", crumb="<b>New deck</b> / Context")


@app.route("/dashboard")
def dashboard():
    return shell(render_template_string(DASHBOARD_BODY, **_dash_stats()),
                 active="home", crumb="<b>Dashboard</b> / Overview", tabs=["Overview", "Activity"])


@app.route("/library")
def library():
    try:
        recs = json.load(open("tagged_library.json", encoding="utf-8"))
    except Exception:
        recs = []
    slides = []
    for r in recs:
        t = r.get("tags", {})
        slides.append({
            "id": r["slide_id"], "title": r.get("title", ""),
            "wt": t.get("work_type", {}).get("value") or "",
            "kind": t.get("kind", {}).get("value") or "",
            "ind": t.get("industry", {}).get("value") or "",
            "fn": t.get("function", {}).get("value") or "",
            "kw": r.get("keywords", [])[:6],
            "search": (r["slide_id"] + " " + r.get("title", "") + " " +
                       " ".join(r.get("keywords", []))).lower(),
        })
    industries = sorted({s["ind"] for s in slides if s["ind"]})
    body = render_template_string(LIBRARY_BODY, slides=slides, industries=industries, total=len(slides))
    return shell(body, active="library", crumb="<b>Library</b> / All slides")


@app.route("/staging")
def staging_page():
    # read-only history, newest first (records have no timestamp pre-this build -> keep order)
    items = list(reversed(staging.all_items()))
    body = render_template_string(STAGING_BODY, items=items)
    return shell(body, active="staging", crumb="<b>AI history</b>")


@app.route("/staging/<sid>/approve", methods=["POST"])
def staging_approve(sid):
    staging.promote(sid)
    return redirect("/staging")


@app.route("/staging/<sid>/discard", methods=["POST"])
def staging_discard(sid):
    staging.discard(sid)
    return redirect("/staging")


@app.route("/templates")
def templates_page():
    items = []
    try:
        for name, slide in slide_generator.list_templates().items():
            markers = set()
            text = ""
            for sh in slide.shapes:
                if sh.has_text_frame:
                    markers.update(re.findall(r"\{\{[A-Z]+\}\}", sh.text_frame.text))
                    text += sh.text_frame.text
            status = "placeholder" if "Generated slide" in text else "active"
            items.append({"name": name, "markers": sorted(markers), "status": status})
    except Exception:
        items = []
    body = render_template_string(TEMPLATES_BODY, items=items)
    return shell(body, active="templates", crumb="<b>Templates</b>")


@app.route("/meetings")
def meetings():
    f_ind = request.args.get("industry", "").strip()
    f_wt = request.args.get("work_type", "").strip()
    f_phase = request.args.get("phase", "").strip()
    rows = meeting_log.all_meetings()          # newest first
    if f_ind:
        rows = [r for r in rows if r.get("industry") == f_ind]
    if f_wt:
        rows = [r for r in rows if f_wt in r.get("work_types", [])]
    if f_phase:
        rows = [r for r in rows if r.get("phase") == f_phase]
    body = render_template_string(MEETINGS_BODY, rows=rows, total=len(rows),
                                  industries=INDUSTRIES, work_types=WORK_TYPES,
                                  phases=PHASES, wt_labels=WT_LABELS,
                                  f_ind=f_ind, f_wt=f_wt, f_phase=f_phase)
    return shell(body, active="meetings", crumb="<b>Deck repository</b> / All created decks")


@app.route("/deck")
def deck_resume():
    """Re-open the deck-in-progress (held in the browser's deck tray). The list and
    context are hydrated client-side from localStorage; the server just supplies
    the slide catalogue."""
    titles = matcher._title_lookup()
    all_slides = sorted(titles.items(), key=lambda kv: matcher._num(kv[0]))
    empty_ctx = {"client_name": "", "industry": "", "transcript": "",
                 "phase": "", "recipient": "", "functions": [], "work_types": []}
    body = render_template_string(BUILD_BODY, ctx=empty_ctx, picks=[], gaps=[],
                                  titles=titles, all_slides=all_slides,
                                  suggestions=[], suggested=[], ai_used=False,
                                  resume=True, build_id="")
    return shell(body, active="new", crumb="<b>New deck</b> / Your deck")


@app.route("/create_ai", methods=["POST"])
def create_ai():
    """The 'Create with AI' button: write a full structured CASE STUDY from a
    free-text brief (strict format + self-review). Stages it and returns the content
    as JSON so the page shows it inline. Added to the deck as a NEW:<id> order item;
    built into THIS deck at finalize (not promoted to the master library)."""
    brief = request.form.get("brief", "").strip()
    if not brief:
        return {"ok": False, "error": "Please describe the slide you want."}, 400
    industry = request.form.get("industry", "")
    client = request.form.get("client_name", "")
    content = slide_generator.draft_case_study(brief, {"industry": industry})
    content["kind"] = "user_created"
    rec = staging.add(content, "", industry, client)
    return {"ok": True, "id": "NEW:" + rec["id"],
            "title": content["title"], "subhead": content["subhead"],
            "challenge": content["challenge"], "solution": content["solution"],
            "capabilities": content["capabilities"], "results": content["results"],
            "review": content["review"]}


@app.route("/build", methods=["POST"])
def build():
    ctx = {
        "client_name": request.form.get("client_name", "Client").strip(),
        "industry": request.form.get("industry", "").strip(),
        "work_types": request.form.getlist("work_types"),
        "functions": request.form.getlist("functions"),
        "phase": request.form.get("phase", "").strip(),
        "recipient": request.form.get("recipient", "").strip(),
        "salesperson": current_salesperson(),
        "transcript": request.form.get("transcript", "").strip(),
    }
    # Backstop for the browser's at-least-one-work-type check (e.g. JS disabled).
    if not ctx["work_types"]:
        try:
            lib_count = len(json.load(open("tagged_library.json", encoding="utf-8")))
        except Exception:
            lib_count = 0
        body = render_template_string(NEW_FORM_BODY, industries=INDUSTRIES, functions=FUNCTIONS,
                                      phases=PHASES, library_count=lib_count,
                                      error="Please select at least one work type.")
        return shell(body, active="new", crumb="<b>New deck</b> / Context")
    result = matcher.plan(ctx, use_ai=True)   # AI is always on now
    # Gaps are FLAGS only now (no inline generation) — nothing to pre-fill here.
    titles = matcher._title_lookup()

    # --- skills slides (PURE Workforce only): auto-add after the standard block,
    #     before the case studies; labeled + removable in the panel ---
    sk = skills.candidates(ctx)
    if sk:
        picks = result["picks"]
        insert_at = next((i for i, p in enumerate(picks) if p["reason"].startswith("case")), None)
        if insert_at is None:
            insert_at = next((i for i, p in enumerate(picks) if p["slide_id"] in matcher.PIN_TO_END),
                             len(picks))
        sk_picks = []
        _chip = {"industry_strength": "IND", "skill_deepdive": "SKL",
                 "company_footprint": "FOOT"}
        for c in sk:
            titles[c["id"]] = c["label"]
            reason = {
                "industry_strength": "auto-added — RFI: industry strength slide",
                "skill_deepdive":    "auto-added — RFI: skills deployed slide",
                "company_footprint": "auto-added — RFI: existing client relationship",
            }.get(c["kind"], "auto-added — RFI data slide")
            sk_picks.append({"slide_id": c["id"], "reason": reason, "skill": True,
                             "tag": _chip.get(c["kind"], "SKL"),
                             "label": c["label"]})
        picks[insert_at:insert_at] = sk_picks
    all_slides = sorted(titles.items(), key=lambda kv: matcher._num(kv[0]))
    body = render_template_string(BUILD_BODY, ctx=ctx, picks=result["picks"],
                                  gaps=result["gaps"], titles=titles, all_slides=all_slides,
                                  suggestions=result.get("suggestions", []),
                                  suggested=result.get("suggested", []),
                                  ai_used=result.get("ai_used", False),
                                  persona_labels=result.get("persona_labels", []),
                                  resume=False, build_id=uuid.uuid4().hex)
    return shell(body, active="new", crumb="<b>New deck</b> / Suggested slides")


def _file_busy_page(err):
    """Friendly page when the output .pptx can't be written (open/locked/syncing)."""
    body = ("<div class='card' style='border-left:5px solid #c0392b;background:#fdecea;"
            "color:#8a2a1e'><h2 class='sec-title' style='color:#8a2a1e'>Couldn't save the deck</h2>"
            "<p style='margin:8px 0 0'>%s</p>"
            "<p style='margin:10px 0 0;font-size:13px'>Tip: if a previous version of this "
            "deck is open in PowerPoint, close it, then go <a href='javascript:history.back()'>back</a> "
            "and click Download again.</p></div>" % err)
    return shell(body, active="new", crumb="<b>New deck</b> / Preview")


@app.route("/download", methods=["POST"])
def download():
    client = request.form.get("client_name", "Client").strip()
    order = request.form.get("order", "")
    ids = [x for x in order.split(",") if x]   # the FINAL drag-reordered order
    if not ids:
        abort(400)
    os.makedirs(OUTPUT_DIR, exist_ok=True)
    filename = f"Tailored_Deck_{_safe(client)}.pptx"
    path = os.path.join(OUTPUT_DIR, filename)
    try:
        assembler.build_deck(ids, out=path)    # build in the user's exact order
        if client:
            editor.replace_tokens(path, {"[CLIENT]": client, "[Client]": client, "[client]": client})
        unverified = _maybe_generate(path)     # add any AI-generated gap slides
    except (PermissionError, BadZipFile) as e:
        return _file_busy_page(e)
    meeting_log.save(                          # auto-log this meeting (no extra step)
        client=client,
        industry=request.form.get("industry", ""),
        functions=request.form.getlist("functions"),
        work_types=request.form.getlist("work_types"),
        phase=request.form.get("phase", ""),
        recipient=request.form.get("recipient", ""),
        salesperson=current_salesperson(),
        slide_ids=ids,
        deck_file=filename,
    )
    count = len(list(Presentation(path).slides))
    body = render_template_string(PREVIEW_BODY, client=client, filename=filename, count=count)
    return shell(body, active="new", crumb="<b>New deck</b> / Done")


@app.route("/output/<path:fname>")
def output_file(fname):
    fname = os.path.basename(fname)
    path = os.path.join(OUTPUT_DIR, fname)
    if not os.path.exists(path):
        abort(404)
    return send_file(path, as_attachment=bool(request.args.get("dl")), download_name=fname)


@app.route("/slide/<sid>/download")
def slide_download(sid):
    sid = sid.upper()
    os.makedirs(OUTPUT_DIR, exist_ok=True)
    path = os.path.join(OUTPUT_DIR, f"Slide_{sid}.pptx")
    kept, _ = assembler.build_deck([sid], out=path)
    if not kept:
        abort(404)
    return send_file(path, as_attachment=True, download_name=f"{sid}.pptx")


@app.route("/review", methods=["POST"])
def review():
    client = request.form.get("client_name", "Client").strip()
    ids = [x for x in request.form.get("order", "").split(",") if x]
    industry = request.form.get("industry", "")
    transcript = request.form.get("transcript", "")
    prs = Presentation(assembler.SOURCE)
    by_id = {read_id(s): s for s in prs.slides if read_id(s)}

    # ---- existing library slides: editable title/subtitle (as before) ----
    cards = []
    for sid in ids:
        slide = by_id.get(sid)
        if not slide:
            continue
        fields = [(idx, label, text.replace("[CLIENT]", client))
                  for idx, label, text in editor.editable_fields(slide)]
        shown = {f[2] for f in fields}
        context = "\n".join(t.replace("[CLIENT]", client)
                            for t in editor.full_text(slide) if t not in shown)
        cards.append({"id": sid, "fields": fields, "context": context[:400]})

    # Gaps are FLAGS only now — they are never auto-written here. AI slides come
    # solely from the deliberate "Create a slide with AI" tool (NEW:<id> items).
    ai_cards, ai_ids = [], []

    body = render_template_string(REVIEW_BODY, client=client, order=",".join(ids),
                                  cards=cards, ai_cards=ai_cards, ai_ids=",".join(ai_ids),
                                  industry=industry, transcript=transcript,
                                  phase=request.form.get("phase", ""),
                                  recipient=request.form.get("recipient", ""),
                                  functions=request.form.getlist("functions"),
                                  work_types=request.form.getlist("work_types"))
    return shell(body, active="new", crumb="<b>New deck</b> / Review &amp; edit")


@app.route("/finalize", methods=["POST"])
def finalize():
    client = request.form.get("client_name", "Client").strip()
    ids = [x for x in request.form.get("order", "").split(",") if x]
    if not ids:
        abort(400)
    os.makedirs(OUTPUT_DIR, exist_ok=True)
    filename = f"Tailored_Deck_{_safe(client)}.pptx"
    path = os.path.join(OUTPUT_DIR, filename)

    # ---- AI slides: apply edits, then ACCEPT (promote -> library) or REJECT ----
    final_ids = list(ids)
    accepted = []
    for gid in [x for x in request.form.get("ai_ids", "").split(",") if x]:
        decision = request.form.get("ai_decision__" + gid, "accept")
        if decision == "reject":
            staging.discard(gid)
            continue
        staging.update_content(
            gid,
            title=request.form.get("ai_title__" + gid),
            keywords=request.form.get("ai_keywords__" + gid),
            bullets=[b.strip() for b in request.form.get("ai_bullets__" + gid, "").splitlines() if b.strip()],
        )
        new_cs = staging.promote(gid)          # full sign-off -> real, client-ready slide
        if new_cs:
            accepted.append(new_cs)
    # Slot the AI slides in BEFORE the closing slides (Next Steps / Let's win together),
    # not at the very end of the deck.
    if accepted:
        insert_at = next((i for i, s in enumerate(final_ids) if s in matcher.PIN_TO_END),
                         len(final_ids))
        final_ids[insert_at:insert_at] = accepted

    # Skills slides ride along in final_ids (SK:/FP: ids); re-derive their data here.
    skills_ctx = {"work_types": request.form.getlist("work_types"),
                  "industry": request.form.get("industry", ""),
                  "transcript": request.form.get("transcript", ""),
                  "client_name": client}
    skills_cands = skills.candidates(skills_ctx)
    # "Create with AI" slides ride as NEW:<staging_id>; build them from the staged
    # case-study content into THIS deck (not promoted to the master library).
    create_items = []
    for oid in final_ids:
        if oid.startswith("NEW:"):
            rec = staging.get(oid[4:])
            if rec:
                create_items.append({"id": oid, "template": "case_study_full", "content": rec})
    try:
        assembler.build_deck(final_ids, out=path)     # builds the CS slides (skills/NEW ids ignored)
        edits = {}
        for key, val in request.form.items():
            if key.startswith("edit__"):
                _, sid, idx = key.split("__")
                edits.setdefault(sid, {})[int(idx)] = val
        if edits:
            editor.apply_edits(path, edits)
        if client:
            editor.replace_tokens(path, {"[CLIENT]": client, "[Client]": client, "[client]": client})
        skills.build_into(path, final_ids, skills_cands + create_items)   # fill + slot extras
    except (PermissionError, BadZipFile) as e:
        return _file_busy_page(e)
    meeting_log.save(                          # auto-log this meeting (no extra step)
        client=client,
        industry=request.form.get("industry", ""),
        functions=request.form.getlist("functions"),
        work_types=request.form.getlist("work_types"),
        phase=request.form.get("phase", ""),
        recipient=request.form.get("recipient", ""),
        salesperson=current_salesperson(),
        slide_ids=final_ids,
        deck_file=filename,
    )
    count = len(list(Presentation(path).slides))
    body = render_template_string(PREVIEW_BODY, client=client, filename=filename, count=count)
    return shell(body, active="new", crumb="<b>New deck</b> / Done")


if __name__ == "__main__":
    app.run(host="0.0.0.0", port=5000, debug=False)
