# -*- coding: utf-8 -*-
"""
build_library.py  --  Box 0, Step 1: build the CONTENT LIBRARY.

Opens the deck and, for every slide, creates one standalone RECORD:
  - slide_id : the stable ID read from the slide's notes (J2W_ID:)
  - title / subtitle / body_text / full_text : the slide's words
  - source   : a POINTER back to the slide (file + 0-based index + slide number)

The records do NOT depend on deck order. The library is saved as library.json.
(No tagging here yet -- that's Step 2.)
"""

import json
import re
import sys
from pptx import Presentation

ID_LINE_RE = re.compile(r"^J2W_ID:\s*(\S+)\s*$", re.MULTILINE)


def _collect_text(shapes, out):
    """Recursively pull text out of normal shapes, grouped shapes, and tables."""
    for sh in shapes:
        # Grouped shapes -> recurse into their children
        if sh.shape_type == 6:  # MSO_SHAPE_TYPE.GROUP
            _collect_text(sh.shapes, out)
            continue
        if sh.has_text_frame:
            t = sh.text_frame.text.strip()
            if t:
                out.append(t)
        if sh.has_table:
            for row in sh.table.rows:
                for cell in row.cells:
                    t = cell.text.strip()
                    if t:
                        out.append(t)


def read_id(slide):
    if not slide.has_notes_slide:
        return None
    m = ID_LINE_RE.search(slide.notes_slide.notes_text_frame.text or "")
    return m.group(1) if m else None


def build(path):
    prs = Presentation(path)
    records = []
    for idx, slide in enumerate(prs.slides):
        blocks = []
        _collect_text(slide.shapes, blocks)

        # Title: prefer the real title placeholder, else first text block.
        title = ""
        try:
            if slide.shapes.title and slide.shapes.title.text.strip():
                title = slide.shapes.title.text.strip()
        except (ValueError, AttributeError):
            pass
        if not title and blocks:
            title = blocks[0]

        subtitle = ""
        for b in blocks:
            if b != title:
                subtitle = b
                break

        title = " ".join(title.split())
        subtitle = " ".join(subtitle.split())
        full_text = "\n".join(blocks)

        records.append({
            "slide_id": read_id(slide),
            "title": title,
            "subtitle": subtitle,
            "body_text": [b for b in blocks if b != title],
            "full_text": full_text,
            "source": {
                "file": path,
                "slide_index": idx,        # 0-based; the pointer the assembler uses
                "slide_number": idx + 1,   # 1-based; human-friendly
            },
        })
    return records


if __name__ == "__main__":
    src = sys.argv[1] if len(sys.argv) > 1 else "WORKING_COPY_Master_Deck.pptx"
    out = sys.argv[2] if len(sys.argv) > 2 else "library.json"
    recs = build(src)
    with open(out, "w", encoding="utf-8") as f:
        json.dump(recs, f, ensure_ascii=False, indent=2)
    missing = [r["source"]["slide_number"] for r in recs if not r["slide_id"]]
    print(f"Built {len(recs)} records -> {out}")
    print(f"Records WITHOUT an ID: {len(missing)} {missing if missing else ''}")
