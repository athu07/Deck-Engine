from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

prs = Presentation(r"C:\Users\E36250417\Downloads\Slides for skills.pptx")
for i, slide in enumerate(prs.slides):
    print(f"\n{'='*60}")
    print(f"SLIDE {i+1}")
    print(f"{'='*60}")
    for sh in slide.shapes:
        l = round(sh.left / 914400, 3)
        t = round(sh.top / 914400, 3)
        w = round(sh.width / 914400, 3)
        h = round(sh.height / 914400, 3)
        txt = ""
        if sh.has_text_frame:
            txt = repr(sh.text_frame.text[:60])
        stype = sh.shape_type
        print(f"  [{stype:2}] {sh.name!r:28} left={l:5.3f} top={t:5.3f} w={w:5.3f} h={h:5.3f}  {txt}")
