# -*- coding: utf-8 -*-
"""
assembler.py  --  Box 2: a list of slide IDs -> a real, tailored .pptx.

How it works (the safe way):
  - Open the working-copy master (all 68 slides, each stamped with J2W_ID).
  - Keep ONLY the requested slides, put them in the requested order, drop the
    rest, and save to a NEW file.
Because we edit the deck in place (delete + reorder) rather than copying slide
internals by hand, every image, font, table, and layout is preserved perfectly.

The original and the working copy are never modified — we always save to a new
output file.

Slides the matcher couldn't find (e.g. a future "NEEDS TO BE CREATED" gap) are
reported, not silently dropped. Generating missing slides is a later milestone.
"""

import os
import sys
import tempfile

from pptx import Presentation
from pptx.oxml.ns import qn

from build_library import read_id   # reuse the J2W_ID reader

SOURCE = "WORKING_COPY_Master_Deck.pptx"


def _atomic_save(prs, out):
    """Save to a temp file in the same folder, then swap it into place. This way
    we never leave a half-written / corrupt .pptx behind — and if the target is
    locked (e.g. open in PowerPoint) the swap fails cleanly with a clear message
    instead of corrupting the file."""
    out_dir = os.path.dirname(os.path.abspath(out)) or "."
    fd, tmp = tempfile.mkstemp(suffix=".pptx", dir=out_dir)
    os.close(fd)
    try:
        prs.save(tmp)
        os.replace(tmp, out)
    except PermissionError:
        if os.path.exists(tmp):
            os.remove(tmp)
        raise PermissionError(
            "Could not write '%s' — it looks like the file is open (e.g. in "
            "PowerPoint) or syncing. Close it and try again." % os.path.basename(out))
    except Exception:
        if os.path.exists(tmp):
            os.remove(tmp)
        raise


def build_deck(slide_ids, source=SOURCE, out="Tailored_Deck.pptx"):
    prs = Presentation(source)

    sld_id_lst = prs.slides._sldIdLst          # the <p:sldIdLst> ordering element
    elements = list(sld_id_lst)                # one <p:sldId> per slide, in order
    slides = list(prs.slides)                  # same order as `elements`

    id_to_elem = {}
    pairs = []                                 # (slide_id, sldId element)
    for slide, elem in zip(slides, elements):
        sid = read_id(slide)
        if sid:
            id_to_elem[sid] = elem
            pairs.append((sid, elem))

    desired_ids = [s for s in slide_ids if s in id_to_elem]
    missing = [s for s in slide_ids if s not in id_to_elem]
    keep = set(desired_ids)

    # DROP unwanted slides for real: remove the presentation->slide relationship
    # AND the order entry. Once a slide part is unreferenced, python-pptx prunes
    # it (and any images only it used) on save -> smaller file.
    # NOTE: iterate ALL slides (not just ID'd ones) so ID-less slides — e.g. the
    # J2W_TEMPLATE skills/footprint slides living in the master — are dropped too,
    # instead of silently leaking into every deck.
    for slide, elem in zip(slides, elements):
        if read_id(slide) not in keep:
            rid = elem.get(qn("r:id"))
            if rid:
                prs.part.drop_rel(rid)
            sld_id_lst.remove(elem)

    # Re-append the kept slides in the requested order (append moves the node).
    for sid in desired_ids:
        sld_id_lst.append(id_to_elem[sid])

    _atomic_save(prs, out)
    return len(desired_ids), missing


if __name__ == "__main__":
    import json
    import matcher

    # Demo: same BFSI context as the matcher, end-to-end into a .pptx.
    context = {
        "client_name": "Acme Bank",
        "industry": "BFSI",
        "work_types": ["AI_POD", "WORKFORCE"],
        "deck_phase": "intro",
        "recipient": "CTO",
    }
    picks = matcher.match(context)
    ids = [p["slide_id"] for p in picks]
    safe_name = context["client_name"].replace(" ", "_")
    out = f"Tailored_Deck_{safe_name}.pptx"

    kept, missing = build_deck(ids, out=out)
    print(f"Requested {len(ids)} slides -> built {kept} into {out}")
    if missing:
        print(f"NOT FOUND in deck (need creating later): {missing}")
