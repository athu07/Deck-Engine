#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
create_case_study_template.py
Creates the J2W branded case-study slide template.

Run:    py create_case_study_template.py
Output: case_study_v2.pptx   (1 slide; open in PowerPoint to inspect)

Markers used (the fill script replaces these with real content):
  Header  : {{TITLE}}  {{CLIENT}}  {{DOMAIN}}
  Cards   : {{CHALLENGE}}  {{SOLUTION}}
  Caps    : {{CAP_1_TITLE}} .. {{CAP_6_TITLE}}
            {{CAP_1_BODY}}  .. {{CAP_6_BODY}}
  Results : {{RESULT_1_PCT}} .. {{RESULT_3_PCT}}
            {{RESULT_1_TEXT}} .. {{RESULT_3_TEXT}}
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_AUTO_SIZE

# ---------------------------------------------------------------------------
# J2W brand colours
# ---------------------------------------------------------------------------
C_BLACK     = RGBColor(0x11, 0x11, 0x10)
C_RED       = RGBColor(0xC0, 0x20, 0x26)   # CRIMSON RED
C_TEAL      = RGBColor(0x3A, 0x8B, 0x82)   # DEEP TEAL
C_WHITE     = RGBColor(0xFF, 0xFF, 0xFF)
C_CARD_BG   = RGBColor(0xF5, 0xF5, 0xF5)
C_CARD_LINE = RGBColor(0xDE, 0xDE, 0xDE)
C_BODY      = RGBColor(0x3E, 0x3E, 0x3E)

FONT = "Calibri"

# ---------------------------------------------------------------------------
# Font sizes (points) — set per owner's spec
# ---------------------------------------------------------------------------
SZ_TITLE      = 24   # main heading
SZ_SUBHEAD    = 15   # CLIENT | DOMAIN line
SZ_BOX_HEAD   = 13   # "The Challenge" / "The Solution" / capability card titles
SZ_CAPS_LABEL = 14   # "Key Capabilities Developed"
SZ_BODY       = 11   # all body paragraph text
SZ_RESULT_PCT = 34   # big stat numbers in the results bar
SZ_RESULT_TXT = 9    # caption under each stat


# ---------------------------------------------------------------------------
# Helpers
# ---------------------------------------------------------------------------

def _in(v):
    return Inches(v)


def _pt(v):
    return Pt(v)


def _rect(slide, l, t, w, h, fill=None, line=None, lw=0.5):
    """Add a rectangle. fill / line are RGBColor; lw is line width in points."""
    shape = slide.shapes.add_shape(1, _in(l), _in(t), _in(w), _in(h))
    if fill:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill
    else:
        shape.fill.background()
    if line:
        shape.line.color.rgb = line
        shape.line.width = _pt(lw)
    else:
        shape.line.fill.background()
    return shape


def _txb(slide, l, t, w, h, text, size,
         bold=False, color=None, align=PP_ALIGN.LEFT, wrap=True, shrink=False):
    """Add a text box with a single run.
    shrink=True -> 'Shrink text on overflow' so long content auto-fits the box
    (PowerPoint computes the scale when the file is opened)."""
    tb = slide.shapes.add_textbox(_in(l), _in(t), _in(w), _in(h))
    tf = tb.text_frame
    tf.word_wrap = wrap
    if shrink:
        tf.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
    tf.margin_top = tf.margin_bottom = tf.margin_left = tf.margin_right = 0
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.name = FONT
    run.font.size = _pt(size)
    run.font.bold = bold
    if color:
        run.font.color.rgb = color
    return tb


# ---------------------------------------------------------------------------
# Build
# ---------------------------------------------------------------------------

def build(out="case_study_v2.pptx"):
    prs = Presentation()
    prs.slide_width  = _in(13.33)
    prs.slide_height = _in(7.50)

    slide = prs.slides.add_slide(prs.slide_layouts[6])   # blank layout

    # ── 0. Top split bar — crimson (left) + teal (right) ───────────────────
    SPLIT = 6.40                                  # where red ends / teal begins
    _rect(slide, 0.00, 0.00, SPLIT,         0.10, fill=C_RED)
    _rect(slide, SPLIT, 0.00, 13.33 - SPLIT, 0.10, fill=C_TEAL)

    # ── 1. Red accent bar (left of title) ──────────────────────────────────
    _rect(slide, 0.25, 0.34, 0.065, 0.55, fill=C_RED)

    # ── 2. Main title ──────────────────────────────────────────────────────
    _txb(slide, 0.44, 0.30, 9.50, 0.60,
         "{{TITLE}}", SZ_TITLE, bold=True, color=C_BLACK)

    # ── 3. CLIENT | DOMAIN subtitle ────────────────────────────────────────
    _txb(slide, 0.44, 0.74, 12.50, 0.34,
         "CLIENT: {{CLIENT}}  |  DOMAIN: {{DOMAIN}}",
         SZ_SUBHEAD, bold=True, color=C_TEAL)

    # ── 4. Challenge card (red left bar + white card) ──────────────────────
    _rect(slide, 0.250, 1.40, 0.065, 2.00, fill=C_RED)
    _rect(slide, 0.315, 1.40, 5.900, 2.00, fill=C_WHITE, line=C_CARD_LINE)
    _txb(slide, 0.45, 1.47, 5.65, 0.32,
         "The Challenge", SZ_BOX_HEAD, bold=True, color=C_BLACK)
    _txb(slide, 0.45, 1.85, 5.65, 1.48,
         "{{CHALLENGE}}", SZ_BODY, color=C_BODY, wrap=True, shrink=True)

    # ── 5. Solution card (teal left bar + white card) ──────────────────────
    SOL_L, SOL_W = 6.635, 6.40
    _rect(slide, SOL_L - 0.065, 1.40, 0.065, 2.00, fill=C_TEAL)      # left bar
    _rect(slide, SOL_L,         1.40, SOL_W, 2.00, fill=C_WHITE, line=C_CARD_LINE)
    _txb(slide, SOL_L + 0.13, 1.47, SOL_W - 0.26, 0.32,
         "The Solution", SZ_BOX_HEAD, bold=True, color=C_BLACK)
    _txb(slide, SOL_L + 0.13, 1.85, SOL_W - 0.26, 1.48,
         "{{SOLUTION}}", SZ_BODY, color=C_BODY, wrap=True, shrink=True)

    # ── 6. "Key Capabilities Developed" label ──────────────────────────────
    _txb(slide, 0.28, 3.50, 7.00, 0.28,
         "Key Capabilities Developed", SZ_CAPS_LABEL, color=C_TEAL)

    # ── 7. Capability cards — 3 columns × 2 rows ───────────────────────────
    CW  = 4.190   # card width
    CH  = 1.000   # card height
    GAP = 0.105   # gap between cards

    COL_X = [0.28,
             0.28 + CW + GAP,
             0.28 + (CW + GAP) * 2]
    ROW_Y = [3.84,
             3.84 + CH + 0.11]

    caps = [
        ("{{CAP_1_TITLE}}", "{{CAP_1_BODY}}"),
        ("{{CAP_2_TITLE}}", "{{CAP_2_BODY}}"),
        ("{{CAP_3_TITLE}}", "{{CAP_3_BODY}}"),
        ("{{CAP_4_TITLE}}", "{{CAP_4_BODY}}"),
        ("{{CAP_5_TITLE}}", "{{CAP_5_BODY}}"),
        ("{{CAP_6_TITLE}}", "{{CAP_6_BODY}}"),
    ]

    n = 0
    for ry in ROW_Y:
        for cx in COL_X:
            ct, cb = caps[n]
            _rect(slide, cx, ry, CW, CH,
                  fill=C_CARD_BG, line=C_CARD_LINE, lw=0.4)
            _txb(slide, cx + 0.12, ry + 0.07,
                 CW - 0.24, 0.26, ct, SZ_BOX_HEAD, bold=True, color=C_BLACK)
            _txb(slide, cx + 0.12, ry + 0.36,
                 CW - 0.24, CH - 0.44, cb, SZ_BODY, color=C_BODY,
                 wrap=True, shrink=True)
            n += 1

    # ── 8. Results bar (full-width deep teal) ──────────────────────────────
    BAR_T = 6.00
    BAR_H = 7.50 - BAR_T          # fills to the bottom of the slide
    COL_W = 13.33 / 3

    _rect(slide, 0.00, BAR_T, 13.33, BAR_H, fill=C_TEAL)

    results = [
        ("{{RESULT_1_PCT}}", "{{RESULT_1_TEXT}}"),
        ("{{RESULT_2_PCT}}", "{{RESULT_2_TEXT}}"),
        ("{{RESULT_3_PCT}}", "{{RESULT_3_TEXT}}"),
    ]
    for idx, (pct, txt) in enumerate(results):
        x = idx * COL_W
        _txb(slide, x + 0.20, BAR_T + 0.14, COL_W - 0.40, 0.66,
             pct, SZ_RESULT_PCT, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)
        _txb(slide, x + 0.15, BAR_T + 0.86, COL_W - 0.30, 0.55,
             txt, SZ_RESULT_TXT, color=C_WHITE, align=PP_ALIGN.CENTER, wrap=True)

    # ── 9. Notes tag (engine reads this to identify the template) ──────────
    slide.notes_slide.notes_text_frame.text = "J2W_TEMPLATE: case_study_v2"

    prs.save(out)
    print(f"Saved  ->  {out}")
    print()
    print("Markers in this template:")
    print("  Header  : {{TITLE}}  {{CLIENT}}  {{DOMAIN}}")
    print("  Cards   : {{CHALLENGE}}  {{SOLUTION}}")
    for n in range(1, 7):
        print(f"  Cap {n}   : {{{{CAP_{n}_TITLE}}}}  {{{{CAP_{n}_BODY}}}}")
    for n in range(1, 4):
        print(f"  Result {n}: {{{{RESULT_{n}_PCT}}}}  {{{{RESULT_{n}_TEXT}}}}")


if __name__ == "__main__":
    build()
