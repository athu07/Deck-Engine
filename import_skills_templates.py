# -*- coding: utf-8 -*-
"""
import_skills_templates.py
Takes the three designed slides from "Slides for skills.pptx", fixes each one
(notes tag, chart placeholder, removes hardcoded/designer-only content), and
saves the result as "skills_templates.pptx" for use by skills.py.

Run ONCE (safe to re-run — overwrites skills_templates.pptx):
    py import_skills_templates.py
"""
import copy
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor

SRC   = r"C:\Users\E36250417\Downloads\Slides for skills.pptx"
OUT   = "skills_templates.pptx"
GREY  = RGBColor(0x6E, 0x6E, 0x69)
_R    = 'http://schemas.openxmlformats.org/officeDocument/2006/relationships'
EMU   = 914400      # English Metric Units per inch


# ------------------------------------------------------------------ #
# Helpers
# ------------------------------------------------------------------ #
def _remove(slide, shape):
    slide.shapes._spTree.remove(shape._element)


def _add_chart_box(slide, left_in, top_in, w_in, h_in):
    tb  = slide.shapes.add_textbox(Inches(left_in), Inches(top_in),
                                   Inches(w_in), Inches(h_in))
    tf  = tb.text_frame
    p   = tf.paragraphs[0]
    run = p.add_run()
    run.text = "{{CHART}}"
    run.font.size = Pt(12)
    run.font.color.rgb = GREY


def _set_notes(slide, tag):
    slide.notes_slide.notes_text_frame.text = tag


def _anchor_top(text_frame):
    """Set text frame vertical anchor to TOP so text starts at the top of the box."""
    text_frame._txBody.bodyPr.set('anchor', 't')


# ------------------------------------------------------------------ #
# Slide 1 — skill_deepdive
#
# Key shapes (from diagnostic, inches):
#   Image 4       : icon next to "DEPLOYMENT SUMMARY" top=1.573 → KEEP
#   Text 8        : "DEPLOYMENT SUMMARY"   top=1.563 → set font 14pt
#   Shape 9       : white left panel       top=1.979 h=4.375 (ends 6.354)
#   Text 12       : "{{SKILL_SUMMARY}}"    top=2.708 h=0.250 → move to 2.020, h=4.20, anchor TOP
#   Shape 13/16/19: 3 bullet circles       w=h=0.125 → REMOVE (empty-text auto-shapes!)
#   Shape 22      : gray callout box       top=5.542 w=4.156 → REMOVE (empty-text auto-shape!)
#   Image 5       : icon inside gray box   top=5.673 left=1.385 → REMOVE
#
# BUG FIXED: in python-pptx, auto-shapes (incl. circles/boxes with no text) still have
# has_text_frame=True; checking the TEXT is how to detect "real" text vs. empty shape.
# ------------------------------------------------------------------ #
_SLIDE1_REMOVE_EXACT = {
    "Skill Distribution Overview",
    "45%", "30%", "25%",
    "Technical Skills", "Domain Expertise", "Soft Skills",
    "Horizontal bar chart showing consultants per skill",
    "Y-axis: Skill names | X-axis: Number of consultants",
    "{{TOTAL_CONSULTANTS}} Consultants",
    "{{ACTIVE_ENGAGEMENTS}} Active Engagements",
    "{{RELATIONSHIP_DURATION}} Partnering",
    "Last updated: {{SNAPSHOT_DATE}}",
    "Data reflects current deployment as of {{SNAPSHOT_DATE}}",
}
_SLIDE1_REMOVE_STARTS = ("J2W_TEMPLATE:",)


def fix_slide1(slide):
    chart_shape = None
    to_remove   = []

    for sh in list(slide.shapes):
        # Get text only if there is meaningful content
        txt = sh.text_frame.text.strip() if sh.has_text_frame else ""

        if txt:
            # ---------- shapes WITH actual text ----------
            if txt in _SLIDE1_REMOVE_EXACT:
                to_remove.append(sh)
            elif txt.startswith(_SLIDE1_REMOVE_STARTS):
                to_remove.append(sh)
            elif "{{CHART}}" in txt:
                chart_shape = sh
            elif txt == "DEPLOYMENT SUMMARY":
                for p in sh.text_frame.paragraphs:
                    for run in p.runs:
                        run.font.size = Pt(14)
            elif "{{SKILL_SUMMARY}}" in txt:
                # Move UP to sit just inside the white panel (panel top=1.979")
                # Set anchor=top so content doesn't float to mid-box
                sh.top    = int(2.020 * EMU)
                sh.height = int(4.200 * EMU)   # fills panel to ~6.22" (panel ends 6.354")
                sh.text_frame.word_wrap = True
                _anchor_top(sh.text_frame)
                for p in sh.text_frame.paragraphs:
                    for run in p.runs:
                        run.font.size = Pt(10)

        elif sh.shape_type == 13:
            # ---------- PICTURE shapes (no text) ----------
            left_in = sh.left  / EMU
            top_in  = sh.top   / EMU
            w_in    = sh.width / EMU
            if top_in > 6.8:
                # Footer images
                to_remove.append(sh)
            elif left_in > 8.5 and 3.0 < top_in < 5.5 and w_in < 1.0:
                # Chart-mockup icon in right panel
                to_remove.append(sh)
            elif left_in < 3.0 and top_in > 5.5:
                # Icon inside the gray callout box (Image 5, top≈5.67)
                to_remove.append(sh)

        else:
            # ---------- AUTO_SHAPEs with EMPTY text (circles, boxes) ----------
            left_in = sh.left   / EMU
            top_in  = sh.top    / EMU
            w_in    = sh.width  / EMU
            h_in    = sh.height / EMU
            # 3 bullet circles: tiny shapes in the left panel mid-area
            if left_in < 5.0 and top_in > 3.0 and w_in <= 0.130 and h_in <= 0.130:
                to_remove.append(sh)
            # Gray callout box: medium-sized at lower-left
            elif left_in < 3.0 and top_in > 5.0 and w_in > 3.5 and h_in > 0.4:
                to_remove.append(sh)

    for sh in to_remove:
        try:
            _remove(slide, sh)
        except Exception:
            pass

    # Position and size the {{CHART}} placeholder inside the right panel
    # Right panel (Shape 3): left=5.979 top=1.979 w=6.500 h=4.396 → fill it
    if chart_shape is not None:
        chart_shape.left   = int(6.15 * EMU)
        chart_shape.top    = int(2.10 * EMU)
        chart_shape.width  = int(6.20 * EMU)
        chart_shape.height = int(4.10 * EMU)
    else:
        _add_chart_box(slide, 6.15, 2.10, 6.20, 4.10)

    _set_notes(slide, "J2W_TEMPLATE: skill_deepdive")
    print("  Slide 1 (skill_deepdive): fixed")


# ------------------------------------------------------------------ #
# Slide 2 — industry_strength
#
# Bottom-right panel (Shape 11): left=6.875 top=3.854 w=5.833 h=2.990 (ends 6.844)
# "Function Distribution" label: top=4.010 h=0.250 (ends 4.260)
# Chart goes right below the label; panel bottom at 6.844" → chart h=2.50"
# ------------------------------------------------------------------ #
def fix_slide2(slide):
    has_chart = any(
        sh.has_text_frame and "{{CHART}}" in sh.text_frame.text
        for sh in slide.shapes
    )
    if not has_chart:
        _add_chart_box(slide, 6.92, 4.30, 5.73, 2.50)
    _set_notes(slide, "J2W_TEMPLATE: industry_strength")
    print("  Slide 2 (industry_strength): fixed")


# ------------------------------------------------------------------ #
# Slide 3 — company_footprint
#
# Key shapes (from diagnostic, inches):
#   Text 15 ({{TOTAL_DEPLOYED}})     top=1.615 h=0.323  ─┐
#   Text 17 ({{NUM_FUNCTIONS_CO}})   top=1.615 h=0.646  ─┤ normalize to h=0.500
#   Text 19 ({{NUM_SKILLS_CO}})      top=1.615 h=0.323  ─┤
#   Text 21 ({{ENGAGEMENT_TYPE}})    top=1.615 h=0.282  ─┘
#   Text 16 (Consultants deployed)   top=2.019  ─┐
#   Text 18 (Functions delivered)    top=2.340  ─┤ align all to 2.200
#   Text 20 (Distinct skills)        top=2.019  ─┤ (1.615 + 0.500 + 0.085 gap)
#   Text 22 (Engagement type)        top=1.973  ─┘
#   Text 24 ({{FUNCTION_BREAKDOWN}}) top=3.594 h=0.240 → move to 3.420, h=3.15, anchor TOP
#   Image 5 (chart mockup)           top=3.542 w=5.167 h=3.125 → REMOVE, add {{CHART}}
#   Text 26 (footer {{SNAPSHOT_DATE}}) → REMOVE
# ------------------------------------------------------------------ #
_VALUE_MARKERS = {"{{TOTAL_DEPLOYED}}", "{{NUM_FUNCTIONS_CO}}", "{{NUM_SKILLS_CO}}", "{{ENGAGEMENT_TYPE}}"}
_LABEL_TEXTS   = {"Consultants deployed", "Functions delivered", "Distinct skills", "Engagement type"}


def fix_slide3(slide):
    to_remove = []
    chart_img = None

    for sh in list(slide.shapes):
        txt = sh.text_frame.text.strip() if sh.has_text_frame else ""

        if txt:
            if "Account snapshot as of" in txt or "{{SNAPSHOT_DATE}}" in txt:
                to_remove.append(sh)
            elif txt in _VALUE_MARKERS:
                # Normalize all 4 metric value boxes to same height (0.50")
                # so labels below them can align on a single horizontal line
                sh.height = int(0.500 * EMU)
                sh.text_frame.word_wrap = False
            elif txt in _LABEL_TEXTS:
                # Align all 4 metric labels to same top
                # = value-box top (1.615) + normalized height (0.500) + gap (0.085)
                sh.top = int(2.200 * EMU)
            elif "{{FUNCTION_BREAKDOWN}}" in txt:
                # Move up close to the heading, expand height, anchor TOP
                # Heading (Text 23) ends at 3.125 + 0.271 = 3.396"
                sh.top    = int(3.430 * EMU)
                sh.height = int(3.150 * EMU)   # fills to 6.580" (panel ends 6.896")
                sh.text_frame.word_wrap = True
                _anchor_top(sh.text_frame)
                for p in sh.text_frame.paragraphs:
                    for run in p.runs:
                        run.font.size = Pt(9)

        elif sh.shape_type == 13:
            top_in = sh.top   / EMU
            w_in   = sh.width / EMU
            # Chart-mockup picture: wide (>4.5") in mid-page (top 3–6")
            if w_in > 4.5 and 3.0 < top_in < 6.5:
                chart_img = sh
                to_remove.append(sh)

    # Capture chart image position before removing
    if chart_img:
        cl = chart_img.left  / EMU
        ct = chart_img.top   / EMU
        cw = chart_img.width / EMU
        ch = chart_img.height / EMU
    else:
        cl, ct, cw, ch = 7.19, 3.54, 5.17, 3.12

    for sh in to_remove:
        try:
            _remove(slide, sh)
        except Exception:
            pass

    _add_chart_box(slide, cl, ct, cw, ch)
    _set_notes(slide, "J2W_TEMPLATE: company_footprint")
    print("  Slide 3 (company_footprint): fixed")


# ------------------------------------------------------------------ #
# Copy a slide including image relationships
# ------------------------------------------------------------------ #
def _blank_layout(prs):
    for layout in prs.slide_layouts:
        if (layout.name or "").lower().strip() == "blank":
            return layout
    return prs.slide_layouts[-1]


def _copy_with_images(dest_prs, src_slide):
    """Deep-copy a slide into dest_prs, including all image relationships."""
    new = dest_prs.slides.add_slide(_blank_layout(dest_prs))
    for shp in list(new.shapes):
        shp._element.getparent().remove(shp._element)

    src_part  = src_slide._part
    dest_part = new._part
    rId_map   = {}
    _REMAP    = {f'{{{_R}}}embed', f'{{{_R}}}link'}

    for rId, rel in list(src_part.rels.items()):
        if rel.is_external:
            continue
        if '/image' in (rel.reltype or ''):
            new_rId = dest_part.relate_to(rel._target, rel.reltype)
            if new_rId != rId:
                rId_map[rId] = new_rId

    def remap(elem):
        for attr, val in list(elem.attrib.items()):
            if attr in _REMAP and val in rId_map:
                elem.attrib[attr] = rId_map[val]
        for child in elem:
            remap(child)

    for shp in src_slide.shapes:
        elem = copy.deepcopy(shp._element)
        if rId_map:
            remap(elem)
        new.shapes._spTree.append(elem)

    if src_slide.has_notes_slide:
        txt = src_slide.notes_slide.notes_text_frame.text
        if txt.strip():
            new.notes_slide.notes_text_frame.text = txt

    return new


if __name__ == "__main__":
    src_prs = Presentation(SRC)
    slides  = list(src_prs.slides)

    if len(slides) != 3:
        raise SystemExit(f"Expected 3 slides in source, found {len(slides)}")

    fix_slide1(slides[0])   # skill_deepdive
    fix_slide2(slides[1])   # industry_strength
    fix_slide3(slides[2])   # company_footprint

    out_prs = Presentation()
    out_prs.slide_width  = src_prs.slide_width
    out_prs.slide_height = src_prs.slide_height

    for slide in slides:
        _copy_with_images(out_prs, slide)

    out_prs.save(OUT)
    print(f"\nSaved {OUT}  ({len(out_prs.slides)} slides)")
    for s in out_prs.slides:
        note = s.notes_slide.notes_text_frame.text.strip()
        print(f"  Notes: {repr(note)}")
