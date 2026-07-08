import os as _os, sys as _sys
_sys.path.insert(0, _os.path.dirname(_os.path.dirname(_os.path.abspath(__file__))))
from deckengine import config  # noqa: E402  (repo-root data paths)
#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
create_case_study_template.py
Creates the J2W branded case-study slide template.

Run:    py create_case_study_template.py
Output: case_study_v2.pptx   (1 slide; open in PowerPoint to inspect)

Markers used (the fill script replaces these with real content):
  Header  : {{TITLE}}  {{CLIENT}}  {{DOMAIN}}
  Cards   : {{CHALLENGE}}  {{SOLUTION}}
  Caps    : {{CAP_1_TITLE}} .. {{CAP_6_TITLE}}
            {{CAP_1_BODY}}  .. {{CAP_6_BODY}}
  Results : {{RESULT_1_PCT}} .. {{RESULT_3_PCT}}
            {{RESULT_1_TEXT}} .. {{RESULT_3_TEXT}}
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_AUTO_SIZE

# ---------------------------------------------------------------------------
# J2W brand colours — matched to the WORKING master deck's own title slide
# (extracted pixel-exact from its top-bar images, 2026-07-07): red #D62839,
# teal #2A9D8F. Every J2W template should share these, not a close approximation.
# ---------------------------------------------------------------------------
C_BLACK     = RGBColor(0x11, 0x11, 0x10)
C_RED       = RGBColor(0xD6, 0x28, 0x39)   # master-deck red
C_TEAL      = RGBColor(0x2A, 0x9D, 0x8F)   # master-deck teal
C_WHITE     = RGBColor(0xFF, 0xFF, 0xFF)
C_CARD_BG   = RGBColor(0xF5, 0xF5, 0xF5)
C_CARD_LINE = RGBColor(0xDE, 0xDE, 0xDE)
C_BODY      = RGBColor(0x3E, 0x3E, 0x3E)

# Fonts (owner's spec, 2026-07-07): heading = Oswald; subheading = Roboto
# Condensed; everything else (body/content) = Raleway.
# (If a font isn't installed on the viewing machine, PowerPoint substitutes a
# similar one but keeps the name — so it renders correctly wherever it IS installed.)
FONT_HEAD    = "Oswald"            # main title only
FONT_SUBHEAD = "Roboto Condensed"  # CLIENT | DOMAIN subheading
FONT_BODY    = "Raleway"           # all other content

# ---------------------------------------------------------------------------
# Font sizes (points) — set per owner's spec
# ---------------------------------------------------------------------------
SZ_TITLE      = 24   # main heading
SZ_SUBHEAD    = 15   # CLIENT | DOMAIN line
SZ_BOX_HEAD   = 14   # "The Challenge" / "The Solution" / capability card titles
SZ_CAPS_LABEL = 14   # "Key Capabilities Developed"
SZ_BODY       = 13   # all body paragraph text (owner 2026-07-07: content = 13-14)
SZ_RESULT_PCT = 25   # big stat numbers in the results bar (owner: 25, was 34)
SZ_RESULT_TXT = 13   # caption under each stat (owner 2026-07-07: content = 13-14)


# ---------------------------------------------------------------------------
# Helpers
# ---------------------------------------------------------------------------

def _in(v):
    return Inches(v)


def _pt(v):
    return Pt(v)


def _rect(slide, l, t, w, h, fill=None, line=None, lw=0.5):
    """Add a rectangle. fill / line are RGBColor; lw is line width in points."""
    shape = slide.shapes.add_shape(1, _in(l), _in(t), _in(w), _in(h))
    if fill:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill
    else:
        shape.fill.background()
    if line:
        shape.line.color.rgb = line
        shape.line.width = _pt(lw)
    else:
        shape.line.fill.background()
    return shape


def _txb(slide, l, t, w, h, text, size,
         bold=False, color=None, align=PP_ALIGN.LEFT, wrap=True, shrink=False,
         font=FONT_BODY):
    """Add a text box with a single run.
    shrink=True -> 'Shrink text on overflow' so long content auto-fits the box
    (PowerPoint computes the scale when the file is opened)."""
    tb = slide.shapes.add_textbox(_in(l), _in(t), _in(w), _in(h))
    tf = tb.text_frame
    tf.word_wrap = wrap
    if shrink:
        tf.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
    tf.margin_top = tf.margin_bottom = tf.margin_left = tf.margin_right = 0
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.name = font
    run.font.size = _pt(size)
    run.font.bold = bold
    if color:
        run.font.color.rgb = color
    return tb


def _title_box(slide, l, t, w, h):
    """The main heading: a red 'CASE STUDY:' label + the black case title, both
    Oswald bold, on one paragraph so they flow (and wrap) together. The title
    text is the {{TITLE}} marker the fill script replaces (keeps it black)."""
    tb = slide.shapes.add_textbox(_in(l), _in(t), _in(w), _in(h))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_top = tf.margin_bottom = tf.margin_left = tf.margin_right = 0
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    label = p.add_run()
    label.text = "CASE STUDY: "
    label.font.name = FONT_HEAD
    label.font.size = _pt(SZ_TITLE)
    label.font.bold = True
    label.font.color.rgb = C_RED
    title = p.add_run()
    title.text = "{{TITLE}}"
    title.font.name = FONT_HEAD
    title.font.size = _pt(SZ_TITLE)
    title.font.bold = True
    title.font.color.rgb = C_BLACK
    return tb


# ---------------------------------------------------------------------------
# Build
# ---------------------------------------------------------------------------

def build(out=config.CASE_TEMPLATE_PPTX):
    prs = Presentation()
    prs.slide_width  = _in(13.33)
    prs.slide_height = _in(7.50)

    slide = prs.slides.add_slide(prs.slide_layouts[6])   # blank layout

    # ── 0. Top split bar — crimson (left) + teal (right) ───────────────────
    # Proportions + thickness match the master deck's own title-slide bar
    # exactly (extracted from its embedded images, 2026-07-07): red covers the
    # first ~2/3, teal the last ~1/3, both only 0.083in tall ("very less").
    SPLIT = 8.89                                  # where red ends / teal begins
    BAR_H_TOP = 0.083
    _rect(slide, 0.00, 0.00, SPLIT,         BAR_H_TOP, fill=C_RED)
    _rect(slide, SPLIT, 0.00, 13.33 - SPLIT, BAR_H_TOP, fill=C_TEAL)

    # ── 1. Red accent bar (left of title, spans the header block) ──────────
    # Header sits a little lower so there's a clear gap under the top split bar
    # (it was congested right beneath it). Still clears the cards at 1.40.
    _rect(slide, 0.25, 0.34, 0.065, 0.80, fill=C_RED)

    # ── 2. Main title — "CASE STUDY:" (red) + title (black), Oswald ─────────
    # Wide box + room for a 2nd line: a short title sits on line 1; a long one
    # wraps to line 2, and the fill script drops the subheading a line to match.
    _title_box(slide, 0.44, 0.34, 12.45, 0.80)

    # ── 3. CLIENT | DOMAIN subtitle (Roboto Condensed, teal) ────────────────
    _txb(slide, 0.44, 0.88, 12.45, 0.30,
         "CLIENT: {{CLIENT}}  |  DOMAIN: {{DOMAIN}}",
         SZ_SUBHEAD, bold=True, color=C_TEAL, font=FONT_SUBHEAD)

    # ── 4. Challenge card (red left bar + white card) ──────────────────────
    _rect(slide, 0.250, 1.40, 0.065, 2.00, fill=C_RED)
    _rect(slide, 0.315, 1.40, 5.900, 2.00, fill=C_WHITE, line=C_CARD_LINE)
    _txb(slide, 0.45, 1.47, 5.65, 0.32,
         "The Challenge", SZ_BOX_HEAD, bold=True, color=C_BLACK)
    _txb(slide, 0.45, 1.85, 5.65, 1.48,
         "{{CHALLENGE}}", SZ_BODY, color=C_BODY, wrap=True, shrink=True)

    # ── 5. Solution card (teal left bar + white card) ──────────────────────
    SOL_L, SOL_W = 6.635, 6.40
    _rect(slide, SOL_L - 0.065, 1.40, 0.065, 2.00, fill=C_TEAL)      # left bar
    _rect(slide, SOL_L,         1.40, SOL_W, 2.00, fill=C_WHITE, line=C_CARD_LINE)
    _txb(slide, SOL_L + 0.13, 1.47, SOL_W - 0.26, 0.32,
         "The Solution", SZ_BOX_HEAD, bold=True, color=C_BLACK)
    _txb(slide, SOL_L + 0.13, 1.85, SOL_W - 0.26, 1.48,
         "{{SOLUTION}}", SZ_BODY, color=C_BODY, wrap=True, shrink=True)

    # ── 6. "Key Capabilities Developed" label ──────────────────────────────
    _txb(slide, 0.28, 3.50, 7.00, 0.28,
         "Key Capabilities Developed", SZ_CAPS_LABEL, color=C_TEAL)

    # ── 7. Capability cards — 3 columns × 2 rows ───────────────────────────
    CW  = 4.190   # card width
    CH  = 1.000   # card height
    GAP = 0.105   # gap between cards

    COL_X = [0.28,
             0.28 + CW + GAP,
             0.28 + (CW + GAP) * 2]
    ROW_Y = [3.84,
             3.84 + CH + 0.11]

    caps = [
        ("{{CAP_1_TITLE}}", "{{CAP_1_BODY}}"),
        ("{{CAP_2_TITLE}}", "{{CAP_2_BODY}}"),
        ("{{CAP_3_TITLE}}", "{{CAP_3_BODY}}"),
        ("{{CAP_4_TITLE}}", "{{CAP_4_BODY}}"),
        ("{{CAP_5_TITLE}}", "{{CAP_5_BODY}}"),
        ("{{CAP_6_TITLE}}", "{{CAP_6_BODY}}"),
    ]

    n = 0
    for ry in ROW_Y:
        for cx in COL_X:
            ct, cb = caps[n]
            _rect(slide, cx, ry, CW, CH,
                  fill=C_CARD_BG, line=C_CARD_LINE, lw=0.4)
            _txb(slide, cx + 0.12, ry + 0.07,
                 CW - 0.24, 0.26, ct, SZ_BOX_HEAD, bold=True, color=C_BLACK)
            _txb(slide, cx + 0.12, ry + 0.36,
                 CW - 0.24, CH - 0.44, cb, SZ_BODY, color=C_BODY,
                 wrap=True, shrink=True)
            n += 1

    # ── 8. Results bar (full-width deep teal) ──────────────────────────────
    # Starts a little lower than the capability cards (which end ~5.95) so there
    # is a clear gap above it; the stats inside are positioned relative to BAR_T,
    # so they move down with the bar.
    BAR_T = 6.20
    BAR_H = 7.50 - BAR_T          # fills to the bottom of the slide
    COL_W = 13.33 / 3

    _rect(slide, 0.00, BAR_T, 13.33, BAR_H, fill=C_TEAL)

    results = [
        ("{{RESULT_1_PCT}}", "{{RESULT_1_TEXT}}"),
        ("{{RESULT_2_PCT}}", "{{RESULT_2_TEXT}}"),
        ("{{RESULT_3_PCT}}", "{{RESULT_3_TEXT}}"),
    ]
    for idx, (pct, txt) in enumerate(results):
        x = idx * COL_W
        _txb(slide, x + 0.20, BAR_T + 0.16, COL_W - 0.40, 0.52,
             pct, SZ_RESULT_PCT, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)
        _txb(slide, x + 0.15, BAR_T + 0.66, COL_W - 0.30, 0.60,
             txt, SZ_RESULT_TXT, color=C_WHITE, align=PP_ALIGN.CENTER, wrap=True)

    # ── 9. Notes tag (engine reads this to identify the template) ──────────
    slide.notes_slide.notes_text_frame.text = "J2W_TEMPLATE: case_study_v2"

    prs.save(out)
    print(f"Saved  ->  {out}")
    print()
    print("Markers in this template:")
    print("  Header  : {{TITLE}}  {{CLIENT}}  {{DOMAIN}}")
    print("  Cards   : {{CHALLENGE}}  {{SOLUTION}}")
    for n in range(1, 7):
        print(f"  Cap {n}   : {{{{CAP_{n}_TITLE}}}}  {{{{CAP_{n}_BODY}}}}")
    for n in range(1, 4):
        print(f"  Result {n}: {{{{RESULT_{n}_PCT}}}}  {{{{RESULT_{n}_TEXT}}}}")


if __name__ == "__main__":
    build()
