# -*- coding: utf-8 -*-
import os as _os, sys as _sys
_sys.path.insert(0, _os.path.dirname(_os.path.dirname(_os.path.abspath(__file__))))
from deckengine import config  # noqa: E402  (repo-root data paths)
"""
set_case_content_font.py  --  set the case-study slide CONTENT font to Raleway.

The branded case_study_v2 template uses Oswald for the heading ('CASE STUDY: ...')
and subheading ('CLIENT: ... | DOMAIN: ...'), and Helvetica for everything in the
body (the section labels, the challenge/solution text, the capabilities and the
results). The owner wants the heading + subheading kept as-is and the body content
in Raleway, so this rewrites every Helvetica run to Raleway and leaves Oswald alone.

Idempotent and re-runnable. Run after any template swap:
    py scripts/set_case_content_font.py
"""
from pptx import Presentation

FROM_FONT = "Helvetica"
TO_FONT = "Raleway"


def main():
    prs = Presentation(config.CASE_TEMPLATE_PPTX)
    changed = 0
    for slide in prs.slides:
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            for para in shape.text_frame.paragraphs:
                for run in para.runs:
                    if run.font.name == FROM_FONT:
                        run.font.name = TO_FONT
                        changed += 1
    prs.save(config.CASE_TEMPLATE_PPTX)
    print(f"Set {changed} content runs from {FROM_FONT} -> {TO_FONT} in {config.CASE_TEMPLATE_PPTX}")


if __name__ == "__main__":
    main()
