# -*- coding: utf-8 -*-
import os as _os, sys as _sys
_sys.path.insert(0, _os.path.dirname(_os.path.dirname(_os.path.abspath(__file__))))
from deckengine import config  # noqa: E402  (repo-root data paths)
"""
swap_master_deck_39slide.py  --  ONE-OFF: replace the master deck with the new
39-slide "checking_Master deck 2.pptx" (per SLIDE_SELECTION_BRIEF.md), fixing
the two structural issues the brief calls out before it can be used:

  1. Slide 3 duplicates slide 2's J2W_ID (both CS02) -- the brief needs them
     distinguishable. Slide 3 (the Intro-only short "Who We Are") gets CS02B.
  2. 12 other slides carry no J2W_ID at all (today's assembler silently DROPS
     any ID-less slide, so these would vanish from every deck). Stamped with
     fresh IDs starting ABOVE the highest number already used anywhere in the
     current registry/library (136) -- not above this file's own internal max
     (68) -- so nothing collides with old registry rows.

Registry rows are added for all 13 new IDs with SAFE, non-guessing defaults:
  - Slides whose rule the brief states unconditionally (dividers that never
    auto-include, on-demand-only slides, the always-before-case-studies
    divider) get their real rule now.
  - Everything else gets an inert placeholder ("PENDING_PHASE3...") that the
    matcher's include_rule parser does not recognise, so it can NEVER be
    auto-included until Phase 3 wires the real stage/work-type/trigger logic.
    (Verified: matcher.plan only auto-includes on `rule.upper()=="ALWAYS"` or
    `rule.startswith("IF work_type includes")` -- anything else is inert.)

Run once:  py scripts/swap_master_deck_39slide.py
Idempotent-ish: safe to re-run against the SAME source deck (already-stamped
slides are left untouched by the stamping pass), but do not re-run after the
registry has been hand-edited for Phase 3 -- it will not clobber existing rows
for these 13 IDs, but check before re-running.
"""

import re
import shutil

import openpyxl
from pptx import Presentation

from deckengine.services.content import build_library
from deckengine.services.matching import tagger

SOURCE_DECK = r"C:\Users\E36250417\Downloads\checking_Master deck 2.pptx"
TARGET_DECK = config.MASTER_DECK
REGISTRY = config.REGISTRY_XLSX

ID_LINE_RE = re.compile(r"^J2W_ID:\s*(\S+)\s*$", re.MULTILINE)

# 1-indexed slide numbers (matches SLIDE_SELECTION_BRIEF.md's own numbering)
# that carry NO J2W_ID in the source file, in document order.
MISSING_ID_SLIDES = [4, 10, 11, 21, 22, 23, 24, 25, 26, 30, 35, 37, 39]

# For each missing slide, the registry defaults to write once it gets its new
# CS number. "PENDING" rows use an include_rule the matcher does not recognise
# (never auto-included) until Phase 3 replaces it with the real stage/trigger
# logic from SLIDE_SELECTION_BRIEF.md.
_PENDING = "PENDING_PHASE3 - stage/work-type/trigger rule not yet wired (see SLIDE_SELECTION_BRIEF.md); never auto-included until then"
SLIDE_DEFAULTS = {
    4:  {"section": "CORE",      "kind": "STANDARD", "std_group": "CORE",
         "work_types": None,      "include_rule": _PENDING},        # GCC trigger - Phase 3
    10: {"section": "CORE",      "kind": "STANDARD", "std_group": "CORE",
         "work_types": None,      "include_rule": _PENDING},        # GCC trigger - Phase 3
    11: {"section": "CORE",      "kind": "STANDARD", "std_group": "CORE",
         "work_types": None,      "include_rule": _PENDING},        # GCC trigger - Phase 3
    21: {"section": "WORKFORCE", "kind": "STANDARD", "std_group": "WORKFORCE",
         "work_types": "WORKFORCE", "include_rule": _PENDING},      # stage-gated - Phase 3
    22: {"section": "WORKFORCE", "kind": "STANDARD", "std_group": "WORKFORCE",
         "work_types": "WORKFORCE", "include_rule": _PENDING},      # stage-gated - Phase 3
    23: {"section": "WORKFORCE", "kind": "STANDARD", "std_group": "WORKFORCE",
         "work_types": "WORKFORCE", "include_rule": _PENDING},      # stage-gated - Phase 3
    24: {"section": "WORKFORCE", "kind": "OPTIONAL", "std_group": "OPTIONAL",
         "work_types": "WORKFORCE",
         "include_rule": "OPTIONAL - manual only, never auto-selected"},   # brief: on-demand only
    25: {"section": "WORKFORCE", "kind": "OPTIONAL", "std_group": "OPTIONAL",
         "work_types": "WORKFORCE",
         "include_rule": "OPTIONAL - manual only, never auto-selected"},   # brief: on-demand only
    26: {"section": "WORKFORCE", "kind": "OPTIONAL", "std_group": "OPTIONAL",
         "work_types": "WORKFORCE",
         "include_rule": "OPTIONAL - manual only, never auto-selected"},   # brief: on-demand only; has {{CLIENT}} marker
    30: {"section": "AI_POD",    "kind": "STANDARD", "std_group": "AI_POD",
         "work_types": "AI_POD",  "include_rule": _PENDING},        # stage-gated - Phase 3
    35: {"section": "CASE_STUDIES", "kind": "DIVIDER", "std_group": "CASE_STUDIES",
         "work_types": None,
         "include_rule": "ALWAYS"},   # brief: ALWAYS, unlike the other dividers - deliberate exception
    37: {"section": "MS",        "kind": "STANDARD", "std_group": "MS",
         "work_types": "MS",     "include_rule": _PENDING},        # stage-gated - Phase 3
    39: {"section": "APPENDIX",  "kind": "OPTIONAL", "std_group": "OPTIONAL",
         "work_types": None,
         "include_rule": "OPTIONAL - manual only, never auto-selected"},   # brief: manual/on-demand only
}

# Slide 3's own row (short "Who We Are") - Intro-only, pure-MSS-only per the
# brief; depends on the new stage+deck_composition machinery -> pending.
SLIDE3_DEFAULT = {"section": "CORE", "kind": "STANDARD", "std_group": "CORE",
                   "work_types": "MS", "include_rule": _PENDING}


def _existing_id(slide):
    if not slide.has_notes_slide:
        return None
    m = ID_LINE_RE.search(slide.notes_slide.notes_text_frame.text or "")
    return m.group(1) if m else None


def _set_id(slide, new_id):
    tf = slide.notes_slide.notes_text_frame
    existing = (tf.text or "").strip()
    # keep any real notes text that isn't itself a stale J2W_ID line
    existing = ID_LINE_RE.sub("", existing).strip()
    tf.text = f"J2W_ID: {new_id}" + (f"\n{existing}" if existing else "")


def _global_max_cs_number():
    wb = openpyxl.load_workbook(REGISTRY, data_only=True)
    ws = wb["Slide Registry"]
    mx = 0
    for r in range(2, ws.max_row + 1):
        sid = str(ws.cell(r, 1).value or "")
        m = re.match(r"CS(\d+)", sid)
        if m:
            mx = max(mx, int(m.group(1)))
    return mx


def main():
    prs = Presentation(SOURCE_DECK)
    slides = list(prs.slides)
    print(f"Source deck: {len(slides)} slides")

    # ---- 1. fix the CS02 duplicate: slide 3 (index 2) -> CS02B ----
    dup_slide = slides[2]   # 1-indexed slide 3
    cur = _existing_id(dup_slide)
    if cur != "CS02":
        raise SystemExit(f"Expected slide 3 to be CS02, found {cur!r} - aborting, re-check the source file.")
    _set_id(dup_slide, "CS02B")
    print("Slide 3: CS02 -> CS02B (own identity, was duplicating slide 2)")

    # ---- 2. stamp the 13 ID-less slides with fresh, non-colliding numbers ----
    start = _global_max_cs_number()   # 136 today; NOT this file's own max (68)
    new_ids = {}   # 1-indexed slide number -> new CS id
    next_num = start
    for sno in MISSING_ID_SLIDES:
        slide = slides[sno - 1]
        if _existing_id(slide) is not None:
            raise SystemExit(f"Slide {sno} unexpectedly already has an ID - re-check MISSING_ID_SLIDES.")
        next_num += 1
        new_id = f"CS{next_num}"
        _set_id(slide, new_id)
        new_ids[sno] = new_id
        print(f"Slide {sno}: (no id) -> {new_id}")

    # ---- 3. save as the new working master (backup already taken by hand) ----
    prs.save(TARGET_DECK)
    print(f"\nSaved -> {TARGET_DECK}")

    # ---- 4. rebuild library.json / tagged_library.json from the new master ----
    recs = build_library.build(TARGET_DECK)
    import json
    json.dump(recs, open(config.LIBRARY_JSON, "w", encoding="utf-8"), ensure_ascii=False, indent=2)
    tagged = tagger.tag_library(recs)
    json.dump(tagged, open(config.TAGGED_LIBRARY_JSON, "w", encoding="utf-8"), ensure_ascii=False, indent=2)
    print(f"Rebuilt library.json + tagged_library.json ({len(recs)} slide records)")

    # ---- 5. add registry rows for CS02B + the 13 new IDs ----
    wb = openpyxl.load_workbook(REGISTRY)
    ws = wb["Slide Registry"]
    col = {c.value: i for i, c in enumerate(ws[1])}
    existing_ids = {ws.cell(r, col["slide_id"] + 1).value for r in range(2, ws.max_row + 1)}

    def _title_of(sno):
        for sh in slides[sno - 1].shapes:
            if sh.has_text_frame and sh.text_frame.text.strip():
                return sh.text_frame.text.strip().replace("\n", " ")[:120]
        return ""

    def _append_row(sid, defaults, title):
        if sid in existing_ids:
            print(f"  (skip {sid} - already in registry)")
            return
        row = [None] * len(col)
        row[col["slide_id"]] = sid
        row[col["section"]] = defaults["section"]
        row[col["kind"]] = defaults["kind"]
        row[col["std_group"]] = defaults["std_group"]
        row[col["include_rule"]] = defaults["include_rule"]
        row[col["work_types"]] = defaults["work_types"]
        row[col["confidence"]] = "MANUAL"
        row[col["title"]] = title
        ws.append(row)
        existing_ids.add(sid)

    _append_row("CS02B", SLIDE3_DEFAULT, _title_of(3))
    for sno, new_id in new_ids.items():
        _append_row(new_id, SLIDE_DEFAULTS[sno], _title_of(sno))

    wb.save(REGISTRY)
    print(f"\nRegistry updated -> {REGISTRY}")
    print("\nDone. New IDs:", {"CS02B": "slide 3", **{v: f"slide {k}" for k, v in new_ids.items()}})


if __name__ == "__main__":
    main()
