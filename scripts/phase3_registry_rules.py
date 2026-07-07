# -*- coding: utf-8 -*-
import os as _os, sys as _sys
_sys.path.insert(0, _os.path.dirname(_os.path.dirname(_os.path.abspath(__file__))))
from deckengine import config  # noqa: E402  (repo-root data paths)
"""
phase3_registry_rules.py  --  ONE-OFF: write the SLIDE_SELECTION_BRIEF.md stage x
work-type x trigger rules into the registry, for the 39-slide master.

Also finishes the row-additions Phase 1 (swap_master_deck_39slide.py) queued but
couldn't write (CS02B + the 13 newly-stamped ids) -- this script does both in one
pass since neither ever landed on disk.

WHAT THIS ADDS to the registry schema (new columns, existing ones untouched):
  include_rule_legacy  -- a FROZEN SNAPSHOT of each row's include_rule as it was
                          BEFORE this script ran. Proposal-stage decks (and any
                          unrecognised phase) are evaluated against THIS column
                          with the ORIGINAL always/work-type logic, byte-for-byte
                          unchanged -- Phase 3 only changes Intro/First/Second.
  stage_intro/first/second -- "Y" or blank; eligibility at that stage (new path).
  composition           -- blank (any) / "pure" / "mixed" -- deck_composition gate.
  trigger               -- blank (none) / "gcc" / "asked" / "manual" / "client_context"
                          ("manual"/"client_context" rows are NEVER auto-included
                          via the registry pass -- on-demand only, or handled by
                          a dedicated module e.g. client_context.py)

Three slide-ids' rules can't be expressed as a flat row at all (see brief Slides
8, 28-29, 18-19) -- those are special-cased BY ID directly in matcher.py /
client_context.py, not encoded here; this script marks their include_rule with a
plain sentinel ("SPECIAL" / "MANUAL - filled via client_context module...") so a
reader of the spreadsheet isn't misled into thinking the registry alone decides
their fate.

Run once (safe to re-run -- fully idempotent, always recomputed from RULES below):
    py scripts/phase3_registry_rules.py
"""

import openpyxl

REGISTRY = config.REGISTRY_XLSX

NEW_COLUMNS = ["include_rule_legacy", "stage_intro", "stage_first", "stage_second",
               "composition", "trigger"]

# ---------------------------------------------------------------------------
# The rule table -- one entry per slide id covered by SLIDE_SELECTION_BRIEF.md.
# Anything NOT listed here is left completely untouched apart from getting its
# include_rule_legacy backfilled (old legacy master-deck rows, CS61/62/68's
# already-correct OPTIONAL rows, CS09/14/17 dividers already hard-excluded via
# matcher.EXCLUDE regardless of registry text).
# ---------------------------------------------------------------------------
# stage = (intro Y/blank, first Y/blank, second Y/blank)
RULES = {
    # slide 1 - Title: unconditional, every stage (unchanged, no touch needed,
    # included here only so include_rule_legacy is explicit)
    "CS01":  dict(rule="ALWAYS"),
    # slide 2 - Who We Are (full): SPECIAL (First/Second any work type; Intro
    # only for AI Pods/Workforce -- owner-confirmed 2026-07-07, see matcher._rule_cs02)
    "CS02":  dict(rule="SPECIAL"),
    # slide 3 - Who We Are (short): SPECIAL (Intro only, pure MS only)
    "CS02B": dict(rule="SPECIAL"),
    # slide 4 - Global Footprint & Industries We Serve: any stage, GCC-triggered
    "CS137": dict(stage=("Y", "Y", "Y"), wt=None, trigger="gcc"),
    # slide 5 - Architects of Transformation (long bios): First/Second only
    "CS03":  dict(stage=("", "Y", "Y"), wt=None, trigger=""),
    # slide 8 - What We Offer: SPECIAL (composition-conditional at Intro only)
    "CS04":  dict(rule="SPECIAL"),
    # slide 9 - What We Own: unconditional, every stage
    "CS06":  dict(rule="ALWAYS"),
    # slides 10/11 - GCC Build Offerings / Greenfield & Brownfield: any stage, GCC
    "CS138": dict(stage=("Y", "Y", "Y"), wt=None, trigger="gcc"),
    "CS139": dict(stage=("Y", "Y", "Y"), wt=None, trigger="gcc"),
    # slide 13 - J2W Operating Model: First/Second only
    "CS05":  dict(stage=("", "Y", "Y"), wt=None, trigger=""),
    # slide 14 - Sourcing Engine: any stage, WF only
    "CS10":  dict(stage=("Y", "Y", "Y"), wt="WORKFORCE", trigger=""),
    # slide 15 - WF How We Deliver: on-demand only, WF
    "CS11":  dict(stage=("", "", ""), wt="WORKFORCE", trigger="manual"),
    # slide 16 - Structural Differentiation: First/Second, WF, asked
    "CS12":  dict(stage=("", "Y", "Y"), wt="WORKFORCE", trigger="asked"),
    # slide 17 - WF What We Deploy: First/Second, WF, asked
    "CS63":  dict(stage=("", "Y", "Y"), wt="WORKFORCE", trigger="asked"),
    # slides 18/19 - Client Context / Tailored Approach: handled entirely by
    # client_context.py, never via the registry pass-1 mechanism
    "CS65":  dict(rule="CLIENT_CONTEXT"),
    "CS66":  dict(rule="CLIENT_CONTEXT"),
    # slide 20 - Why J2W: First/Second, WF, asked
    "CS13":  dict(stage=("", "Y", "Y"), wt="WORKFORCE", trigger="asked"),
    # slide 21 - Our Engagement Model: First/Second, WF, asked
    "CS140": dict(stage=("", "Y", "Y"), wt="WORKFORCE", trigger="asked"),
    # slide 22 - Contract Process Outsourcing Operating Model: First/Second, WF, asked
    "CS141": dict(stage=("", "Y", "Y"), wt="WORKFORCE", trigger="asked"),
    # slide 23 - The Strategic Hiring Imperative: First/Second, WF, asked
    "CS142": dict(stage=("", "Y", "Y"), wt="WORKFORCE", trigger="asked"),
    # slides 24/25/26 - Lifecycle / Engagement & Retention / Competitive Advantages:
    # on-demand only, WF (26 also carries the {{CLIENT}} token, unrelated to this)
    "CS143": dict(stage=("", "", ""), wt="WORKFORCE", trigger="manual"),
    "CS144": dict(stage=("", "", ""), wt="WORKFORCE", trigger="manual"),
    "CS145": dict(stage=("", "", ""), wt="WORKFORCE", trigger="manual"),
    # slides 28/29 - AI-First Pod Model + AI Research Wing: SPECIAL (pure-AIP-only
    # at Intro/First, ALWAYS at Second whenever AI Pods is in the deck)
    "CS15":  dict(rule="SPECIAL"),
    "CS16":  dict(rule="SPECIAL"),
    # slide 30 - AI Pods + Interview-as-a-Service: First/Second, AIP, asked
    "CS146": dict(stage=("", "Y", "Y"), wt="AI_POD", trigger="asked"),
    # slide 31 - AI Pod Accountability Metrics: on-demand only, AIP (already
    # OPTIONAL/manual-equivalent today; re-stated here for clarity)
    "CS68":  dict(stage=("", "", ""), wt="AI_POD", trigger="manual"),
    # slide 33 - Managed Service Solutions overview: First/Second only (owner-
    # confirmed, NOT Intro), MS
    "CS18":  dict(stage=("", "Y", "Y"), wt="MS", trigger=""),
    # slide 34 - Strategic Managed Services Offerings: First/Second, MS
    "CS19":  dict(stage=("", "Y", "Y"), wt="MS", trigger=""),
    # slide 35 - Case Studies divider: ALWAYS, the deliberate exception among dividers
    "CS147": dict(rule="ALWAYS"),
    # slide 36 - Next Steps: First/Second, WF or AI Pods (not pure MSS)
    "CS07":  dict(stage=("", "Y", "Y"), wt="WORKFORCE,AI_POD", trigger=""),
    # slide 37 - How We Move Forward: First/Second, MS
    "CS148": dict(stage=("", "Y", "Y"), wt="MS", trigger=""),
    # slide 38 - "Let's win together": unconditional, every stage
    "CS08":  dict(rule="ALWAYS"),
    # slide 39 - Appendix: manual/on-demand only
    "CS149": dict(stage=("", "", ""), wt=None, trigger="manual"),
}

# Rows that don't exist yet (Phase 1's swap queued these but the write failed) --
# same defaults as swap_master_deck_39slide.py used, PLUS the real Phase-3 rule
# from RULES above (no more "PENDING_PHASE3" placeholder -- this replaces it).
NEW_ROWS = {
    "CS02B": {"section": "CORE", "kind": "STANDARD", "std_group": "CORE", "title": ""},
    "CS137": {"section": "CORE", "kind": "STANDARD", "std_group": "CORE", "title": ""},
    "CS138": {"section": "CORE", "kind": "STANDARD", "std_group": "CORE", "title": ""},
    "CS139": {"section": "CORE", "kind": "STANDARD", "std_group": "CORE", "title": ""},
    "CS140": {"section": "WORKFORCE", "kind": "STANDARD", "std_group": "WORKFORCE", "title": ""},
    "CS141": {"section": "WORKFORCE", "kind": "STANDARD", "std_group": "WORKFORCE", "title": ""},
    "CS142": {"section": "WORKFORCE", "kind": "STANDARD", "std_group": "WORKFORCE", "title": ""},
    "CS143": {"section": "WORKFORCE", "kind": "OPTIONAL", "std_group": "OPTIONAL", "title": ""},
    "CS144": {"section": "WORKFORCE", "kind": "OPTIONAL", "std_group": "OPTIONAL", "title": ""},
    "CS145": {"section": "WORKFORCE", "kind": "OPTIONAL", "std_group": "OPTIONAL", "title": ""},
    "CS146": {"section": "AI_POD", "kind": "STANDARD", "std_group": "AI_POD", "title": ""},
    "CS147": {"section": "CASE_STUDIES", "kind": "DIVIDER", "std_group": "CASE_STUDIES", "title": ""},
    "CS148": {"section": "MS", "kind": "STANDARD", "std_group": "MS", "title": ""},
    "CS149": {"section": "APPENDIX", "kind": "OPTIONAL", "std_group": "OPTIONAL", "title": ""},
}


def _title_of(deck_path, slide_id):
    from pptx import Presentation
    from deckengine.services.content.build_library import read_id
    prs = Presentation(deck_path)
    for s in prs.slides:
        if read_id(s) == slide_id:
            for sh in s.shapes:
                if sh.has_text_frame and sh.text_frame.text.strip():
                    return sh.text_frame.text.strip().replace("\n", " ")[:120]
    return ""


def main():
    wb = openpyxl.load_workbook(REGISTRY)
    ws = wb["Slide Registry"]
    hdr = [c.value for c in ws[1]]
    col = {h: i + 1 for i, h in enumerate(hdr)}   # 1-indexed for openpyxl cells

    # ---- add any missing new columns at the end ----
    next_col = len(hdr) + 1
    for name in NEW_COLUMNS:
        if name not in col:
            ws.cell(1, next_col).value = name
            col[name] = next_col
            next_col += 1

    # ---- backfill include_rule_legacy for EVERY existing row (Proposal safety net) ----
    for r in range(2, ws.max_row + 1):
        if ws.cell(r, col["include_rule_legacy"]).value in (None, ""):
            ws.cell(r, col["include_rule_legacy"]).value = ws.cell(r, col["include_rule"]).value

    # ---- append the rows Phase 1 queued but never wrote ----
    existing_ids = {ws.cell(r, col["slide_id"]).value for r in range(2, ws.max_row + 1)}
    for sid, defaults in NEW_ROWS.items():
        if sid in existing_ids:
            continue
        row = [None] * len(hdr)
        row[col["slide_id"] - 1] = sid
        row[col["section"] - 1] = defaults["section"]
        row[col["kind"] - 1] = defaults["kind"]
        row[col["std_group"] - 1] = defaults["std_group"]
        row[col["confidence"] - 1] = "MANUAL"
        row[col["title"] - 1] = _title_of(config.MASTER_DECK, sid)
        row[col["include_rule"] - 1] = "PENDING (set below)"
        ws.append(row)
        existing_ids.add(sid)
        print(f"Added missing row: {sid}")

    # rebuild sid -> row index now that new rows exist
    row_of = {ws.cell(r, col["slide_id"]).value: r for r in range(2, ws.max_row + 1)}

    # ---- write the Phase-3 rule for every id in RULES ----
    for sid, rule in RULES.items():
        r = row_of.get(sid)
        if r is None:
            print(f"WARNING: {sid} not found in registry -- skipped (check the id).")
            continue

        if rule.get("rule") == "ALWAYS":
            ws.cell(r, col["include_rule"]).value = "ALWAYS"
            continue
        if rule.get("rule") == "SPECIAL":
            ws.cell(r, col["include_rule"]).value = "SPECIAL - see matcher._SPECIAL_STAGE_RULES"
            continue
        if rule.get("rule") == "CLIENT_CONTEXT":
            ws.cell(r, col["include_rule"]).value = "MANUAL - filled via client_context module, not registry-driven"
            continue

        s_i, s_f, s_s = rule["stage"]
        ws.cell(r, col["stage_intro"]).value = s_i or None
        ws.cell(r, col["stage_first"]).value = s_f or None
        ws.cell(r, col["stage_second"]).value = s_s or None
        ws.cell(r, col["composition"]).value = rule.get("composition") or None
        ws.cell(r, col["trigger"]).value = rule.get("trigger") or None
        if rule.get("wt") is not None:
            ws.cell(r, col["work_types"]).value = rule["wt"]
        trig = rule.get("trigger") or "none"
        ws.cell(r, col["include_rule"]).value = (
            f"STAGE_RULE (stage={''.join(x or '-' for x in rule['stage'])}, "
            f"wt={rule.get('wt') or ws.cell(r, col['work_types']).value}, "
            f"trigger={trig})"
        )

    wb.save(REGISTRY)
    print(f"\nSaved -> {REGISTRY}")
    print(f"Rows with a Phase-3 rule written: {len(RULES)}")


if __name__ == "__main__":
    main()
