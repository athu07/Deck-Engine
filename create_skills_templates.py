# -*- coding: utf-8 -*-
"""
create_skills_templates.py
Builds skills_templates.pptx from scratch — no source PPT needed.

The 3 data-driven "capability" slides, drawn programmatically with the EXACT
{{PLACEHOLDER}} names + J2W_TEMPLATE notes tags the engine (skills.py) expects.

Design language matches the case-study template (case_study_v2.pptx):
  - red + teal split bar across the very top
  - red accent bar to the left of a 24pt black title
  - teal subtitle line carrying the dynamic name (like the case study's
    CLIENT | DOMAIN line)
  - white, hairline-bordered cards with a red or teal left/top accent
  - the J2W brand palette throughout, plain hyphens only (no em-dashes)

Run once (safe to re-run — overwrites skills_templates.pptx):
    py create_skills_templates.py
"""

from pptx import Presentation
from pptx.util import Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.oxml.ns import qn
from lxml import etree

OUT = "skills_templates.pptx"

# ── J2W brand palette (identical to create_case_study_template.py) ────────────
C_BLACK = RGBColor(0x11, 0x11, 0x10)
C_RED   = RGBColor(0xC0, 0x20, 0x26)   # CRIMSON RED
C_TEAL  = RGBColor(0x3A, 0x8B, 0x82)   # DEEP TEAL
C_WHITE = RGBColor(0xFF, 0xFF, 0xFF)
C_CARD  = RGBColor(0xF5, 0xF5, 0xF5)   # light card fill
C_LINE  = RGBColor(0xDE, 0xDE, 0xDE)   # hairline border
C_BODY  = RGBColor(0x3E, 0x3E, 0x3E)   # body text
C_MUTE  = RGBColor(0x6E, 0x6E, 0x69)   # muted labels

FONT = "Calibri"

# ── Slide dimensions (13.33" x 7.50" widescreen) ─────────────────────────────
SW, SH = 13.33, 7.50
EMU    = 914400
I      = lambda v: int(v * EMU)

# ── Shared layout constants ───────────────────────────────────────────────────
SPLIT   = 6.40    # where the red top bar ends / teal begins
MARGIN  = 0.30
PAD     = 0.16


# ── Low-level helpers ─────────────────────────────────────────────────────────
def _blank_layout(prs):
    for layout in prs.slide_layouts:
        if (layout.name or "").lower().strip() == "blank":
            return layout
    return prs.slide_layouts[-1]


def _no_border(shape):
    spPr = shape._element.spPr
    ln = spPr.find(qn('a:ln'))
    if ln is None:
        ln = etree.SubElement(spPr, qn('a:ln'))
    for c in list(ln):
        ln.remove(c)
    etree.SubElement(ln, qn('a:noFill'))


def bar(slide, l, t, w, h, fill):
    """A flat, border-less filled rectangle (bars / accents)."""
    s = slide.shapes.add_shape(1, I(l), I(t), I(w), I(h))
    s.fill.solid()
    s.fill.fore_color.rgb = fill
    _no_border(s)
    s.shadow.inherit = False
    return s


def card(slide, l, t, w, h, fill=C_WHITE, line=C_LINE, lw=0.75):
    """A flat card with a hairline border."""
    s = slide.shapes.add_shape(1, I(l), I(t), I(w), I(h))
    s.fill.solid()
    s.fill.fore_color.rgb = fill
    s.line.color.rgb = line
    s.line.width = Pt(lw)
    s.shadow.inherit = False
    return s


def txt(slide, l, t, w, h, text, size, color,
        bold=False, align=PP_ALIGN.LEFT, italic=False,
        anchor='t', wrap=True):
    tb = slide.shapes.add_textbox(I(l), I(t), I(w), I(h))
    tf = tb.text_frame
    tf.word_wrap = wrap
    tf._txBody.bodyPr.set('anchor', anchor)
    tf.margin_top = tf.margin_bottom = 0
    p = tf.paragraphs[0]
    p.alignment = align
    r = p.add_run()
    r.text = text
    r.font.name = FONT
    r.font.size = Pt(size)
    r.font.color.rgb = color
    r.font.bold = bold
    r.font.italic = italic
    return tb


# ── Shared chrome (top split bar + title + subtitle) ─────────────────────────
def header(slide, title, subtitle):
    # top split bar: crimson (left) + teal (right)
    bar(slide, 0.00, 0.00, SPLIT,      0.10, C_RED)
    bar(slide, SPLIT, 0.00, SW - SPLIT, 0.10, C_TEAL)
    # red accent bar left of the title
    bar(slide, 0.25, 0.30, 0.065, 0.52, C_RED)
    # title + subtitle
    txt(slide, 0.44, 0.26, 12.4, 0.52, title, 24, C_BLACK, bold=True)
    txt(slide, 0.44, 0.80, 12.6, 0.34, subtitle, 14, C_TEAL, bold=True)
    # slim teal footer strip for a finished, framed look
    bar(slide, 0.00, SH - 0.10, SW, 0.10, C_TEAL)


def section_label(slide, l, t, w, label):
    txt(slide, l + PAD, t + PAD, w - PAD * 2, 0.28, label, 11, C_TEAL, bold=True)
    bar(slide, l + PAD, t + PAD + 0.30, w - PAD * 2, 0.022, C_TEAL)


def text_panel(slide, l, t, w, h, accent, label, marker, body_size=10):
    """White card + left accent bar + section label + a marker body box."""
    card(slide, l, t, w, h)
    bar(slide, l, t, 0.06, h, accent)                 # left accent
    section_label(slide, l, t, w, label)
    txt(slide, l + PAD, t + PAD + 0.44, w - PAD * 2, h - PAD - 0.52,
        marker, body_size, C_BODY, wrap=True, anchor='t')


def chart_panel(slide, l, t, w, h, label):
    """White card + teal top accent + section label + the {{CHART}} box."""
    card(slide, l, t, w, h)
    bar(slide, l, t, w, 0.06, C_TEAL)                 # top accent
    section_label(slide, l, t, w, label)
    # the {{CHART}} placeholder box — skills.py reads its position, removes it,
    # and drops the native chart in its place.
    txt(slide, l + PAD, t + PAD + 0.44, w - PAD * 2, h - PAD - 0.52,
        "{{CHART}}", 10, C_MUTE, anchor='t')


def metric_tile(slide, l, t, w, h, value_marker, label):
    """White stat card + teal top accent + big teal number + muted label."""
    card(slide, l, t, w, h)
    bar(slide, l, t, w, 0.05, C_TEAL)
    txt(slide, l + 0.08, t + 0.20, w - 0.16, 0.66,
        value_marker, 27, C_TEAL, bold=True, align=PP_ALIGN.CENTER)
    txt(slide, l + 0.08, t + 0.90, w - 0.16, h - 0.96,
        label, 9, C_MUTE, align=PP_ALIGN.CENTER, anchor='t')


def set_notes(slide, tag):
    slide.notes_slide.notes_text_frame.text = tag


# ── geometry shared by the two-tile slides ───────────────────────────────────
def _four_tiles(slide, tiles, top):
    gap = 0.18
    tw = (SW - 2 * MARGIN - 3 * gap) / 4
    th = 1.28
    for i, (marker, label) in enumerate(tiles):
        tl = MARGIN + i * (tw + gap)
        metric_tile(slide, tl, top, tw, th, marker, label)
    return top + th


def _lower_split():
    """(left_l, left_w, right_l, right_w) for the lower two-panel row."""
    left_w = 5.55
    gap = 0.24
    right_l = MARGIN + left_w + gap
    right_w = SW - right_l - MARGIN
    return MARGIN, left_w, right_l, right_w


# ── SLIDE 1 — skill_deepdive ──────────────────────────────────────────────────
#   markers: {{SKILLS_HEADER}} {{SKILL_SUMMARY}} {{CHART}}
def build_slide1(prs):
    slide = prs.slides.add_slide(_blank_layout(prs))
    header(slide, "WHERE J2W HAS DELIVERED",
           "SKILLS DEPLOYED:  {{SKILLS_HEADER}}")

    top = 1.32
    h = SH - top - 0.28
    ll, lw, rl, rw = _lower_split()
    text_panel(slide, ll, top, lw, h, C_RED, "DEPLOYMENT SUMMARY", "{{SKILL_SUMMARY}}")
    chart_panel(slide, rl, top, rw, h, "SKILL DISTRIBUTION")

    set_notes(slide, "J2W_TEMPLATE: skill_deepdive")
    print("  Slide 1 (skill_deepdive): done")


# ── SLIDE 2 — industry_strength ───────────────────────────────────────────────
#   markers: {{INDUSTRY_NAME}} {{TOTAL_CONSULTANTS}} {{NUM_COMPANIES}}
#            {{NUM_FUNCTIONS}} {{NUM_SKILLS}} {{TOP_SKILL_1/2/3}} {{CHART}}
def build_slide2(prs):
    slide = prs.slides.add_slide(_blank_layout(prs))
    header(slide, "OUR DELIVERY FOOTPRINT",
           "INDUSTRY:  {{INDUSTRY_NAME}}")

    tiles = [
        ("{{TOTAL_CONSULTANTS}}", "Consultants deployed"),
        ("{{NUM_COMPANIES}}",     "Companies served"),
        ("{{NUM_FUNCTIONS}}",     "Functions delivered"),
        ("{{NUM_SKILLS}}",        "Distinct skills"),
    ]
    tiles_bottom = _four_tiles(slide, tiles, 1.30)

    top = tiles_bottom + 0.22
    h = SH - top - 0.28
    ll, lw, rl, rw = _lower_split()

    # left: top-3 hired skills, ranked cards
    card(slide, ll, top, lw, h)
    bar(slide, ll, top, 0.06, h, C_RED)
    section_label(slide, ll, top, lw, "TOP HIRED SKILLS - BY COMPANIES SERVED")
    card_t0 = top + PAD + 0.46
    avail = h - PAD - 0.54
    cgap = 0.10
    ch = (avail - 2 * cgap) / 3
    badge_w = 0.46
    for j, marker in enumerate(["{{TOP_SKILL_1}}", "{{TOP_SKILL_2}}", "{{TOP_SKILL_3}}"]):
        ct = card_t0 + j * (ch + cgap)
        card(slide, ll + PAD, ct, lw - PAD * 2, ch, fill=C_CARD)
        bar(slide, ll + PAD, ct, badge_w, ch, C_RED)
        txt(slide, ll + PAD, ct, badge_w, ch, f"#{j+1}", 14, C_WHITE,
            bold=True, align=PP_ALIGN.CENTER, anchor='ctr')
        txt(slide, ll + PAD + badge_w + 0.12, ct + 0.06, lw - PAD * 2 - badge_w - 0.18,
            ch - 0.12, marker, 11, C_BODY, wrap=True, anchor='ctr')

    # right: function distribution pie
    chart_panel(slide, rl, top, rw, h, "FUNCTION DISTRIBUTION")

    set_notes(slide, "J2W_TEMPLATE: industry_strength")
    print("  Slide 2 (industry_strength): done")


# ── SLIDE 3 — company_footprint ───────────────────────────────────────────────
#   markers: {{COMPANY_NAME}} {{TOTAL_DEPLOYED}} {{NUM_FUNCTIONS_CO}}
#            {{NUM_SKILLS_CO}} {{ENGAGEMENT_TYPE}} {{FUNCTION_BREAKDOWN}} {{CHART}}
def build_slide3(prs):
    slide = prs.slides.add_slide(_blank_layout(prs))
    header(slide, "OUR PARTNERSHIP OVERVIEW",
           "CLIENT:  {{COMPANY_NAME}}   |   EXISTING J2W RELATIONSHIP")

    tiles = [
        ("{{TOTAL_DEPLOYED}}",   "Consultants deployed"),
        ("{{NUM_FUNCTIONS_CO}}", "Functions delivered"),
        ("{{NUM_SKILLS_CO}}",    "Distinct skills"),
        ("{{ENGAGEMENT_TYPE}}",  "Engagement type"),
    ]
    tiles_bottom = _four_tiles(slide, tiles, 1.30)

    top = tiles_bottom + 0.22
    h = SH - top - 0.28
    ll, lw, rl, rw = _lower_split()
    text_panel(slide, ll, top, lw, h, C_RED, "FUNCTION BREAKDOWN", "{{FUNCTION_BREAKDOWN}}")
    chart_panel(slide, rl, top, rw, h, "FUNCTION DISTRIBUTION")

    set_notes(slide, "J2W_TEMPLATE: company_footprint")
    print("  Slide 3 (company_footprint): done")


# ── Main ──────────────────────────────────────────────────────────────────────
if __name__ == "__main__":
    prs = Presentation()
    prs.slide_width  = I(SW)
    prs.slide_height = I(SH)

    build_slide1(prs)
    build_slide2(prs)
    build_slide3(prs)

    prs.save(OUT)
    print(f"\nSaved: {OUT}  ({len(prs.slides)} slides)")
    for i, s in enumerate(prs.slides, 1):
        note = s.notes_slide.notes_text_frame.text.strip()
        print(f"  Slide {i} notes: {repr(note)}")
