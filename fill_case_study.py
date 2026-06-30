# -*- coding: utf-8 -*-
"""
fill_case_study.py  --  Pour ONE content-store row into the case-study template
and save a finished, single-slide .pptx.

This is the heart of the content-referencing model: a case study lives as DATA
in case_study_content_store.json, and a real slide is built on demand by
dropping that data into case_study_v2.pptx (the J2W branded template).

What it does:
  * opens the template (1 slide, all {{MARKERS}})
  * builds a value for every marker from the content row
  * smart-splits each narrative RESULT into a headline metric + caption
    (e.g. "13 engineers added in the initial phase" -> "13" / "engineers added…")
  * splits a CAPABILITY into title/body when it carries a separator, else
    shows it as a title-only card
  * replaces markers IN PLACE, preserving the template's fonts/colours
  * any marker with no data is blanked (never left as "{{…}}" on the slide)

Use as a library:
    import fill_case_study as fcs
    fcs.fill_row(row_dict, "out.pptx", client_name="Acme Bank")

Or from the CLI (demo on one stored case):
    py fill_case_study.py                 # fills MSS001 -> output/<id>.pptx
    py fill_case_study.py AIP007 "Acme"   # fill a chosen id, with a client name
"""

import json
import os
import re
import sys

from pptx import Presentation

TEMPLATE = "case_study_v2.pptx"
STORE = "case_study_content_store.json"
OUT_DIR = "output"

MARKER_RE = re.compile(r"\{\{([A-Z0-9_]+)\}\}")

# headline metric inside a results sentence: 70%, 2.5x, $4M, 13, 24x7, 100+
_METRIC_RE = re.compile(
    r"(\d+(?:\.\d+)?\s*%"            # 70%  12.5 %
    r"|\d+(?:\.\d+)?\s*[xX]\d*"      # 2x  24x7  3.5x
    r"|\$\s*\d[\d,]*\.?\d*\s*[KMB]?" # $4M  $1,200
    r"|\b\d[\d,]*\+?)"              # 13  1,200  100+
)

_CAP_SEPS = [" — ", " – ", " - ", ": "]

# connectives to trim off the edges of a caption once the highlight is pulled
_LEAD_RE = re.compile(r"^(?:of|by|in|to|the|a|an|and|with|via|for|on|at|from)\s+", re.I)
_TRAIL_RE = re.compile(r"\s+(?:of|by|in|to|the|a|an|and|with|via|for|on|at|from)$", re.I)

# CONFIDENTIALITY: a case-study slide never names a real client/company — only
# J2W is named anywhere. The CLIENT line shows a generic anonymised descriptor,
# matching J2W's own deck convention ("Leading Manufacturing Institution").
_CLIENT_SUFFIX = {
    "BFSI": "Institution", "PRIVATE_EQUITY": "Firm", "AVIATION": "Operator",
    "ENERGY": "Utility", "TELECOM": "Operator",
}


def anon_client(row):
    """A generic, name-free client descriptor derived from the domain.
    Never returns a real company name."""
    domain = (row.get("domain") or "").strip()
    if not domain:
        return "Global Enterprise"
    suffix = _CLIENT_SUFFIX.get(row.get("industry"), "Enterprise")
    # don't echo a word the domain already carries ("Leading Retail Retailer")
    domain_words = {w.lower() for w in re.findall(r"[A-Za-z]+", domain)}
    if suffix.lower() in domain_words:
        suffix = "Enterprise"
    if suffix.lower() in domain_words:        # domain literally contains "Enterprise"
        return f"Leading {domain}"
    return f"Leading {domain} {suffix}"


def split_result(sentence):
    """Return (highlight, caption) for one result.

    The highlight is the MOST relevant short bit shown big on top — a number
    if there is one (70%, 24x7, $4M), otherwise the leading key phrase
    ("Automated", "Single platform"). The rest becomes the caption below."""
    s = (sentence or "").strip().rstrip(".")
    if not s:
        return "", ""

    m = _METRIC_RE.search(s)
    if m:
        top = re.sub(r"\s+", "", m.group(1))
        bottom = (s[:m.start()] + " " + s[m.end():])
    else:
        words = s.split()
        first = words[0]
        # short modifier (Single/Zero/Real) reads better paired with the next word
        if len(first) <= 7 and len(words) > 1:
            top = first + " " + words[1]
            bottom = " ".join(words[2:])
        else:
            top = first
            bottom = " ".join(words[1:])
        top = top[0].upper() + top[1:]

    bottom = _LEAD_RE.sub("", bottom.strip(" ,–—-"))
    bottom = _TRAIL_RE.sub("", bottom).strip(" ,–—-")
    bottom = re.sub(r"\s{2,}", " ", bottom)
    return top, bottom


def split_capability(cap):
    """Normalise one capability to (title, body).
    Accepts a {title, body} dict (enriched store) or a plain string (splits on
    a separator if present, else title-only)."""
    if isinstance(cap, dict):
        return (cap.get("title") or "").strip(), (cap.get("body") or "").strip()
    c = (cap or "").strip()
    for sep in _CAP_SEPS:
        if sep in c:
            t, b = c.split(sep, 1)
            return t.strip(), b.strip()
    return c, ""


def _pad(seq, n):
    seq = list(seq)[:n]
    return seq + [""] * (n - len(seq))


def build_mapping(row):
    """Marker -> replacement string, for one content-store row."""
    caps = _pad(row.get("capabilities", []), 6)
    results = _pad(row.get("results", []), 3)

    mapping = {
        "TITLE": row.get("title", ""),
        "CLIENT": anon_client(row),
        "DOMAIN": row.get("domain", "") or "-",
        "CHALLENGE": row.get("challenge", ""),
        "SOLUTION": row.get("solution", ""),
    }
    for i, cap in enumerate(caps, 1):
        t, b = split_capability(cap)
        mapping[f"CAP_{i}_TITLE"] = t
        mapping[f"CAP_{i}_BODY"] = b
    for i, res in enumerate(results, 1):
        pct, txt = split_result(res)
        mapping[f"RESULT_{i}_PCT"] = pct
        mapping[f"RESULT_{i}_TEXT"] = txt
    # owner rule: no em/en dash ever reaches a slide — always a plain hyphen
    return {k: (v or "").replace("—", "-").replace("–", "-") for k, v in mapping.items()}


def _apply(text, mapping):
    """Replace every {{MARKER}}; unknown markers become empty string."""
    return MARKER_RE.sub(lambda m: mapping.get(m.group(1), ""), text)


def fill_row(row, out_path, template=TEMPLATE):
    """Build a finished slide from one content row. Returns out_path.

    Note: there is deliberately no client-name parameter — a case-study slide
    never carries a real client/company name (only J2W is named)."""
    prs = Presentation(template)
    slide = prs.slides[0]
    mapping = build_mapping(row)

    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue
        for para in shape.text_frame.paragraphs:
            for run in para.runs:
                if "{{" in run.text:
                    run.text = _apply(run.text, mapping)

    out_dir = os.path.dirname(out_path)
    if out_dir:
        os.makedirs(out_dir, exist_ok=True)
    prs.save(out_path)
    return out_path, mapping


def _load_store():
    with open(STORE, encoding="utf-8") as f:
        return json.load(f)


def fill_by_id(case_id, out_path=None):
    recs = {r["id"]: r for r in _load_store()}
    if case_id not in recs:
        raise KeyError(f"{case_id} not in {STORE}")
    out_path = out_path or os.path.join(OUT_DIR, f"{case_id}.pptx")
    return fill_row(recs[case_id], out_path)


if __name__ == "__main__":
    try:                                   # so the ✓ stat prints on a cp1252 console
        sys.stdout.reconfigure(encoding="utf-8")
    except (AttributeError, ValueError):
        pass
    cid = sys.argv[1] if len(sys.argv) > 1 else "MSS001"
    out, mapping = fill_by_id(cid)

    print(f"Filled {cid} -> {out}")
    print(f"Client: {mapping['CLIENT']}   Domain: {mapping['DOMAIN']}")
    print(f"Title : {mapping['TITLE']}")
    print("\nCapabilities:")
    for i in range(1, 7):
        t, b = mapping[f"CAP_{i}_TITLE"], mapping[f"CAP_{i}_BODY"]
        print(f"  {i}. {t}" + (f"  —  {b}" if b else ""))
    print("\nResults:")
    for i in range(1, 4):
        print(f"  {mapping[f'RESULT_{i}_PCT']:>6}  {mapping[f'RESULT_{i}_TEXT']}")

    # safety: no leftover markers anywhere on the slide
    leftover = [v for v in mapping.values() if "{{" in str(v)]
    print("\nLeftover markers in values:", leftover or "none")
