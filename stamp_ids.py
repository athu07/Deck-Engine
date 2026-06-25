# -*- coding: utf-8 -*-
"""
stamp_ids.py  --  give every slide a permanent, stable ID (Prompt D).

Writes a machine-readable line into each slide's SPEAKER NOTES:

    J2W_ID: CS01

Why notes? The ID then lives INSIDE the slide and travels with it no matter
where the slide moves, so nothing depends on position anymore.

STABLE-BY-DESIGN assignment:
  - A slide that ALREADY has a J2W_ID is left completely untouched.
  - Only slides WITHOUT an ID get a new one, numbered (highest existing) + 1,
    in document order.
  => Inserting / reordering slides and re-running NEVER renumbers existing
     slides. New slides simply get the next free IDs. Re-running is a no-op
     once everything is stamped.

Always run this on the COPY, never the original.
"""

import re
import sys
from pptx import Presentation

ID_PREFIX = "J2W_ID:"
ID_LINE_RE = re.compile(r"^J2W_ID:\s*(\S+)\s*$", re.MULTILINE)
ID_NUM_RE = re.compile(r"^CS(\d+)$")


def _existing_id(slide):
    """Return the slide's current ID, or None. Does NOT create a notes part."""
    if not slide.has_notes_slide:
        return None
    m = ID_LINE_RE.search(slide.notes_slide.notes_text_frame.text or "")
    return m.group(1) if m else None


def stamp(path):
    prs = Presentation(path)
    slides = list(prs.slides)

    # Pass 1: read what's already there and find the highest CS number in use.
    current = []
    max_num = 0
    for slide in slides:
        cur = _existing_id(slide)
        current.append(cur)
        if cur:
            m = ID_NUM_RE.match(cur)
            if m:
                max_num = max(max_num, int(m.group(1)))

    # Pass 2: assign IDs ONLY to slides that don't have one yet.
    kept, assigned = [], []
    next_num = max_num
    for slide, cur in zip(slides, current):
        if cur:
            kept.append(cur)
            continue
        next_num += 1
        new_id = f"CS{next_num:02d}"
        tf = slide.notes_slide.notes_text_frame   # creates the notes part now
        existing = tf.text or ""
        if existing.strip():
            tf.text = f"{ID_PREFIX} {new_id}\n{existing}"   # keep real notes
        else:
            tf.text = f"{ID_PREFIX} {new_id}"
        assigned.append(new_id)

    prs.save(path)
    return len(slides), kept, assigned


if __name__ == "__main__":
    target = sys.argv[1] if len(sys.argv) > 1 else "WORKING_COPY_Master_Deck.pptx"
    total, kept, assigned = stamp(target)
    print(f"File: {target}")
    print(f"  slides total      : {total}")
    print(f"  already had an ID  : {len(kept)} (left untouched)")
    print(f"  newly assigned     : {len(assigned)} {assigned if assigned else ''}")
