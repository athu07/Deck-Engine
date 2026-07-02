# -*- coding: utf-8 -*-
"""
slide_generator.py  --  Step 04: turn a "needs to be created" gap into a real
slide, built from a TEMPLATE.

Design = pluggable templates:
  - Templates live in templates.pptx. Each template is one slide, tagged in its
    notes with  J2W_TEMPLATE: <name>  (e.g. case_study), and contains marker
    tokens in its text:  {{TITLE}}  {{KEYWORDS}}  {{BULLETS}}
  - To generate a slide: the AI writes the words, then we COPY the chosen
    template slide into the deck and REPLACE the markers with those words.
  - Add another template later = add another tagged slide to templates.pptx.
    No code change. Swap the temporary template for the real J2W design anytime.

The template made by create_temp_template() is a PLACEHOLDER — plain text boxes,
text-only (images in a template need extra work). Replace it with the real
J2W-designed template slide when ready; keep the same marker tokens + tag.
"""

import copy
import json

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

import editor
from build_library import read_id  # noqa: F401  (kept for parity / future use)

MASTER = "WORKING_COPY_Master_Deck.pptx"
TEMPLATES_FILE = "templates.pptx"
MODEL = "gpt-4o-mini"
TEMPLATE_TAG = "J2W_TEMPLATE:"


# --------------------------------------------------------------------------- #
# Template file
# --------------------------------------------------------------------------- #
def create_temp_template(path=TEMPLATES_FILE):
    """Create a temporary, text-only 'case_study' template slide."""
    master = Presentation(MASTER)
    tp = Presentation()
    tp.slide_width = master.slide_width
    tp.slide_height = master.slide_height
    slide = tp.slides.add_slide(tp.slide_layouts[6])      # 6 = Blank
    slide.notes_slide.notes_text_frame.text = f"{TEMPLATE_TAG} case_study"

    def textbox(left, top, width, height, marker, size, color, bold=False, italic=False):
        tb = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
        p = tb.text_frame.paragraphs[0]
        run = p.add_run()
        run.text = marker
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.italic = italic
        run.font.color.rgb = color
        return tb

    textbox(0.6, 0.5, 12.0, 1.1, "{{TITLE}}", 28, RGBColor(0x14, 0x28, 0x50), bold=True)
    textbox(0.6, 1.6, 12.0, 0.5, "{{KEYWORDS}}", 13, RGBColor(0x0F, 0x6E, 0x56), italic=True)
    textbox(0.6, 2.3, 12.0, 4.0, "{{BULLETS}}", 16, RGBColor(0x33, 0x33, 0x33))
    textbox(0.6, 6.9, 12.0, 0.4, "Generated slide - replace with the real J2W template",
            10, RGBColor(0xAA, 0xAA, 0xAA), italic=True)
    tp.save(path)
    return path


def list_templates(path=TEMPLATES_FILE):
    """{template_name: slide} for every tagged slide in templates.pptx."""
    prs = Presentation(path)
    out = {}
    for slide in prs.slides:
        if slide.has_notes_slide:
            txt = slide.notes_slide.notes_text_frame.text or ""
            for line in txt.splitlines():
                if line.strip().startswith(TEMPLATE_TAG):
                    out[line.split(":", 1)[1].strip()] = slide
    return out


# --------------------------------------------------------------------------- #
# AI content
# --------------------------------------------------------------------------- #
def default_brief(work_type, industry, transcript, topic=""):
    """A starting 'what should this slide cover' brief, pre-filled from the notes.
    The salesperson edits it before generating. No AI call.
    If `topic` is given (a specific client ask with no slide, e.g. 'ADAS'), the
    brief is centred on demonstrating J2W's capability in that exact topic."""
    labels = {"WORKFORCE": "Workforce", "AI_POD": "AI Pod", "MS": "Managed Services"}
    t = (transcript or "").strip()
    snippet = (t[:180].rstrip() + "…") if len(t) > 180 else t
    topic = (topic or "").strip()
    if topic:
        base = "Slide demonstrating J2W's capability in %s" % topic
        if industry:
            base += " for the %s industry" % industry.replace("_", " ").title()
        base += " — the client specifically asked about this"
        base += (", in context: " + snippet) if snippet else ""
        return base.strip() + "."
    base = "%s slide" % labels.get(work_type, (work_type or "").replace("_", " ").title())
    if industry:
        base += " for the %s industry" % industry.replace("_", " ").title()
    base += (", addressing: " + snippet) if snippet else \
            ", covering J2W's relevant capability and the outcomes it delivers"
    return base.strip() + "."


def _similar_slides(work_type, query, n=3):
    """The closest real J2W case-study slides — used as a FORMAT/style example so a
    generated slide matches the existing deck instead of being invented from nothing."""
    try:
        lib = json.load(open("tagged_library.json", encoding="utf-8"))
    except Exception:
        return []
    q = (query or "").lower()
    scored = []
    for r in lib:
        tags = r.get("tags", {})
        if (tags.get("kind", {}) or {}).get("value") != "CASE_STUDY":
            continue
        kws = r.get("keywords", []) or []
        score = sum(1 for k in kws if k and k.lower() in q)
        if (tags.get("work_type", {}) or {}).get("value") == work_type:
            score += 1
        scored.append((score, r))
    scored.sort(key=lambda x: -x[0])
    return [{"title": r.get("title", ""), "keywords": " · ".join((r.get("keywords") or [])[:8])}
            for _, r in scored[:n]]


def draft(gap, context):
    """Ask the LLM to write a gap slide's content, GUIDED by a brief and grounded in
    the format of similar real J2W slides. Returns {title, keywords, bullets:[...]}.
    Falls back to a stub on any error."""
    wt = gap.get("work_type", "")
    topic = (gap.get("topic") or context.get("topic") or "").strip()
    industry = context.get("industry", "")
    transcript = (context.get("transcript") or "")[:3000]
    brief = (context.get("brief") or "").strip()
    examples = _similar_slides(wt, topic or brief or transcript)
    ex_text = "\n".join("  - %s (keywords: %s)" % (e["title"], e["keywords"]) for e in examples) or "  (none found)"
    focus = (f"The CLIENT specifically asked about \"{topic}\" and the deck has no "
             f"slide on it. Write a slide that demonstrates J2W's capability in "
             f"\"{topic}\".\n\n") if topic else ""
    prompt = (
        f"You are writing ONE slide for a J2W sales deck. Work type: {wt}; "
        f"industry: {industry or 'the client'}.\n\n"
        f"{focus}"
        f"WHAT THIS SLIDE SHOULD COVER (follow this brief):\n"
        f"{brief or '(no brief given — infer it from the meeting notes below)'}\n\n"
        f"MEETING NOTES (context):\n\"\"\"\n{transcript}\n\"\"\"\n\n"
        f"FORMAT — follow the style of these existing J2W slides:\n{ex_text}\n\n"
        "Write the slide so it satisfies the brief. It can be a CASE STUDY that "
        "follows the format of the examples above, OR a different client-specific "
        "slide if the brief calls for that. Keep every claim credible and do NOT "
        "invent a specific real client name. Return ONLY JSON: "
        '{"title": "...", "keywords": "A · B · C · D", "bullets": ["...", "...", "..."]}. '
        "3-4 short bullets, each a concrete outcome or capability."
    )
    try:
        from secrets_loader import load_env
        load_env()
        from openai import OpenAI
        resp = OpenAI().chat.completions.create(
            model=MODEL, temperature=0.4,
            response_format={"type": "json_object"},
            messages=[
                {"role": "system", "content": "You write concise B2B sales slide "
                                              "content. Reply with one JSON object only."},
                {"role": "user", "content": prompt},
            ],
        )
        data = json.loads(resp.choices[0].message.content)
        return {
            "title": str(data.get("title", "Proposed case study")),
            "keywords": str(data.get("keywords", "")),
            "bullets": [str(b) for b in (data.get("bullets") or [])][:4],
            "template": "case_study",
        }
    except Exception:
        return {
            "title": f"{wt} CASE STUDY - {industry or 'CLIENT'} (TO BE CREATED)",
            "keywords": "Draft · placeholder · replace",
            "bullets": ["Content could not be generated - add details manually."],
            "template": "case_study",
        }


# --------------------------------------------------------------------------- #
# Full structured CASE STUDY (the "Create with AI" feature) — strict format
# --------------------------------------------------------------------------- #
CASE_STUDY_RULES = (
    "NON-NEGOTIABLE RULES:\n"
    "- Exactly 6 capabilities, exactly 3 results, always.\n"
    "- No em dashes anywhere. All numbers as numerals. No real company names.\n"
    "- Capability names never reference technology labels (no LLM, RAG, GPT, ML, API, NLP); "
    "name each by business function.\n"
    "- Infer realistic metrics from industry benchmarks if none are provided.\n"
    "- Tone: straight, professional, executive.\n"
    "RESULTS RULES:\n"
    "- Result 1: most impactful metric (percentage, number, or time contrast), one punchy line.\n"
    "- Result 2: operational or financial outcome, one line.\n"
    "- Result 3: qualitative shift in decision quality, confidence, or leverage, one line.\n"
    "- No result exceeds 15 words. No filler words.\n"
    "- Only use a time contrast when the gap is genuinely significant and specific; "
    "never a throwaway 'hours to minutes'.\n"
)


def _clean(s):
    return str("" if s is None else s).replace("—", "-").replace("–", "-").strip()


def _normalize_case_study(data, industry=""):
    caps = [_clean(c) for c in (data.get("capabilities") or []) if _clean(c)][:6]
    while len(caps) < 6:
        caps.append("Capability: to be defined.")
    res = [_clean(r) for r in (data.get("results") or []) if _clean(r)][:3]
    while len(res) < 3:
        res.append("Result to be defined.")
    rev = data.get("review") or {}
    subhead = _clean(data.get("subhead")) or (
        "Client: Leading %s | Domain: %s | Function: " % ((industry or "Enterprise").title(), industry))
    return {
        "template": "case_study_full",
        "title": _clean(data.get("title")) or "Proposed Case Study",
        "subhead": subhead,
        "challenge": _clean(data.get("challenge")),
        "solution": _clean(data.get("solution")),
        "capabilities": caps,
        "results": res,
        "review": {"quality": _clean(rev.get("quality")) or "Needs Revision",
                   "weakest": _clean(rev.get("weakest")), "fix": _clean(rev.get("fix"))},
    }


def draft_case_study(brief, context=None):
    """Generate ONE full case study from a salesperson's free-text brief, in the
    strict format + self-review. Returns the structured fields (+ 'review'). The
    human supplies the facts in the brief; the model just writes them up."""
    context = context or {}
    industry = context.get("industry", "")
    recipient = context.get("recipient", "")
    function = context.get("function", "")
    notes = context.get("notes", "")
    prompt = (
        "Write ONE proof-point case study for a J2W sales meeting. It must feel "
        "SPECIFIC and OPERATIONAL to THIS account — never generic, never padded.\n\n"
        "ACCOUNT CONTEXT — tailor the case to this exact situation:\n"
        + (f"- Industry / domain: {industry}\n" if industry else "")
        + (f"- Stakeholder we are meeting: {recipient}\n" if recipient else "")
        + (f"- Their function / remit: {function}\n" if function else "")
        + (f"- Meeting notes / research:\n\"\"\"\n{notes[:1800]}\n\"\"\"\n" if notes else "")
        + "\nThe capability / use case to prove:\n\"\"\"\n" + (brief or "")[:1500] + "\n\"\"\"\n\n"
        "Write it as a REAL J2W engagement with an ANONYMISED client in the SAME "
        "domain as this account, whose situation mirrors what this stakeholder "
        "personally owns, solving exactly the capability above. Ground every claim "
        "in that scenario — no boilerplate.\n\n"
        "Follow these rules exactly:\n" + CASE_STUDY_RULES
        + "\nEach field:\n"
        "- title: a specific, concrete case study title (name the capability, not a slogan).\n"
        "- subhead: 'Client: <generic descriptor, e.g. Leading Manufacturing Enterprise> | "
        "Domain: <this account's domain> | Function: <the stakeholder's business function>'.\n"
        "- challenge: 3-4 sentences, plain and operational; who the client is (NEVER a real name) "
        "and what was breaking, specific to this domain. No solution language. Max 100 words.\n"
        "- solution: 3-4 sentences; what J2W deployed, how it works operationally, what the client "
        "can now do. No bullets. No hype. Max 100 words.\n"
        "- capabilities: EXACTLY 6, each 'Capability Name: one line max 18 words' (name = business function).\n"
        "- results: EXACTLY 3, following the RESULTS RULES.\n"
        "Then SELF-REVIEW what you wrote (quality verdict, weakest part, fix).\n"
        'Return ONLY JSON: {"title":"...","subhead":"...","challenge":"...","solution":"...",'
        '"capabilities":["six strings"],"results":["three strings"],'
        '"review":{"quality":"Strong or Needs Revision","weakest":"one sentence or None",'
        '"fix":"one sentence or None"}}'
    )
    try:
        from secrets_loader import load_env
        load_env()
        from openai import OpenAI
        resp = OpenAI().chat.completions.create(
            model=MODEL, temperature=0.5,
            response_format={"type": "json_object"},
            messages=[
                {"role": "system", "content": "You are a B2B case study writer and sales reviewer "
                 "for an enterprise AI and technology services company. Reply with one JSON object only."},
                {"role": "user", "content": prompt},
            ],
        )
        data = json.loads(resp.choices[0].message.content)
    except Exception:
        data = {}
    return _normalize_case_study(data, industry)


def fill_case_study(slide, content):
    """Fill a case_study_full template slide from structured content."""
    singles = {"{{TITLE}}": content.get("title", ""), "{{SUBHEAD}}": content.get("subhead", ""),
               "{{CHALLENGE}}": content.get("challenge", ""), "{{SOLUTION}}": content.get("solution", "")}
    for sh in slide.shapes:
        if not sh.has_text_frame:
            continue
        txt = sh.text_frame.text
        hit = next((m for m in singles if m in txt), None)
        if hit:
            editor.set_text(sh, singles[hit])
        elif "{{CAPABILITIES}}" in txt:
            _set_bullets(sh, content.get("capabilities", []))
        elif "{{RESULTS}}" in txt:
            _set_bullets(sh, content.get("results", []))


# --------------------------------------------------------------------------- #
# Build the slide into a deck
# --------------------------------------------------------------------------- #
def _blank_layout(prs):
    for layout in prs.slide_layouts:
        if (layout.name or "").lower().strip() == "blank":
            return layout
    return prs.slide_layouts[-1]


def _copy_slide(dest_prs, src_slide):
    """Copy a template slide into dest_prs as a new slide, including image parts
    so branded picture shapes (logos, backgrounds, icons) render correctly."""
    new = dest_prs.slides.add_slide(_blank_layout(dest_prs))
    for shp in list(new.shapes):
        shp._element.getparent().remove(shp._element)

    # Copy image relationships: for each image part in the source slide,
    # relate it to the new slide so r:embed rIds resolve in the destination.
    src_part  = src_slide._part
    dest_part = new._part
    rId_map   = {}
    _R = 'http://schemas.openxmlformats.org/officeDocument/2006/relationships'
    for rId, rel in list(src_part.rels.items()):
        if rel.is_external:
            continue
        if '/image' in (rel.reltype or ''):
            new_rId = dest_part.relate_to(rel._target, rel.reltype)
            if new_rId != rId:
                rId_map[rId] = new_rId

    _REMAP = {f'{{{_R}}}embed', f'{{{_R}}}link'}

    def _remap(elem):
        for attr, val in list(elem.attrib.items()):
            if attr in _REMAP and val in rId_map:
                elem.attrib[attr] = rId_map[val]
        for child in elem:
            _remap(child)

    for shp in src_slide.shapes:
        elem = copy.deepcopy(shp._element)
        if rId_map:
            _remap(elem)
        new.shapes._spTree.append(elem)

    return new


def _set_bullets(shape, bullets):
    tf = shape.text_frame
    tf.clear()
    for i, b in enumerate(bullets):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = "• " + b


def _fill(slide, content):
    for sh in slide.shapes:
        if not sh.has_text_frame:
            continue
        txt = sh.text_frame.text
        if "{{TITLE}}" in txt:
            editor.set_text(sh, content.get("title", ""))
        elif "{{KEYWORDS}}" in txt:
            editor.set_text(sh, content.get("keywords", ""))
        elif "{{BULLETS}}" in txt:
            _set_bullets(sh, content.get("bullets", []))


def _add_verify_banner(slide, slide_width):
    """Stamp a loud red bar across the top of a slide: this AI-written slide has
    not been checked by an expert, so it must not reach a client as-is. The banner
    travels with the slide in the .pptx (visible in preview and in PowerPoint)."""
    bar = slide.shapes.add_textbox(0, 0, slide_width, Inches(0.5))
    bar.fill.solid()
    bar.fill.fore_color.rgb = RGBColor(0xC0, 0x39, 0x2B)      # red
    bar.line.fill.background()
    tf = bar.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = "⚠ AI-GENERATED - NEEDS EXPERT VERIFICATION - NOT CLIENT-READY"
    run.font.size = Pt(13)
    run.font.bold = True
    run.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)


def append_generated(deck_path, items):
    """items = [{template, title, keywords, bullets, verified}] -> add each as a
    new slide at the end of the deck. Unverified (verified != True) slides get a
    loud 'needs expert verification' banner stamped on them."""
    templates = list_templates()
    dest = Presentation(deck_path)
    for it in items:
        src = templates.get(it.get("template", "case_study"))
        if src is None and templates:
            src = next(iter(templates.values()))
        if src is None:
            continue
        new_slide = _copy_slide(dest, src)
        _fill(new_slide, it)
        if not it.get("verified"):
            _add_verify_banner(new_slide, dest.slide_width)
    dest.save(deck_path)
    return len(items)


if __name__ == "__main__":
    print("Created template:", create_temp_template())
    print("Templates found:", list(list_templates().keys()))
