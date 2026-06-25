# -*- coding: utf-8 -*-
"""
staging.py  --  the self-learning loop (#8), with a safety gate.

When the engine writes a new slide to fill a gap, the CONTENT is saved here as
'pending'. Two things then happen:
  - REUSE: the next time the same slot (work type + industry) is needed, we
    reuse the stored content instead of paying for a fresh AI draft. The engine
    gets better/cheaper the more it's used.
  - APPROVE: a human reviews pending slides on /staging. Approving one PROMOTES
    it into the real master deck + library + registry (so it becomes a normal,
    matchable slide). Until then it never pollutes the trusted library.

We store CONTENT (title / keywords / bullets), not rendered slides — the slide
is built from the template on demand. Keeps it simple and swap-safe.
"""

import json
import os
import re
from copy import deepcopy
from datetime import datetime

from pptx import Presentation

import slide_generator
import tagger
from build_library import read_id

STAGE_DIR = "staging"
STAGE_JSON = os.path.join(STAGE_DIR, "staging.json")
MASTER = "WORKING_COPY_Master_Deck.pptx"
REGISTRY = "J2W_CaseStudy_Portfolio_Metadata.xlsx"

# The loud, human-facing status for an AI-written slide that no expert has checked.
# A 'pending' record == this status; approving clears it (becomes client-ready).
VERIFY_STATUS = "NEEDS EXPERT VERIFICATION - not client-ready"


def _load():
    try:
        return json.load(open(STAGE_JSON, encoding="utf-8"))
    except Exception:
        return []


def _save(items):
    os.makedirs(STAGE_DIR, exist_ok=True)
    json.dump(items, open(STAGE_JSON, "w", encoding="utf-8"), ensure_ascii=False, indent=2)


def find(work_type, industry):
    """Reuse a prior generated slide for this slot (approved first, then pending)."""
    cands = [it for it in _load()
             if it["work_type"] == work_type
             and (it.get("industry") or "") == (industry or "")
             and it["status"] in ("pending", "approved")]
    cands.sort(key=lambda it: 0 if it["status"] == "approved" else 1)
    return cands[0] if cands else None


def add(content, work_type, industry, client=""):
    items = _load()
    rec = {
        "id": "G%03d" % (len(items) + 1),
        "status": "pending",                                  # pending == needs expert verification
        "work_type": work_type,
        "industry": industry or "",
        "title": content.get("title", ""),
        "keywords": content.get("keywords", ""),
        "bullets": content.get("bullets", []),
        "template": content.get("template", "case_study"),
        "for_client": client,
        "created_at": datetime.now().strftime("%Y-%m-%d %H:%M"),
    }
    items.append(rec)
    _save(items)
    return rec


def get(stg_id):
    return next((it for it in _load() if it["id"] == stg_id), None)


def update_content(stg_id, title=None, keywords=None, bullets=None):
    """Save a reviewer's edits onto a pending slide before it's accepted/promoted."""
    items = _load()
    for it in items:
        if it["id"] == stg_id:
            if title is not None:
                it["title"] = title
            if keywords is not None:
                it["keywords"] = keywords
            if bullets is not None:
                it["bullets"] = bullets
    _save(items)


def pending():
    return [it for it in _load() if it["status"] == "pending"]


def all_items():
    return _load()


def _set_status(stg_id, status):
    items = _load()
    for it in items:
        if it["id"] == stg_id:
            it["status"] = status
    _save(items)


def discard(stg_id):
    _set_status(stg_id, "discarded")


def _stamp_id(prs, slide, id_text):
    """Write the J2W_ID into a slide's notes. A freshly-copied slide has an empty
    notes slide (no placeholder), so borrow a notes placeholder from an existing
    slide first, then write."""
    notes = slide.notes_slide
    if notes.notes_text_frame is not None:
        notes.notes_text_frame.text = id_text
        return
    for s in prs.slides:
        if s is slide or not s.has_notes_slide:
            continue
        ph = s.notes_slide.notes_placeholder
        if ph is not None:
            notes.shapes._spTree.append(deepcopy(ph._element))
            notes.notes_text_frame.text = id_text
            return


def promote(stg_id):
    """Approve: build the slide into the master deck and register it as a real,
    matchable slide (source = generated)."""
    rec = get(stg_id)
    if not rec or rec["status"] == "approved":
        return None

    prs = Presentation(MASTER)
    nums = [int(s[2:]) for s in (read_id(sl) for sl in prs.slides) if s and s[2:].isdigit()]
    new_id = "CS%02d" % (max(nums) + 1 if nums else 1)

    templates = slide_generator.list_templates()
    src = templates.get(rec["template"]) or (next(iter(templates.values())) if templates else None)
    if src is None:
        return None
    new_slide = slide_generator._copy_slide(prs, src)
    slide_generator._fill(new_slide, rec)
    _stamp_id(prs, new_slide, "J2W_ID: " + new_id)
    prs.save(MASTER)

    _add_to_library_and_registry(new_id, rec)
    _set_status(stg_id, "approved")
    rec["promoted_id"] = new_id
    items = _load()
    for it in items:
        if it["id"] == stg_id:
            it["promoted_id"] = new_id
    _save(items)
    return new_id


def _add_to_library_and_registry(new_id, rec):
    import build_library
    import openpyxl

    # rebuild library + tags from the (now larger) master
    recs = build_library.build(MASTER)
    json.dump(recs, open("library.json", "w", encoding="utf-8"), ensure_ascii=False, indent=2)
    recs = tagger.tag_library(recs)
    json.dump(recs, open("tagged_library.json", "w", encoding="utf-8"), ensure_ascii=False, indent=2)

    # infer a function tag from the content for the registry row
    pseudo = {"title": rec["title"], "subtitle": rec["keywords"],
              "full_text": rec["title"] + " " + rec["keywords"] + " " + " ".join(rec["bullets"])}
    tagger.tag_record(pseudo)
    function = pseudo["tags"]["function"]["value"]

    wb = openpyxl.load_workbook(REGISTRY)
    ws = wb["Slide Registry"]
    col = {c.value: i for i, c in enumerate(ws[1])}
    existing = {ws.cell(r, col["slide_id"] + 1).value for r in range(2, ws.max_row + 1)}
    if new_id not in existing:
        row = [None] * len(col)
        row[col["slide_id"]] = new_id
        row[col["section"]] = rec["work_type"]
        row[col["kind"]] = "CASE_STUDY"
        row[col["include_rule"]] = "IF industry/function/keywords match context or transcript"
        row[col["work_types"]] = rec["work_type"]
        row[col["primary_industry"]] = rec["industry"] or None
        row[col["primary_function"]] = function
        row[col["keywords"]] = rec["keywords"]
        row[col["confidence"]] = "AUTO"
        row[col["title"]] = rec["title"]
        ws.append(row)
        wb.save(REGISTRY)
