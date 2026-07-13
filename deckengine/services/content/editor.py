# -*- coding: utf-8 -*-
"""
editor.py  --  Box / Step 05 helper: read a slide's editable text, and write
edits back into a .pptx while preserving formatting.

We let the user edit the two text fields that matter most per slide — the
TITLE and the SUBTITLE/headline — because those are single-line and edit
cleanly. Body text is shown for context but not edited here (editing richly
formatted multi-run body text safely is a later step).
"""

from pptx import Presentation

from deckengine.services.content.build_library import read_id


def _text_shapes(slide):
    """(index, shape) for every shape on the slide that has non-empty text."""
    out = []
    for i, sh in enumerate(slide.shapes):
        if sh.has_text_frame and sh.text_frame.text.strip():
            out.append((i, sh))
    return out


def editable_fields(slide):
    """Return [(shape_index, label, current_text)] for the Title + Subtitle only.
    Kept for callers that want just the two headline fields."""
    fields = all_editable_fields(slide)
    return [(i, lbl, txt) for i, lbl, txt, _multi in fields[:2]]


def all_editable_fields(slide):
    """Every text-bearing shape on the slide, as (shape_index, label, current_text,
    is_multiline), in reading order. Lets the review page edit a library slide's WHOLE
    text -- title, subtitle, and every bullet block -- not just its first two lines
    (owner's spec, 2026-07-13). is_multiline is True when the shape holds more than one
    paragraph, so the editor shows a textarea and the writer preserves the bullets."""
    shapes = _text_shapes(slide)

    title_idx = None
    try:
        title_shape = slide.shapes.title
    except (ValueError, AttributeError):
        title_shape = None
    if title_shape is not None and title_shape.text.strip():
        for i, sh in enumerate(slide.shapes):
            if sh is title_shape:
                title_idx = i
                break
    if title_idx is None and shapes:
        title_idx = shapes[0][0]

    # title first, then the rest in reading order
    ordered = ([(i, sh) for i, sh in shapes if i == title_idx]
               + [(i, sh) for i, sh in shapes if i != title_idx])
    out = []
    body_n = 0
    for i, sh in ordered:
        tf = sh.text_frame
        text = tf.text.strip()
        multiline = len([p for p in tf.paragraphs if p.text.strip()]) > 1
        if i == title_idx:
            label = "Title"
        elif body_n == 0:
            label = "Subtitle"; body_n += 1
        else:
            body_n += 1
            label = "Text %d" % (body_n - 1)
        out.append((i, label, text, multiline))
    return out


def set_text(shape, text):
    """Replace a shape's text with `text`, keeping the FIRST run's formatting
    (font, size, colour). Extra runs/paragraphs are removed."""
    tf = shape.text_frame
    paras = tf.paragraphs
    first = paras[0]
    if first.runs:
        first.runs[0].text = text
        for r in first.runs[1:]:
            r._r.getparent().remove(r._r)
    else:
        first.text = text
    for p in list(paras[1:]):
        p._p.getparent().remove(p._p)


def full_text(slide):
    """All of a slide's text, for read-only context in the review screen."""
    return [sh.text_frame.text.strip() for _, sh in _text_shapes(slide)]


def set_text_multiline(shape, text):
    """Replace a shape's text, keeping its BULLETS: each line becomes its own paragraph,
    reusing the formatting of the paragraph that was there (or the first). Used when the
    review page edits a multi-paragraph library shape -- set_text would collapse the
    bullets to one line."""
    tf = shape.text_frame
    lines = (text or "").split("\n")
    paras = list(tf.paragraphs)
    template_p = paras[0]._p                       # clone this paragraph's formatting
    import copy
    # write the first line into the first paragraph (via set_text's run-safe path)
    _set_paragraph_text(paras[0], lines[0] if lines else "")
    # remove any extra existing paragraphs
    for p in paras[1:]:
        p._p.getparent().remove(p._p)
    # append one paragraph per remaining line, cloned from the first for consistent bullets
    body = tf._txBody
    for line in lines[1:]:
        new_p = copy.deepcopy(template_p)
        body.append(new_p)
        from pptx.text.text import _Paragraph
        _set_paragraph_text(_Paragraph(new_p, tf), line)


def _set_paragraph_text(paragraph, text):
    """Set one paragraph's text, keeping its first run's formatting."""
    if paragraph.runs:
        paragraph.runs[0].text = text
        for r in paragraph.runs[1:]:
            r._r.getparent().remove(r._r)
    else:
        paragraph.text = text


def apply_edits(path, edits):
    """edits = {slide_id: {shape_index(int): new_text}} -> write into the deck.
    A shape whose text has newlines keeps them as separate paragraphs (bullets)."""
    prs = Presentation(path)
    for slide in prs.slides:
        sid = read_id(slide)
        if sid in edits:
            shapes = list(slide.shapes)
            for idx, text in edits[sid].items():
                if 0 <= idx < len(shapes) and shapes[idx].has_text_frame:
                    if "\n" in (text or ""):
                        set_text_multiline(shapes[idx], text)
                    else:
                        set_text(shapes[idx], text)
    prs.save(path)


def replace_tokens(path, tokens):
    """Replace literal tokens (e.g. {'[CLIENT]': 'Acme Bank'}) in every run.
    Run-level replace keeps formatting intact."""
    prs = Presentation(path)
    for slide in prs.slides:
        for sh in slide.shapes:
            if not sh.has_text_frame:
                continue
            for p in sh.text_frame.paragraphs:
                for r in p.runs:
                    for old, new in tokens.items():
                        if old in r.text:
                            r.text = r.text.replace(old, new)
    prs.save(path)
