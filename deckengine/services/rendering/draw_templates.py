# -*- coding: utf-8 -*-
"""
draw_templates.py  --  the ten slide shapes from the owner's style guide.

Each one is DRAWN programmatically onto a copied header-only template slide (title +
subtitle + the two-tone top bar live in skills_templates.pptx; everything below the
header is drawn here, because none of these shapes has a fixed slot count a static
template could carry).

Kept out of skills.py, which was already 1,100 lines and doing three jobs. skills.py
now dispatches into DRAWERS below.

## The design language (Slide_Template_Style_Guide.docx, 2026-07-10)

Composition, fixed for every shape:
    title block (red accent bar + bold 24pt title + teal subtitle)
      -> 3-6 soft cards in a grid, row, or split
      -> optional colored callout strip at the bottom stating the one takeaway

Two decisions the owner made where the sources disagreed:

  * The two-tone top bar runs RED LEFT, TEAL RIGHT. The guide's prose says the
    opposite, and four of the ten screenshots show it flipped -- but every existing
    J2W slide, the master deck's own title slide, and six of the ten screenshots run
    red-left. Consistency with the 100+ slides already in the library wins.

  * Icon chips are drawn EMPTY for BUILDER-made slides -- the app ships no icon assets
    of its own (the web UI pulls Tabler from a CDN; skills_templates.pptx contains no
    images), and a unicode glyph renders inconsistently through LibreOffice, so the chip
    just carries the colour rhythm. But Recreate-with-AI now carries the SOURCE slide's
    own icons across into the chips (owner's spec, 2026-07-13); see `_chip`/`_pop_icon`
    and recreate._slide_icons.

Accents alternate TEAL, RED per card -- the guide is explicit ("one accent color per
card/number, alternating rather than random"), and amber, which appears once in the
screenshots, is not in the palette.
"""

from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

# ── palette, sampled from the style guide ────────────────────────────────────
RED = RGBColor(0xD6, 0x28, 0x39)     # primary red
TEAL = RGBColor(0x2A, 0x9D, 0x8F)    # primary teal
NAVY = RGBColor(0x1E, 0x29, 0x3B)    # governance banners, dark callouts
LABEL = RGBColor(0x26, 0x31, 0x42)   # card titles, stat labels
BODY = RGBColor(0x3A, 0x3F, 0x47)    # description copy
CARD = RGBColor(0xF8, 0xFA, 0xFC)    # the default soft-card fill
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
BORDER = RGBColor(0xD9, 0xDC, 0xE1)  # hairline card borders
GHOST = RGBColor(0xE4, 0xE7, 0xEB)   # the oversized background numbers
MUTED = RGBColor(0x8A, 0x90, 0x99)   # small uppercase labels

PALE_TEAL = RGBColor(0xE8, 0xF3, 0xF1)
PALE_RED = RGBColor(0xFD, 0xEC, 0xEA)

# ── canvas ───────────────────────────────────────────────────────────────────
SW, SH = 13.33, 7.50
MARGIN = 0.30
TOP = 1.32              # first line below the header block
BOTTOM = SH - 0.30      # above the slim teal footer strip
GAP = 0.20
PAD = 0.16

FONT_HEAD = "Oswald"
FONT_BODY = "Raleway"

SZ_CARD_TITLE = 13
SZ_BODY = 10.5
SZ_SMALL = 9
SZ_STAT = 28


def _accent(i):
    """Alternating card accent. Never random -- the guide is explicit."""
    return TEAL if i % 2 == 0 else RED


def _pale(colour):
    return PALE_TEAL if colour == TEAL else PALE_RED


# ── primitives ───────────────────────────────────────────────────────────────
def _bar(slide, l, t, w, h, fill):
    s = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(l), Inches(t), Inches(w), Inches(h))
    s.fill.solid(); s.fill.fore_color.rgb = fill
    s.line.fill.background(); s.shadow.inherit = False
    return s


def _card(slide, l, t, w, h, fill=CARD, line=BORDER, rounded=False):
    shape = MSO_SHAPE.ROUNDED_RECTANGLE if rounded else MSO_SHAPE.RECTANGLE
    s = slide.shapes.add_shape(shape, Inches(l), Inches(t), Inches(w), Inches(h))
    s.fill.solid(); s.fill.fore_color.rgb = fill
    if line is None:
        s.line.fill.background()
    else:
        s.line.color.rgb = line; s.line.width = Pt(0.75)
    s.shadow.inherit = False
    if rounded:                       # a gentle radius, not a pill
        try:
            s.adjustments[0] = 0.12
        except (IndexError, AttributeError):
            pass
    return s


def _is_light_icon(icon):
    """True if the icon's own pixels are mostly light/white -- such an icon is designed
    for a DARK background and vanishes on a pale chip (owner-reported, 2026-07-13: the
    white stat icons were invisible). Fails safe to False."""
    try:
        from PIL import Image
        import io as _io
        im = Image.open(_io.BytesIO(icon)).convert("RGBA")
        im.thumbnail((32, 32))
        px = [p for p in im.getdata() if p[3] > 40]      # opaque-ish pixels only
        if not px:
            return False
        lum = sum(0.299 * r + 0.587 * g + 0.114 * b for r, g, b, _a in px) / len(px)
        return lum > 200
    except Exception:
        return False


def _chip(slide, l, t, size, colour, icon=None):
    """The icon chip: a rounded square. `icon` is optional PNG/JPEG bytes -- when Recreate
    carries the source slide's own icon across (owner's spec, 2026-07-13) it's dropped in,
    centred and inset. A WHITE/light icon gets a DARK chip (the accent colour at full
    saturation) so it's visible; a coloured icon keeps the pale tint. An empty chip stays
    pale, as the builder-drawn slides leave it (the app ships no icon set of its own)."""
    fill = _pale(colour)
    if icon and _is_light_icon(icon):
        fill = colour                    # white icon on the full accent -- good contrast
    s = _card(slide, l, t, size, size, fill=fill, line=None, rounded=True)
    if icon:
        try:
            import io as _io
            pad = size * 0.20
            slide.shapes.add_picture(_io.BytesIO(icon), Inches(l + pad), Inches(t + pad),
                                     Inches(size - 2 * pad), Inches(size - 2 * pad))
        except Exception:
            pass                    # a bad blob just leaves the chip empty
    return s


def draw_score_ring(slide, cx, cy, r, value, outof, colour=TEAL):
    """A filled circular score gauge -- a big value on a coloured disc with '/outof'
    under it. Recreate draws this where the source showed a score as a filled circle and
    Recreate would otherwise have flattened it to a bullet (owner's spec, 2026-07-13)."""
    d = 2 * r
    disc = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(cx - r), Inches(cy - r),
                                  Inches(d), Inches(d))
    disc.fill.solid(); disc.fill.fore_color.rgb = colour
    disc.line.fill.background(); disc.shadow.inherit = False
    _text(slide, cx - r, cy - r * 0.58, d, r * 0.82, str(value), int(23 * r / 0.62),
          WHITE, bold=True, align=PP_ALIGN.CENTER, font=FONT_HEAD, anchor="ctr")
    _text(slide, cx - r, cy + r * 0.30, d, r * 0.5, "/ %s" % outof, int(11 * r / 0.62),
          WHITE, align=PP_ALIGN.CENTER, anchor="ctr")


def _pop_icon(data):
    """Next source icon for a chip, or None. Recreate seeds data['_icons'] with the
    source slide's icons in reading order; drawers pull them as they draw chips, so the
    Nth chip gets the Nth source icon. Anything not from Recreate has no icons -> None."""
    icons = (data or {}).get("_icons")
    if icons:
        return icons.pop(0)
    return None


def _text(slide, l, t, w, h, text, size, colour, bold=False, align=PP_ALIGN.LEFT,
          font=FONT_BODY, anchor=None, italic=False):
    tb = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = Pt(0)
    tf.margin_top = tf.margin_bottom = Pt(1)
    if anchor:
        tf._txBody.bodyPr.set("anchor", anchor)
    p = tf.paragraphs[0]
    p.alignment = align
    r = p.add_run(); r.text = text or ""
    r.font.name = font; r.font.size = Pt(size)
    r.font.color.rgb = colour; r.font.bold = bold; r.font.italic = italic
    return tb


def _rich(slide, l, t, w, h, parts, size, align=PP_ALIGN.LEFT, font=FONT_BODY):
    """One paragraph, several runs: [(text, bold, colour), ...]. For the bold lead-in
    the guardrail and comparison layouts use ("Approved model registry: with version...")."""
    tb = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = Pt(0)
    tf.margin_top = tf.margin_bottom = Pt(1)
    p = tf.paragraphs[0]
    p.alignment = align
    for text, bold, colour in parts:
        r = p.add_run(); r.text = text
        r.font.name = font; r.font.size = Pt(size)
        r.font.color.rgb = colour; r.font.bold = bold
    return tb


def _tick(slide, l, t, colour):
    """The small square that stands in for the tick / chevron marker."""
    _bar(slide, l, t + 0.05, 0.055, 0.055, colour)


def _label(slide, l, t, w, text, colour=MUTED, size=SZ_SMALL):
    """A small uppercase label -- OPPORTUNITY, METRIC, ARCHITECTURE."""
    return _text(slide, l, t, w, 0.20, (text or "").upper(), size, colour, bold=True)


def _grid(n, cols, top=TOP, bottom=BOTTOM, gap=GAP, left=MARGIN, right=MARGIN,
          max_row_h=None):
    """Positions for n cards in `cols` columns.

    By default the rows FILL the body box. Pass `max_row_h` to cap each row's height at
    the tallest its content needs and CENTRE the grid vertically -- otherwise a slide with
    one or two short cards stretches them to full slide height, which is the white-space
    the owner reported (2026-07-13). A grid whose natural height already exceeds the cap
    fills as before."""
    rows = -(-n // cols)
    w = (SW - left - right - (cols - 1) * gap) / cols
    natural = (bottom - top - (rows - 1) * gap) / rows
    h = natural
    grid_top = top
    if max_row_h is not None and natural > max_row_h:
        h = max_row_h
        used = rows * h + (rows - 1) * gap
        # sit the grid a THIRD of the way down the leftover space, not centred: a small
        # top margin plus a slightly larger bottom one reads as deliberate padding, where
        # an equal split leaves an obvious hole under the header.
        grid_top = top + max(0.0, (bottom - top - used)) / 3
    for i in range(n):
        r, c = divmod(i, cols)
        yield i, left + c * (w + gap), grid_top + r * (h + gap), w, h


def _footer_bar(slide, text, fill=NAVY, height=0.62):
    """The 'so what' strip: a full-width coloured band at the bottom of the slide."""
    t = BOTTOM - height
    _bar(slide, MARGIN, t, SW - 2 * MARGIN, height, fill)
    _bar(slide, MARGIN, t, 0.06, height, RED)
    _text(slide, MARGIN + 0.30, t, SW - 2 * MARGIN - 0.60, height, text,
          SZ_CARD_TITLE, WHITE, bold=True, align=PP_ALIGN.CENTER,
          font=FONT_HEAD, anchor="ctr")
    return t


# ── header markers ───────────────────────────────────────────────────────────
def _head(data):
    return {"TITLE": (data.get("title") or "").upper(),
            "SUBHEAD": data.get("subhead") or ""}


def _head_intro(data):
    m = _head(data)
    m["INTRO"] = data.get("intro") or ""
    return m


# ═════════════════════════════════════════════════════════════════════════════
# 1. Pain-point list -- vertical rows: accent bar, chip, bold label, divider, body
# ═════════════════════════════════════════════════════════════════════════════
def draw_pain_point_list(slide, data):
    rows = (data.get("rows") or [])[:6]
    if not rows:
        return
    n = len(rows)
    # cap the row height: four rows stretched over the whole body look like empty boxes
    h = min(1.12, (BOTTOM - TOP - (n - 1) * GAP) / n)
    w = SW - 2 * MARGIN
    top = TOP + max(0.0, (BOTTOM - TOP - (n * h + (n - 1) * GAP)) / 2)   # centre the stack
    for i, row in enumerate(rows):
        t = top + i * (h + GAP)
        colour = _accent(i)
        _card(slide, MARGIN, t, w, h, CARD, BORDER)
        _bar(slide, MARGIN, t, 0.055, h, colour)                 # coloured left edge
        _chip(slide, MARGIN + 0.24, t + (h - 0.34) / 2, 0.34, colour, _pop_icon(data))
        _text(slide, MARGIN + 0.70, t, 2.55, h, row.get("label", ""),
              SZ_CARD_TITLE, LABEL, bold=True, font=FONT_HEAD, anchor="ctr")
        _bar(slide, MARGIN + 3.35, t + 0.18, 0.008, h - 0.36, BORDER)   # hairline divider
        _text(slide, MARGIN + 3.60, t, w - 3.90, h, row.get("body", ""),
              SZ_BODY, BODY, anchor="ctr")


# ═════════════════════════════════════════════════════════════════════════════
# 2. Platform overview -- stat row + capability chip row + dark footer band
# ═════════════════════════════════════════════════════════════════════════════
def draw_platform_overview(slide, data):
    stats = (data.get("stats") or [])[:4]
    caps = (data.get("capabilities") or [])[:6]
    footer_title = data.get("footer_title") or ""
    footer_items = data.get("footer_items") or []

    # Distribute the body vertically so it fills the slide rather than clustering at the
    # top with a gap above the footer (the white space the owner reported, 2026-07-13).
    footer_h = 0.72 if (footer_title or footer_items) else 0.0
    body_bottom = BOTTOM - (footer_h + (0.24 if footer_h else 0))
    # the stat row and the capability row share the space above the footer, gap between
    stat_h = 1.70 if stats else 0.0
    cap_label_h = 0.34 if caps else 0.0
    if stats and caps:
        gap_mid = 0.34
        cap_h = body_bottom - TOP - stat_h - gap_mid - cap_label_h
        cap_h = max(0.9, min(cap_h, 1.9))              # a chip + a label, not a banner
    else:
        cap_h = body_bottom - TOP - cap_label_h

    if stats:
        w = (SW - 2 * MARGIN - (len(stats) - 1) * GAP) / len(stats)
        # a source stat card leads with an icon, then the number, then the label; mirror
        # that. Drawing a chip here also consumes the source's stat icons in order, so the
        # capability chips below get the CAPABILITY icons and not the stat ones.
        has_icons = bool((data or {}).get("_icons"))
        icon_h = 0.42 if has_icons else 0.0
        block_h = icon_h + 0.90
        vy = TOP + (stat_h - block_h) / 2
        for i, s in enumerate(stats):
            l = MARGIN + i * (w + GAP)
            colour = _accent(i)
            _card(slide, l, TOP, w, stat_h, WHITE, BORDER)
            _bar(slide, l, TOP, w, 0.05, colour)                 # coloured top rule
            if has_icons:
                _chip(slide, l + (w - icon_h) / 2, vy, icon_h, colour, _pop_icon(data))
            _text(slide, l, vy + icon_h, w, 0.62, s.get("value", ""), SZ_STAT, colour,
                  bold=True, align=PP_ALIGN.CENTER, font=FONT_HEAD)
            _text(slide, l, vy + icon_h + 0.64, w, 0.28, (s.get("label", "") or "").upper(),
                  SZ_SMALL, LABEL, bold=True, align=PP_ALIGN.CENTER)

    cap_top = TOP + stat_h + (0.34 if stats else 0)
    if caps:
        _label(slide, MARGIN, cap_top, 4.0, "Core capabilities", LABEL, 11)
        cw = (SW - 2 * MARGIN - (len(caps) - 1) * GAP) / len(caps)
        ct = cap_top + cap_label_h
        for i, c in enumerate(caps):
            l = MARGIN + i * (cw + GAP)
            _card(slide, l, ct, cw, cap_h, CARD, BORDER)
            _chip(slide, l + (cw - 0.36) / 2, ct + (cap_h - 0.60) / 2 - 0.02, 0.36, _accent(i), _pop_icon(data))
            _text(slide, l + 0.06, ct + (cap_h - 0.60) / 2 + 0.40, cw - 0.12, 0.30, c,
                  SZ_SMALL, LABEL, bold=True, align=PP_ALIGN.CENTER)

    if footer_title or footer_items:
        h = footer_h
        t = BOTTOM - h
        _bar(slide, MARGIN, t, SW - 2 * MARGIN, h, NAVY)
        _bar(slide, MARGIN, t, 0.06, h, RED)
        _text(slide, MARGIN + 0.28, t, 3.0, h, (footer_title or "").upper(),
              12, TEAL, bold=True, font=FONT_HEAD, anchor="ctr")
        if footer_items:
            _text(slide, MARGIN + 3.30, t, SW - 2 * MARGIN - 3.60, h,
                  "     |     ".join(footer_items), SZ_BODY, WHITE, anchor="ctr")


# ═════════════════════════════════════════════════════════════════════════════
# 3. Before/after split -- two workflow lanes, then the questions they raise
# ═════════════════════════════════════════════════════════════════════════════
def _stage_lane(slide, l, t, w, h, title, stages, colour):
    _card(slide, l, t, w, h, WHITE, BORDER)
    _chip(slide, l + PAD, t + PAD, 0.28, colour)   # a lane header, not a per-card icon
    _text(slide, l + PAD + 0.40, t + PAD - 0.02, w - 0.80, 0.32, (title or "").upper(),
          11, LABEL, bold=True, font=FONT_HEAD)
    stages = (stages or [])[:4]
    if not stages:
        return
    sw = (w - 2 * PAD - (len(stages) - 1) * 0.10) / len(stages)
    st = t + 0.66
    sh = h - 0.66 - PAD
    for i, s in enumerate(stages):
        sl = l + PAD + i * (sw + 0.10)
        tag = (s.get("tag") or "").strip()
        is_agent = "agent" in tag.lower() or "ai" in tag.lower()
        tone = RED if is_agent else (TEAL if colour == TEAL else None)
        fill = _pale(tone) if tone else CARD
        _card(slide, sl, st, sw, sh, fill, BORDER)
        if tag:
            _bar(slide, sl + 0.06, st + 0.08, sw - 0.12, 0.22, tone or BORDER)
            _text(slide, sl + 0.06, st + 0.08, sw - 0.12, 0.22, tag.upper(), 7.5,
                  WHITE if tone else LABEL, bold=True, align=PP_ALIGN.CENTER, anchor="ctr")
        _text(slide, sl + 0.05, st + sh - 0.62, sw - 0.10, 0.56,
              (s.get("label", "") or "").upper(), 8.5, LABEL, bold=True,
              align=PP_ALIGN.CENTER, anchor="ctr")


def draw_before_after_split(slide, data):
    questions = (data.get("questions") or [])[:5]
    intro = (data.get("intro") or "").strip()

    # The intro is drawn HERE, not through the template's {{INTRO}} marker: that marker
    # sits at 1.00in and the teal subtitle runs to 1.14in, so the two collided. Drawing
    # it also lets it be red and bold, as the reference design has it.
    lane_top = TOP
    if intro:
        _text(slide, MARGIN, TOP - 0.04, SW - 2 * MARGIN, 0.32, intro, 11.5, RED, bold=True)
        lane_top = TOP + 0.40

    lane_h = 2.30 if questions else (BOTTOM - lane_top - 0.10)
    lane_w = (SW - 2 * MARGIN - 0.55) / 2

    _stage_lane(slide, MARGIN, lane_top, lane_w,
                lane_h, data.get("before_title") or "Before",
                data.get("before_stages"), TEAL)
    _stage_lane(slide, MARGIN + lane_w + 0.55, lane_top, lane_w, lane_h,
                data.get("after_title") or "After", data.get("after_stages"), RED)

    # the arrow between the lanes: a chip with a chevron drawn as two bars
    ax = MARGIN + lane_w + 0.155
    ay = lane_top + lane_h / 2 - 0.12
    _card(slide, ax, ay, 0.24, 0.24, PALE_RED, line=None, rounded=True)
    _bar(slide, ax + 0.09, ay + 0.09, 0.07, 0.055, RED)

    if not questions:
        return
    qt = lane_top + lane_h + 0.42
    _label(slide, MARGIN, qt - 0.32, 6.0, "The questions this raises", LABEL, 11)
    qw = (SW - 2 * MARGIN - (len(questions) - 1) * 0.14) / len(questions)
    qh = min(1.80, BOTTOM - qt)          # a question is two lines, not half a slide
    for i, q in enumerate(questions):
        l = MARGIN + i * (qw + 0.14)
        colour = _accent(i)
        _card(slide, l, qt, qw, qh, CARD, BORDER)
        _bar(slide, l, qt, qw, 0.04, colour)
        _text(slide, l + 0.12, qt + 0.14, qw - 0.24, 0.30,
              "%d. %s" % (i + 1, (q.get("title", "") or "").upper()),
              SZ_SMALL, colour, bold=True)
        _text(slide, l + 0.12, qt + 0.48, qw - 0.24, qh - 0.60,
              q.get("body", ""), SZ_SMALL, BODY)


# ═════════════════════════════════════════════════════════════════════════════
# 4. Comparison split -- capability cards left, a metric table right
# ═════════════════════════════════════════════════════════════════════════════
def draw_comparison_split(slide, data):
    features = (data.get("features") or [])[:4]
    rows = (data.get("rows") or [])[:6]
    left_w = 6.05
    right_l = MARGIN + left_w + 0.45
    right_w = SW - MARGIN - right_l

    _text(slide, MARGIN, TOP, left_w, 0.30, data.get("panel_title") or "Capability overview",
          SZ_CARD_TITLE, LABEL, bold=True, font=FONT_HEAD)
    _text(slide, MARGIN, TOP + 0.32, left_w, 0.28, data.get("panel_intro") or "",
          SZ_SMALL, BODY)
    if features:
        ft = TOP + 0.72
        fh = (BOTTOM - ft - (len(features) - 1) * 0.14) / len(features)
        for i, f in enumerate(features):
            t = ft + i * (fh + 0.14)
            colour = _accent(i)
            _card(slide, MARGIN, t, left_w, fh, WHITE, BORDER)
            _bar(slide, MARGIN, t, 0.055, fh, colour)
            _chip(slide, MARGIN + 0.22, t + 0.14, 0.24, colour, _pop_icon(data))
            _text(slide, MARGIN + 0.56, t + 0.10, left_w - 0.75, 0.28,
                  f.get("heading", ""), 12, LABEL, bold=True, font=FONT_HEAD)
            _text(slide, MARGIN + 0.22, t + 0.42, left_w - 0.44, fh - 0.52,
                  f.get("body", ""), SZ_SMALL, BODY)

    # the table
    _bar(slide, right_l, TOP, right_w, 0.52, NAVY)
    _text(slide, right_l, TOP, right_w, 0.52, data.get("table_title") or "Comparison",
          14, WHITE, bold=True, align=PP_ALIGN.CENTER, font=FONT_HEAD, anchor="ctr")
    col = [right_w * 0.40, right_w * 0.30, right_w * 0.30]
    hy = TOP + 0.52
    _bar(slide, right_l, hy, right_w, 0.34, CARD)
    _label(slide, right_l + 0.14, hy + 0.07, col[0], "Metric")
    _label(slide, right_l + col[0], hy + 0.07, col[1], data.get("col_a") or "Option A", TEAL)
    _label(slide, right_l + col[0] + col[1], hy + 0.07, col[2], data.get("col_b") or "Option B", RED)

    ry = hy + 0.34
    takeaway = data.get("takeaway") or ""
    avail = (BOTTOM - ry - (0.62 if takeaway else 0))
    rh = min(0.46, avail / max(len(rows), 1)) if rows else 0
    for i, r in enumerate(rows):
        t = ry + i * rh
        if i % 2 == 1:
            _bar(slide, right_l, t, right_w, rh, CARD)
        _text(slide, right_l + 0.14, t, col[0] - 0.14, rh, r.get("metric", ""),
              SZ_BODY, LABEL, bold=True, anchor="ctr")
        _text(slide, right_l + col[0], t, col[1], rh, r.get("a", ""), SZ_BODY, MUTED, anchor="ctr")
        _text(slide, right_l + col[0] + col[1], t, col[2], rh, r.get("b", ""),
              SZ_BODY, LABEL, bold=True, anchor="ctr")
    if takeaway:
        _text(slide, right_l, BOTTOM - 0.56, right_w, 0.52, takeaway, SZ_SMALL, BODY,
              italic=True, align=PP_ALIGN.CENTER, anchor="ctr")


# ═════════════════════════════════════════════════════════════════════════════
# 5. Pillar grid -- numbered capability pillars, 2x2 or 1x3
# ═════════════════════════════════════════════════════════════════════════════
def draw_pillar_grid(slide, data):
    pillars = (data.get("pillars") or [])[:6]
    if not pillars:
        return
    n = len(pillars)
    cols = 2 if n in (2, 4) else min(3, n)
    # the height the fullest card actually needs. Body wraps to ~2 lines; every point is
    # kept (owner: never drop content), so budget for the true max, not a cap.
    max_pts = max((len(p.get("points") or []) for p in pillars), default=0)
    card_h = 1.90 + max_pts * 0.34 + 0.24
    for i, l, t, w, h in _grid(n, cols, max_row_h=card_h):
        p = pillars[i]
        colour = _accent(i)
        _card(slide, l, t, w, h, WHITE, BORDER)
        _bar(slide, l, t, w, 0.045, colour)
        _chip(slide, l + PAD, t + 0.22, 0.40, colour, _pop_icon(data))
        _text(slide, l + w - 0.95, t + 0.16, 0.80, 0.52, "%02d" % (i + 1),
              26, GHOST, bold=True, align=PP_ALIGN.RIGHT, font=FONT_HEAD)
        _text(slide, l + PAD, t + 0.74, w - 2 * PAD, 0.34, p.get("heading", ""),
              SZ_CARD_TITLE, LABEL, bold=True, font=FONT_HEAD)
        _text(slide, l + PAD, t + 1.10, w - 2 * PAD, 0.62, p.get("body", ""), SZ_BODY, BODY)
        pts = p.get("points") or []       # keep every point; no cap
        if pts:
            # fit ALL points in the remaining card height, shrinking the step (and, when
            # it's tight, the text) rather than dropping any
            py0 = t + 1.80
            avail = (t + h - 0.12) - py0
            step = max(0.19, min(0.52, avail / len(pts)))
            pt_sz = SZ_SMALL if step >= 0.24 else 8
            for j, pt in enumerate(pts):
                py = py0 + j * step
                _tick(slide, l + PAD, py, colour)
                _text(slide, l + PAD + 0.20, py - 0.03, w - 2 * PAD - 0.20, step, pt, pt_sz, BODY)


# ═════════════════════════════════════════════════════════════════════════════
# 6. Option columns -- three architecture options, compared row by row
# ═════════════════════════════════════════════════════════════════════════════
def draw_option_columns(slide, data):
    options = (data.get("options") or [])[:4]
    if not options:
        return
    rec = data.get("recommendation") or ""
    bottom = BOTTOM - (0.86 if rec else 0)
    n = len(options)
    w = (SW - 2 * MARGIN - (n - 1) * GAP) / n
    head_h, tag_h = 0.52, 0.34

    labels = []
    for o in options:
        labels = labels or [r.get("label", "") for r in (o.get("rows") or [])]
    for i, o in enumerate(options):
        l = MARGIN + i * (w + GAP)
        colour = _accent(i)
        _card(slide, l, TOP, w, bottom - TOP, WHITE, BORDER)
        _bar(slide, l, TOP, w, head_h, colour)
        _text(slide, l, TOP, w, head_h, o.get("name", ""), 14, WHITE, bold=True,
              align=PP_ALIGN.CENTER, font=FONT_HEAD, anchor="ctr")
        _bar(slide, l, TOP + head_h, w, tag_h, colour)
        _text(slide, l, TOP + head_h, w, tag_h, (o.get("tag", "") or "").upper(), 9.5,
              WHITE, bold=True, align=PP_ALIGN.CENTER, anchor="ctr")

        rows = (o.get("rows") or [])[:6]
        ry = TOP + head_h + tag_h + 0.10
        rh = (bottom - ry - 0.10) / max(len(rows), 1)
        for j, r in enumerate(rows):
            t = ry + j * rh
            _label(slide, l + 0.14, t + 0.04, w - 0.28, r.get("label", ""))
            _text(slide, l + 0.14, t + 0.24, w - 0.28, rh - 0.30, r.get("value", ""),
                  SZ_BODY, BODY)
            if j < len(rows) - 1:
                _bar(slide, l + 0.14, t + rh - 0.02, w - 0.28, 0.006, BORDER)

    if rec:
        t = BOTTOM - 0.74
        _card(slide, MARGIN, t, SW - 2 * MARGIN, 0.74, CARD, BORDER)
        _chip(slide, MARGIN + 0.20, t + 0.20, 0.34, RED)
        _rich(slide, MARGIN + 0.70, t + 0.10, SW - 2 * MARGIN - 0.95, 0.56,
              [("Recommendation:  ", True, LABEL), (rec, False, BODY)], SZ_BODY)


# ═════════════════════════════════════════════════════════════════════════════
# 7. Agent architecture -- a card grid of agents, each with a metric pill
# ═════════════════════════════════════════════════════════════════════════════
def draw_agent_architecture(slide, data):
    agents = (data.get("agents") or [])[:6]
    if not agents:
        return
    footer = data.get("footer") or ""
    bottom = BOTTOM - (0.78 if footer else 0)

    # 3 over 2, as the reference design: a full top row, a centred second row
    top_n = 3 if len(agents) > 3 else len(agents)
    rows = [agents[:top_n], agents[top_n:]]
    rows = [r for r in rows if r]
    row_h = (bottom - TOP - (len(rows) - 1) * 0.26) / len(rows)

    idx = 0
    for ri, row in enumerate(rows):
        w = (SW - 2 * MARGIN - (top_n - 1) * GAP) / top_n
        total = len(row) * w + (len(row) - 1) * GAP
        l0 = (SW - total) / 2                       # centre a short row
        t = TOP + ri * (row_h + 0.26)
        for j, a in enumerate(row):
            l = l0 + j * (w + GAP)
            colour = _accent(idx); idx += 1
            _card(slide, l, t, w, row_h, WHITE, BORDER)
            _bar(slide, l, t, w, 0.045, colour)
            _chip(slide, l + PAD, t + 0.20, 0.28, colour, _pop_icon(data))
            _text(slide, l + PAD + 0.40, t + 0.18, w - 0.90, 0.32, a.get("name", ""),
                  12, LABEL, bold=True, font=FONT_HEAD)
            _bar(slide, l + PAD, t + 0.58, w - 2 * PAD, 0.008, BORDER)
            _text(slide, l + PAD, t + 0.70, w - 2 * PAD, row_h - 1.40, a.get("body", ""),
                  SZ_SMALL, BODY)
            badge = a.get("badge") or ""
            if badge:
                bw = min(w - 2 * PAD, 0.085 * len(badge) + 0.34)
                _card(slide, l + PAD, t + row_h - 0.52, bw, 0.32, _pale(colour),
                      line=None, rounded=True)
                _text(slide, l + PAD, t + row_h - 0.52, bw, 0.32, badge, 8.5, colour,
                      bold=True, align=PP_ALIGN.CENTER, anchor="ctr")
    if footer:
        _footer_bar(slide, footer, NAVY, 0.66)


# ═════════════════════════════════════════════════════════════════════════════
# 8. Governance timeline -- chips on a vertical spine, each with a card
# ═════════════════════════════════════════════════════════════════════════════
def draw_governance_list(slide, data):
    items = (data.get("items") or [])[:5]
    if not items:
        return
    n = len(items)
    h = (BOTTOM - TOP - (n - 1) * GAP) / n
    spine_x = 0.72
    _bar(slide, spine_x, TOP + h / 2, 0.012, (n - 1) * (h + GAP), BORDER)
    for i, it in enumerate(items):
        t = TOP + i * (h + GAP)
        colour = _accent(i)
        _chip(slide, spine_x - 0.19, t + h / 2 - 0.19, 0.38, colour, _pop_icon(data))
        l = 1.15
        w = SW - MARGIN - l
        _card(slide, l, t, w, h, WHITE, BORDER)
        _bar(slide, l, t, 0.055, h, colour)
        _text(slide, l + 0.24, t + 0.12, w - 0.44, 0.32, (it.get("heading", "") or "").upper(),
              12, LABEL, bold=True, font=FONT_HEAD)
        _text(slide, l + 0.24, t + 0.46, w - 0.44, h - 0.58, it.get("body", ""), SZ_BODY, BODY)


# ═════════════════════════════════════════════════════════════════════════════
# 9. Guardrail columns -- three cards of bold-lead bullets + a quick-win callout
# ═════════════════════════════════════════════════════════════════════════════
def draw_guardrail_columns(slide, data):
    columns = (data.get("columns") or [])[:4]
    if not columns:
        return
    callout_label = data.get("callout_label") or ""
    callout_body = data.get("callout_body") or ""
    bottom = BOTTOM - (0.86 if (callout_label or callout_body) else 0)

    n = len(columns)
    w = (SW - 2 * MARGIN - (n - 1) * GAP) / n
    for i, c in enumerate(columns):
        l = MARGIN + i * (w + GAP)
        colour = _accent(i)
        _card(slide, l, TOP, w, bottom - TOP, CARD, BORDER)
        _bar(slide, l, TOP, w, 0.045, colour)
        _chip(slide, l + PAD, TOP + 0.24, 0.40, colour, _pop_icon(data))
        _text(slide, l + PAD, TOP + 0.80, w - 2 * PAD, 0.36, c.get("heading", ""),
              SZ_CARD_TITLE, LABEL, bold=True, font=FONT_HEAD)
        py = TOP + 1.26
        for p in (c.get("points") or [])[:4]:
            if py + 0.40 > bottom - 0.10:
                break
            _tick(slide, l + PAD, py, colour)
            lead = (p.get("lead") or "").strip()
            body = (p.get("body") or "").strip()
            parts = []
            if lead:
                parts.append((lead + (": " if body else ""), True, LABEL))
            if body:
                parts.append((body, False, BODY))
            _rich(slide, l + PAD + 0.20, py - 0.02, w - 2 * PAD - 0.20, 0.80,
                  parts or [("", False, BODY)], SZ_SMALL)
            py += 0.86

    if callout_label or callout_body:
        t = BOTTOM - 0.74
        _card(slide, MARGIN, t, SW - 2 * MARGIN, 0.74, CARD, BORDER)
        _chip(slide, MARGIN + 0.20, t + 0.20, 0.34, RED)
        _rich(slide, MARGIN + 0.70, t + 0.10, SW - 2 * MARGIN - 0.95, 0.56,
              [((callout_label or "").upper() + "   ", True, RED),
               (callout_body, False, BODY)], SZ_BODY)


# ═════════════════════════════════════════════════════════════════════════════
# 10. Opportunity cards -- numbered columns, each an opportunity and its outcome
# ═════════════════════════════════════════════════════════════════════════════
def draw_opportunity_cards(slide, data):
    cards = (data.get("cards") or [])[:4]
    if not cards:
        return
    n = len(cards)
    w = (SW - 2 * MARGIN - (n - 1) * GAP) / n
    h = BOTTOM - TOP
    for i, c in enumerate(cards):
        l = MARGIN + i * (w + GAP)
        colour = _accent(i)
        _card(slide, l, TOP, w, h, WHITE, BORDER)
        _bar(slide, l, TOP, w, 0.045, colour)
        _chip(slide, l + PAD, TOP + 0.20, 0.40, colour, _pop_icon(data))
        _text(slide, l + w - 1.10, TOP + 0.14, 0.95, 0.66, "%02d" % (i + 1),
              30, GHOST, bold=True, align=PP_ALIGN.RIGHT, font=FONT_HEAD)
        _text(slide, l + PAD, TOP + 0.72, w - 2 * PAD, 0.60, (c.get("heading", "") or "").upper(),
              12, LABEL, bold=True, font=FONT_HEAD)

        inner_w = w - 2 * PAD
        box_h = (h - 1.56) / 2
        ot = TOP + 1.42
        _card(slide, l + PAD, ot, inner_w, box_h, CARD, BORDER)
        _bar(slide, l + PAD, ot, 0.045, box_h, MUTED)
        _label(slide, l + PAD + 0.20, ot + 0.10, inner_w - 0.36, "Opportunity")
        _text(slide, l + PAD + 0.20, ot + 0.34, inner_w - 0.36, box_h - 0.44,
              c.get("opportunity", ""), SZ_SMALL, BODY)

        ut = ot + box_h + 0.14
        _card(slide, l + PAD, ut, inner_w, box_h, _pale(colour), line=None)
        _bar(slide, l + PAD, ut, 0.045, box_h, colour)
        _label(slide, l + PAD + 0.20, ut + 0.10, inner_w - 0.36, "Outcome", colour)
        _text(slide, l + PAD + 0.20, ut + 0.34, inner_w - 0.36, box_h - 0.44,
              c.get("outcome", ""), SZ_SMALL, BODY)


# ═════════════════════════════════════════════════════════════════════════════
# 11. Cover brief -- an intro line, then a row of highlight chips
# ═════════════════════════════════════════════════════════════════════════════
def draw_cover_brief(slide, data):
    intro = (data.get("intro") or "").strip()
    chips = (data.get("chips") or [])[:6]
    y = TOP
    if intro:
        _text(slide, MARGIN, y, SW - 2 * MARGIN, 0.85, intro, 13, BODY)
        y += 1.05
    if not chips:
        return
    n = len(chips)
    w = (SW - 2 * MARGIN - (n - 1) * GAP) / n
    h = min(1.15, BOTTOM - y)
    y += max(0.0, (BOTTOM - y - h) / 3)   # a little breathing room, not dead-centred
    for i, c in enumerate(chips):
        l = MARGIN + i * (w + GAP)
        colour = _accent(i)
        _card(slide, l, y, w, h, WHITE, BORDER)
        _bar(slide, l, y, 0.06, h, colour)
        _text(slide, l + 0.22, y, w - 0.36, h, c, SZ_CARD_TITLE, LABEL, bold=True, anchor="ctr")


# ═════════════════════════════════════════════════════════════════════════════
# 12. Pain vs. advantage split -- two full-sentence lanes + a closing stat strip
# ═════════════════════════════════════════════════════════════════════════════
def _lane_rows(slide, l, t, w, h, title, rows, colour):
    _card(slide, l, t, w, h, WHITE, BORDER)
    _bar(slide, l, t, w, 0.05, colour)
    _text(slide, l + PAD, t + 0.16, w - 2 * PAD, 0.30, (title or "").upper(),
          12, LABEL, bold=True, font=FONT_HEAD)
    rows = (rows or [])[:5]
    if not rows:
        return
    ry = t + 0.56
    rh = (h - 0.56 - PAD) / len(rows)
    for i, r in enumerate(rows):
        rt = ry + i * rh
        _tick(slide, l + PAD, rt + 0.06, colour)
        lead = (r.get("lead") or "").strip()
        body = (r.get("body") or "").strip()
        parts = []
        if lead:
            parts.append((lead + (": " if body else ""), True, LABEL))
        if body:
            parts.append((body, False, BODY))
        _rich(slide, l + PAD + 0.20, rt, w - 2 * PAD - 0.20, rh - 0.06,
              parts or [("", False, BODY)], SZ_SMALL)


def draw_pain_advantage_split(slide, data):
    intro = (data.get("intro") or "").strip()
    lane_top = TOP
    if intro:
        _text(slide, MARGIN, TOP - 0.04, SW - 2 * MARGIN, 0.32, intro, 11, BODY)
        lane_top = TOP + 0.38

    strip_stats = (data.get("strip_stats") or [])[:4]
    strip_h = 1.05 if strip_stats else 0.0
    lane_bottom = BOTTOM - (strip_h + 0.30 if strip_stats else 0)
    lane_w = (SW - 2 * MARGIN - 0.55) / 2

    _lane_rows(slide, MARGIN, lane_top, lane_w, lane_bottom - lane_top,
               data.get("left_title") or "Current state", data.get("left_rows"), RED)
    _lane_rows(slide, MARGIN + lane_w + 0.55, lane_top, lane_w, lane_bottom - lane_top,
               data.get("right_title") or "J2W advantage", data.get("right_rows"), TEAL)

    # the VS badge between the two lanes
    cx = MARGIN + lane_w + 0.275
    cy = lane_top + (lane_bottom - lane_top) / 2
    r = 0.26
    badge = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(cx - r), Inches(cy - r),
                                   Inches(2 * r), Inches(2 * r))
    badge.fill.solid(); badge.fill.fore_color.rgb = NAVY
    badge.line.fill.background(); badge.shadow.inherit = False
    _text(slide, cx - r, cy - r * 0.62, 2 * r, r * 1.2, "VS", 13, WHITE, bold=True,
          align=PP_ALIGN.CENTER, font=FONT_HEAD, anchor="ctr")

    if not strip_stats:
        return
    st = BOTTOM - strip_h
    _card(slide, MARGIN, st, SW - 2 * MARGIN, strip_h, CARD, BORDER)
    label = data.get("strip_label") or ""
    seg_l = MARGIN + PAD
    if label:
        _label(slide, MARGIN + PAD, st + 0.10, 2.2, label, LABEL, 10)
        seg_l = MARGIN + 2.35
    n = len(strip_stats)
    seg_w = (SW - MARGIN - PAD - seg_l - (n - 1) * 0.02) / n
    for i, s in enumerate(strip_stats):
        l = seg_l + i * seg_w
        if i:
            _bar(slide, l - 0.02, st + 0.14, 0.01, strip_h - 0.28, BORDER)
        _text(slide, l, st + 0.10, seg_w, 0.44, s.get("value", ""), 20, _accent(i),
              bold=True, align=PP_ALIGN.CENTER, font=FONT_HEAD)
        _text(slide, l, st + 0.58, seg_w, 0.30, (s.get("label", "") or "").upper(),
              8.5, MUTED, align=PP_ALIGN.CENTER)


# ═════════════════════════════════════════════════════════════════════════════
# 13. Stat + table + callout -- headline stats, a 3-column data table, a callout
# ═════════════════════════════════════════════════════════════════════════════
def draw_stat_table_callout(slide, data):
    stats = (data.get("stats") or [])[:4]
    cols = (data.get("col_labels") or ["", "", ""])[:3]
    while len(cols) < 3:
        cols.append("")
    rows = (data.get("rows") or [])[:8]
    callout_label = data.get("callout_label") or ""
    callout_body = data.get("callout_body") or ""
    closing = (data.get("closing") or "").strip()

    y = TOP
    if stats:
        n = len(stats)
        w = (SW - 2 * MARGIN - (n - 1) * GAP) / n
        sh = 1.05
        for i, s in enumerate(stats):
            l = MARGIN + i * (w + GAP)
            colour = _accent(i)
            _card(slide, l, y, w, sh, WHITE, BORDER)
            _bar(slide, l, y, w, 0.05, colour)
            _text(slide, l, y + 0.14, w, 0.50, s.get("value", ""), SZ_STAT - 4, colour,
                  bold=True, align=PP_ALIGN.CENTER, font=FONT_HEAD)
            _text(slide, l, y + 0.66, w, 0.30, (s.get("label", "") or "").upper(),
                  SZ_SMALL, LABEL, bold=True, align=PP_ALIGN.CENTER)
        y += sh + 0.24

    callout_h = 0.66 if (callout_label or callout_body) else 0.0
    closing_h = 0.30 if closing else 0.0
    table_bottom = (BOTTOM - callout_h - closing_h
                    - (0.14 if callout_h else 0) - (0.10 if closing_h else 0))

    if rows:
        body_w = SW - 2 * MARGIN
        col_w = [body_w * 0.30, body_w * 0.24, body_w * 0.46]
        col_l = [MARGIN, MARGIN + col_w[0], MARGIN + col_w[0] + col_w[1]]
        head_h = 0.30
        for cl, cw, label in zip(col_l, col_w, cols):
            _text(slide, cl, y, cw, head_h, label, SZ_BODY, MUTED, bold=True)
        _bar(slide, MARGIN, y + head_h, body_w, 0.015, BORDER)
        ry = y + head_h + 0.10
        rh = (table_bottom - ry) / len(rows)
        for i, r in enumerate(rows):
            rt = ry + i * rh
            if i % 2 == 1:
                _bar(slide, MARGIN, rt, body_w, rh, CARD)
            _text(slide, col_l[0] + 0.06, rt, col_w[0] - 0.12, rh, r.get("col1", ""),
                  SZ_BODY, LABEL, bold=True, anchor="ctr")
            _text(slide, col_l[1] + 0.06, rt, col_w[1] - 0.12, rh, r.get("col2", ""),
                  SZ_BODY, TEAL, bold=True, anchor="ctr")
            _text(slide, col_l[2] + 0.06, rt, col_w[2] - 0.12, rh, r.get("col3", ""),
                  SZ_SMALL, BODY, anchor="ctr")
            if i < len(rows) - 1:
                _bar(slide, MARGIN, rt + rh - 0.01, body_w, 0.008, BORDER)

    if callout_label or callout_body:
        t = BOTTOM - callout_h - closing_h - (0.10 if closing_h else 0)
        _card(slide, MARGIN, t, SW - 2 * MARGIN, callout_h, PALE_TEAL, line=None)
        _bar(slide, MARGIN, t, 0.06, callout_h, TEAL)
        _rich(slide, MARGIN + 0.30, t + (callout_h - 0.30) / 2, SW - 2 * MARGIN - 0.55, 0.30,
              [((callout_label or "").upper() + ":  ", True, TEAL),
               (callout_body, False, BODY)], SZ_SMALL)
    if closing:
        _text(slide, MARGIN, BOTTOM - closing_h, SW - 2 * MARGIN, closing_h, closing,
              SZ_SMALL, BODY, italic=True, align=PP_ALIGN.CENTER)


# ── the registry skills.build_into dispatches through ────────────────────────
# kind -> (header-marker mapping, body drawer)
DRAWERS = {
    "pain_point_list":     (_head, draw_pain_point_list),
    "platform_overview":   (_head, draw_platform_overview),
    "before_after_split":  (_head, draw_before_after_split),
    "comparison_split":    (_head, draw_comparison_split),
    "pillar_grid":         (_head, draw_pillar_grid),
    "option_columns":      (_head, draw_option_columns),
    "agent_architecture":  (_head, draw_agent_architecture),
    "governance_list":     (_head, draw_governance_list),
    "guardrail_columns":   (_head, draw_guardrail_columns),
    "opportunity_cards":   (_head, draw_opportunity_cards),
    "cover_brief":            (_head, draw_cover_brief),
    "pain_advantage_split":   (_head, draw_pain_advantage_split),   # draws its own intro (see draw_before_after_split)
    "stat_table_callout":     (_head_intro, draw_stat_table_callout),
}
