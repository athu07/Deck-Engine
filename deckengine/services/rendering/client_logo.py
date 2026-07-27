# -*- coding: utf-8 -*-
"""
client_logo.py  --  the client's logo, background-removed, stamped onto the
title slide (CS01) side by side with the J2W wordmark ("Client x JoulesToWatts",
owner's spec, 2026-07-08).

Two ways in: fetch_via_search() finds it automatically from a company name (+
optional domain to disambiguate) using OpenAI's web_search tool -- no upload
needed. save() also still accepts a manually-uploaded file, kept as a fallback
(see below for why this matters -- auto-fetch is real but not reliable enough
to be the ONLY path). Both end up going through the same storage/background-
removal pipeline.

On auto-fetch reliability (owner asked for name+domain auto-fetch, 2026-07-08;
this is what was actually tested and found, live, before building):
Clearbit's free logo-by-domain endpoint -- the obvious first choice -- turned
out to be dead (logo.clearbit.com no longer even resolves, sunset after the
HubSpot acquisition); logo.dev and Brandfetch's public endpoints both require
an API key/client ID now, so neither is a keyless option. What DOES work with
zero new keys is asking the model to search the web and return a direct image
URL, but real testing showed a meaningful chunk of those URLs 404 or point to
a non-public staging subdomain the model can't tell apart from a real one
(e.g. it confidently returned "webdev.waaree.com" for one real company --
looks like a real asset path, isn't publicly reachable). So: fetch_via_search()
ACTUALLY FETCHES AND VALIDATES the URL it's given (opens it as a real image
with Pillow) before accepting it, rather than trusting the model's URL at face
value -- and returns None on any failure, same fail-safe contract as every
other lookup in this app. Because it doesn't always find one, the manual
upload path stays in the UI as a fallback, and the web UI shows a preview
before the logo is used, so the salesperson always knows whether it actually
found the right thing before it ends up in a client-facing deck.

Storage: one PNG per client name (save() overwrites any previous logo for the
same client -- always the latest one). Background removal is intentionally
simple, not a full ML cutout: most real brand logos are already transparent
PNGs, so this only needs to handle the remaining case -- a flat/near-solid
background colour (a screenshot, a scan, a JPEG export) -- by flood-filling
that colour to transparent from the four corners. A photographic or gradient
background is left alone rather than risk eating into the logo's own artwork.

Placement: CS01's own J2W wordmark sits at approximately (5.06in, 0.49in,
3.21in x 1.55in) on the master deck's 13.33x7.5in canvas -- pixel-measured
2026-07-08. The space to its left (down to the title text starting at
y=2.47in) is empty in the master's own design, so the client logo + a small
"x" divider are placed there, right-aligned up to the wordmark, matching its
vertical centre. Scaled down slightly from the wordmark's own height so the
client logo reads as a guest mark, not the dominant one.
"""

import io
import os
import re

from PIL import Image, UnidentifiedImageError
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

from deckengine import config
from deckengine.services.content.build_library import read_id

CS01 = "CS01"

# J2W wordmark geometry on CS01 (pixel-measured against the real master deck).
_WORDMARK_LEFT = 5.06
_WORDMARK_TOP = 0.49
_WORDMARK_H = 1.55
_WORDMARK_CENTER_Y = _WORDMARK_TOP + _WORDMARK_H / 2

_GAP_W = 0.5             # reserved for the "x" divider, immediately left of the wordmark
_LOGO_MAX_W = 2.0
_LOGO_MAX_H = 1.2         # slightly under the wordmark's own 1.55in -- a guest mark, not the lead

_BG_TOLERANCE = 28        # per-channel colour distance treated as "same as background"
_MAX_FLOODFILL_PIXELS = 2_000_000   # skip removal on anything larger (too slow, unlikely for a logo)


def _slug(name):
    return re.sub(r"[^A-Za-z0-9]+", "_", (name or "").strip()).strip("_") or "client"


def _remove_solid_background(img):
    """If this image's four corners are a consistent, near-solid colour, flood-
    fill that colour to transparent from each corner (4-connected -- so it only
    clears the connected background region, never a similar colour trapped
    INSIDE the logo's own artwork). Left untouched if the corners disagree (a
    photo/gradient background) or the image is too large to process quickly."""
    img = img.convert("RGBA")
    w, h = img.size
    if w * h > _MAX_FLOODFILL_PIXELS or w < 2 or h < 2:
        return img
    px = img.load()
    corners = [px[0, 0], px[w - 1, 0], px[0, h - 1], px[w - 1, h - 1]]

    def _dist(a, b):
        return abs(a[0] - b[0]) + abs(a[1] - b[1]) + abs(a[2] - b[2])

    base = corners[0]
    if any(_dist(base, c) > _BG_TOLERANCE for c in corners[1:]):
        return img          # corners disagree -- not a uniform background, leave it as-is

    seen = bytearray(w * h)
    stack = [(0, 0), (w - 1, 0), (0, h - 1), (w - 1, h - 1)]
    for x, y in stack:
        seen[y * w + x] = 1
    while stack:
        x, y = stack.pop()
        r, g, b, a = px[x, y]
        if _dist((r, g, b), base) > _BG_TOLERANCE:
            continue
        px[x, y] = (r, g, b, 0)
        for nx, ny in ((x + 1, y), (x - 1, y), (x, y + 1), (x, y - 1)):
            if 0 <= nx < w and 0 <= ny < h and not seen[ny * w + nx]:
                seen[ny * w + nx] = 1
                stack.append((nx, ny))
    return img


_MAX_LOGO_BYTES = 8 * 1024 * 1024   # a real logo is small; refuse anything absurd


def _fetch_url(url):
    """GET a URL and return its bytes only if they're small AND actually open
    as a real image -- this is the real safety net, not the URL itself. The
    model can be confidently wrong about a URL being fetchable (a non-public
    staging subdomain, a stale link); trusting it at face value is how a
    build would end up with a broken image or, worse, silently no logo with
    no way to tell why."""
    try:
        import urllib.request
        req = urllib.request.Request(url, headers={"User-Agent": "Mozilla/5.0"})
        with urllib.request.urlopen(req, timeout=8) as resp:
            if resp.status != 200:
                return None
            data = resp.read(_MAX_LOGO_BYTES + 1)
            if len(data) > _MAX_LOGO_BYTES:
                return None
    except Exception:
        return None
    try:
        img = Image.open(io.BytesIO(data))
        img.load()                    # actually decode it -- a real image, not an error page
    except (UnidentifiedImageError, Exception):
        return None
    # re-encode as PNG -- the fetched bytes could be JPEG/WEBP/whatever the source
    # served; normalising here means every caller (the data-URI preview, save())
    # can safely assume "these bytes are a PNG" instead of re-detecting format.
    out = io.BytesIO()
    try:
        img.convert("RGBA").save(out, "PNG")
    except Exception:
        return None
    return out.getvalue()


def fetch_via_search(client_name, domain=""):
    """Best-effort logo fetch from a company NAME (+ optional domain to help
    disambiguate) via OpenAI's web_search tool -- no API key or vendor beyond
    the OPENAI_API_KEY already configured for matching/generation. Every
    candidate URL is actually fetched and verified as a real image (see
    _fetch_url) before being accepted, since the model's URL is a guess, not
    a guarantee. Returns the raw image bytes, or None if nothing verifiably
    fetchable was found -- fails safe, same as every other lookup in this
    app; the caller just ends up with no logo, never a broken build."""
    client_name = (client_name or "").strip()
    if not client_name:
        return None
    prompt = (
        f'Find a direct URL to a RASTER image file (.png or .jpg ONLY, never .svg) of the '
        f'official logo for the company "{client_name}"'
        + (f' (domain: {domain})' if domain else '') + '. Prefer the company\'s own public '
        'website (a press/media/brand-assets page, not an internal staging subdomain), or a '
        'Wikipedia Commons /commons/ full-resolution URL (NOT a /thumb/ URL with a manually '
        'chosen pixel width -- those are often rejected). Reply with ONLY the raw URL and '
        'nothing else -- no explanation, no markdown.'
    )
    try:
        from deckengine.services.infra import load_env
        load_env()
        from openai import OpenAI
        resp = OpenAI().responses.create(
            model="gpt-4o-mini", tools=[{"type": "web_search"}], input=prompt)
        url = (resp.output_text or "").strip().strip('"').strip("'")
    except Exception:
        return None
    if not url.startswith("http") or " " in url:
        return None
    return _fetch_url(url)


def save(client_name, data):
    """Store `data` (raw uploaded image bytes) as this client's logo, background-
    removed if it looks like a solid fill. Returns the saved path, or None if
    `data` isn't a readable image (fails safe -- a bad upload just means no logo,
    never a broken build)."""
    try:
        img = Image.open(io.BytesIO(data))
        img.load()
    except Exception:
        return None
    img = _remove_solid_background(img)
    os.makedirs(config.CLIENT_LOGOS_DIR, exist_ok=True)
    path = os.path.join(config.CLIENT_LOGOS_DIR, _slug(client_name) + ".png")
    try:
        img.save(path, "PNG")
    except Exception:
        return None
    return path


def path_for(client_name):
    p = os.path.join(config.CLIENT_LOGOS_DIR, _slug(client_name) + ".png")
    return p if os.path.exists(p) else None


def _add_x_divider(slide, logo_right):
    divider = slide.shapes.add_textbox(Inches(logo_right), Inches(_WORDMARK_CENTER_Y - 0.25),
                                       Inches(_GAP_W), Inches(0.5))
    tf = divider.text_frame
    tf.word_wrap = False
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    r.text = "×"          # "x" -- Client x JoulesToWatts
    r.font.size = Pt(28)
    r.font.bold = True
    r.font.color.rgb = RGBColor(0x6B, 0x6B, 0x6B)


def stamp_into(deck_path, client_name):
    """After assembler.build_deck() has already copied CS01 into the deck, add
    this client's logo beside the J2W wordmark, plus a small "x" divider -- only
    when a real logo was auto-fetched or uploaded for this client name. If none
    was found, CS01 is left exactly as it is in the master deck (owner-spec,
    2026-07-23: reverted the dashed placeholder box from 2026-07-20 -- no logo
    means no visible cue at all, not a stand-in graphic). Returns True if a real
    logo was stamped."""
    prs = Presentation(deck_path)
    slide = next((s for s in prs.slides if read_id(s) == CS01), None)
    if slide is None:
        return False

    logo_path = path_for(client_name)
    if not logo_path:
        return False
    logo_right = _WORDMARK_LEFT - _GAP_W

    with Image.open(logo_path) as im:
        iw, ih = im.size
    scale = min(_LOGO_MAX_W / iw, _LOGO_MAX_H / ih)
    logo_w, logo_h = iw * scale, ih * scale
    logo_left = logo_right - logo_w
    logo_top = _WORDMARK_CENTER_Y - logo_h / 2
    slide.shapes.add_picture(logo_path, Inches(logo_left), Inches(logo_top),
                             width=Inches(logo_w), height=Inches(logo_h))
    _add_x_divider(slide, logo_right)

    prs.save(deck_path)
    return True
