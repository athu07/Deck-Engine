# -*- coding: utf-8 -*-
"""
client_context.py  --  the "Client Context" (CS65) + "Tailored Approach" (CS66)
slide pair, per SLIDE_SELECTION_BRIEF.md.

Gate: Workforce-only, First/Second stage only (never Intro). Unlike the other
Workforce slides, these aren't picked from the registry's stage/trigger columns
at all -- they need a whole AI extraction (client name, stats, up to 4 real
talent challenges + solutions) to fill their `[BRACKET]`-style markers, so
inclusion is decided the same way skills.py decides its data-driven slides:
candidates() runs the extraction; if it comes back empty (fails closed -- see
ai_matcher.extract_client_context), there is NO slide, never a half-filled one.

CS65/CS66 are real slides already in the master deck (not a separate template
file), so assembler.build_deck() copies them in like any other CSxx once their
ids are in the final order. This module's job is the FILL step that runs after
assembly: replace every `[BRACKET]` marker + illustrative "e.g. ..." example
paragraph with the extracted content, and BLANK (never leave a literal bracket
or an "e.g." example) anything the notes didn't supply.
"""

import re

from pptx import Presentation

from deckengine.services.matching import ai_matcher
from deckengine.services.content.build_library import read_id

CS_CONTEXT = "CS65"    # Client Context
CS_APPROACH = "CS66"   # Tailored Approach

_CLIENT_NAME_RE = re.compile(r"\[CLIENT NAME\]", re.I)
_CLIENT_RE = re.compile(r"\[CLIENT\]", re.I)
_DATE_RE = re.compile(r"\[DATE\]", re.I)
_CHALLENGE_TITLE_RE = re.compile(r"^\[Challenge \d Title\]$")
_SOLUTION_TITLE_RE = re.compile(r"^\[Solution \d Title\]$")


def gate(context):
    from deckengine.services.matching.matcher import stage_of
    wts = {str(w).upper() for w in (context.get("work_types") or [])}
    stage = stage_of(context.get("phase"))
    return "WORKFORCE" in wts and stage in ("first", "second")


def candidates(context):
    """[] unless the gate passes AND the AI extraction returns enough real
    content (fails closed -- see ai_matcher.extract_client_context)."""
    if not gate(context):
        return []
    data = ai_matcher.extract_client_context(context.get("transcript", ""))
    if not data:
        return []
    label = f"Client context - {data['client_name']}"
    return [
        {"id": CS_CONTEXT, "kind": "client_context", "label": label, "data": data},
        {"id": CS_APPROACH, "kind": "tailored_approach", "label": "Tailored approach", "data": data},
    ]


def _replace_bracket(shape, pattern, value):
    """Substring-replace a bracket pattern within a shape's text, run-by-run
    (keeps formatting) -- same convention as editor.replace_tokens."""
    for p in shape.text_frame.paragraphs:
        for r in p.runs:
            if pattern.search(r.text):
                r.text = pattern.sub(value, r.text)


def _blank(shape):
    """Blank a shape's whole text (never leave a literal bracket or an
    illustrative 'e.g. ...' example in a client-facing deck)."""
    tf = shape.text_frame
    if tf.paragraphs and tf.paragraphs[0].runs:
        tf.paragraphs[0].runs[0].text = ""
        for r in tf.paragraphs[0].runs[1:]:
            r.text = ""
        for p in tf.paragraphs[1:]:
            for r in p.runs:
                r.text = ""


def _set_whole(shape, value):
    tf = shape.text_frame
    if tf.paragraphs and tf.paragraphs[0].runs:
        tf.paragraphs[0].runs[0].text = value
        for r in tf.paragraphs[0].runs[1:]:
            r.text = ""
        for p in tf.paragraphs[1:]:
            for r in p.runs:
                r.text = ""


def _fill_context_slide(slide, data):
    challenges = data.get("challenges") or []
    scalar_subs = [
        (re.compile(r"\[X\]%\s*Offer Drop Rate", re.I), "offer_drop_pct", "{}% Offer Drop Rate"),
        (re.compile(r"\[X\]\s*people strong", re.I), "org_size", "{} people strong"),
        (re.compile(r"\[X[–—-]X\]\s*hires in\s*\[YEAR\]", re.I), None, None),  # handled separately
        (re.compile(r"\[X\]%\s*junior workforce", re.I), "junior_pct", "{}% junior workforce"),
        (re.compile(r"\[City\]\s*extension of\s*\[HQ\]", re.I), None, None),             # handled separately
    ]
    challenge_i = 0
    for sh in slide.shapes:
        if not sh.has_text_frame:
            continue
        text = sh.text_frame.text

        if _CLIENT_NAME_RE.search(text):
            _replace_bracket(sh, _CLIENT_NAME_RE, data["client_name"])
            continue
        if "[DATE]" in text:
            if data.get("date"):
                _replace_bracket(sh, _DATE_RE, data["date"])
            else:
                _blank(sh)
            continue
        if _CHALLENGE_TITLE_RE.match(text.strip()):
            if challenge_i < len(challenges):
                _set_whole(sh, challenges[challenge_i]["title"])
            else:
                _blank(sh)
            continue
        if text.strip().startswith("e.g. "):
            if challenge_i < len(challenges):
                _set_whole(sh, challenges[challenge_i]["body"])
                challenge_i += 1
            else:
                _blank(sh)
                challenge_i += 1
            continue
        if re.search(r"\[X\]%\s*Offer Drop Rate", text, re.I):
            if data.get("offer_drop_pct"):
                _set_whole(sh, f"{data['offer_drop_pct']} Offer Drop Rate")
            else:
                _blank(sh)
            continue
        if re.search(r"\[X\]\s*people strong", text, re.I):
            if data.get("org_size"):
                _set_whole(sh, f"{data['org_size']} people strong")
            else:
                _blank(sh)
            continue
        if re.search(r"hires in", text, re.I):
            if data.get("hiring_range") and data.get("hiring_year"):
                _set_whole(sh, f"{data['hiring_range']} hires in {data['hiring_year']}")
            else:
                _blank(sh)
            continue
        if re.search(r"junior workforce", text, re.I):
            if data.get("junior_pct"):
                _set_whole(sh, f"{data['junior_pct']} junior workforce")
            else:
                _blank(sh)
            continue
        if re.search(r"extension of", text, re.I):
            if data.get("city") and data.get("hq"):
                _set_whole(sh, f"{data['city']} extension of {data['hq']}")
            else:
                _blank(sh)
            continue


def _fill_approach_slide(slide, data):
    solutions = data.get("solutions") or []
    solution_i = 0
    for sh in slide.shapes:
        if not sh.has_text_frame:
            continue
        text = sh.text_frame.text

        if _CLIENT_RE.search(text):
            _replace_bracket(sh, _CLIENT_RE, data["client_name"])
            continue
        if _SOLUTION_TITLE_RE.match(text.strip()):
            if solution_i < len(solutions):
                _set_whole(sh, solutions[solution_i]["title"])
            else:
                _blank(sh)
            continue
        if text.strip().startswith("e.g. "):
            if solution_i < len(solutions):
                _set_whole(sh, solutions[solution_i]["body"])
            else:
                _blank(sh)
            continue
        if text.strip().startswith("Solves:"):
            if solution_i < len(solutions) and solutions[solution_i].get("solves"):
                _set_whole(sh, "Solves: " + solutions[solution_i]["solves"])
                solution_i += 1
            elif solution_i < len(solutions):
                _blank(sh)
                solution_i += 1
            else:
                _blank(sh)
            continue


def fill_into(deck_path, final_ids, cands):
    """After assembler.build_deck() has already copied CS65/CS66 into the deck
    (because their ids were in final_ids), replace their bracket markers with
    the extracted content. No-op if neither id is in final_ids or cands is empty."""
    cand_by_id = {c["id"]: c for c in cands if c["id"] in final_ids}
    if not cand_by_id:
        return 0

    prs = Presentation(deck_path)
    filled = 0
    for slide in prs.slides:
        sid = read_id(slide)
        cand = cand_by_id.get(sid)
        if not cand:
            continue
        if sid == CS_CONTEXT:
            _fill_context_slide(slide, cand["data"])
        elif sid == CS_APPROACH:
            _fill_approach_slide(slide, cand["data"])
        filled += 1

    if filled:
        prs.save(deck_path)
    return filled
