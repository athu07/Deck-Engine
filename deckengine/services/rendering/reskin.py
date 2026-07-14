# -*- coding: utf-8 -*-
"""
reskin.py  --  rebrand an uploaded PowerPoint into the J2W look WITHOUT changing
any of its content or slide count (beyond the two bookend slides described below).

Approach: restyle the ORIGINAL deck in place, slide by slide, box by box. We open
the uploaded file and keep every shape exactly where it is and every word exactly
as written -- all text, tables, images (including vector WMF/EMF), charts and the
slide size are untouched -- and only:
  * swap fonts + EXACT sizes to the J2W system (owner's spec, 2026-07-07): heading
    = Oswald 24 ALL-CAPS + always black, subheading = Roboto Condensed 15, Title
    Case + always teal, body = Raleway 13 (shrink-to-fit enabled as an overflow
    safety net, since arbitrary uploaded boxes weren't sized for our exact numbers),
  * force every slide's background to plain white, no matter what the original deck
    used -- including a full-bleed shape standing in for a background. If the
    original was dark, any body text authored light-on-dark is flipped to ink so it
    stays readable once the background goes white (owner's spec, 2026-07-07),
  * add a small red vertical accent bar to the left of each slide's own heading,
  * recolour "box" shapes (solid-filled containers with no text of their own --
    card backgrounds, icon/badge holders) to our own card convention: small,
    roughly-square ones become a red accent badge (the same red-badge style our
    own ranked-skill cards use); larger ones become our light card fill + hairline
    border (the same style every capability-card grid in our own templates uses),
  * add a J2W brand overlay to every original slide (the same thin red|teal top
    split bar every other J2W template uses),
  * prepend OUR OWN title slide and append OUR OWN "Let's win together" closing
    slide (pulled verbatim from the real master deck, CS01 / CS08) to every
    reskinned deck, scaled to the uploaded deck's own canvas size. These two are
    already fully on-brand by construction so, unlike the uploaded deck's own
    slides, they are copied as-is -- no restyling, no extra overlay.

Because nothing is rebuilt, nothing can be dropped — the exact tradeoff the owner
asked for ("every word, table, image, heading preserved"). The preview is produced
by rendering the restyled deck to real page images via LibreOffice (render_pngs), so
what you see is exactly what downloads.
"""

import glob
import os
import re
import subprocess

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.dml import MSO_COLOR_TYPE, MSO_FILL_TYPE
from pptx.enum.shapes import MSO_SHAPE, MSO_SHAPE_TYPE, PP_PLACEHOLDER
from pptx.enum.text import MSO_AUTO_SIZE

from deckengine import config
from deckengine.services.content.build_library import read_id
from deckengine.services.rendering.slide_generator import _copy_slide

# ── J2W identity — matched to the WORKING master deck's own title slide
# (pixel-extracted, 2026-07-07): red #D62839, teal #2A9D8F. Every J2W surface
# (case studies, skills slides, and this re-skinner) now shares one palette.
TEAL = RGBColor(0x2A, 0x9D, 0x8F)
RED = RGBColor(0xD6, 0x28, 0x39)
INK = RGBColor(0x11, 0x11, 0x10)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
SOFT = RGBColor(0xE7, 0xF0, 0xEE)
CARD_BG = RGBColor(0xF5, 0xF5, 0xF5)    # same card fill as case_study_v2 / skills templates
CARD_LINE = RGBColor(0xDE, 0xDE, 0xDE)  # same hairline border
HEAD_FONT = "Oswald"            # matches the case-study heading font
SUB_FONT = "Roboto Condensed"   # matches the case-study subheading font
BODY_FONT = "Raleway"           # matches the case-study content font
SZ_HEAD = 24
SZ_SUB = 15
SZ_BODY = 11        # owner's spec, 2026-07-08: plain body/content text
SZ_BODY_HEAD = 13   # a bold mini-heading NESTED inside a card (still Raleway, not
                     # the page-level Oswald heading) -- e.g. a card's own title
                     # line above its paragraph. Owner's spec: 13, content: 11.
# A run only qualifies for SZ_BODY_HEAD if it was originally >= this size AND
# bold. Without the size floor, a tiny bold eyebrow LABEL (e.g. "FOR YOUR
# CONFIRMATION" at 7pt in the real deck that surfaced this bug) would also get
# bumped all the way to 13 in a box sized for ~7pt text -- exactly the kind of
# overflow this whole distinction exists to avoid. A real card mini-heading is
# already sized noticeably larger than its own tiny caption labels.
SZ_BODY_HEAD_MIN_ORIG = 10

# ── contrast rule (owner's spec, 2026-07-07) ──────────────────────────────────
# Red-on-teal and teal-on-red are NOT legible -- both are mid-brightness,
# similarly-saturated colours with too little contrast against each other.
# White reads cleanly against either. This is what caught the badge-number bug:
# a step badge coloured red by _restyle_box, with its number label forced teal
# by the (correct, on its own) "subheadings are always teal" rule -- red bg +
# teal text is exactly the illegible pairing above.
#
# Checked with real WCAG-style contrast ratios (relative luminance, gamma-
# corrected): red-vs-teal = 1.50 (illegible), red-vs-white = 4.97, teal-vs-
# white = 3.32, white-vs-our-card-grey = 1.09 (also illegible -- the same bug
# shape shows up any time original light/white text is left on a box we've
# since recoloured light). _MIN_CONTRAST sits comfortably between those two
# groups. Any text box found sitting on top of a coloured box gets checked
# against this floor and, if the pairing is under it, flipped to a legible
# colour -- see _fix_text_on_box_contrast().
_MIN_CONTRAST = 2.2
# Named override: on our own red or teal brand fills specifically, the owner's
# rule is always WHITE text (not "whichever of white/ink technically scores
# higher" -- for teal, ink actually scores marginally higher, but white-on-
# red-or-teal is the house look). Everything else (e.g. our light card grey)
# falls through to the general best-of-white-or-ink pick below.
_NAMED_SAFE_TEXT = {str(RED): WHITE, str(TEAL): WHITE}


def _srgb_linear(c):
    c = c / 255.0
    return c / 12.92 if c <= 0.03928 else ((c + 0.055) / 1.055) ** 2.4


def _relative_luminance(rgb):
    r, g, b = rgb
    return 0.2126 * _srgb_linear(r) + 0.7152 * _srgb_linear(g) + 0.0722 * _srgb_linear(b)


def _contrast_ratio(rgb_a, rgb_b):
    la, lb = _relative_luminance(rgb_a), _relative_luminance(rgb_b)
    lighter, darker = max(la, lb), min(la, lb)
    return (lighter + 0.05) / (darker + 0.05)


def _legible_replacement(box_rgb):
    key = str(RGBColor(*box_rgb))
    if key in _NAMED_SAFE_TEXT:
        return _NAMED_SAFE_TEXT[key]
    c_white = _contrast_ratio(box_rgb, (255, 255, 255))
    c_ink = _contrast_ratio(box_rgb, (0x11, 0x11, 0x10))
    return WHITE if c_white >= c_ink else INK

_HEAD_PH = (PP_PLACEHOLDER.TITLE, PP_PLACEHOLDER.CENTER_TITLE)


# ── 1) restyle the original deck in place ─────────────────────────────────────
def restyle_deck(data, out_path):
    """Rebrand the uploaded pptx bytes to J2W and save to out_path (content intact)."""
    prs = Presentation(_as_stream(data))
    sw, sh = prs.slide_width, prs.slide_height
    for slide in prs.slides:
        size_ranks = _slide_sizes(slide.shapes)   # this slide's own size hierarchy
        dark_bg = _slide_is_dark(slide, slide.shapes, sw, sh)   # BEFORE we whiten it
        _restyle_shapes(slide.shapes, size_ranks, sw, sh, dark_bg)
        _fix_text_on_box_contrast(slide.shapes)
        _fix_picture_contrast(slide.shapes)
        _resolve_bottom_banner_overlap(slide.shapes, sw, sh)
        heading = _first_heading_shape(slide.shapes, size_ranks)
        if heading is not None:
            _add_heading_bar(slide, heading)
        _whiten_slide_background(slide)
        _brand_slide(slide, sw, sh)
    _append_bookend_slides(prs, sw, sh)
    out_dir = os.path.dirname(out_path)
    if out_dir:
        os.makedirs(out_dir, exist_ok=True)
    prs.save(out_path)
    return out_path


def _as_stream(data):
    import io
    return io.BytesIO(data) if isinstance(data, (bytes, bytearray)) else data


def _max_run_size(sh):
    """The largest font size (pt) among this shape's runs, or None if unset."""
    try:
        sizes = [r.font.size.pt for p in sh.text_frame.paragraphs
                for r in p.runs if r.text.strip() and r.font.size]
        return max(sizes) if sizes else None
    except Exception:
        return None


def _slide_sizes(shapes, acc=None):
    """Every text shape's max run size on this slide (recursing into groups) --
    the raw material for the size-based role fallback below."""
    acc = [] if acc is None else acc
    for sh in shapes:
        try:
            if sh.shape_type == MSO_SHAPE_TYPE.GROUP:
                _slide_sizes(sh.shapes, acc)
                continue
        except Exception:
            pass
        try:
            if sh.has_text_frame and sh.text_frame.text.strip():
                sz = _max_run_size(sh)
                if sz:
                    acc.append(sz)
        except Exception:
            pass
    return acc


def _role(shape, size_ranks=None):
    """'head' | 'sub' | 'body' — how a text shape should be styled."""
    name = (shape.name or "").lower()
    try:
        if shape.is_placeholder:
            t = shape.placeholder_format.type
            if t in _HEAD_PH:
                return "head"
            if t == PP_PLACEHOLDER.SUBTITLE:
                return "sub"
    except Exception:
        pass
    if "subtitle" in name:
        return "sub"
    if "title" in name:
        return "head"
    # FALLBACK for decks with no native placeholders (common for decks built
    # outside PowerPoint's template system -- everything is a generic "Text N"
    # box). Conservative by design: a slide with no real heading (all similar,
    # small sizes) never gets one invented, since the >=20pt threshold has to
    # be met. Deliberately NOT restricted to the single biggest run on the
    # slide -- a slide can legitimately carry a real section heading (e.g. 32pt
    # "The problem") alongside separately-large stat numbers (e.g. 36pt "100%")
    # that outrank it; both are heading-scale text and both should get the
    # heading treatment, not just whichever happens to be a point or two bigger.
    if size_ranks:
        sz = _max_run_size(shape)
        distinct = sorted(set(size_ranks), reverse=True)
        if sz is not None and distinct:
            if sz >= 20:
                return "head"
            if len(distinct) > 1 and sz == distinct[1] and 13 <= sz < 20:
                return "sub"
    return "body"


def _restyle_shapes(shapes, size_ranks=None, slide_w=None, slide_h=None, dark_bg=False):
    for sh in shapes:
        try:
            if sh.shape_type == MSO_SHAPE_TYPE.GROUP:
                _restyle_shapes(sh.shapes, size_ranks, slide_w, slide_h, dark_bg)
                continue
        except Exception:
            pass
        try:
            if sh.has_table:                       # tables: J2W palette, data untouched
                _restyle_table(sh.table)
                continue
        except Exception:
            pass
        try:
            has_text = bool(sh.has_text_frame and sh.text_frame.text.strip())
        except Exception:
            has_text = False
        if not has_text:
            # a solid-filled shape with no text of its own -- a card container
            # or icon/badge holder (never touched before; every such shape was
            # simply skipped, so uploaded "boxes" kept their original colours
            # even though every other J2W surface uses the same card/badge look)
            _restyle_box(sh, slide_w, slide_h)
            continue
        role = _role(sh, size_ranks)
        font = HEAD_FONT if role == "head" else SUB_FONT if role == "sub" else BODY_FONT
        if role == "body":
            try:
                sh.text_frame.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
            except Exception:
                pass
        for para in sh.text_frame.paragraphs:
            for run in para.runs:
                orig_size = run.font.size.pt if run.font.size else None   # BEFORE we overwrite it
                run.font.name = font
                if role == "head":
                    run.font.size = Pt(SZ_HEAD)
                    run.text = run.text.upper()      # owner's spec: headings are always caps
                    _force_color(run, INK)            # owner's spec: headings are always black
                elif role == "sub":
                    run.font.size = Pt(SZ_SUB)
                    run.text = _title_case(run.text)  # owner's spec: subheadings are Title Case
                    _force_color(run, TEAL)            # owner's spec: subheadings are always teal
                else:
                    # a card can nest its own bold mini-heading above its paragraph
                    # (e.g. "Share an open innovation challenge" sitting above a
                    # description) -- that gets the bigger SZ_BODY_HEAD size, plain
                    # paragraph/list text gets SZ_BODY. A size floor on the ORIGINAL
                    # run size keeps a tiny bold caption label (e.g. "FOR YOUR
                    # CONFIRMATION") from also being bumped to the bigger size in a
                    # box that was never sized to hold it (see SZ_BODY_HEAD_MIN_ORIG).
                    is_mini_head = bool(run.font.bold) and (
                        orig_size is None or orig_size >= SZ_BODY_HEAD_MIN_ORIG)
                    run.font.size = Pt(SZ_BODY_HEAD if is_mini_head else SZ_BODY)
                    _ink_if_dark_bg(run, dark_bg)
        if role == "body":
            _fit_body_text(sh)


_AVG_CHAR_WIDTH_EM = 0.50    # Raleway's rough average glyph width as a fraction
                             # of its point size -- the same estimate PowerPoint's
                             # own autofit uses internally, good enough to size a
                             # shrink from, not meant to be typographically exact.
_LINE_HEIGHT_EM = 1.22
_FIT_SCALE_FLOOR = 0.6       # never shrink body text below 60% of its set size --
                             # a badly-undersized box should still read as "too
                             # much text, shrunk a lot" rather than vanish entirely.


def _fit_body_text(sh):
    """auto_size = TEXT_TO_FIT_SHAPE (set just before this runs) only marks a
    shape as "should shrink text to fit" -- python-pptx does not compute the
    actual shrink, and testing against a real reskinned deck showed the shrink
    doesn't reliably get (re)computed by every renderer either: text can render
    at its full literal size and visibly overflow into the next shape below it
    -- the exact bug the owner reported (screenshot: a card's body paragraph
    overlapping its own "FOR YOUR CONFIRMATION" label and the checklist below
    it). This estimates the wrapped line count from character count vs. the
    box's own width (same rough heuristic PowerPoint's autofit uses) and, if
    that's taller than the box, bakes a real fontScale into the XML -- so the
    shrink is part of the saved file, not dependent on a renderer recalculating
    it live."""
    try:
        w_in = (sh.width or 0) / 914400
        h_in = (sh.height or 0) / 914400
        tf = sh.text_frame
        if w_in <= 0 or h_in <= 0:
            return
        avail_w = max(0.1, w_in - (tf.margin_left + tf.margin_right) / 914400)
        avail_h = max(0.1, h_in - (tf.margin_top + tf.margin_bottom) / 914400)
    except Exception:
        return
    total_lines = 0.0
    line_h_in = 0.0
    for para in tf.paragraphs:
        text = "".join(r.text for r in para.runs)
        sizes = [r.font.size.pt for r in para.runs if r.font.size]
        size = max(sizes) if sizes else SZ_BODY
        line_h_in = _LINE_HEIGHT_EM * size / 72.0
        if not text.strip():
            total_lines += 1
            continue
        char_w_in = _AVG_CHAR_WIDTH_EM * size / 72.0
        chars_per_line = max(1, int(avail_w / char_w_in))
        total_lines += -(-len(text) // chars_per_line)   # ceil division
    if total_lines <= 0 or line_h_in <= 0:
        return
    required_h = total_lines * line_h_in
    if required_h <= avail_h:
        return
    scale = max(_FIT_SCALE_FLOOR, avail_h / required_h)
    _set_font_scale(sh, scale)


def _set_font_scale(sh, scale):
    """Write a real, explicit normAutofit fontScale (+ a matching line-spacing
    reduction, same as PowerPoint's own autofit does) so the shrink is baked
    into the file rather than left for a renderer to (maybe) recompute."""
    try:
        from pptx.oxml.ns import qn
        bodyPr = sh.text_frame._txBody.bodyPr
        norm = bodyPr.find(qn("a:normAutofit"))
        if norm is None:
            norm = bodyPr.makeelement(qn("a:normAutofit"), {})
            bodyPr.append(norm)
        norm.set("fontScale", str(int(round(scale * 100000))))
        norm.set("lnSpcReduction", str(int(round(min(0.2, 1 - scale) * 100000))))
    except Exception:
        pass


def _restyle_box(sh, slide_w, slide_h):
    """A solid-filled shape with no text of its own -- either a card container
    or a small icon/badge holder, judged by size. Recoloured to our own card/
    badge convention (see module docstring). A near-full-slide shape is treated
    as a fake background (a common export pattern) and forced to plain white,
    same as the slide's own real background -- the owner's rule is the
    background is white no matter what, not just left at its original colour."""
    try:
        if sh.fill.type is None:
            return
    except Exception:
        return
    w = (sh.width or 0) / 914400
    h = (sh.height or 0) / 914400
    if w <= 0 or h <= 0:
        return
    if slide_w and slide_h:
        sw_in, sh_in = slide_w / 914400, slide_h / 914400
        if w >= sw_in * 0.92 and h >= sh_in * 0.92:
            try:
                sh.fill.solid()
                sh.fill.fore_color.rgb = WHITE
                sh.line.fill.background()
                sh.shadow.inherit = False
            except Exception:
                pass
            return
    is_badge = w <= 1.2 and h <= 1.2 and abs(w - h) < 0.35
    try:
        sh.fill.solid()
        if is_badge:
            sh.fill.fore_color.rgb = RED           # matches our ranked-badge style
            sh.line.fill.background()
        else:
            sh.fill.fore_color.rgb = CARD_BG       # matches our capability-card style
            sh.line.color.rgb = CARD_LINE
            sh.line.width = Pt(0.75)
        sh.shadow.inherit = False
    except Exception:
        pass


def _rect_of(sh):
    try:
        l, t, w, h = sh.left, sh.top, sh.width, sh.height
        if None in (l, t, w, h):
            return None
        return (l, t, l + w, t + h)
    except Exception:
        return None


def _overlaps(a, b, thresh=0.5):
    """True if >= thresh of rect a's area is covered by rect b."""
    ax0, ay0, ax1, ay1 = a
    bx0, by0, bx1, by1 = b
    iw = max(0, min(ax1, bx1) - max(ax0, bx0))
    ih = max(0, min(ay1, by1) - max(ay0, by0))
    area_a = max(1, (ax1 - ax0) * (ay1 - ay0))
    return (iw * ih) / area_a >= thresh


def _fix_text_on_box_contrast(shapes):
    """Second pass, run AFTER box fills and role text-colours are both already
    set: a text shape can visually sit on top of a coloured box (e.g. a number
    label centred over a badge circle -- a very common step/stat design, and
    exactly the bug the owner flagged: a red badge with its number forced teal
    by the subheading rule), and that combination needs to be checked against
    what's actually behind it, not just the text's own role colour. Any run
    whose contrast against its underlying box falls below _MIN_CONTRAST is
    flipped to a legible replacement (see _legible_replacement above). Top-
    level shapes only: a shape inside a group has coordinates relative to the
    group, not the slide, so overlap math against top-level box rects would be
    meaningless for it."""
    boxes = []
    for sh in shapes:
        try:
            if sh.shape_type == MSO_SHAPE_TYPE.GROUP:
                continue
            if sh.has_text_frame and sh.text_frame.text.strip():
                continue
            if sh.fill.type is None:
                continue
            fc = sh.fill.fore_color
            if fc.type != MSO_COLOR_TYPE.RGB:
                continue
        except Exception:
            continue
        rect = _rect_of(sh)
        if rect:
            boxes.append((rect, tuple(fc.rgb)))
    if not boxes:
        return
    for sh in shapes:
        try:
            if sh.shape_type == MSO_SHAPE_TYPE.GROUP:
                continue
            if not (sh.has_text_frame and sh.text_frame.text.strip()):
                continue
        except Exception:
            continue
        rect = _rect_of(sh)
        if not rect:
            continue
        # a text shape can overlap MULTIPLE boxes at once -- e.g. a small badge
        # nested inside a larger card behind it -- so take the SMALLEST-area
        # match (the immediately-underlying shape), not just the first one in
        # document order, which would pick the outer card instead of the badge
        # actually sitting under the text.
        covering = [(rgb, (brect[2] - brect[0]) * (brect[3] - brect[1]))
                    for brect, rgb in boxes if _overlaps(rect, brect)]
        if not covering:
            continue
        box_rgb = min(covering, key=lambda c: c[1])[0]
        for para in sh.text_frame.paragraphs:
            for run in para.runs:
                try:
                    col = run.font.color
                    cur_rgb = tuple(col.rgb) if col.type == MSO_COLOR_TYPE.RGB else None
                except Exception:
                    cur_rgb = None
                if cur_rgb is None:
                    continue
                try:
                    if _contrast_ratio(box_rgb, cur_rgb) < _MIN_CONTRAST:
                        run.font.color.rgb = _legible_replacement(box_rgb)
                except Exception:
                    pass


def _recolor_picture(sh, target_rgb):
    """Recolour a picture's own glyph pixels to target_rgb, keeping each
    pixel's original alpha untouched. Icon images in a real uploaded deck
    (owner's file, 2026-07-08) were confirmed to be single-colour glyphs on a
    transparent PNG background -- a uniform new RGB with the SAME alpha mask
    reproduces the identical glyph shape in a new colour, not a solid blob.
    A shape's fill can't reach into embedded image pixels, so this is the one
    place in the module that edits raw image bytes rather than an XML
    property. Best-effort: a picture that isn't a plain RGBA glyph (rare) is
    left untouched rather than risking a corrupted image."""
    try:
        from PIL import Image
        import io as _io
        im = Image.open(_io.BytesIO(sh.image.blob)).convert("RGBA")
        # Recolour ONLY a monochrome glyph on a transparent background. A photo or any
        # OPAQUE image (or a colourful logo) has no transparent backdrop, so painting every
        # pixel one colour turns the whole rectangle into a solid block -- the 'black box'
        # bug on real uploaded decks (owner-reported, 2026-07-13). Leave those as-authored.
        small = im.copy()
        small.thumbnail((48, 48))
        px = list(small.getdata())
        if not px:
            return
        opaque = [p for p in px if p[3] > 200]
        if len(opaque) / len(px) > 0.70:
            return                       # mostly opaque -> a picture, not a cut-out glyph
        if opaque:
            chroma = sum(max(p[:3]) - min(p[:3]) for p in opaque) / len(opaque)
            if chroma > 60:
                return                   # multi-colour art -> recolouring would destroy it
        alpha = im.split()[3]
        solid = Image.new("RGBA", im.size, tuple(target_rgb) + (255,))
        solid.putalpha(alpha)
        buf = _io.BytesIO()
        solid.save(buf, format="PNG")
        _img_part, rId = sh.part.get_or_add_image_part(buf)
        sh._element.blipFill.blip.rEmbed = rId
    except Exception:
        pass


def _fix_picture_contrast(shapes):
    """Icon PICTURES get the same contrast rule as text (see
    _fix_text_on_box_contrast) -- white on a red/teal badge, ink otherwise --
    but a picture's colour is baked into its pixels, not a settable property,
    so this recolours the glyph itself (_recolor_picture) instead of a font
    colour. Must run AFTER _restyle_box, so the box colour checked behind each
    icon is the NEW one, not the uploaded deck's original. A picture with no
    box behind it sits directly on the (about-to-be-forced-white) slide
    background, so it defaults to ink -- owner's spec, 2026-07-08: 'every page
    is on a white background' so an icon with nothing coloured behind it
    should read black. Top-level shapes only, same reasoning as the text
    version (a grouped shape's coordinates aren't slide-relative)."""
    boxes = []
    for sh in shapes:
        try:
            if sh.shape_type == MSO_SHAPE_TYPE.GROUP:
                continue
            if sh.shape_type == MSO_SHAPE_TYPE.PICTURE:
                continue
            if sh.has_text_frame and sh.text_frame.text.strip():
                continue
            if sh.fill.type is None:
                continue
            fc = sh.fill.fore_color
            if fc.type != MSO_COLOR_TYPE.RGB:
                continue
        except Exception:
            continue
        rect = _rect_of(sh)
        if rect:
            boxes.append((rect, tuple(fc.rgb)))
    for sh in shapes:
        try:
            if sh.shape_type != MSO_SHAPE_TYPE.PICTURE:
                continue
        except Exception:
            continue
        rect = _rect_of(sh)
        if not rect:
            continue
        covering = [(rgb, (brect[2] - brect[0]) * (brect[3] - brect[1]))
                    for brect, rgb in boxes if _overlaps(rect, brect)]
        box_rgb = min(covering, key=lambda c: c[1])[0] if covering else (255, 255, 255)
        _recolor_picture(sh, _legible_replacement(box_rgb))


def _resolve_bottom_banner_overlap(shapes, slide_w, slide_h, margin=Inches(0.05)):
    """A closing/summary banner near the bottom of a slide can overlap the
    last content row above it -- confirmed on a real uploaded deck, 2026-07-08
    (a 'Command Centre' banner's own top sat above the last list-row's bottom
    edge by ~0.27in, baked into the SOURCE file's own shape coordinates --
    nothing to do with reskinning or font size, since body text is already
    forced small). Deliberately narrow: only shapes at least half the slide's
    width (sequential row/banner content) anchored in the bottom 40% of the
    slide are candidates, so small icon/badge shapes nested inside their own
    card (which routinely overlap their card by DESIGN) are never touched."""
    top_level = [(sh, _rect_of(sh)) for sh in shapes]
    top_level = [(sh, r) for sh, r in top_level if r is not None]
    wide = [(sh, r) for sh, r in top_level if (r[2] - r[0]) >= 0.5 * slide_w]
    wide.sort(key=lambda t: t[1][1])   # by top

    # Cluster into BANDS first: two shapes authored at (near) the same top are
    # ONE visual element -- e.g. a coloured bar plus its own text sitting on
    # top of it, both at an identical position -- never a sequential-content
    # pair to check for overlap against each other. Comparing raw shapes
    # (rather than bands) caused exactly that: a banner's own fill-shape and
    # its text label "overlapped" each other and each triggered its own
    # shift, double-moving the banner (caught in real testing, 2026-07-08).
    bands = []   # [(top, bottom, left, right, [shapes])]
    for sh, r in wide:
        if bands and abs(r[1] - bands[-1][0]) <= Inches(0.02):
            top, bottom, left, right, band_shapes = bands[-1]
            bands[-1] = (top, max(bottom, r[3]), min(left, r[0]), max(right, r[2]),
                        band_shapes + [sh])
        else:
            bands.append((r[1], r[3], r[0], r[2], [sh]))

    for i in range(1, len(bands)):
        b_top, b_bottom, b_left, b_right, b_shapes = bands[i]
        if b_top < 0.6 * slide_h:
            continue   # not a trailing/closing element -- leave it alone
        a_top, a_bottom, a_left, a_right, _a_shapes = bands[i - 1]
        if a_bottom <= b_top:
            continue   # no real overlap
        # horizontal overlap check -- only shift if they actually collide, not
        # just happen to be stacked with unrelated x-ranges
        ox = min(a_right, b_right) - max(a_left, b_left)
        if ox <= 0:
            continue
        delta = (a_bottom - b_top) + margin
        for sh in b_shapes:
            if sh.top is not None:
                sh.top = sh.top + delta


def _restyle_table(table):
    """Recolour a table to the J2W palette (teal header, striped body, Raleway text)
    WITHOUT touching any cell text — only fills and fonts change."""
    for ri, row in enumerate(table.rows):
        header = (ri == 0)
        for cell in row.cells:
            try:
                cell.fill.solid()
                cell.fill.fore_color.rgb = TEAL if header else (SOFT if ri % 2 == 0 else WHITE)
            except Exception:
                pass
            for para in cell.text_frame.paragraphs:
                for run in para.runs:
                    run.font.name = BODY_FONT
                    try:
                        run.font.color.rgb = WHITE if header else INK
                        if header:
                            run.font.bold = True
                    except Exception:
                        pass


def _force_color(run, rgb):
    """Set a run's colour unconditionally (used for heading/subheading, which are
    always one fixed colour regardless of what the original deck authored)."""
    try:
        run.font.color.rgb = rgb
    except Exception:
        pass


def _is_dark_rgb(rgb):
    try:
        return (rgb[0] + rgb[1] + rgb[2]) < 460
    except Exception:
        return False


def _ink_if_dark_bg(run, dark_bg):
    """Body text only: if this slide's background used to be dark, any run that
    was authored light (or inherits a theme colour, which on a dark deck almost
    always resolves light) is flipped to ink so it stays readable once the
    background is forced to white. Left untouched on an originally light-
    background slide, where the original body colour is already readable."""
    if not dark_bg:
        return
    try:
        col = run.font.color
        if col.type == MSO_COLOR_TYPE.RGB:
            if not _is_dark_rgb(col.rgb):
                run.font.color.rgb = INK
        else:
            run.font.color.rgb = INK
    except Exception:
        pass


def _slide_is_dark(slide, shapes, slide_w, slide_h):
    """True if this slide's ORIGINAL background reads as dark -- checked before
    any restyling touches it, since the background is about to be forced white
    and any light/white text authored for a dark background would otherwise
    become invisible. Checks both the real pptx slide-background property and a
    full-bleed shape standing in for one (a common export pattern)."""
    try:
        fill = slide.background.fill
        if fill.type == MSO_FILL_TYPE.SOLID and fill.fore_color.type == MSO_COLOR_TYPE.RGB:
            if _is_dark_rgb(fill.fore_color.rgb):
                return True
    except Exception:
        pass
    if slide_w and slide_h:
        sw_in, sh_in = slide_w / 914400, slide_h / 914400
        for sh in shapes:
            try:
                w = (sh.width or 0) / 914400
                h = (sh.height or 0) / 914400
                if w >= sw_in * 0.92 and h >= sh_in * 0.92 and sh.fill.type == MSO_FILL_TYPE.SOLID \
                        and sh.fill.fore_color.type == MSO_COLOR_TYPE.RGB \
                        and _is_dark_rgb(sh.fill.fore_color.rgb):
                    return True
            except Exception:
                pass
    return False


def _whiten_slide_background(slide):
    """Force the slide's own background to plain white -- owner's rule, no
    matter what the original deck used."""
    try:
        slide.background.fill.solid()
        slide.background.fill.fore_color.rgb = WHITE
    except Exception:
        pass


_WORD_RE = re.compile(r"([^A-Za-z0-9']+)")


def _title_case(text):
    """Each word's first letter capitalised, rest lower-cased -- except an
    all-caps token of 2+ letters (an acronym: AI, ROI, ERP, J2W...) is left as-is
    rather than being flattened to one capital letter."""
    def cap(tok):
        if len(tok) > 1 and tok.isupper():
            return tok
        if tok[:1].isalpha():
            return tok[0].upper() + tok[1:].lower()
        return tok
    parts = _WORD_RE.split(text)
    return "".join(cap(p) if p else p for p in parts)


def _first_heading_shape(shapes, size_ranks):
    """The slide's primary heading -- the first top-level text shape (in document
    order) that resolves to the 'head' role. Deliberately the FIRST one, not
    every shape that happens to be >=20pt (a slide can carry a real heading
    alongside separately-large stat numbers -- see _role -- and only the actual
    heading should get the accent bar, not every big number). Grouped shapes are
    skipped: their coordinates live in the group's own child coordinate space,
    not the slide's, so a bar placed using them directly would be misplaced."""
    for sh in shapes:
        try:
            if sh.shape_type == MSO_SHAPE_TYPE.GROUP:
                continue
        except Exception:
            pass
        try:
            if sh.has_text_frame and sh.text_frame.text.strip() and _role(sh, size_ranks) == "head":
                return sh
        except Exception:
            pass
    return None


def _add_heading_bar(slide, heading_shape):
    """Small red vertical accent bar to the left of the slide's heading (owner's
    spec) -- the same accent-bar convention used on every other J2W template."""
    try:
        left, top, height = heading_shape.left, heading_shape.top, heading_shape.height
    except Exception:
        return
    if left is None or top is None:
        return
    height = height or Inches(0.5)
    bar_x = max(0, left - Inches(0.15))
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, bar_x, top, Inches(0.065), height)
    _flat_fill(bar, RED)


def _brand_slide(slide, sw, sh):
    """Add the J2W overlay: the same top split bar every other J2W template uses
    -- red covers the first ~2/3, teal the last ~1/3, both very thin (0.083in,
    matching the master deck's own title-slide bar exactly). The split is
    PROPORTIONAL (not a fixed inch count) so it looks identical whether the
    uploaded deck is our standard 13.33in-wide or a different size (this deck
    is 10in) -- matters here since a fixed segment width would read as a
    different proportion on a narrower slide. No bottom-right wordmark (owner's
    spec, 2026-07-08: the top bar alone is the brand mark on a reskinned slide)."""
    bar_h = Inches(0.083)
    split = int(sw * (8.89 / 13.33))
    red = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, split, bar_h)
    _flat_fill(red, RED)
    teal = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, split, 0, sw - split, bar_h)
    _flat_fill(teal, TEAL)


def _flat_fill(shape, rgb):
    shape.fill.solid()
    shape.fill.fore_color.rgb = rgb
    shape.line.fill.background()
    try:
        shape.shadow.inherit = False
    except Exception:
        pass


# ── 1b) bookend every reskinned deck with OUR OWN title + closing slide ───────
def _append_bookend_slides(prs, dest_sw, dest_sh):
    """Prepend the master deck's own title slide (CS01) and append its closing
    'Let's win together' slide (CS08) -- owner's spec: every reskinned deck
    should open and close on our own slides, not the uploaded deck's own cover.
    Pulled verbatim via the same slide-copy machinery the rest of the app uses
    (skills slides, AI-drafted slides), then scaled to the uploaded deck's own
    canvas size so nothing is cropped or mis-sized on a non-standard aspect."""
    try:
        master = Presentation(config.MASTER_DECK)
    except Exception:
        return
    title_src = winback_src = None
    for slide in master.slides:
        sid = read_id(slide)
        if sid == "CS01":
            title_src = slide
        elif sid == "CS08":
            winback_src = slide
    if title_src is None and winback_src is None:
        return
    msw, msh = master.slide_width, master.slide_height
    sx = (dest_sw / msw) if msw else 1.0
    sy = (dest_sh / msh) if msh else 1.0
    if title_src is not None:
        new_title = _copy_slide(prs, title_src)
        _scale_shapes(new_title.shapes, sx, sy)
        _move_slide_to_front(prs, new_title)
    if winback_src is not None:
        new_winback = _copy_slide(prs, winback_src)
        _scale_shapes(new_winback.shapes, sx, sy)


def _scale_shapes(shapes, sx, sy):
    """Rescale a freshly-copied slide's TOP-LEVEL shapes to fit a differently-
    sized destination canvas. Group shapes only need their own outer left/top/
    width/height scaled -- PowerPoint renders a group's children (and their
    text) relative to that outer extent automatically. A non-group shape's box
    can be resized the same way, but its text does NOT auto-scale with the box,
    so its runs are scaled by hand too."""
    is_uniform_ish = abs(sx - sy) < 1e-6
    for sh in shapes:
        try:
            if sh.left is not None:
                sh.left = int(round(sh.left * sx))
            if sh.top is not None:
                sh.top = int(round(sh.top * sy))
            if sh.width is not None:
                sh.width = int(round(sh.width * sx))
            if sh.height is not None:
                sh.height = int(round(sh.height * sy))
        except Exception:
            pass
        try:
            is_group = sh.shape_type == MSO_SHAPE_TYPE.GROUP
        except Exception:
            is_group = False
        if is_group:
            continue     # children scale automatically with the group's own extent
        try:
            if sh.has_text_frame:
                s = sx if is_uniform_ish else (sx + sy) / 2.0
                for p in sh.text_frame.paragraphs:
                    for r in p.runs:
                        if r.font.size:
                            r.font.size = Pt(r.font.size.pt * s)
        except Exception:
            pass


def _move_slide_to_front(prs, slide):
    sld_id_lst = prs.slides._sldIdLst
    elements = list(sld_id_lst)
    slides = list(prs.slides)
    idx = slides.index(slide)
    elem = elements[idx]
    sld_id_lst.remove(elem)
    sld_id_lst.insert(0, elem)


# ── 2) render the restyled deck to page images (true preview) ─────────────────
def _on_path(name):
    from shutil import which
    return which(name) is not None


def render_pngs(pptx_path, out_dir, dpi=120):
    """Render every slide of pptx_path to a PNG in out_dir via LibreOffice (pptx->pdf)
    then poppler (pdf->png). Returns the ordered list of PNG paths, or [] if the tools
    aren't available (the caller then falls back to a download-only preview)."""
    soffice = "libreoffice" if _on_path("libreoffice") else ("soffice" if _on_path("soffice") else None)
    if not soffice or not _on_path("pdftoppm"):
        return []
    os.makedirs(out_dir, exist_ok=True)
    profile = "file://" + os.path.join(out_dir, ".loprofile")
    try:
        subprocess.run(
            [soffice, "-env:UserInstallation=" + profile, "--headless", "--convert-to",
             "pdf", "--outdir", out_dir, pptx_path],
            timeout=150, check=True, stdout=subprocess.DEVNULL, stderr=subprocess.DEVNULL)
    except Exception:
        return []
    base = os.path.splitext(os.path.basename(pptx_path))[0]
    pdf = os.path.join(out_dir, base + ".pdf")
    if not os.path.exists(pdf):
        return []
    try:
        subprocess.run(["pdftoppm", "-png", "-r", str(dpi), pdf, os.path.join(out_dir, "slide")],
                       timeout=150, check=True, stdout=subprocess.DEVNULL, stderr=subprocess.DEVNULL)
    except Exception:
        return []
    return sorted(glob.glob(os.path.join(out_dir, "slide*.png")))
