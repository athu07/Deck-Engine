# -*- coding: utf-8 -*-
"""
fill_case_study.py  --  Pour ONE content-store row into the case-study template
and save a finished, single-slide .pptx.

This is the heart of the content-referencing model: a case study lives as DATA
in case_study_content_store.json, and a real slide is built on demand by
dropping that data into case_study_v2.pptx (the J2W branded template).

What it does:
  * opens the template (1 slide, all {{MARKERS}})
  * builds a value for every marker from the content row
  * smart-splits each narrative RESULT into a headline metric + caption
    (e.g. "13 engineers added in the initial phase" -> "13" / "engineers added…")
  * splits a CAPABILITY into title/body when it carries a separator, else
    shows it as a title-only card
  * replaces markers IN PLACE, preserving the template's fonts/colours
  * any marker with no data is blanked (never left as "{{…}}" on the slide)

Use as a library:
    import fill_case_study as fcs
    fcs.fill_row(row_dict, "out.pptx", client_name="Acme Bank")

Or from the CLI (demo on one stored case):
    py fill_case_study.py                 # fills MSS001 -> output/<id>.pptx
    py fill_case_study.py AIP007 "Acme"   # fill a chosen id, with a client name
"""

import json
import math
import os
import re
import sys

from pptx import Presentation
from pptx.util import Inches

from deckengine import config

TEMPLATE = config.CASE_TEMPLATE_PPTX
STORE = config.CONTENT_STORE_JSON
OUT_DIR = config.OUTPUT_DIR

# Item 4 — heading fit. The heading reads "CASE STUDY: <title>" in Oswald 24pt
# bold across a ~12.2in usable box. If it's too long for one line it wraps to a
# second line, and we drop the CLIENT|DOMAIN subheading one line so it doesn't
# collide. (Heuristic: python-pptx can't measure rendered text, so we estimate
# characters-per-line. Oswald is condensed, hence the generous count.)
TITLE_CHARS_PER_LINE = 92    # ~chars of the prefixed title that fit on one line
SUBHEAD_LINE_DROP_IN = 0.34  # how far to move the subheading down per extra line

MARKER_RE = re.compile(r"\{\{([A-Z0-9_]+)\}\}")

# The headline metric inside a results sentence. Ordered most-specific first so
# a COMPOUND metric (a range, a ratio, an "X out of Y") is captured WHOLE and put
# on top as one unit — otherwise the caption below would repeat part of it
# (e.g. "3 out of 3 ..." must not become "3" on top + "out of 3 ..." below).
_METRIC_RE = re.compile(
    r"("
    r"\d+(?:\.\d+)?\s*(?:out\s+of|of)\s*\d+(?:\.\d+)?"   # 3 out of 3, 3 of 3
    r"|\d+\s*/\s*\d+"                                      # 3/3
    r"|\d+(?:\.\d+)?\s*:\s*\d+(?:\.\d+)?"                 # 8:1
    r"|\d+(?:\.\d+)?\s*-\s*\d+(?:\.\d+)?\s*%"             # 11-13%
    r"|\d+(?:\.\d+)?\s*%"                                  # 70%
    r"|\$\s*\d[\d,]*\.?\d*\s*[KMBkmb]?"                   # $4M
    r"|\d+(?:\.\d+)?\s*[xX]\s*\d+"                        # 24x7
    r"|\d+(?:\.\d+)?[xX]\b"                               # 2x, 3.5x
    r"|\d+(?:\.\d+)?\s*-\s*\d+(?:\.\d+)?"                 # 6-9 (range; unit added below)
    r"|\d+(?:\.\d+)?\s*to\s*\d+(?:\.\d+)?"               # 3.0 to 4.0, 10 to 5
    r"|\d[\d,]*\+"                                         # 40+
    r"|\b\d[\d,]*(?:\.\d+)?\b"                            # 13, 1,200, 3.0
    r")"
)

# a qualifier immediately BEFORE the number belongs WITH it ("under 5", ">8:1",
# "Rs 2 crore")
_QUALIFIER_RE = re.compile(
    r"(?:>=?|<=?|~|₹)\s*$"
    r"|(?:\brs\.?|\binr)\s*$"
    r"|\b(?:up\s+to|greater\s+than|less\s+than|more\s+than|at\s+least|"
    r"under|over|nearly|almost|about|around|approx\.?|approximately)\s+$",
    re.I,
)

# a unit word immediately AFTER the number belongs WITH it ("5 minutes", "6-9
# points", "2 crore")
_TRAIL_UNIT_RE = re.compile(
    r"\s*(?:percentage\s+points?|points?|pts?|minutes?|mins?|hours?|hrs?|"
    r"days?|weeks?|months?|years?|yrs?|crores?|lakhs?|lacs?)\b",
    re.I,
)

# an "or more"/"+" tail belongs WITH the number ("40 or more", "90%+")
# (the \b binds only to the word branch — a bare "+" has no following word boundary)
_TRAIL_MORE_RE = re.compile(
    r"\s*(?:\+|or\s+(?:more|less|fewer|higher|lower|above|below|greater)\b)",
    re.I,
)

_CAP_SEPS = [" — ", " – ", " - ", ": "]

# connectives to trim off the edges of a caption once the highlight is pulled
_LEAD_RE = re.compile(r"^(?:of|by|in|to|the|a|an|and|with|via|for|on|at|from)\s+", re.I)
_TRAIL_RE = re.compile(r"\s+(?:of|by|in|to|the|a|an|and|with|via|for|on|at|from)$", re.I)

# CONFIDENTIALITY: a case-study slide never names a real client/company — only
# J2W is named anywhere. The CLIENT line shows a generic anonymised descriptor,
# matching J2W's own deck convention ("Leading Manufacturing Institution").
_CLIENT_SUFFIX = {
    "BFSI": "Institution", "PRIVATE_EQUITY": "Firm", "AVIATION": "Operator",
    "ENERGY": "Utility", "TELECOM": "Operator",
}


def anon_client(row):
    """A generic, name-free client descriptor derived from the domain.
    Never returns a real company name."""
    domain = (row.get("domain") or "").strip()
    if not domain:
        return "Global Enterprise"
    suffix = _CLIENT_SUFFIX.get(row.get("industry"), "Enterprise")
    # don't echo a word the domain already carries ("Leading Retail Retailer")
    domain_words = {w.lower() for w in re.findall(r"[A-Za-z]+", domain)}
    if suffix.lower() in domain_words:
        suffix = "Enterprise"
    if suffix.lower() in domain_words:        # domain literally contains "Enterprise"
        return f"Leading {domain}"
    return f"Leading {domain} {suffix}"


def split_result(sentence):
    """Return (highlight, caption) for one result.

    The highlight (shown big + bold on top) is the WHOLE meaningful metric —
    a complete number expression if there is one (70%, 11-13%, 3 out of 3, 8:1,
    Under 5 minutes, $4M), otherwise the leading key phrase ("Automated"). The
    ENTIRE metric span is removed from the sentence, so the caption below carries
    only the descriptive rest — never a piece of the number repeated."""
    s = (sentence or "").strip().rstrip(".")
    if not s:
        return "", ""

    m = _METRIC_RE.search(s)
    if m:
        start, end = m.start(1), m.end(1)
        # pull a leading qualifier ("under", "greater than", ">") into the metric
        qm = _QUALIFIER_RE.search(s[:start])
        if qm:
            start = qm.start()
        # pull a trailing unit word ("minutes", "points") into the metric
        um = _TRAIL_UNIT_RE.match(s[end:])
        if um:
            end += um.end()
        # pull a trailing "or more" / "+" into the metric
        mm = _TRAIL_MORE_RE.match(s[end:])
        if mm:
            end += mm.end()
        top = s[start:end].strip()
        # drop a connective that merely introduced the metric ("...improved from"
        # 3.0 to 4.0 -> caption shouldn't keep the dangling "from")
        left = re.sub(r"\s+(?:from|to|of|by|between)\s*$", "", s[:start], flags=re.I)
        bottom = left + " " + s[end:]
        # tidy the highlight: no space before %, no space after $
        top = re.sub(r"\s+%", "%", top)
        top = re.sub(r"\$\s+", "$", top)
        top = re.sub(r"\s{2,}", " ", top).strip()
    else:
        words = s.split()
        first = words[0]
        # short modifier (Single/Zero/Real) reads better paired with the next word
        if len(first) <= 7 and len(words) > 1:
            top = first + " " + words[1]
            bottom = " ".join(words[2:])
        else:
            top = first
            bottom = " ".join(words[1:])
        top = top[0].upper() + top[1:]

    bottom = _LEAD_RE.sub("", bottom.strip(" ,–—-"))
    bottom = _TRAIL_RE.sub("", bottom).strip(" ,–—-")
    bottom = re.sub(r"\s{2,}", " ", bottom).strip()
    return top, bottom


def split_capability(cap):
    """Normalise one capability to (title, body).
    Accepts a {title, body} dict (enriched store) or a plain string (splits on
    a separator if present, else title-only)."""
    if isinstance(cap, dict):
        return (cap.get("title") or "").strip(), (cap.get("body") or "").strip()
    c = (cap or "").strip()
    for sep in _CAP_SEPS:
        if sep in c:
            t, b = c.split(sep, 1)
            return t.strip(), b.strip()
    return c, ""


def _pad(seq, n):
    seq = list(seq)[:n]
    return seq + [""] * (n - len(seq))


def build_mapping(row):
    """Marker -> replacement string, for one content-store row."""
    caps = _pad(row.get("capabilities", []), 6)
    results = _pad(row.get("results", []), 3)

    mapping = {
        "TITLE": row.get("title", ""),
        "CLIENT": anon_client(row),
        "DOMAIN": row.get("domain", "") or "-",
        "CHALLENGE": row.get("challenge", ""),
        "SOLUTION": row.get("solution", ""),
    }
    for i, cap in enumerate(caps, 1):
        t, b = split_capability(cap)
        mapping[f"CAP_{i}_TITLE"] = t
        mapping[f"CAP_{i}_BODY"] = b
    for i, res in enumerate(results, 1):
        pct, txt = split_result(res)
        mapping[f"RESULT_{i}_PCT"] = pct
        mapping[f"RESULT_{i}_TEXT"] = txt
    # owner rule: no em/en dash ever reaches a slide — always a plain hyphen
    return {k: (v or "").replace("—", "-").replace("–", "-") for k, v in mapping.items()}


def _apply(text, mapping):
    """Replace every {{MARKER}}; unknown markers become empty string."""
    return MARKER_RE.sub(lambda m: mapping.get(m.group(1), ""), text)


def _reflow_subhead(subhead_shape, full_title):
    """Item 4: if the prefixed heading wraps past one line, move the CLIENT|DOMAIN
    subheading down so it clears the extra title line(s)."""
    if subhead_shape is None:
        return
    lines = max(1, math.ceil(len(full_title) / TITLE_CHARS_PER_LINE))
    extra = min(lines - 1, 1)      # cap the drop at one line (box budget)
    if extra:
        subhead_shape.top = subhead_shape.top + Inches(SUBHEAD_LINE_DROP_IN * extra)


def fill_row(row, out_path, template=TEMPLATE):
    """Build a finished slide from one content row. Returns out_path.

    Note: there is deliberately no client-name parameter — a case-study slide
    never carries a real client/company name (only J2W is named)."""
    prs = Presentation(template)
    slide = prs.slides[0]
    mapping = build_mapping(row)

    subhead_shape = None
    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue
        if "CLIENT:" in shape.text_frame.text:
            subhead_shape = shape
        for para in shape.text_frame.paragraphs:
            for run in para.runs:
                if "{{" in run.text:
                    run.text = _apply(run.text, mapping)

    # item 4: drop the subheading a line when the heading wraps
    _reflow_subhead(subhead_shape, "CASE STUDY: " + mapping.get("TITLE", ""))

    out_dir = os.path.dirname(out_path)
    if out_dir:
        os.makedirs(out_dir, exist_ok=True)
    prs.save(out_path)
    return out_path, mapping


def _load_store():
    with open(STORE, encoding="utf-8") as f:
        return json.load(f)


def fill_by_id(case_id, out_path=None):
    recs = {r["id"]: r for r in _load_store()}
    if case_id not in recs:
        raise KeyError(f"{case_id} not in {STORE}")
    out_path = out_path or os.path.join(OUT_DIR, f"{case_id}.pptx")
    return fill_row(recs[case_id], out_path)


if __name__ == "__main__":
    try:                                   # so the ✓ stat prints on a cp1252 console
        sys.stdout.reconfigure(encoding="utf-8")
    except (AttributeError, ValueError):
        pass
    cid = sys.argv[1] if len(sys.argv) > 1 else "MSS001"
    out, mapping = fill_by_id(cid)

    print(f"Filled {cid} -> {out}")
    print(f"Client: {mapping['CLIENT']}   Domain: {mapping['DOMAIN']}")
    print(f"Title : {mapping['TITLE']}")
    print("\nCapabilities:")
    for i in range(1, 7):
        t, b = mapping[f"CAP_{i}_TITLE"], mapping[f"CAP_{i}_BODY"]
        print(f"  {i}. {t}" + (f"  —  {b}" if b else ""))
    print("\nResults:")
    for i in range(1, 4):
        print(f"  {mapping[f'RESULT_{i}_PCT']:>6}  {mapping[f'RESULT_{i}_TEXT']}")

    # safety: no leftover markers anywhere on the slide
    leftover = [v for v in mapping.values() if "{{" in str(v)]
    print("\nLeftover markers in values:", leftover or "none")
