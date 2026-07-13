# -*- coding: utf-8 -*-
"""
slide_generator.py  --  Step 04: turn a "needs to be created" gap into a real
slide, built from a TEMPLATE.

Design = pluggable templates:
  - Templates live in templates.pptx. Each template is one slide, tagged in its
    notes with  J2W_TEMPLATE: <name>  (e.g. case_study), and contains marker
    tokens in its text:  {{TITLE}}  {{KEYWORDS}}  {{BULLETS}}
  - To generate a slide: the AI writes the words, then we COPY the chosen
    template slide into the deck and REPLACE the markers with those words.
  - Add another template later = add another tagged slide to templates.pptx.
    No code change. Swap the temporary template for the real J2W design anytime.

The template made by create_temp_template() is a PLACEHOLDER — plain text boxes,
text-only (images in a template need extra work). Replace it with the real
J2W-designed template slide when ready; keep the same marker tokens + tag.
"""

import copy
import json

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

from deckengine import config
from deckengine.services.content import editor
from deckengine.services.content.build_library import read_id  # noqa: F401  (kept for parity / future use)

MASTER = config.MASTER_DECK
TEMPLATES_FILE = config.TEMPLATES_PPTX
MODEL = "gpt-4o-mini"
TEMPLATE_TAG = "J2W_TEMPLATE:"


# --------------------------------------------------------------------------- #
# Template file
# --------------------------------------------------------------------------- #
def create_temp_template(path=TEMPLATES_FILE):
    """Create a temporary, text-only 'case_study' template slide."""
    master = Presentation(MASTER)
    tp = Presentation()
    tp.slide_width = master.slide_width
    tp.slide_height = master.slide_height
    slide = tp.slides.add_slide(tp.slide_layouts[6])      # 6 = Blank
    slide.notes_slide.notes_text_frame.text = f"{TEMPLATE_TAG} case_study"

    def textbox(left, top, width, height, marker, size, color, bold=False, italic=False):
        tb = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
        p = tb.text_frame.paragraphs[0]
        run = p.add_run()
        run.text = marker
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.italic = italic
        run.font.color.rgb = color
        return tb

    textbox(0.6, 0.5, 12.0, 1.1, "{{TITLE}}", 28, RGBColor(0x14, 0x28, 0x50), bold=True)
    textbox(0.6, 1.6, 12.0, 0.5, "{{KEYWORDS}}", 13, RGBColor(0x0F, 0x6E, 0x56), italic=True)
    textbox(0.6, 2.3, 12.0, 4.0, "{{BULLETS}}", 16, RGBColor(0x33, 0x33, 0x33))
    textbox(0.6, 6.9, 12.0, 0.4, "Generated slide - replace with the real J2W template",
            10, RGBColor(0xAA, 0xAA, 0xAA), italic=True)
    tp.save(path)
    return path


def list_templates(path=TEMPLATES_FILE):
    """{template_name: slide} for every tagged slide in templates.pptx."""
    prs = Presentation(path)
    out = {}
    for slide in prs.slides:
        if slide.has_notes_slide:
            txt = slide.notes_slide.notes_text_frame.text or ""
            for line in txt.splitlines():
                if line.strip().startswith(TEMPLATE_TAG):
                    out[line.split(":", 1)[1].strip()] = slide
    return out


# --------------------------------------------------------------------------- #
# AI content
# --------------------------------------------------------------------------- #
def default_brief(work_type, industry, transcript, topic=""):
    """A starting 'what should this slide cover' brief, pre-filled from the notes.
    The salesperson edits it before generating. No AI call.
    If `topic` is given (a specific client ask with no slide, e.g. 'ADAS'), the
    brief is centred on demonstrating J2W's capability in that exact topic."""
    labels = {"WORKFORCE": "Workforce", "AI_POD": "AI Pod", "MS": "Managed Services"}
    t = (transcript or "").strip()
    snippet = (t[:180].rstrip() + "…") if len(t) > 180 else t
    topic = (topic or "").strip()
    if topic:
        base = "Slide demonstrating J2W's capability in %s" % topic
        if industry:
            base += " for the %s industry" % industry.replace("_", " ").title()
        base += " — the client specifically asked about this"
        base += (", in context: " + snippet) if snippet else ""
        return base.strip() + "."
    base = "%s slide" % labels.get(work_type, (work_type or "").replace("_", " ").title())
    if industry:
        base += " for the %s industry" % industry.replace("_", " ").title()
    base += (", addressing: " + snippet) if snippet else \
            ", covering J2W's relevant capability and the outcomes it delivers"
    return base.strip() + "."


def _similar_slides(work_type, query, n=3):
    """The closest real J2W case-study slides — used as a FORMAT/style example so a
    generated slide matches the existing deck instead of being invented from nothing."""
    try:
        lib = json.load(open(config.TAGGED_LIBRARY_JSON, encoding="utf-8"))
    except Exception:
        return []
    q = (query or "").lower()
    scored = []
    for r in lib:
        tags = r.get("tags", {})
        if (tags.get("kind", {}) or {}).get("value") != "CASE_STUDY":
            continue
        kws = r.get("keywords", []) or []
        score = sum(1 for k in kws if k and k.lower() in q)
        if (tags.get("work_type", {}) or {}).get("value") == work_type:
            score += 1
        scored.append((score, r))
    scored.sort(key=lambda x: -x[0])
    return [{"title": r.get("title", ""), "keywords": " · ".join((r.get("keywords") or [])[:8])}
            for _, r in scored[:n]]


def draft(gap, context):
    """Ask the LLM to write a gap slide's content, GUIDED by a brief and grounded in
    the format of similar real J2W slides. Returns {title, keywords, bullets:[...]}.
    Falls back to a stub on any error."""
    wt = gap.get("work_type", "")
    topic = (gap.get("topic") or context.get("topic") or "").strip()
    industry = context.get("industry", "")
    transcript = (context.get("transcript") or "")[:3000]
    brief = (context.get("brief") or "").strip()
    examples = _similar_slides(wt, topic or brief or transcript)
    ex_text = "\n".join("  - %s (keywords: %s)" % (e["title"], e["keywords"]) for e in examples) or "  (none found)"
    focus = (f"The CLIENT specifically asked about \"{topic}\" and the deck has no "
             f"slide on it. Write a slide that demonstrates J2W's capability in "
             f"\"{topic}\".\n\n") if topic else ""
    prompt = (
        f"You are writing ONE slide for a J2W sales deck. Work type: {wt}; "
        f"industry: {industry or 'the client'}.\n\n"
        f"{focus}"
        f"WHAT THIS SLIDE SHOULD COVER (follow this brief):\n"
        f"{brief or '(no brief given — infer it from the meeting notes below)'}\n\n"
        f"MEETING NOTES (context):\n\"\"\"\n{transcript}\n\"\"\"\n\n"
        f"FORMAT — follow the style of these existing J2W slides:\n{ex_text}\n\n"
        "Write the slide so it satisfies the brief. It can be a CASE STUDY that "
        "follows the format of the examples above, OR a different client-specific "
        "slide if the brief calls for that. Keep every claim credible and do NOT "
        "invent a specific real client name. Return ONLY JSON: "
        '{"title": "...", "keywords": "A · B · C · D", "bullets": ["...", "...", "..."]}. '
        "3-4 short bullets, each a concrete outcome or capability."
    )
    try:
        from deckengine.services.infra import load_env
        load_env()
        from openai import OpenAI
        resp = OpenAI().chat.completions.create(
            model=MODEL, temperature=0.4,
            response_format={"type": "json_object"},
            messages=[
                {"role": "system", "content": "You write concise B2B sales slide "
                                              "content. Reply with one JSON object only."},
                {"role": "user", "content": prompt},
            ],
        )
        data = json.loads(resp.choices[0].message.content)
        return {
            "title": str(data.get("title", "Proposed case study")),
            "keywords": str(data.get("keywords", "")),
            "bullets": [str(b) for b in (data.get("bullets") or [])][:4],
            "template": "case_study",
        }
    except Exception:
        return {
            "title": f"{wt} CASE STUDY - {industry or 'CLIENT'} (TO BE CREATED)",
            "keywords": "Draft · placeholder · replace",
            "bullets": ["Content could not be generated - add details manually."],
            "template": "case_study",
        }


# --------------------------------------------------------------------------- #
# Full structured CASE STUDY (the "Create with AI" feature) — strict format
# --------------------------------------------------------------------------- #
CASE_STUDY_RULES = (
    "NON-NEGOTIABLE RULES:\n"
    "- Exactly 6 capabilities, exactly 3 results, always.\n"
    "- No em dashes anywhere. All numbers as numerals. No real company names.\n"
    "- Capability names never reference technology labels (no LLM, RAG, GPT, ML, API, NLP); "
    "name each by business function.\n"
    "- Capability names must be SPECIFIC and CONCRETE — name the actual asset, mechanism, or "
    "deliverable (e.g. 'Pre-Trained Bilingual Talent Pipeline', 'Proven Delivery Model "
    "Lift-and-Scale', 'End-to-End CRM and CTI Integration'). NEVER a generic label like "
    "'Facilities Management', 'Resource Optimization', 'Data-Driven Insights', 'Operational "
    "Efficiency', 'Compliance Monitoring', 'Risk Management' or 'X Management'.\n"
    "- ALIGNMENT: the SOLUTION must directly answer EVERY problem raised in the CHALLENGE — "
    "each pain point in the challenge has a specific, matching response in the solution.\n"
    "- Every claim is concrete: name specific scale numbers, phases, sites, tools, workflows "
    "or outcomes. Never vague ('comprehensive framework', 'streamlined processes', 'enhanced "
    "efficiency', 'improved operations'). Invent realistic, specific details when the input "
    "is thin, the way a real proof point reads.\n"
    "- Infer realistic metrics from industry benchmarks if none are provided.\n"
    "- Tone: straight, professional, executive.\n"
    "RESULTS RULES:\n"
    "- Result 1: most impactful metric (percentage, number, or time contrast), one punchy line.\n"
    "- Result 2: operational or financial outcome, one line.\n"
    "- Result 3: qualitative shift in decision quality, confidence, or leverage, one line.\n"
    "- No result exceeds 15 words. No filler words.\n"
    "- Only use a time contrast when the gap is genuinely significant and specific; "
    "never a throwaway 'hours to minutes'.\n"
)

# The output-field contract shared by every case-study generator (draft / from-content).
CASE_FIELDS_SPEC = (
    "\nEach field:\n"
    "- title: a specific, concrete case study title (name the capability, not a slogan).\n"
    "- subhead: 'Client: <generic descriptor, e.g. Leading Manufacturing Enterprise> | "
    "Domain: <this account's domain> | Function: <the stakeholder's business function>'.\n"
    "- challenge: 3-4 sentences, plain and operational; who the client is (NEVER a real name) "
    "and what was breaking, specific to this domain. No solution language. Max 100 words.\n"
    "- solution: 3-4 sentences; what J2W deployed, how it works operationally, what the client "
    "can now do. No bullets. No hype. Max 100 words.\n"
    "- capabilities: EXACTLY 6, each 'Capability Name: one line max 18 words' (name = business function).\n"
    "- results: EXACTLY 3, following the RESULTS RULES.\n"
    "Then SELF-REVIEW what you wrote (quality verdict, weakest part, fix).\n"
    'Return ONLY JSON: {"title":"...","subhead":"...","challenge":"...","solution":"...",'
    '"capabilities":["six strings"],"results":["three strings"],'
    '"review":{"quality":"Strong or Needs Revision","weakest":"one sentence or None",'
    '"fix":"one sentence or None"}}'
)


def case_from_content(content, context=None):
    """Turn a user's COMPLETE case-study content (pasted text or an uploaded document)
    into ONE J2W case-study record — preserving THEIR facts, client situation and
    numbers, only reformatting into our structure. Returns the structured fields
    (+ 'review'); fails safe to a placeholder record."""
    context = context or {}
    industry = context.get("industry", "")
    prompt = (
        "Below is the COMPLETE content for ONE case study, provided by the user. Restructure "
        "it into J2W's case-study format. PRESERVE their facts: the client's situation, what "
        "was delivered, and every real number/metric. Do NOT invent facts, and do NOT drop "
        "content that fits a field — only reformat, tighten and organise it. Keep the client "
        "anonymised (never a real company name; only J2W is named).\n\n"
        + (f"Account domain (use for the Domain if the content doesn't state one): {industry}\n\n" if industry else "")
        + "CASE STUDY CONTENT:\n\"\"\"\n" + (content or "")[:12000] + "\n\"\"\"\n\n"
        "Follow these rules exactly:\n" + CASE_STUDY_RULES + CASE_FIELDS_SPEC
    )
    try:
        from deckengine.services.infra import load_env
        load_env()
        from openai import OpenAI
        resp = OpenAI().chat.completions.create(
            model=config.GEN_MODEL, temperature=0.3,
            response_format={"type": "json_object"},
            messages=[
                {"role": "system", "content": "You restructure a client's case-study content "
                 "into a strict format without losing facts. Reply with one JSON object only."},
                {"role": "user", "content": prompt},
            ],
        )
        data = json.loads(resp.choices[0].message.content)
    except Exception:
        data = {}
    return _normalize_case_study(data, industry)


def classify_content(content):
    """Which built-in template shape best fits this pasted content -- picks from
    CONTENT_TEMPLATES (the registry at the bottom of this file), so adding a new
    shape later means adding ONE entry there, not rewriting this prompt by hand.
    Owner's spec, 2026-07-08 (extended 2026-07-08 to more than two shapes):
    pasted content isn't always a case study; auto-detect which shape it
    actually is (like Gen Spark does), instead of always forcing the case-study
    mould. Fails safe to the registry's first entry ('case_study') -- the
    safer, more full-featured existing default when the AI call is unavailable
    or the reply doesn't parse to a real choice."""
    content = (content or "").strip()
    default_key = CONTENT_TEMPLATES[0]["key"]
    if not content:
        return default_key
    letters = "ABCDEFGHIJKLMNOPQRSTUVWXYZ"   # one per CONTENT_TEMPLATES entry (18 and counting)
    choices = "\n".join(f"({letters[i]}) {t['classify_desc']}"
                        for i, t in enumerate(CONTENT_TEMPLATES))
    key_by_letter = {letters[i]: t["key"] for i, t in enumerate(CONTENT_TEMPLATES)}
    prompt = (
        "Does the content below read as one of these template shapes?\n"
        + choices + "\n\nIf genuinely unsure, or the content doesn't clearly fit "
        "any one shape well, answer (A).\n\nCONTENT:\n\"\"\"\n" + content[:4000] + "\n\"\"\"\n\n"
        'Reply with ONLY this JSON: {"choice":"A"} -- the single letter of your pick.'
    )
    try:
        from deckengine.services.infra import load_env
        load_env()
        from openai import OpenAI
        resp = OpenAI().chat.completions.create(
            model=MODEL, temperature=0,
            response_format={"type": "json_object"},
            messages=[
                {"role": "system", "content": "You classify pasted slide content into "
                 "one of several template shapes. Reply with one JSON object only."},
                {"role": "user", "content": prompt},
            ],
        )
        data = json.loads(resp.choices[0].message.content)
        letter = str(data.get("choice", "")).strip().upper()[:1]
        return key_by_letter.get(letter, default_key)
    except Exception:
        return default_key


FOUR_BOX_FIELDS_SPEC = (
    "Return ONLY this JSON: {\"title\":\"...\",\"subhead\":\"...\","
    "\"boxes\":[{\"heading\":\"...\",\"body\":\"...\"},...]}\n"
    "- title: a short slide title (3-7 words) naming what the four boxes are.\n"
    "- subhead: one line of context/framing under the title (can be empty).\n"
    "- boxes: EXACTLY 4 items, in the content's own order. Each heading is 2-5 "
    "words; each body is 1-3 sentences, preserving the source's own facts/numbers.\n"
)


def four_box_from_content(content, context=None):
    """Turn pasted content that reads as a 4-way breakdown into J2W's four_box
    format: one short heading + one short paragraph per box. PRESERVE the
    user's own facts and numbers; only reformat into the fixed 4-box
    structure. Fails safe to a placeholder record (mirrors case_from_content)."""
    context = context or {}
    industry = context.get("industry", "")
    prompt = (
        "Below is content the user pasted that reads as FOUR roughly-parallel "
        "sections (pillars, categories, steps, or findings) -- not a single-"
        "client case study. Restructure it into J2W's four-box format. PRESERVE "
        "their facts and numbers; only reformat, tighten and organise -- do not "
        "invent content that isn't there, and if there are more or fewer than 4 "
        "natural sections, group or split them sensibly into exactly 4.\n\n"
        + (f"Account domain (context only, doesn't need to appear verbatim): {industry}\n\n" if industry else "")
        + "CONTENT:\n\"\"\"\n" + (content or "")[:12000] + "\n\"\"\"\n\n"
        + FOUR_BOX_FIELDS_SPEC
    )
    try:
        from deckengine.services.infra import load_env
        load_env()
        from openai import OpenAI
        resp = OpenAI().chat.completions.create(
            model=config.GEN_MODEL, temperature=0.3,
            response_format={"type": "json_object"},
            messages=[
                {"role": "system", "content": "You restructure pasted content into a "
                 "strict 4-box format without losing facts. Reply with one JSON "
                 "object only."},
                {"role": "user", "content": prompt},
            ],
        )
        data = json.loads(resp.choices[0].message.content)
    except Exception:
        data = {}
    return _normalize_four_box(data)


def _normalize_four_box(data):
    boxes = [b for b in (data.get("boxes") or []) if isinstance(b, dict)][:4]
    out_boxes = []
    for b in boxes:
        out_boxes.append({"heading": _clean(b.get("heading")) or "Untitled",
                          "body": _clean(b.get("body")) or "Content to be defined."})
    while len(out_boxes) < 4:
        out_boxes.append({"heading": "Untitled", "body": "Content to be defined."})
    return {"content_type": "four_box", "template": "four_box",
            "title": _clean(data.get("title")) or "Untitled",
            "subhead": _clean(data.get("subhead")),
            "boxes": out_boxes}


ROADMAP_FIELDS_SPEC = (
    "Return ONLY this JSON: {\"title\":\"...\",\"subhead\":\"...\",\"intro\":\"...\","
    "\"columns\":[{\"name\":\"...\",\"tag\":\"...\",\"items\":[\"...\"]}],"
    "\"legend\":[{\"tag\":\"...\",\"note\":\"...\"}],"
    "\"footer_title\":\"...\",\"footer_body\":\"...\"}\n"
    "- title: a short slide title (3-8 words).\n"
    "- subhead: one short line under the title, e.g. a section label (can be empty).\n"
    "- intro: a 1-3 sentence lead-in paragraph giving context for the plan/board, "
    "preserving the source's own facts (can be empty).\n"
    "- columns: one per distinct category/lane/function in the content, in the "
    "LITERAL reading order the content lists them in -- do NOT group, sort, or "
    "reorder columns by their tag/phase (e.g. do not move every 'Phase 1' column "
    "to the front); keep each column exactly where it appears in the source, even "
    "if that interleaves different tags. name = 2-4 words, the content's own "
    "names; tag = the short group/phase/stage label this column belongs to (e.g. "
    "'Phase 1'), copied EXACTLY as written elsewhere in the content -- every "
    "column in the same group must use the IDENTICAL tag string, character for "
    "character; items = its own bullet points, in the content's own order and "
    "wording, tightened to a few words each.\n"
    "- legend: one entry per DISTINCT tag used above, in the order each tag first "
    "appears. note = the short descriptor accompanying that tag in the source, if "
    "any (e.g. 'the anchors'), else empty.\n"
    "- footer_title / footer_body: a closing summary/callout banner if the content "
    "has one (e.g. a shared foundation or common thread across every column); both "
    "empty if the content has no such closing summary.\n"
)


def roadmap_from_content(content, context=None):
    """Turn pasted content that reads as a phased roadmap or board -- several
    categories/lanes each tagged with a phase/stage and its own bullet items --
    into J2W's roadmap_board format. PRESERVE the user's own facts and structure
    -- this shape, unlike four_box, has NO fixed slot count; as many columns as
    the content actually has (rendering draws them programmatically, see
    skills.py's _draw_roadmap_columns). Fails safe to a placeholder record
    (mirrors case_from_content / four_box_from_content)."""
    context = context or {}
    industry = context.get("industry", "")
    prompt = (
        "Below is content the user pasted that reads as a PHASED ROADMAP OR BOARD "
        "-- several categories/lanes/functions, each tagged with a phase, stage, "
        "or status, each listing its own items -- not a single-client case study "
        "and not a simple four-way breakdown. Restructure it into J2W's roadmap "
        "format. PRESERVE their facts, structure, and however many columns the "
        "content actually has; do not invent content that isn't there.\n\n"
        + (f"Account domain (context only, doesn't need to appear verbatim): {industry}\n\n" if industry else "")
        + "CONTENT:\n\"\"\"\n" + (content or "")[:12000] + "\n\"\"\"\n\n"
        + ROADMAP_FIELDS_SPEC
    )
    try:
        from deckengine.services.infra import load_env
        load_env()
        from openai import OpenAI
        resp = OpenAI().chat.completions.create(
            model=config.GEN_MODEL, temperature=0.3,
            response_format={"type": "json_object"},
            messages=[
                {"role": "system", "content": "You restructure pasted content into a "
                 "phased-roadmap/board format without losing facts or columns. "
                 "Reply with one JSON object only."},
                {"role": "user", "content": prompt},
            ],
        )
        data = json.loads(resp.choices[0].message.content)
    except Exception:
        data = {}
    return _normalize_roadmap(data)


def _extract_shape(content, context, instruction, fields_spec, system_note):
    """Shared AI-extraction call for the box_grid/pillar_deepdive/scored_list/
    stat_overview/data_table shapes below -- returns the raw parsed JSON dict
    (or {} on any failure); each shape's own _normalize_* turns that into a
    safe, complete record. Factored out since these 5 shapes are structurally
    identical calls (prompt + schema + one JSON response) -- case_from_content/
    four_box_from_content/roadmap_from_content predate this and are left as-is
    (working, tested, not worth touching for this)."""
    context = context or {}
    industry = context.get("industry", "")
    prompt = (
        instruction
        + (f"\n\nAccount domain (context only, doesn't need to appear verbatim): {industry}\n" if industry else "")
        + "\n\nCONTENT:\n\"\"\"\n" + (content or "")[:12000] + "\n\"\"\"\n\n"
        + fields_spec
    )
    try:
        from deckengine.services.infra import load_env
        load_env()
        from openai import OpenAI
        resp = OpenAI().chat.completions.create(
            model=config.GEN_MODEL, temperature=0.3,
            response_format={"type": "json_object"},
            messages=[
                {"role": "system", "content": system_note},
                {"role": "user", "content": prompt},
            ],
        )
        return json.loads(resp.choices[0].message.content)
    except Exception:
        return {}


BOX_GRID_FIELDS_SPEC = (
    "Return ONLY this JSON: {\"title\":\"...\",\"subhead\":\"...\","
    "\"boxes\":[{\"heading\":\"...\",\"body\":\"...\"},...]}\n"
    "- title: a short slide title (3-8 words).\n"
    "- subhead: one line of context/framing under the title (can be empty).\n"
    "- boxes: one per distinct section/category/pillar/finding, in the "
    "content's own order -- as many as the content naturally has (2 to 8), do "
    "NOT force a fixed count. Each heading is 2-5 words; each body is 1-3 "
    "sentences, preserving the source's own facts/numbers.\n"
)


def box_grid_from_content(content, context=None):
    """A generalised four_box for N sections (2-8) instead of exactly 4 --
    used by the per-slide 'Recreate with AI' pipeline (see recreate.py) where
    a source slide's own grid might have 3, 5, or 6 items, not always 4."""
    data = _extract_shape(
        content, context,
        "Below is content that reads as several roughly-parallel sections, "
        "categories, pillars, or findings -- not a single-client case study. "
        "Restructure it into J2W's box-grid format. PRESERVE the content's own "
        "facts and numbers; only reformat, tighten and organise -- use as many "
        "boxes as the content naturally has, do not force a fixed count.",
        BOX_GRID_FIELDS_SPEC,
        "You restructure pasted content into a box-grid format without losing "
        "facts. Reply with one JSON object only.")
    return _normalize_box_grid(data)


def _normalize_box_grid(data):
    boxes_in = [b for b in (data.get("boxes") or []) if isinstance(b, dict)][:8]
    boxes = []
    for b in boxes_in:
        heading = _clean(b.get("heading"))
        body = _clean(b.get("body"))
        if heading or body:
            boxes.append({"heading": heading or "Untitled",
                          "body": body or "Content to be defined."})
    while len(boxes) < 2:
        boxes.append({"heading": "Untitled", "body": "Content to be defined."})
    return {"content_type": "box_grid", "template": "box_grid",
            "title": _clean(data.get("title")) or "Untitled",
            "subhead": _clean(data.get("subhead")),
            "boxes": boxes}


PILLAR_FIELDS_SPEC = (
    "Return ONLY this JSON: {\"eyebrow\":\"...\",\"title\":\"...\","
    "\"blocks\":[{\"heading\":\"...\",\"body\":\"...\",\"subpoints\":[\"...\"]}]}\n"
    "- eyebrow: a short label above the title (e.g. 'PILLAR 01'), can be empty.\n"
    "- title: the capability/pillar's own name (2-5 words).\n"
    "- blocks: one per distinct feature/component within this pillar, in the "
    "content's own order -- as many as it naturally has (2 to 4). heading = "
    "2-5 words (the feature's own name); body = 1 sentence describing it; "
    "subpoints = the feature's own supporting details/specifics, 1 short line "
    "each (2 to 4 subpoints), preserving the source's own facts/numbers.\n"
)


def pillar_deepdive_from_content(content, context=None):
    """ONE capability broken into a few feature blocks, each with its own
    supporting sub-points -- matches the 'PILLAR 0N' deep-dive layout pattern
    (owner's reference deck, 2026-07-09), used by 'Recreate with AI' for a
    source slide that goes deep on a single capability rather than surveying
    several in parallel (that's box_grid instead)."""
    data = _extract_shape(
        content, context,
        "Below is content that reads as ONE capability/pillar broken into "
        "several features or components, each with its own supporting details "
        "-- restructure it into J2W's pillar-deepdive format. PRESERVE the "
        "content's own facts and numbers.",
        PILLAR_FIELDS_SPEC,
        "You restructure pasted content into a pillar-deepdive format without "
        "losing facts. Reply with one JSON object only.")
    return _normalize_pillar(data)


def _normalize_pillar(data):
    blocks_in = [b for b in (data.get("blocks") or []) if isinstance(b, dict)][:6]
    blocks = []
    for b in blocks_in:
        subs = [_clean(s) for s in (b.get("subpoints") or []) if _clean(s)][:10]
        heading = _clean(b.get("heading"))
        if not (heading or subs):
            continue
        blocks.append({"heading": heading or "Untitled",
                       "body": _clean(b.get("body")), "subpoints": subs})
    if not blocks:
        blocks = [{"heading": "Untitled", "body": "", "subpoints": ["Content to be defined."]}]
    return {"content_type": "pillar_deepdive", "template": "pillar_deepdive",
            "eyebrow": _clean(data.get("eyebrow")),
            "title": _clean(data.get("title")) or "Untitled",
            "blocks": blocks}


SCORED_LIST_FIELDS_SPEC = (
    "Return ONLY this JSON: {\"title\":\"...\",\"subhead\":\"...\","
    "\"rows\":[{\"name\":\"...\",\"description\":\"...\",\"stat\":\"...\"}]}\n"
    "- title: a short slide title (3-8 words).\n"
    "- subhead: one line of context under the title (can be empty).\n"
    "- rows: one per distinct named item (agent, step, component, metric "
    "owner), in the content's own order -- as many as it naturally has (2 to "
    "8). name = 2-4 words; description = 1 short sentence; stat = a short "
    "figure/score/label associated with that row if the content has one (e.g. "
    "'85% confidence', '112 markets'), else empty.\n"
)


def scored_list_from_content(content, context=None):
    """A named-row list with an optional per-row stat chip -- matches the
    'Agent Architecture' style layout (owner's reference deck, 2026-07-09):
    several named things, each with a short description and a figure."""
    data = _extract_shape(
        content, context,
        "Below is content that reads as a list of named items (agents, steps, "
        "components), each with a short description and often a figure/score "
        "attached -- restructure it into J2W's scored-list format. PRESERVE "
        "the content's own facts and numbers.",
        SCORED_LIST_FIELDS_SPEC,
        "You restructure pasted content into a scored-list format without "
        "losing facts. Reply with one JSON object only.")
    return _normalize_scored_list(data)


def _normalize_scored_list(data):
    rows_in = [r for r in (data.get("rows") or []) if isinstance(r, dict)][:14]
    rows = []
    for r in rows_in:
        name = _clean(r.get("name"))
        if not name:
            continue
        rows.append({"name": name, "description": _clean(r.get("description")),
                     "stat": _clean(r.get("stat"))})
    if not rows:
        rows = [{"name": "Untitled", "description": "Content to be defined.", "stat": ""}]
    return {"content_type": "scored_list", "template": "scored_list",
            "title": _clean(data.get("title")) or "Untitled",
            "subhead": _clean(data.get("subhead")),
            "rows": rows}


STAT_OVERVIEW_FIELDS_SPEC = (
    "Return ONLY this JSON: {\"title\":\"...\",\"subhead\":\"...\",\"intro\":\"...\","
    "\"stats\":[{\"value\":\"...\",\"label\":\"...\"}],\"items\":[\"...\"],"
    "\"footer_title\":\"...\",\"footer_body\":\"...\"}\n"
    "- title: a short slide title (3-8 words).\n"
    "- subhead: one line under the title (can be empty).\n"
    "- intro: a 1-2 sentence lead-in paragraph (can be empty).\n"
    "- stats: the content's own headline numbers/metrics, in the content's own "
    "order -- as many as it naturally has (2 to 6). value = the number/figure "
    "AS WRITTEN (e.g. '112', '85%', '47s'); label = 1-3 words naming it.\n"
    "- items: a short list of named components/capabilities the content "
    "enumerates alongside the stats, if any (0 to 6), else empty.\n"
    "- footer_title / footer_body: a closing summary banner if the content "
    "has one, else both empty.\n"
)


def stat_overview_from_content(content, context=None):
    """A headline-numbers overview -- matches the 'What is X?' style slide
    (owner's reference deck, 2026-07-09): a handful of big stats, maybe a row
    of named components, maybe a closing summary strip."""
    data = _extract_shape(
        content, context,
        "Below is content that reads as a headline STATS OVERVIEW -- a "
        "handful of key numbers/metrics, maybe alongside a short list of "
        "named components -- restructure it into J2W's stat-overview format. "
        "PRESERVE the content's own facts and numbers exactly as written.",
        STAT_OVERVIEW_FIELDS_SPEC,
        "You restructure pasted content into a stat-overview format without "
        "losing facts. Reply with one JSON object only.")
    return _normalize_stat_overview(data)


def _normalize_stat_overview(data):
    stats_in = [s for s in (data.get("stats") or []) if isinstance(s, dict)][:6]
    stats = []
    for s in stats_in:
        value = _clean(s.get("value"))
        if not value:
            continue
        stats.append({"value": value, "label": _clean(s.get("label")) or "-"})
    if not stats:
        stats = [{"value": "-", "label": "Untitled"}]
    items = [_clean(i) for i in (data.get("items") or []) if _clean(i)][:6]
    return {"content_type": "stat_overview", "template": "stat_overview",
            "title": _clean(data.get("title")) or "Untitled",
            "subhead": _clean(data.get("subhead")),
            "intro": _clean(data.get("intro")),
            "stats": stats, "items": items,
            "footer_title": _clean(data.get("footer_title")),
            "footer_body": _clean(data.get("footer_body"))}


DATA_TABLE_FIELDS_SPEC = (
    "Return ONLY this JSON: {\"title\":\"...\",\"subhead\":\"...\",\"intro\":\"...\","
    "\"col_labels\":[\"...\",\"...\",\"...\"],"
    "\"rows\":[{\"label\":\"...\",\"tag\":\"...\",\"value\":\"...\"}]}\n"
    "- title: a short slide title (3-8 words).\n"
    "- subhead: one line under the title (can be empty).\n"
    "- intro: a 1-3 sentence lead-in paragraph describing what the table shows "
    "(can be empty).\n"
    "- col_labels: EXACTLY 3 short column headers for the table (e.g. "
    "['Market','Filing Type','Count']), matching what the content's own table/"
    "list actually tracks.\n"
    "- rows: one per data row, in the content's own order -- as many as it "
    "naturally has (2 to 10). label = the row's own name (e.g. a market/"
    "entity); tag = the short categorical value for that row (e.g. a type/"
    "status), copied EXACTLY as written -- rows sharing a category use the "
    "IDENTICAL tag string; value = the row's own number/figure.\n"
)


def data_table_from_content(content, context=None):
    """A narrative panel + a small data table with a colour-coded category
    per row -- matches the 'Market Heatmap' style layout (owner's reference
    deck, 2026-07-09): real structured data, not prose to summarise."""
    data = _extract_shape(
        content, context,
        "Below is content that reads as a DATA TABLE or structured list of "
        "rows, each with a category/type and a figure -- restructure it into "
        "J2W's data-table format. PRESERVE the content's own facts and "
        "numbers exactly, and the content's own row order.",
        DATA_TABLE_FIELDS_SPEC,
        "You restructure pasted content into a data-table format without "
        "losing facts. Reply with one JSON object only.")
    return _normalize_data_table(data)


def _normalize_data_table(data):
    rows_in = [r for r in (data.get("rows") or []) if isinstance(r, dict)][:10]
    rows = []
    for r in rows_in:
        label = _clean(r.get("label"))
        if not label:
            continue
        rows.append({"label": label, "tag": _clean(r.get("tag")),
                     "value": _clean(r.get("value"))})
    if not rows:
        rows = [{"label": "Untitled", "tag": "", "value": ""}]
    cols = [_clean(c) for c in (data.get("col_labels") or []) if _clean(c)][:3]
    defaults = ["Item", "Category", "Value"]
    while len(cols) < 3:
        cols.append(defaults[len(cols)])
    return {"content_type": "data_table", "template": "data_table",
            "title": _clean(data.get("title")) or "Untitled",
            "subhead": _clean(data.get("subhead")),
            "intro": _clean(data.get("intro")),
            "col_labels": cols, "rows": rows}


def _normalize_roadmap(data):
    # 16 is a sanity ceiling (a slide can only be so wide), not a design target --
    # unlike four_box, real column counts are whatever the content has.
    cols_in = [c for c in (data.get("columns") or []) if isinstance(c, dict)][:16]
    columns = []
    for c in cols_in:
        items = [_clean(i) for i in (c.get("items") or []) if _clean(i)][:6]
        if not items:
            continue
        columns.append({"name": _clean(c.get("name")) or "Untitled",
                        "tag": _clean(c.get("tag")) or "",
                        "items": items})
    if not columns:
        columns = [{"name": "Untitled", "tag": "", "items": ["Content to be defined."]}]
    legend_in = [l for l in (data.get("legend") or []) if isinstance(l, dict)][:8]
    legend = [{"tag": _clean(l.get("tag")), "note": _clean(l.get("note"))}
             for l in legend_in if _clean(l.get("tag"))]
    return {"content_type": "roadmap_board", "template": "roadmap_board",
            "title": _clean(data.get("title")) or "Untitled",
            "subhead": _clean(data.get("subhead")),
            "intro": _clean(data.get("intro")),
            "columns": columns,
            "legend": legend,
            "footer_title": _clean(data.get("footer_title")),
            "footer_body": _clean(data.get("footer_body"))}


def _clean(s):
    return str("" if s is None else s).replace("—", "-").replace("–", "-").strip()


def _normalize_case_study(data, industry=""):
    caps = [_clean(c) for c in (data.get("capabilities") or []) if _clean(c)][:6]
    while len(caps) < 6:
        caps.append("Capability: to be defined.")
    res = [_clean(r) for r in (data.get("results") or []) if _clean(r)][:3]
    while len(res) < 3:
        res.append("Result to be defined.")
    rev = data.get("review") or {}
    subhead = _clean(data.get("subhead")) or (
        "Client: Leading %s | Domain: %s | Function: " % ((industry or "Enterprise").title(), industry))
    return {
        "template": "case_study_full",
        "title": _clean(data.get("title")) or "Proposed Case Study",
        "subhead": subhead,
        "challenge": _clean(data.get("challenge")),
        "solution": _clean(data.get("solution")),
        "capabilities": caps,
        "results": res,
        "review": {"quality": _clean(rev.get("quality")) or "Needs Revision",
                   "weakest": _clean(rev.get("weakest")), "fix": _clean(rev.get("fix"))},
    }


def draft_case_study(brief, context=None):
    """Generate ONE full case study in the strict format + self-review, SYNTHESISED
    from the real client context — the deep-research brief, stakeholder profile and
    full transcript — not just a definition of the topic. Returns the structured
    fields (+ 'review'). context keys: industry, recipient, function, notes
    (full transcript), research (deep-research brief), profile (stakeholder bio)."""
    context = context or {}
    industry = context.get("industry", "")
    recipient = context.get("recipient", "")
    function = context.get("function", "")
    notes = context.get("notes", "")
    research = context.get("research", "")
    profile = context.get("profile", "")
    prompt = (
        "Write ONE proof-point case study for a J2W sales meeting, grounded in the REAL "
        "client context below. It must be SPECIFIC and OPERATIONAL to THIS account — a "
        "SYNTHESIS of the provided research, profile and notes, never a generic "
        "definition of the topic.\n\n"
        "ACCOUNT CONTEXT:\n"
        + (f"- Industry / domain: {industry}\n" if industry else "")
        + (f"- Stakeholder we are meeting: {recipient}\n" if recipient else "")
        + (f"- Their function / remit: {function}\n" if function else "")
        + (f"\nDEEP RESEARCH BRIEF (background, domain specifics, priorities):\n\"\"\"\n"
           f"{research[:6000]}\n\"\"\"\n" if research else "")
        + (f"\nSTAKEHOLDER PROFILE (tone, remit, what they personally own):\n\"\"\"\n"
           f"{profile[:4000]}\n\"\"\"\n" if profile else "")
        + (f"\nMEETING TRANSCRIPT / NOTES (prior context, expressed needs):\n\"\"\"\n"
           f"{notes[:6000]}\n\"\"\"\n" if notes else "")
        + "\nThe capability / use case to prove (fill THIS gap):\n\"\"\"\n"
        + (brief or "")[:1500] + "\n\"\"\"\n\n"
        "SYNTHESISE the sources above into a case that mirrors what THIS stakeholder "
        "personally owns and the domain described in the research — an anonymised J2W "
        "engagement in the SAME domain, solving exactly this capability. Ground every "
        "claim in the provided context. Only where a specific metric is genuinely "
        "absent from the sources, use a realistic industry benchmark — never leave a "
        "field vague or generic.\n\n"
        "QUALITY BAR — match this level of specificity and note how the solution answers "
        "EACH part of the challenge (this is a DIFFERENT topic; do NOT copy its content):\n"
        "CHALLENGE: A leading global technology hardware company needed to stand up a "
        "scalable, high-performance Customer Engagement Center in Cairo for multilingual EMEA "
        "support at strict SLAs, replicating a proven 900 seat India center in a new geography "
        "with zero execution risk and a fixed November 1st launch.\n"
        "SOLUTION: A lift and scale of the proven India model into Cairo, applying the "
        "identical playbook, tools, workflows and governance already delivering for the "
        "client. A prime Cairo site was secured with phased capacity from 50 to 400 seats, a "
        "pre-trained bilingual pool of 50 resources ready day 1, and end to end IT including "
        "CRM and CTI integration.\n"
        "CAPABILITIES: Proven Delivery Model Lift-and-Scale; Pre-Trained Bilingual Talent "
        "Pipeline; End-to-End IT and Infrastructure Setup; Joint Client-J2W Governance; GDPR "
        "and ISO 27001 Compliance; Innovation Roadmap Integration.\n"
        "RESULTS: 98% SLA adherence sustained to the India blueprint; 50 resources live day 1 "
        "scaling to 400 seats; 5 year contractual partnership established.\n\n"
        "Follow these rules exactly:\n" + CASE_STUDY_RULES + CASE_FIELDS_SPEC
    )
    try:
        from deckengine.services.infra import load_env
        load_env()
        from openai import OpenAI
        resp = OpenAI().chat.completions.create(
            model=config.GEN_MODEL, temperature=0.5,
            response_format={"type": "json_object"},
            messages=[
                {"role": "system", "content": "You are a B2B case study writer and sales reviewer "
                 "for an enterprise AI and technology services company. Reply with one JSON object only."},
                {"role": "user", "content": prompt},
            ],
        )
        data = json.loads(resp.choices[0].message.content)
    except Exception:
        data = {}
    return _normalize_case_study(data, industry)


def fill_case_study(slide, content):
    """Fill a case_study_full template slide from structured content."""
    singles = {"{{TITLE}}": content.get("title", ""), "{{SUBHEAD}}": content.get("subhead", ""),
               "{{CHALLENGE}}": content.get("challenge", ""), "{{SOLUTION}}": content.get("solution", "")}
    for sh in slide.shapes:
        if not sh.has_text_frame:
            continue
        txt = sh.text_frame.text
        hit = next((m for m in singles if m in txt), None)
        if hit:
            editor.set_text(sh, singles[hit])
        elif "{{CAPABILITIES}}" in txt:
            _set_bullets(sh, content.get("capabilities", []))
        elif "{{RESULTS}}" in txt:
            _set_bullets(sh, content.get("results", []))


# --------------------------------------------------------------------------- #
# Build the slide into a deck
# --------------------------------------------------------------------------- #
def _blank_layout(prs):
    for layout in prs.slide_layouts:
        if (layout.name or "").lower().strip() == "blank":
            return layout
    return prs.slide_layouts[-1]


def _copy_slide(dest_prs, src_slide):
    """Copy a template slide into dest_prs as a new slide, including image parts
    so branded picture shapes (logos, backgrounds, icons) render correctly."""
    new = dest_prs.slides.add_slide(_blank_layout(dest_prs))
    for shp in list(new.shapes):
        shp._element.getparent().remove(shp._element)

    # Copy image relationships: for each image part in the source slide, add
    # its BYTES as a new image part in the DESTINATION package (get_or_add_
    # image_part) rather than relating directly to the source's part object.
    # The latter used to leave a part that still belongs to the source
    # package's bookkeeping attached to the destination -- harmless for a
    # single copy, but when the destination package independently numbers its
    # OWN image parts too (e.g. it already has images from other slides),
    # both can be assigned the same sequential filename (image2.png etc.),
    # producing a genuine duplicate-name zip entry on save (silently broken
    # or wrong images, depending on the reader). get_or_add_image_part()
    # registers a real part IN the destination package, so its filename is
    # correctly unique there.
    import io as _io
    src_part  = src_slide._part
    dest_part = new._part
    rId_map   = {}
    _R = 'http://schemas.openxmlformats.org/officeDocument/2006/relationships'
    for rId, rel in list(src_part.rels.items()):
        if rel.is_external:
            continue
        if '/image' in (rel.reltype or ''):
            image_part, new_rId = dest_part.get_or_add_image_part(_io.BytesIO(rel._target.blob))
            if new_rId != rId:
                rId_map[rId] = new_rId

    _REMAP = {f'{{{_R}}}embed', f'{{{_R}}}link'}

    def _remap(elem):
        for attr, val in list(elem.attrib.items()):
            if attr in _REMAP and val in rId_map:
                elem.attrib[attr] = rId_map[val]
        for child in elem:
            _remap(child)

    for shp in src_slide.shapes:
        elem = copy.deepcopy(shp._element)
        if rId_map:
            _remap(elem)
        new.shapes._spTree.append(elem)

    return new


def _set_bullets(shape, bullets):
    tf = shape.text_frame
    tf.clear()
    for i, b in enumerate(bullets):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = "• " + b


def _fill(slide, content):
    for sh in slide.shapes:
        if not sh.has_text_frame:
            continue
        txt = sh.text_frame.text
        if "{{TITLE}}" in txt:
            editor.set_text(sh, content.get("title", ""))
        elif "{{KEYWORDS}}" in txt:
            editor.set_text(sh, content.get("keywords", ""))
        elif "{{BULLETS}}" in txt:
            _set_bullets(sh, content.get("bullets", []))


def _add_verify_banner(slide, slide_width):
    """Stamp a loud red bar across the top of a slide: this AI-written slide has
    not been checked by an expert, so it must not reach a client as-is. The banner
    travels with the slide in the .pptx (visible in preview and in PowerPoint)."""
    bar = slide.shapes.add_textbox(0, 0, slide_width, Inches(0.5))
    bar.fill.solid()
    bar.fill.fore_color.rgb = RGBColor(0xC0, 0x39, 0x2B)      # red
    bar.line.fill.background()
    tf = bar.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = "⚠ AI-GENERATED - NEEDS EXPERT VERIFICATION - NOT CLIENT-READY"
    run.font.size = Pt(13)
    run.font.bold = True
    run.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)


def append_generated(deck_path, items):
    """items = [{template, title, keywords, bullets, verified}] -> add each as a
    new slide at the end of the deck. Unverified (verified != True) slides get a
    loud 'needs expert verification' banner stamped on them."""
    templates = list_templates()
    dest = Presentation(deck_path)
    for it in items:
        src = templates.get(it.get("template", "case_study"))
        if src is None and templates:
            src = next(iter(templates.values()))
        if src is None:
            continue
        new_slide = _copy_slide(dest, src)
        _fill(new_slide, it)
        if not it.get("verified"):
            _add_verify_banner(new_slide, dest.slide_width)
    dest.save(deck_path)
    return len(items)



# ═══════════════════════════════════════════════════════════════════════════════
# The ten style-guide shapes (owner's designs, 2026-07-10).
#
# Each is: a JSON field spec the model fills, and a normalizer that guarantees the
# record the renderer (services/rendering/draw_templates.py) can actually draw --
# right number of cards, no missing keys, no None. Adding an eleventh means one spec,
# one normalizer, one CONTENT_TEMPLATES entry, one SCHEMA entry, one editor branch.
#
# They all share _restructure(): the OpenAI boilerplate is identical for every shape,
# and ten copies of it would rot independently. Every one fails safe to a placeholder
# record, exactly as case_from_content and four_box_from_content do.
# ═══════════════════════════════════════════════════════════════════════════════

_RESTRUCTURE_SYSTEM = ("You restructure pasted content into one strict slide format "
                       "without losing facts or inventing them. Reply with one JSON "
                       "object only.")


def _restructure(content, context, spec, normalize, what):
    """Pasted content -> one shape's record. `what` names the shape for the prompt."""
    context = context or {}
    industry = context.get("industry", "")
    prompt = (
        "Below is content the user pasted that reads as %s. Restructure it into that "
        "format. PRESERVE their facts, numbers and wording where you can -- only "
        "reformat, tighten and organise. Do NOT invent content that isn't there; if a "
        "field has no source material, leave it empty.\n\n" % what
        + (f"Account domain (context only, need not appear verbatim): {industry}\n\n" if industry else "")
        + "CONTENT:\n\"\"\"\n" + (content or "")[:12000] + "\n\"\"\"\n\n"
        + spec
    )
    try:
        from deckengine.services.infra import load_env
        load_env()
        from openai import OpenAI
        resp = OpenAI().chat.completions.create(
            model=config.GEN_MODEL, temperature=0.3,
            response_format={"type": "json_object"},
            messages=[{"role": "system", "content": _RESTRUCTURE_SYSTEM},
                      {"role": "user", "content": prompt}],
        )
        data = json.loads(resp.choices[0].message.content)
    except Exception:
        data = {}
    return normalize(data)


def _dicts(data, key, cap):
    return [d for d in (data.get(key) or []) if isinstance(d, dict)][:cap]


def _strs(data, key, cap):
    return [_clean(x) for x in (data.get(key) or []) if _clean(x)][:cap]


def _head_fields(data, content_type):
    return {"content_type": content_type, "template": content_type,
            "title": _clean(data.get("title")) or "Untitled",
            "subhead": _clean(data.get("subhead"))}


# ── 1. pain_point_list ────────────────────────────────────────────────────────
def _normalize_pain_point_list(data):
    rows = [{"label": _clean(r.get("label")) or "Untitled",
             "body": _clean(r.get("body")) or "Content to be defined."}
            for r in _dicts(data, "rows", 16)]
    if not rows:
        rows = [{"label": "Untitled", "body": "Content to be defined."}]
    return dict(_head_fields(data, "pain_point_list"), rows=rows)


def pain_point_list_from_content(content, context=None):
    return _restructure(content, context,
        'Return ONLY this JSON: {"title":"...","subhead":"...",'
        '"rows":[{"label":"...","body":"..."},...]}\n'
        "- rows: 3-6 problems/pain points, in the content's own order. label is 2-4 "
        "words; body is ONE sentence naming the concrete consequence.\n",
        _normalize_pain_point_list, "a list of problems or pain points")


# ── 2. platform_overview ──────────────────────────────────────────────────────
def _normalize_platform_overview(data):
    stats = [{"value": _clean(s.get("value")) or "-",
              "label": _clean(s.get("label")) or "Untitled"}
             for s in _dicts(data, "stats", 4)]
    return dict(_head_fields(data, "platform_overview"),
                stats=stats or [{"value": "-", "label": "Untitled"}],
                capabilities=_strs(data, "capabilities", 6),
                footer_title=_clean(data.get("footer_title")),
                footer_items=_strs(data, "footer_items", 5))


def platform_overview_from_content(content, context=None):
    return _restructure(content, context,
        'Return ONLY this JSON: {"title":"...","subhead":"...",'
        '"stats":[{"value":"112","label":"Markets indexed"},...],'
        '"capabilities":["...",...],"footer_title":"...","footer_items":["...",...]}\n'
        "- stats: up to 4 headline numbers. value is the figure ONLY (112, 85%, 47s).\n"
        "- capabilities: up to 6 named capabilities, 1-3 words each.\n"
        "- footer_title + footer_items: the one cross-cutting layer and what it covers "
        "(empty if the content names none).\n",
        _normalize_platform_overview, "a platform overview: headline stats plus named capabilities")


# ── 3. before_after_split ─────────────────────────────────────────────────────
def _stages(data, key):
    return [{"tag": _clean(s.get("tag")), "label": _clean(s.get("label")) or "Stage"}
            for s in _dicts(data, key, 4)]


def _normalize_before_after_split(data):
    qs = [{"title": _clean(q.get("title")) or "Question",
           "body": _clean(q.get("body"))} for q in _dicts(data, "questions", 6)]
    return dict(_head_fields(data, "before_after_split"),
                intro=_clean(data.get("intro")),
                before_title=_clean(data.get("before_title")) or "Before",
                before_stages=_stages(data, "before_stages"),
                after_title=_clean(data.get("after_title")) or "After",
                after_stages=_stages(data, "after_stages"),
                questions=qs)


def before_after_split_from_content(content, context=None):
    return _restructure(content, context,
        'Return ONLY this JSON: {"title":"...","subhead":"...","intro":"...",'
        '"before_title":"...","before_stages":[{"tag":"HUMAN","label":"..."},...],'
        '"after_title":"...","after_stages":[{"tag":"AI AGENT","label":"..."},...],'
        '"questions":[{"title":"...","body":"..."},...]}\n'
        "- each lane has up to 4 stages, in order. tag says WHO does that stage "
        "(HUMAN, AI AGENT, ...); label is 2-4 words.\n"
        "- questions: up to 5 questions the change raises (empty list if none).\n",
        _normalize_before_after_split, "a before/after workflow transformation")


# ── 4. comparison_split ───────────────────────────────────────────────────────
def _normalize_comparison_split(data):
    feats = [{"heading": _clean(f.get("heading")) or "Untitled",
              "body": _clean(f.get("body"))} for f in _dicts(data, "features", 4)]
    rows = [{"metric": _clean(r.get("metric")) or "Metric",
             "a": _clean(r.get("a")), "b": _clean(r.get("b"))}
            for r in _dicts(data, "rows", 6)]
    return dict(_head_fields(data, "comparison_split"),
                panel_title=_clean(data.get("panel_title")) or "Capability overview",
                panel_intro=_clean(data.get("panel_intro")),
                features=feats or [{"heading": "Untitled", "body": ""}],
                table_title=_clean(data.get("table_title")) or "Comparison",
                col_a=_clean(data.get("col_a")) or "Option A",
                col_b=_clean(data.get("col_b")) or "Option B",
                rows=rows, takeaway=_clean(data.get("takeaway")))


def comparison_split_from_content(content, context=None):
    return _restructure(content, context,
        'Return ONLY this JSON: {"title":"...","subhead":"...","panel_title":"...",'
        '"panel_intro":"...","features":[{"heading":"...","body":"..."},...],'
        '"table_title":"...","col_a":"...","col_b":"...",'
        '"rows":[{"metric":"...","a":"...","b":"..."},...],"takeaway":"..."}\n'
        "- features: up to 4 capability cards for the left panel.\n"
        "- rows: up to 6 metrics compared across the two options.\n"
        "- takeaway: one italic closing line, or empty.\n",
        _normalize_comparison_split, "capabilities on one side and a two-option comparison on the other")


# ── 5. pillar_grid ────────────────────────────────────────────────────────────
def _normalize_pillar_grid(data):
    pillars = [{"heading": _clean(p.get("heading")) or "Untitled",
                "body": _clean(p.get("body")),
                "points": _strs(p, "points", 10)} for p in _dicts(data, "pillars", 8)]
    if not pillars:
        pillars = [{"heading": "Untitled", "body": "", "points": []}]
    return dict(_head_fields(data, "pillar_grid"), pillars=pillars)


def pillar_grid_from_content(content, context=None):
    return _restructure(content, context,
        'Return ONLY this JSON: {"title":"...","subhead":"...",'
        '"pillars":[{"heading":"...","body":"...","points":["...",...]},...]}\n'
        "- pillars: 3-6 parallel capability pillars, in the content's own order. "
        "heading 2-5 words; body 1-2 sentences; points up to 4 short supporting items.\n",
        _normalize_pillar_grid, "several parallel capability pillars, each with supporting points")


# ── 6. option_columns ─────────────────────────────────────────────────────────
def _normalize_option_columns(data):
    opts = []
    for o in _dicts(data, "options", 4):
        rows = [{"label": _clean(r.get("label")) or "Item",
                 "value": _clean(r.get("value"))} for r in _dicts(o, "rows", 6)]
        opts.append({"name": _clean(o.get("name")) or "Option",
                     "tag": _clean(o.get("tag")), "rows": rows})
    if not opts:
        opts = [{"name": "Option", "tag": "", "rows": []}]
    return dict(_head_fields(data, "option_columns"), options=opts,
                recommendation=_clean(data.get("recommendation")))


def option_columns_from_content(content, context=None):
    return _restructure(content, context,
        'Return ONLY this JSON: {"title":"...","subhead":"...",'
        '"options":[{"name":"...","tag":"...","rows":[{"label":"...","value":"..."},...]},...],'
        '"recommendation":"..."}\n'
        "- options: 2-4 options being compared. name is short; tag is its one-line "
        "characterisation (e.g. Cloud native).\n"
        "- rows: the SAME labels in the SAME order for every option, so they line up.\n",
        _normalize_option_columns, "two to four options compared across the same dimensions")


# ── 7. agent_architecture ─────────────────────────────────────────────────────
def _normalize_agent_architecture(data):
    agents = [{"name": _clean(a.get("name")) or "Agent",
               "body": _clean(a.get("body")),
               "badge": _clean(a.get("badge"))} for a in _dicts(data, "agents", 6)]
    if not agents:
        agents = [{"name": "Agent", "body": "", "badge": ""}]
    return dict(_head_fields(data, "agent_architecture"), agents=agents,
                footer=_clean(data.get("footer")))


def agent_architecture_from_content(content, context=None):
    return _restructure(content, context,
        'Return ONLY this JSON: {"title":"...","subhead":"...",'
        '"agents":[{"name":"...","body":"...","badge":"85% confidence"},...],'
        '"footer":"..."}\n'
        "- agents: 3-6 components/agents. body is one line on what it does; badge is "
        "its headline metric, or empty if none is stated.\n"
        "- footer: the orchestrating layer beneath them, or empty.\n",
        _normalize_agent_architecture, "a set of components or agents, each with a metric")


# ── 8. governance_list ────────────────────────────────────────────────────────
def _normalize_governance_list(data):
    items = [{"heading": _clean(i.get("heading")) or "Untitled",
              "body": _clean(i.get("body"))} for i in _dicts(data, "items", 5)]
    if not items:
        items = [{"heading": "Untitled", "body": ""}]
    return dict(_head_fields(data, "governance_list"), items=items)


def governance_list_from_content(content, context=None):
    return _restructure(content, context,
        'Return ONLY this JSON: {"title":"...","subhead":"...",'
        '"items":[{"heading":"...","body":"..."},...]}\n'
        "- items: 3-5 layers/stages, in order. heading 1-3 words; body 1-3 sentences.\n",
        _normalize_governance_list, "a sequence of layers, stages, or governance mechanisms")


# ── 9. guardrail_columns ──────────────────────────────────────────────────────
def _normalize_guardrail_columns(data):
    cols = []
    for c in _dicts(data, "columns", 4):
        pts = [{"lead": _clean(p.get("lead")) or "Point",
                "body": _clean(p.get("body"))} for p in _dicts(c, "points", 8)]
        cols.append({"heading": _clean(c.get("heading")) or "Untitled", "points": pts})
    if not cols:
        cols = [{"heading": "Untitled", "points": []}]
    return dict(_head_fields(data, "guardrail_columns"), columns=cols,
                callout_label=_clean(data.get("callout_label")),
                callout_body=_clean(data.get("callout_body")))


def guardrail_columns_from_content(content, context=None):
    return _restructure(content, context,
        'Return ONLY this JSON: {"title":"...","subhead":"...",'
        '"columns":[{"heading":"...","points":[{"lead":"...","body":"..."},...]},...],'
        '"callout_label":"Quick win","callout_body":"..."}\n'
        "- columns: 2-4 themes. Each point has a bold lead-in (2-4 words) and the rest "
        "of the sentence as body.\n"
        "- callout: the single takeaway strip at the bottom, or empty.\n",
        _normalize_guardrail_columns, "several themes, each a list of points with a bold lead-in")


# ── 10. opportunity_cards ─────────────────────────────────────────────────────
def _normalize_opportunity_cards(data):
    cards = [{"heading": _clean(c.get("heading")) or "Untitled",
              "opportunity": _clean(c.get("opportunity")),
              "outcome": _clean(c.get("outcome"))} for c in _dicts(data, "cards", 4)]
    if not cards:
        cards = [{"heading": "Untitled", "opportunity": "", "outcome": ""}]
    return dict(_head_fields(data, "opportunity_cards"), cards=cards)


def opportunity_cards_from_content(content, context=None):
    return _restructure(content, context,
        'Return ONLY this JSON: {"title":"...","subhead":"...",'
        '"cards":[{"heading":"...","opportunity":"...","outcome":"..."},...]}\n'
        "- cards: 2-4 opportunities. heading names it; opportunity states the problem "
        "and its scale; outcome states what we would deliver and the result.\n",
        _normalize_opportunity_cards, "several opportunities, each paired with its outcome")


# ── CONTENT_TEMPLATES registry ────────────────────────────────────────────────
# The single source of truth for every built-in "quick content" slide shape --
# classify_content() builds its prompt FROM this list, and callers (decks.py's
# /from_content route, new_form.html's template_hint dropdown) iterate it
# instead of hardcoding each key. Adding a new shape = add one entry here (a
# builder + a normalize function above) -- nothing else needs to change.
CONTENT_TEMPLATES = [
    {"key": "case_study", "label": "Case study",
     "classify_desc": "a CASE STUDY -- a specific client's situation, what was "
                      "delivered to solve it, and the outcome/results",
     "builder": case_from_content},
    {"key": "four_box", "label": "Four-box section",
     "classify_desc": "a FOUR-WAY BREAKDOWN -- roughly four parallel sections, "
                      "categories, pillars, steps, or findings, with no single "
                      "client-situation narrative",
     "builder": four_box_from_content},
    {"key": "roadmap_board", "label": "Phased roadmap / board",
     "classify_desc": "a PHASED ROADMAP OR BOARD -- several categories/lanes/"
                      "functions, each tagged with a phase, stage, or status, "
                      "each listing its own bullet items (e.g. a plan sequenced "
                      "across Phase 1/2/3, or a kanban-style board)",
     "builder": roadmap_from_content},
    {"key": "box_grid", "label": "Box grid (N sections)",
     "classify_desc": "a BOX GRID -- several (anywhere from 2 to 8) roughly-"
                      "parallel sections/categories/pillars/findings, like a "
                      "four-way breakdown but not limited to exactly four",
     "builder": box_grid_from_content},
    {"key": "pillar_deepdive", "label": "Capability deep-dive",
     "classify_desc": "a CAPABILITY DEEP-DIVE -- exactly ONE capability/pillar, "
                      "broken into a few features or components, each with its own "
                      "supporting details. If the content covers SEVERAL parallel "
                      "capabilities or themes side by side, this is the WRONG choice",
     "builder": pillar_deepdive_from_content},
    {"key": "scored_list", "label": "Named list with stats",
     "classify_desc": "a NAMED LIST WITH STATS -- several named items (agents, "
                      "steps, components), each with a short description and "
                      "often a figure/score attached to it",
     "builder": scored_list_from_content},
    {"key": "stat_overview", "label": "Headline stats overview",
     "classify_desc": "a HEADLINE STATS OVERVIEW -- a handful of key numbers/"
                      "metrics presented together, maybe alongside a short "
                      "list of named components",
     "builder": stat_overview_from_content},
    {"key": "data_table", "label": "Data table",
     "classify_desc": "a DATA TABLE -- structured rows of real data (each "
                      "with a category/type and a figure), not prose to "
                      "summarise",
     "builder": data_table_from_content},

    # ── the ten style-guide shapes (owner's designs, 2026-07-10) ──────────────
    {"key": "pain_point_list", "label": "Problem / pain-point list",
     "classify_desc": "a PROBLEM LIST -- several things that are broken or painful "
                      "today, each a short label and one line on its consequence; "
                      "the 'why this is hard' slide, with no solution in it",
     "builder": pain_point_list_from_content},
    {"key": "platform_overview", "label": "Platform overview (stats + capabilities)",
     "classify_desc": "a PLATFORM OVERVIEW -- a handful of headline numbers ABOUT a "
                      "named product/platform, plus the list of capabilities it is "
                      "built from, and often one cross-cutting layer beneath them",
     "builder": platform_overview_from_content},
    {"key": "before_after_split", "label": "Before / after workflow",
     "classify_desc": "a BEFORE AND AFTER -- the same workflow shown twice, as it "
                      "runs today and as it would run afterwards, stage by stage; "
                      "often noting who or what performs each stage",
     "builder": before_after_split_from_content},
    {"key": "comparison_split", "label": "Capabilities + comparison table",
     "classify_desc": "CAPABILITIES BESIDE A COMPARISON -- a few capability "
                      "descriptions AND a metric-by-metric comparison of exactly two "
                      "options, together on one slide",
     "builder": comparison_split_from_content},
    {"key": "pillar_grid", "label": "Numbered capability pillars",
     "classify_desc": "PARALLEL CAPABILITY PILLARS -- three to six capabilities "
                      "surveyed side by side, each with a prose SENTENCE describing it "
                      "and then a short checklist of plain supporting items (tools, "
                      "techniques, deliverables). The checklist items are bare phrases, "
                      "NOT 'Term: definition' pairs",
     "builder": pillar_grid_from_content},
    {"key": "option_columns", "label": "Option columns (A / B / C)",
     "classify_desc": "OPTIONS COMPARED IN COLUMNS -- two to four named options "
                      "(architectures, vendors, approaches) described against the "
                      "SAME set of dimensions, usually with a recommendation",
     "builder": option_columns_from_content},
    {"key": "agent_architecture", "label": "Component / agent architecture",
     "classify_desc": "AN ARCHITECTURE OF COMPONENTS -- several named agents, "
                      "services or modules, each with what it does and often a "
                      "headline metric, orchestrated by something beneath them",
     "builder": agent_architecture_from_content},
    {"key": "governance_list", "label": "Governance / timeline list",
     "classify_desc": "A SEQUENCE OF LAYERS OR STAGES -- three to five named "
                      "mechanisms in order, each with a paragraph; a governance "
                      "model, an escalation path, a layered system",
     "builder": governance_list_from_content},
    {"key": "guardrail_columns", "label": "Themed columns with lead-in bullets",
     "classify_desc": "THEMED COLUMNS OF DEFINITIONS -- two to four themes side by "
                      "side, and crucially each bullet under them reads as "
                      "'Term: what that term means' (a named control, policy or "
                      "mechanism, then its definition). Governance, guardrails, "
                      "controls and compliance slides look like this. Often closes on "
                      "a 'quick win' line. Pick this rather than capability pillars "
                      "when the bullets DEFINE things rather than list them",
     "builder": guardrail_columns_from_content},
    {"key": "opportunity_cards", "label": "Opportunity / outcome cards",
     "classify_desc": "OPPORTUNITIES PAIRED WITH OUTCOMES -- two to four numbered "
                      "opportunities, each stating the problem and its scale, then "
                      "separately what would be delivered and the result",
     "builder": opportunity_cards_from_content},
]


def classify_content_many(slides):
    """Which template shape fits each pasted slide -- for a WHOLE document, in ONE call.

    `slides` is a list of {"heading", "content"} (heading may be ""). Used by
    /builder/parse to resolve every slide whose category the salesperson didn't name
    outright: they wrote a heading ("How we think before we build"), or a phrase we
    don't recognise, or nothing at all. The heading is passed as a HINT, never as the
    answer -- the content decides.

    One call for the whole document: N slides must not cost N round-trips, and the
    shapes are better decided together (a deck's slides inform each other -- three
    consecutive client stories read as case studies; the slide between them that lists
    three principles does not).

    Differs from classify_content() in one deliberate way: that function is told to
    answer (A) -- case_study -- whenever it is unsure, which is right for a single
    unlabelled paste but wrong here. It made every slide of a real 9-slide deck come
    back a case study. Here the model must choose the BEST fit for each slide, and only
    a hard failure falls back to case_study.

    Returns a list of CONTENT_TEMPLATES keys, positionally aligned with `slides`.
    Fails safe: on any error every slide gets the registry's first key, exactly as the
    single-slide classifier does.
    """
    slides = list(slides or [])
    default_key = CONTENT_TEMPLATES[0]["key"]
    if not slides:
        return []
    fallback = [default_key] * len(slides)

    letters = "ABCDEFGHIJKLMNOPQRSTUVWXYZ"   # one per CONTENT_TEMPLATES entry (18 and counting)
    choices = "\n".join(f"({letters[i]}) {t['classify_desc']}"
                        for i, t in enumerate(CONTENT_TEMPLATES))
    key_by_letter = {letters[i]: t["key"] for i, t in enumerate(CONTENT_TEMPLATES)}
    blocks = "\n\n".join(
        'SLIDE %d%s:\n"""\n%s\n"""' % (
            i + 1,
            (' (the author titled it "%s")' % s.get("heading")) if s.get("heading") else "",
            (s.get("content") or "").strip()[:2500])
        for i, s in enumerate(slides))
    prompt = (
        "Below are the slides of one presentation. For EACH slide, pick the template "
        "shape that best fits its content.\n\n" + choices + "\n\n" + blocks + "\n\n"
        "Judge each slide on its own content. A title is a hint, not the answer.\n\n"
        "(A), the case study, is the most misused choice. Pick it ONLY when the slide "
        "tells ONE specific client's story and contains all three of: the client's "
        "situation, what was delivered for them, and the outcome that followed. If any "
        "of those three is missing, it is NOT a case study -- pick the shape that "
        "matches how the content is actually STRUCTURED. Guiding principles, a list of "
        "questions, an engagement model, a set of reassurances, a capability blueprint "
        "and an appendix are all not case studies, however much they discuss clients.\n\n"
        'Reply with ONLY this JSON: {"choices":["A","D",...]} -- one letter per slide, '
        "in order, exactly %d of them." % len(slides)
    )
    try:
        from deckengine.services.infra import load_env
        load_env()
        from openai import OpenAI
        resp = OpenAI().chat.completions.create(
            model=MODEL, temperature=0,
            response_format={"type": "json_object"},
            messages=[
                {"role": "system", "content": "You classify each slide of a pasted "
                 "document into one of several template shapes. Reply with one JSON "
                 "object only."},
                {"role": "user", "content": prompt},
            ],
        )
        picks = json.loads(resp.choices[0].message.content).get("choices") or []
    except Exception:
        return fallback
    if len(picks) != len(slides):
        return fallback                 # a short/long list means it lost track -- don't guess
    return [key_by_letter.get(str(p).strip().upper()[:1], default_key) for p in picks]


def template_keys():
    """Every valid template_hint value (plus 'auto', handled by the caller)."""
    return {t["key"] for t in CONTENT_TEMPLATES}


def build_content_slide(content, industry="", template_hint="auto"):
    """Pasted content -> one branded slide record, ready for staging.add().

    `template_hint` is a CONTENT_TEMPLATES key, or "auto" (anything unrecognised) to
    let classify_content() pick the shape. Returns (record, template_def) so the caller
    can show which shape was chosen. Shared by /from_content and the Custom Slide
    Builder, so both classify and build identically."""
    key = template_hint if template_hint in template_keys() else classify_content(content)
    tdef = next((t for t in CONTENT_TEMPLATES if t["key"] == key), CONTENT_TEMPLATES[0])
    rec = tdef["builder"](content, {"industry": industry})
    rec["kind"] = "user_created"
    return rec, tdef


if __name__ == "__main__":
    print("Created template:", create_temp_template())
    print("Templates found:", list(list_templates().keys()))
