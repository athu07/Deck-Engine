# -*- coding: utf-8 -*-
"""
recreate.py  --  "Recreate with AI": rebuild an uploaded deck from scratch in
J2W's own visual language, slide by slide, instead of restyling the original
shapes in place (see reskin.py for that -- kept as a separate, faster, always-
guaranteed-fidelity option; owner's spec, 2026-07-09).

Owner's reference case (2026-07-09): a GenSpark-produced version of an
uploaded capability deck kept every fact/number from the source verbatim but
gave each slide a layout matched to what that slide's content actually is --
a stat-tile overview, a 3-feature "pillar" deep-dive, a named list with score
chips, a small data table -- instead of a font/colour swap of the original
shapes. This module reproduces that: for each ORIGINAL slide, extract its
text, classify which content SHAPE it is (from the SAME registry
slide_generator.CONTENT_TEMPLATES the "Already have the content?" flow uses,
plus one local escape hatch for slides that don't fit), extract structured
fields for that shape, and render a FRESH J2W-styled slide from those fields.

A slide that doesn't fit any known shape (a title/cover, a closing/contact
slide, or something genuinely data-heavy/visual a text template would lose
meaning from) falls back to reskin's in-place restyle for THAT slide only --
never forced into a mould that doesn't fit.

Bookended the same way reskin.py is: our own title slide (CS01) first, our
own closing slide (CS08) last, pulled verbatim from the master deck. The
FIRST and LAST slide of the uploaded deck are skipped outright (almost always
the source's own cover/closing, made redundant by our bookends); every slide
in between is recreated or restyled.
"""

import io
import json
import os
import string

from pptx import Presentation

from deckengine import config
from deckengine.services.rendering import draw_templates, reskin, skills, slide_generator
from deckengine.services.rendering.slide_generator import CONTENT_TEMPLATES

_NONE_KEY = "_none"

# key -> (mapping fn for the header markers, draw fn for the programmatic
# body). case_study is handled separately (its own template file + the
# active-learned-template precedence -- see _render_case_study).
_MAPPING_FNS = {
    "four_box": skills._mapping_four_box,
    "roadmap_board": skills._mapping_roadmap_head,
    "box_grid": skills._mapping_box_grid_head,
    "pillar_deepdive": skills._mapping_pillar_head,
    "scored_list": skills._mapping_scored_list_head,
    "stat_overview": skills._mapping_stat_overview_head,
    "data_table": skills._mapping_data_table_head,
}
_DRAW_FNS = {
    "roadmap_board": skills._draw_roadmap_columns,
    "box_grid": skills._draw_box_grid,
    "pillar_deepdive": skills._draw_pillar_blocks,
    "scored_list": skills._draw_scored_rows,
    "stat_overview": skills._draw_stat_overview,
    "data_table": skills._draw_data_table,
}

# The ten style-guide shapes register themselves in draw_templates.DRAWERS -- fold them
# in rather than maintaining a second list here. Without this, Recreate-with-AI could
# classify a slide as (say) governance_list and then render only its header.
for _key, (_map_fn, _draw_fn) in draw_templates.DRAWERS.items():
    _MAPPING_FNS.setdefault(_key, _map_fn)
    _DRAW_FNS.setdefault(_key, _draw_fn)


def _as_stream(data):
    return io.BytesIO(data) if isinstance(data, (bytes, bytearray)) else data


def _shape_text(sh):
    """The readable content of ONE shape, recursing into groups and pulling the real
    data out of tables and charts. A plain-text-only reader (the old version) saw only
    the heading of a table slide or a chart slide, so those slides reached the classifier
    as a bare title, were misjudged, and fell to restyle-in-place -- which is exactly why
    Recreate produced Reskin's output on any data-heavy deck (owner-reported, 2026-07-13)."""
    from pptx.enum.shapes import MSO_SHAPE_TYPE
    out = []
    try:
        if sh.shape_type == MSO_SHAPE_TYPE.GROUP:
            for child in sh.shapes:
                t = _shape_text(child)
                if t:
                    out.append(t)
            return "\n".join(out)
    except Exception:
        pass

    try:
        if sh.has_table:
            tbl = sh.table
            for row in tbl.rows:
                cells = [c.text.strip() for c in row.cells]
                if any(cells):
                    out.append(" | ".join(cells))       # a row the classifier can read
            return "\n".join(out)
    except Exception:
        pass

    try:
        if sh.has_chart:
            chart = sh.chart
            try:
                cats = [str(c) for c in chart.plots[0].categories]
            except Exception:
                cats = []
            title = ""
            try:
                if chart.has_title and chart.chart_title.text_frame.text.strip():
                    title = chart.chart_title.text_frame.text.strip()
            except Exception:
                pass
            if title:
                out.append("Chart: " + title)
            for series in chart.series:
                vals = [("" if v is None else str(v)) for v in series.values]
                pairs = ", ".join("%s %s" % (c, v) for c, v in zip(cats, vals) if v != "")
                name = getattr(series, "name", "") or "series"
                out.append("%s: %s" % (name, pairs) if pairs else name)
            return "\n".join(out)
    except Exception:
        pass

    try:
        if sh.has_text_frame and sh.text_frame.text.strip():
            return sh.text_frame.text.strip()
    except Exception:
        pass
    return ""


def _slide_text(slide):
    """All of a slide's readable content, in document order -- text boxes, and now the
    contents of tables, charts and groups too."""
    lines = []
    for sh in slide.shapes:
        t = _shape_text(sh)
        if t:
            lines.append(t)
    return "\n".join(lines)


def _slide_layout_hint(slide, sw, sh):
    """A one-line description of how the source slide is ARRANGED, so the classifier can
    weigh visual structure, not just wording (owner-reported, 2026-07-13: a 2x2 icon grid
    was rebuilt as a vertical list, and a real table was flattened to bullets).

    Reports: whether a real table is present, and how the content cards are laid out --
    a grid (R x C), a row of columns, or a vertical stack -- inferred from the positions
    of the text/picture shapes."""
    from pptx.enum.shapes import MSO_SHAPE_TYPE
    if not sw or not sh:
        return ""

    # The precise row x column count is too noisy to trust (title/subtitle and per-bullet
    # boxes inflate it). The signal that actually decides list-vs-grid is coarse and
    # robust: how many distinct COLUMNS the content below the header occupies. Ignore the
    # header band (top ~22%) so the title/subtitle don't count as a row of content.
    has_table = False
    header_cut = 0.22 * sh
    xs = []                                # left-x of each content box below the header
    for shp in slide.shapes:
        try:
            if shp.has_table:
                has_table = True
                continue
            if shp.shape_type == MSO_SHAPE_TYPE.PICTURE:
                continue                   # icons/photos aren't layout cards
            if not (shp.has_text_frame and shp.text_frame.text.strip()):
                continue
            w, h = (shp.width or 0), (shp.height or 0)
            if w > 0.7 * sw and h > 0.5 * sh:
                continue                   # a full-slide background text box
            if (shp.top or 0) + h / 2 < header_cut:
                continue                   # the title / subtitle
            xs.append(shp.left or 0)
        except Exception:
            continue

    parts = []
    if has_table:
        parts.append("contains a real data TABLE")

    # distinct left-edges, at a coarse tolerance = distinct columns
    ncols = len(_cluster(sorted(xs), 0.14 * sw))
    if ncols >= 3:
        parts.append("content is laid out in %d COLUMNS (a grid or columns, not a single "
                     "vertical list)" % ncols)
    elif ncols == 2:
        parts.append("content is laid out in 2 side-by-side COLUMNS, not a single "
                     "vertical list")
    elif ncols == 1 and len(xs) >= 3:
        parts.append("content is a single vertical LIST")
    return "; ".join(parts)


def _cluster(sorted_vals, tol):
    """Group near-equal coordinates into bands -- a crude 1-D clustering for counting
    rows/columns. Returns the list of band representatives."""
    bands = []
    for v in sorted_vals:
        if bands and v - bands[-1] <= tol:
            continue
        bands.append(v)
    return bands


def _slide_icons(slide, sw, sh):
    """The source slide's own ICONS, in reading order, as (PNG/JPEG) bytes -- so Recreate
    can carry them into the rebuilt slide's chips (owner's spec, 2026-07-13).

    An icon is a SMALL, roughly-square picture: a big banner/photo or a full-bleed
    background is not one, and dragging it into a chip would be wrong. Recurses into
    groups. Sorted top-to-bottom, then left-to-right, to match the order the drawers
    place their chips."""
    from pptx.enum.shapes import MSO_SHAPE_TYPE
    _BLIP = "{http://schemas.openxmlformats.org/drawingml/2006/main}blip"
    _EMBED = "{http://schemas.openxmlformats.org/officeDocument/2006/relationships}embed"
    max_side = 0.20 * (sw or 1)            # <= ~20% of slide width reads as an icon
    found = []

    def walk(shapes, ox=0, oy=0):
        for sh in shapes:
            try:
                if sh.shape_type == MSO_SHAPE_TYPE.GROUP:
                    walk(sh.shapes, ox + (sh.left or 0), oy + (sh.top or 0))
                    continue
                w, h = (sh.width or 0), (sh.height or 0)
                if not w or not h:
                    continue
                if w > max_side or h > max_side:
                    continue               # too big to be an icon
                if max(w, h) / min(w, h) > 2.2:
                    continue               # a bar/line, not a square-ish icon
                # An icon may be a PICTURE shape OR an autoshape whose FILL is a picture
                # (blipFill in spPr) -- both carry an <a:blip>. Reading only PICTUREs missed
                # every fill-icon deck (owner-reported, 2026-07-13). image_normalize has
                # already made each blip a raster, so the blob is readable.
                blip = next(iter(sh._element.iter(_BLIP)), None)
                if blip is None:
                    continue
                rid = blip.get(_EMBED)
                if not rid:
                    continue
                blob = sh.part.related_part(rid).blob
                found.append((oy + (sh.top or 0), ox + (sh.left or 0), blob))
            except Exception:
                continue

    walk(slide.shapes)
    found.sort(key=lambda t: (round(t[0] / 100000), t[1]))   # reading order
    return [blob for _t, _l, blob in found]


def classify_slides(texts, hints=None):
    """Classify EVERY slide of the uploaded deck in ONE call, seeing them together.
    `hints[i]` is an optional one-line description of slide i's visual ARRANGEMENT
    (see _slide_layout_hint) -- a 2x2 grid, side-by-side columns, a real table -- so the
    layout is chosen from structure too, not wording alone.

    Two changes from the old per-slide classifier, both deliberate:

      * Whole-deck context. A slide is judged alongside its neighbours (three
        consecutive stories read as case studies; the list of questions between them
        does not), the same batch approach the Slide Builder already uses. It's also one
        AI call for the deck instead of one per slide.

      * NONE is a LAST resort, not a tie-break. The old prompt said "if unsure, answer
        NONE", which -- with content-poor extraction -- pushed most slides to restyle-in-
        place, so Recreate produced Reskin's output. Now NONE is reserved for a genuine
        cover / closing / pure-visual slide; a real content slide is rebuilt.

    Returns a list of CONTENT_TEMPLATES keys (or _NONE_KEY), aligned with `texts`.
    Fail-safe: on any error every slide is _NONE_KEY (restyle in place), the old default.
    """
    texts = list(texts or [])
    if not texts:
        return []
    fallback = [_NONE_KEY] * len(texts)

    letters = string.ascii_uppercase
    choices = "\n".join(f"({letters[i]}) {t['classify_desc']}"
                        for i, t in enumerate(CONTENT_TEMPLATES))
    none_letter = letters[len(CONTENT_TEMPLATES)]
    key_by_letter = {letters[i]: t["key"] for i, t in enumerate(CONTENT_TEMPLATES)}
    key_by_letter[none_letter] = _NONE_KEY

    hints = list(hints or [])
    blocks = []
    for i, t in enumerate(texts):
        hint = hints[i] if i < len(hints) else ""
        head = "SLIDE %d" % (i + 1)
        if hint:
            head += " [layout: %s]" % hint
        blocks.append('%s:\n"""\n%s\n"""' % (head, (t or "").strip()[:2500] or "(no readable text)"))
    blocks = "\n\n".join(blocks)
    prompt = (
        "Below are the slides of one uploaded presentation, each with the text, table "
        "rows and chart data read off it, and a [layout: ...] note describing how the "
        "slide is visually ARRANGED. For EACH slide, pick the template shape that best "
        "fits how its content is structured, using BOTH the content and the layout.\n\n"
        + choices +
        f"\n({none_letter}) NONE OF THESE -- ONLY for a title/cover slide, a closing/"
        "contact/thank-you slide, or a slide that is purely a diagram or image with no "
        "real textual content to rebuild.\n\n"
        "Rules that matter:\n"
        "- If the layout note says a real TABLE is present, choose (H) data_table -- keep "
        "it a table, do not flatten it to a list.\n"
        "- If the note says a GRID (e.g. 2x2), prefer a grid shape (box_grid, four_box, "
        "or numbered pillars) over a vertical list, even when the words read as problems "
        "or points.\n"
        "- If the note says side-by-side COLUMNS, prefer a columns shape (option_columns, "
        "guardrail_columns, comparison, or a before/after) over a single stack.\n"
        "- Most content slides DO fit a shape; only answer NONE when there is genuinely "
        "nothing rebuildable. Do not answer NONE merely because a slide is unusual.\n\n"
        + blocks +
        '\n\nReply with ONLY this JSON: {"choices":["A","H",...]} -- one letter per '
        "slide, in order, exactly %d of them." % len(texts)
    )
    try:
        from deckengine.services.infra import load_env
        load_env()
        from openai import OpenAI
        resp = OpenAI().chat.completions.create(
            model=slide_generator.MODEL, temperature=0,
            response_format={"type": "json_object"},
            messages=[
                {"role": "system", "content": "You classify each slide of an uploaded "
                 "deck into one of several template shapes, or 'none' if it is a "
                 "cover/closing/pure-visual slide. Reply with one JSON object only."},
                {"role": "user", "content": prompt},
            ],
        )
        picks = json.loads(resp.choices[0].message.content).get("choices") or []
    except Exception:
        return fallback
    if len(picks) != len(texts):
        return fallback                 # lost track of the count -> don't guess
    return [key_by_letter.get(str(p).strip().upper()[:1], _NONE_KEY) for p in picks]


def _restyle_original_slide(dest, src_slide, sw, sh, dest_sw, dest_sh):
    """Fallback for a slide that doesn't fit any content shape: copy the
    ORIGINAL slide and run reskin's own in-place restyle machinery on it
    (role detection, heading bar, brand overlay -- all proportioned to the
    slide's OWN original canvas, sw/sh), so it still gets J2W fonts/colours/
    branding even though its own layout is kept as-authored. THEN rescale the
    whole result from that original canvas to the STANDARD J2W canvas
    (dest_sw/dest_sh) the rest of this deck is being built at (see
    recreate_deck) -- every template-rendered slide is natively authored at
    that standard size already; only this restyle-in-place fallback path
    keeps the source's own (possibly different) coordinates and needs
    scaling. A no-op when the source already IS that standard size."""
    new = slide_generator._copy_slide(dest, src_slide)
    size_ranks = reskin._slide_sizes(new.shapes)
    dark_bg = reskin._slide_is_dark(new, new.shapes, sw, sh)
    reskin._restyle_shapes(new.shapes, size_ranks, sw, sh, dark_bg)
    reskin._recolor_to_palette(new.shapes, sw, sh)   # source brand colours -> J2W palette
    reskin._fix_text_on_box_contrast(new.shapes)
    reskin._fix_picture_contrast(new.shapes)
    reskin._resolve_bottom_banner_overlap(new.shapes, sw, sh)
    heading = reskin._first_heading_shape(new.shapes, size_ranks)
    if heading is not None:
        reskin._add_heading_bar(new, heading)
    reskin._whiten_slide_background(new)
    reskin._brand_slide(new, sw, sh)
    if sw and sh and (sw, sh) != (dest_sw, dest_sh):
        reskin._scale_shapes(new.shapes, dest_sw / sw, dest_sh / sh)
    return new


def _render_case_study(dest, record):
    """Mirrors skills.build_into's own case_study_v2 dispatch exactly: the
    owner's ACTIVE learned template wins if one exists, else the built-in
    case_study_v2 -- kept consistent with every other AI-created case study
    in the app, not a parallel reimplementation."""
    from deckengine.services.rendering import templatize as _templatize
    from deckengine.services.rendering import fill_case_study as _fcs
    active = _templatize.active_template()
    if active:
        _templatize.fill_into(dest, active, record)
        return
    case_tpl = skills.find_template(Presentation(config.CASE_TEMPLATE_PPTX), "case_study_v2")
    if case_tpl is None:
        return
    new = slide_generator._copy_slide(dest, case_tpl)
    _fcs.apply_markers(new, _fcs.build_mapping(record))


def _detect_score(text):
    """A headline score shown as 'N/100' (or /10, /5) in the source -- returned as
    (value, outof) so Recreate can redraw it as a filled ring instead of flattening it to
    a bullet (owner's spec, 2026-07-13). Only a round denominator, value <= denominator,
    to avoid catching 'IA/IB' or a stray fraction."""
    import re
    for m in re.finditer(r"\b(\d{1,3})\s*/\s*(100|10|5)\b", text or ""):
        val, out = int(m.group(1)), int(m.group(2))
        if val <= out:
            return val, out
    return None


_FIGURE_RESERVE = 3.7          # inches of right-column reserved for a carried source figure


def _prominent_image(slide, sw, sh):
    """The source slide's own main illustration/photo -- NOT a small icon (already carried
    into chips), NOT a full-bleed background -- as image bytes, so Recreate can carry it into
    the redrawn slide as a figure (owner's spec, 2026-07-13). Else None. Picks the largest
    picture OR picture-fill whose area is 12%-85% of the slide. image_normalize has already
    made every such image a readable raster."""
    from pptx.enum.shapes import MSO_SHAPE_TYPE
    BLIP = "{http://schemas.openxmlformats.org/drawingml/2006/main}blip"
    EMBED = "{http://schemas.openxmlformats.org/officeDocument/2006/relationships}embed"
    best = [None, 0.0]                        # [blob, area]

    def walk(shapes):
        for s in shapes:
            try:
                if s.shape_type == MSO_SHAPE_TYPE.GROUP:
                    walk(s.shapes)
                    continue
                blip = next(iter(s._element.iter(BLIP)), None)
                if blip is None:
                    continue
                w, h = (s.width or 0), (s.height or 0)
                if not w or not h:
                    continue
                area = (w / sw) * (h / sh)
                if area < 0.12 or area > 0.85 or area <= best[1]:
                    continue
                rid = blip.get(EMBED)
                if not rid:
                    continue
                best[0], best[1] = s.part.related_part(rid).blob, area
            except Exception:
                continue

    walk(slide.shapes)
    return best[0]


def _place_figure(slide, blob, reserve):
    """Drop the carried source image into the right column the drawers left free, fitted to
    that column and centred, aspect preserved."""
    from PIL import Image
    from pptx.util import Inches
    try:
        iw, ih = Image.open(io.BytesIO(blob)).size
    except Exception:
        return
    SW, SH = draw_templates.SW, draw_templates.SH
    bx, bw = SW - reserve + 0.10, reserve - 0.40
    by = draw_templates.TOP
    bh = SH - 0.45 - by
    ar = (iw / ih) if ih else 1.0
    w = bw
    h = w / ar
    if h > bh:
        h, w = bh, bh * ar
    x = bx + (bw - w) / 2.0
    y = by + (bh - h) / 2.0
    try:
        slide.shapes.add_picture(io.BytesIO(blob), Inches(x), Inches(y), Inches(w), Inches(h))
    except Exception:
        pass


def _render_shape(dest, tfile, key, record):
    template_slide = skills.find_template(tfile, record.get("template", key))
    if template_slide is None:
        return False
    new = slide_generator._copy_slide(dest, template_slide)
    mapping_fn = _MAPPING_FNS.get(key)
    if mapping_fn:
        skills.fill_markers(new, mapping_fn(record))
    # A prominent source image is carried in as a figure: reserve a right column so the
    # template lays its content into the left, then drop the image into the freed column
    # (owner's spec, 2026-07-13). No figure -> reserve 0, i.e. the full-width layout.
    figure = record.get("_figure")
    reserve = _FIGURE_RESERVE if figure else 0.0
    draw_fn = _DRAW_FNS.get(key)
    if draw_fn:
        # Two drawer families read slide width differently: the skills.py drawers hardcode
        # it (they subtract skills._RIGHT_RESERVE), the draw_templates.py drawers read the
        # module's SW. Narrow BOTH so whichever draws this key leaves the right column free.
        skills._RIGHT_RESERVE = reserve
        _full_sw = draw_templates.SW
        draw_templates.SW = _full_sw - reserve
        try:
            draw_fn(new, record)
        finally:
            skills._RIGHT_RESERVE = 0.0
            draw_templates.SW = _full_sw
    # a headline score in the source was a filled circle INSIDE a panel; redraw it inside
    # the layout (top-right of the body, over the last card's ghost-number corner), not
    # floating in the header (owner-reported, 2026-07-13).
    score = record.get("_score")
    if score:
        draw_templates.draw_score_ring(
            new, draw_templates.SW - reserve - draw_templates.MARGIN - 0.58,
            draw_templates.TOP + 0.52, 0.46, score[0], score[1], draw_templates.TEAL)
    if figure:
        _place_figure(new, figure, reserve)
    return True


def recreate_deck(data, out_path, industry=""):
    """Rebuild the uploaded pptx bytes into a FRESH J2W-styled deck, slide by
    slide, and save to out_path. Returns (out_path, stats) where stats =
    {"recreated": n, "restyled": n, "skipped": n} for the caller to report."""
    src_prs = Presentation(_as_stream(data))

    # Make every picture a RASTER first, whatever its source format -- SVG icons (often the
    # PRIMARY image, with no PNG fallback), WDP/HD-Photo effect layers, EMF/WMF metafiles.
    # Pillow (and so python-pptx) can't read those, and they were silently lost -- icons
    # vanished and primary-SVG pictures went blank (owner-reported, 2026-07-13). After this,
    # icon extraction and slide copy below see ordinary rasters and miss nothing. Fail-safe.
    try:
        from deckengine.services.rendering import image_normalize
        image_normalize.normalize_deck(src_prs)
    except Exception:
        pass

    src_slides = list(src_prs.slides)
    sw, sh = src_prs.slide_width, src_prs.slide_height   # the UPLOADED deck's own canvas

    # The output canvas is ALWAYS the standard J2W size (matches the master
    # deck / skills_templates.pptx / case_study_v2, all 13.33x7.5in) -- NOT
    # the uploaded deck's own size. Every template-rendered slide below is
    # natively authored at that standard size, so it needs no scaling; only
    # the restyle-in-place fallback (which keeps the source's own shape
    # coordinates) needs rescaling to fit, same as the bookend slides
    # already do. Preserving the source's own (possibly different) canvas
    # size instead was the root cause of a real bug: on a source deck
    # smaller than 13.33x7.5, every template slide overflowed past the
    # actual slide edge and every font rendered oversized relative to it
    # (owner-reported, 2026-07-09, with the actual output file attached).
    try:
        master = Presentation(config.MASTER_DECK)
        dest_sw, dest_sh = master.slide_width, master.slide_height
    except Exception:
        dest_sw, dest_sh = sw, sh   # fail-safe: fall back to the source's own size

    dest = Presentation()
    dest.slide_width, dest.slide_height = dest_sw, dest_sh
    tfile = Presentation(config.SKILLS_TEMPLATES_PPTX)

    stats = {"recreated": 0, "restyled": 0, "skipped": 0, "report": []}
    middle = src_slides[1:-1] if len(src_slides) > 2 else src_slides
    stats["skipped"] = len(src_slides) - len(middle)

    # classify the WHOLE deck in one call, seeing every slide together (see
    # classify_slides). One AI call for the deck, not one per slide, and with
    # cross-slide context.
    texts = [_slide_text(s) for s in middle]
    hints = [_slide_layout_hint(s, sw, sh) for s in middle]
    keys = classify_slides(texts, hints)

    # Deterministic guardrail on top of the model: a source laid out in MULTIPLE COLUMNS
    # must not be rebuilt as a single vertical LIST. The model still picks pain_point_list
    # for a 2x2 grid of problems because the WORDS read as problems; the owner wants the
    # grid preserved (2026-07-13). Remap the three single-column list shapes to box_grid,
    # which lays the same {heading, body} cards out in a grid.
    _LIST_SHAPES = {"pain_point_list", "governance_list", "scored_list"}
    for i, (key, hint) in enumerate(zip(keys, hints)):
        if key in _LIST_SHAPES and "COLUMN" in (hint or ""):
            keys[i] = "box_grid"

    def _label(key):
        if key == _NONE_KEY:
            return "as-authored"
        return next((t["label"] for t in CONTENT_TEMPLATES if t["key"] == key), key)

    for pos, (slide, text, key) in enumerate(zip(middle, texts, keys)):
        n = pos + 2                       # source slide number (1 = skipped cover)
        tdef = next((t for t in CONTENT_TEMPLATES if t["key"] == key), None)
        if key == _NONE_KEY or tdef is None:
            _restyle_original_slide(dest, slide, sw, sh, dest_sw, dest_sh)
            stats["restyled"] += 1
            stats["report"].append({"slide": n, "outcome": "restyled", "shape": "as-authored",
                                    "preview": (text[:60] or "(no text)")})
            continue
        record = tdef["builder"](text, {"industry": industry})
        # carry the source slide's own icons into the rebuilt slide's chips
        record["_icons"] = _slide_icons(slide, sw, sh)
        record["_score"] = _detect_score(text)         # a headline N/100 -> a ring
        record["_figure"] = _prominent_image(slide, sw, sh)   # a big illustration/photo -> figure
        outcome = "recreated"
        if key == "case_study":
            from deckengine.services.rendering.deck_build import ai_to_store_record
            _render_case_study(dest, ai_to_store_record(record, industry))
            stats["recreated"] += 1
        elif _render_shape(dest, tfile, key, record):
            stats["recreated"] += 1
        else:
            _restyle_original_slide(dest, slide, sw, sh, dest_sw, dest_sh)
            stats["restyled"] += 1
            outcome = "restyled"
        stats["report"].append({"slide": n, "outcome": outcome,
                                "shape": _label(key) if outcome == "recreated" else "as-authored",
                                "preview": (text[:60] or "(no text)")})

    _append_bookends(dest, dest_sw, dest_sh)
    out_dir = os.path.dirname(out_path)
    if out_dir:
        os.makedirs(out_dir, exist_ok=True)
    dest.save(out_path)
    return out_path, stats


def _append_bookends(dest, sw, sh):
    """Our own title slide (CS01) first, our own closing slide (CS08) last --
    pulled verbatim from the master deck, same source reskin.py's bookends
    use, scaled to the working canvas size."""
    from deckengine.services.content.build_library import read_id
    try:
        master = Presentation(config.MASTER_DECK)
    except Exception:
        return
    title_src = winback_src = None
    for slide in master.slides:
        sid = read_id(slide)
        if sid == "CS01":
            title_src = slide
        elif sid == "CS08":
            winback_src = slide
    msw, msh = master.slide_width, master.slide_height
    sx = (sw / msw) if msw else 1.0
    sy = (sh / msh) if msh else 1.0
    if title_src is not None:
        new_title = slide_generator._copy_slide(dest, title_src)
        reskin._scale_shapes(new_title.shapes, sx, sy)
        reskin._move_slide_to_front(dest, new_title)
    if winback_src is not None:
        new_close = slide_generator._copy_slide(dest, winback_src)
        reskin._scale_shapes(new_close.shapes, sx, sy)
