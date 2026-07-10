# -*- coding: utf-8 -*-
"""
skills.py  --  Four data-driven slides for Workforce-only decks.

Source data = J2W_Delivery_Footprint_Organized_Latest.xlsx
  Sheet: "Clean - Co x Func x Skill"
  Columns: Industry | Company | Function Delivered | Normalized Skill | Count

Gate: candidates() returns [] unless work type = pure Workforce (no other type
selected). Slides 1-3 additionally require the transcript/notes to mention "RFI"
or "request for information"; slide 4 has no RFI requirement.

Four slide types (templates live in skills_templates.pptx):
  industry_strength     -- overview of J2W's presence in the deck's industry (RFI)
  skill_deepdive        -- one combined slide for all skills matched in the notes (RFI)
  company_footprint     -- existing relationship if client fuzzy-matches a company (RFI)
  target_skill_profile  -- AI-categorised 3-column skill profile (domain expertise /
                           technical stack / academic & professional) from whatever
                           skills the notes name — no RFI requirement
"""

import re
from collections import defaultdict

import openpyxl
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION
from pptx.dml.color import RGBColor

from deckengine import config

EXCEL            = config.DELIVERY_FOOTPRINT_XLSX
SHEET            = "Clean - Co x Func x Skill"
SKILLS_TEMPLATES = config.SKILLS_TEMPLATES_PPTX   # branded template slides (from import_skills_templates.py)
CASE_TEMPLATE    = config.CASE_TEMPLATE_PPTX      # one shared template for content-store case studies

BRAND = [
    RGBColor(0x3A, 0x8B, 0x82),   # DEEP TEAL — matches the case-study / skills template
    RGBColor(0x7F, 0xB2, 0xA9),
    RGBColor(0xCF, 0xE7, 0xE2),
    RGBColor(0x11, 0x11, 0x10),
    RGBColor(0x4A, 0x9E, 0x94),
    RGBColor(0xA5, 0xCC, 0xC6),
    RGBColor(0x1E, 0x4D, 0x47),
]

# Distinct multi-color palette for the industry pie chart (no teal/green)
CHART_COLORS_MULTI = [
    RGBColor(0xE0, 0x6C, 0x1F),  # burnt orange
    RGBColor(0x7B, 0x4E, 0xA5),  # purple
    RGBColor(0x2A, 0x7E, 0xBC),  # steel blue
    RGBColor(0xB8, 0x2E, 0x2E),  # crimson red
    RGBColor(0xF5, 0xA6, 0x23),  # amber yellow
    RGBColor(0x8B, 0x5C, 0x2A),  # warm brown
    RGBColor(0xD4, 0x5E, 0x9A),  # raspberry pink
]

# Words to strip when normalizing company names for fuzzy matching
_CO_STRIP = re.compile(
    r'\b(pvt|ltd|llp|llc|inc|corp|co|pty|pte|sdn|bhd|gmbh|ag|nv|sa|ab|'
    r'consulting|advisory|payroll|ops|coe|technologies|technology|'
    r'software|solutions|services|group|india|global|international|'
    r'c2h|routing|sgp|benz|mercedes)\b',
    re.IGNORECASE
)

# Form industry codes -> Excel Industry column values
_IND_MAP = {
    "BFSI":             "Banking & Financial Services",   # the code the FORM sends
    "BANKING":          "Banking & Financial Services",
    "BANKING_FINANCE":  "Banking & Financial Services",
    "FINANCE":          "Banking & Financial Services",
    "TECH_IT":          "IT Services & Consulting",
    "TECH":             "IT Services & Consulting",
    "IT":               "IT Services & Consulting",
    "HEALTHCARE":       "Healthcare & Life Sciences",
    "HEALTH":           "Healthcare & Life Sciences",
    "INSURANCE":        "Insurance",
    "ENERGY":           "Energy & Utilities",
    "UTILITIES":        "Energy & Utilities",
    "RETAIL":           "E-commerce & Retail",
    "ECOMMERCE":        "E-commerce & Retail",
    "E_COMMERCE":       "E-commerce & Retail",
    "MANUFACTURING":    "Manufacturing & Materials",
    "TELECOM":          "Telecommunications",
    "AUTOMOTIVE":       "Automotive & Industrial",
    "INDUSTRIAL":       "Automotive & Industrial",
    "AEROSPACE":        "Aerospace & Defense",
    "AVIATION":         "Aerospace & Defense",   # the code the FORM sends
    "DEFENSE":          "Aerospace & Defense",
    "REALESTATE":       "Real Estate & Professional Services",
    "REAL_ESTATE":      "Real Estate & Professional Services",
    "CONSULTING":       "Consulting & Professional Services",
    "PROFESSIONAL":     "Consulting & Professional Services",
    "SEMICONDUCTORS":   "Semiconductors",
    "SEMICONDUCTOR":    "Semiconductors",
    "SOFTWARE":         "Software & Cloud",
    "CLOUD":            "Software & Cloud",
    "CYBERSECURITY":    "Cybersecurity",
    "CYBER":            "Cybersecurity",
    "HARDWARE":         "Technology Hardware",
    "TECHNOLOGY_HARDWARE": "Technology Hardware",
}

# Generic words that match too many skills; skip them in keyword matching
_SKILL_STOPWORDS = {
    'with', 'that', 'this', 'from', 'have', 'will', 'been', 'test', 'data',
    'back', 'your', 'tech', 'work', 'code', 'team', 'tion', 'ment', 'able',
    'ware', 'base', 'ness', 'ding', 'over', 'into', 'ting', 'ring', 'port',
}


# ------------------------------------------------------------------ #
# Data loading (module-level cache)
# ------------------------------------------------------------------ #
_cache = None


def _load():
    global _cache
    if _cache is not None:
        return _cache
    wb = openpyxl.load_workbook(EXCEL, data_only=True)
    ws = wb[SHEET]
    rows = []
    for r in ws.iter_rows(min_row=2, values_only=True):
        industry, company, function, skill, count = (r + (None,) * 5)[:5]
        if not company:
            continue
        rows.append({
            "industry": (industry or "").strip(),
            "company":  (company  or "").strip(),
            "function": (function or "").strip(),
            "skill":    (skill    or "").strip(),
            "count":    int(count or 0),
        })
    _cache = rows
    return rows


# ------------------------------------------------------------------ #
# Helpers
# ------------------------------------------------------------------ #
def _is_rfi(transcript):
    t = (transcript or "").lower()
    return bool(re.search(r'\brfi\b|request\s+for\s+information', t))


def _industry_label(deck_industry):
    code = (deck_industry or "").upper().strip()
    code = re.sub(r'[\s\-&]+', '_', code)
    if code in _IND_MAP:
        return _IND_MAP[code]
    # Partial match fallback
    for k, v in _IND_MAP.items():
        if code.startswith(k) or k.startswith(code):
            return v
    return deck_industry  # pass through as-is for an exact Excel match


def _normalize_co(name):
    n = (name or "").lower()
    n = _CO_STRIP.sub("", n)
    n = re.sub(r'[^a-z0-9]', '', n)
    return n.strip()


def _match_companies(client_name, all_companies):
    cn = _normalize_co(client_name)
    if len(cn) < 3:
        return []
    matched = []
    for co in all_companies:
        co_n = _normalize_co(co)
        if not co_n:
            continue
        if cn in co_n or co_n in cn:
            matched.append(co)
    return matched


def _match_skills(transcript, all_skills):
    """Return skills whose name (or significant words) appear in the transcript."""
    notes = (transcript or "").lower()
    matched = []
    for skill in all_skills:
        sk_lower = skill.lower()
        if sk_lower in notes:
            matched.append(skill)
            continue
        words = [w for w in re.findall(r'[a-z]{4,}', sk_lower)
                 if w not in _SKILL_STOPWORDS]
        if len(words) >= 2:
            if all(re.search(r'\b' + re.escape(w) + r'\b', notes) for w in words[:2]):
                matched.append(skill)
        elif len(words) == 1:
            if re.search(r'\b' + re.escape(words[0]) + r'\b', notes):
                matched.append(skill)
    seen = set()
    out = []
    for s in matched:
        if s not in seen:
            seen.add(s)
            out.append(s)
    return out


# ------------------------------------------------------------------ #
# Slide data builders
# ------------------------------------------------------------------ #
def _industry_slide_data(rows, industry_label):
    ind_rows = [r for r in rows
                if r["industry"].lower() == (industry_label or "").lower()]
    if not ind_rows:
        return None

    companies = {r["company"] for r in ind_rows}
    functions = {r["function"] for r in ind_rows}
    skills    = {r["skill"]   for r in ind_rows}
    total     = sum(r["count"] for r in ind_rows)

    # Top 3 frequently hired: sort by # distinct companies, then total headcount
    sk_cos = defaultdict(set)
    sk_cnt = defaultdict(int)
    for r in ind_rows:
        sk_cos[r["skill"]].add(r["company"])
        sk_cnt[r["skill"]] += r["count"]
    top3 = sorted(sk_cos.keys(), key=lambda s: (-len(sk_cos[s]), -sk_cnt[s]))[:3]

    fn_totals = defaultdict(int)
    for r in ind_rows:
        fn_totals[r["function"]] += r["count"]
    fn_sorted = sorted(fn_totals.items(), key=lambda x: -x[1])

    return {
        "industry":      industry_label,
        "total":         total,
        "num_companies": len(companies),
        "num_functions": len(functions),
        "num_skills":    len(skills),
        "top3":          top3,
        "top3_cos":      [len(sk_cos[s]) for s in top3],
        "fn_chart":      fn_sorted,
    }


def _skill_slide_data(rows, matched_skills):
    result = []
    for skill in matched_skills:
        sk_rows = [r for r in rows if r["skill"].lower() == skill.lower()]
        if not sk_rows:
            continue
        total = sum(r["count"] for r in sk_rows)
        cos   = defaultdict(int)
        inds  = set()
        for r in sk_rows:
            cos[r["company"]] += r["count"]
            inds.add(r["industry"])
        result.append({
            "skill":      skill,
            "total":      total,
            "companies":  sorted(cos.items(), key=lambda x: -x[1]),
            "industries": sorted(inds),
        })
    return result


def _company_slide_data(rows, client_name):
    all_cos = list({r["company"] for r in rows})
    matched = _match_companies(client_name, all_cos)
    if not matched:
        return None

    co_rows  = [r for r in rows if r["company"] in set(matched)]
    total    = sum(r["count"] for r in co_rows)
    functions = {r["function"] for r in co_rows}
    skills    = {r["skill"]   for r in co_rows}

    fn_totals = defaultdict(int)
    for r in co_rows:
        fn_totals[r["function"]] += r["count"]
    fn_sorted = sorted(fn_totals.items(), key=lambda x: -x[1])

    # Use the shortest matched name as the display name (most likely the clean variant)
    display_name = min(matched, key=len)

    return {
        "company":       display_name,
        "total":         total,
        "num_functions": len(functions),
        "num_skills":    len(skills),
        "fn_breakdown":  fn_sorted,
    }


# ------------------------------------------------------------------ #
# Public API
# ------------------------------------------------------------------ #
def candidates(context):
    """Return slide candidates. All are Workforce-only. Slides 1-3 (industry
    strength / skill deep-dive / company footprint, sourced from the delivery-
    footprint Excel) additionally require the RFI gate. Slide 4 (target skill
    profile, AI-categorised from the notes) has no RFI requirement — it triggers
    whenever the notes name skills/requirements to hire against."""
    wts = {str(w).upper() for w in (context.get("work_types") or [])}
    if wts != {"WORKFORCE"}:
        return []
    transcript = context.get("transcript", "")
    out = []

    if _is_rfi(transcript):
        rows          = _load()
        industry_lbl  = _industry_label(context.get("industry", ""))
        client        = (context.get("client_name", "") or "").strip()

        # Slide 1: Industry Strength
        ind_data = _industry_slide_data(rows, industry_lbl)
        if ind_data:
            out.append({
                "id":       "SK:industry",
                "kind":     "industry_strength",
                "label":    f"Industry strength - {industry_lbl}",
                "template": "industry_strength",
                "data":     ind_data,
                "stale":    False,
            })

        # Slide 2: Combined Skill Deep-dive
        all_skills = list({r["skill"] for r in rows})
        matched    = _match_skills(transcript, all_skills)
        if matched:
            sk_data = _skill_slide_data(rows, matched)
            if sk_data:
                names = [s["skill"] for s in sk_data[:3]]
                label = "Skills deployed - " + "  ·  ".join(names)
                if len(sk_data) > 3:
                    label += f"  +{len(sk_data)-3} more"
                out.append({
                    "id":       "SK:skills",
                    "kind":     "skill_deepdive",
                    "label":    label,
                    "template": "skill_deepdive",
                    "data":     sk_data,
                    "stale":    False,
                })

        # Slide 3: Company Relationship
        if client:
            co_data = _company_slide_data(rows, client)
            if co_data:
                out.append({
                    "id":       f"FP:{co_data['company']}",
                    "kind":     "company_footprint",
                    "label":    f"Client relationship - {co_data['company']}",
                    "template": "company_footprint",
                    "data":     co_data,
                    "stale":    False,
                })

    # Slide 4: Target Skill Profile — AI-categorised from the notes, no RFI gate
    tsp = _target_skill_profile_candidate(transcript)
    if tsp:
        out.append(tsp)

    return out


def _target_skill_profile_candidate(transcript):
    """Whenever the notes name skills/requirements to hire against, AI-categorise
    them into the 3-column Target Skill Profile slide. Returns None if the notes
    name nothing (never a placeholder slide) or the AI call fails (fail-safe)."""
    if not (transcript or "").strip():
        return None
    from deckengine.services.matching import ai_matcher
    try:
        profile = ai_matcher.extract_skill_profile(transcript)
    except Exception:
        return None
    if not profile:
        return None
    total = sum(len(profile.get(k) or []) for k in
                ("domain_expertise", "technical_stack", "academic_professional"))
    return {
        "id":       "TSP:skills",
        "kind":     "target_skill_profile",
        "label":    f"Target skill profile - {total} skill" + ("s" if total != 1 else ""),
        "template": "target_skill_profile",
        "data":     profile,
        "stale":    False,
    }


def by_id(context, sid):
    return next((c for c in candidates(context) if c["id"] == sid), None)


# ------------------------------------------------------------------ #
# Marker filling
# ------------------------------------------------------------------ #
def fill_markers(slide, mapping):
    for sh in slide.shapes:
        if not sh.has_text_frame:
            continue
        for p in sh.text_frame.paragraphs:
            full = "".join(run.text for run in p.runs)
            if "{{" not in full:
                continue
            new = full
            for m, v in mapping.items():
                new = new.replace("{{" + m + "}}", v)
            if new != full and p.runs:
                p.runs[0].text = new
                for run in p.runs[1:]:
                    run.text = ""


def _mapping_industry(data):
    top3     = data["top3"]
    top3_cos = data["top3_cos"]
    def _sk(i):
        if i < len(top3):
            return f"{top3[i]}  ({top3_cos[i]} {'company' if top3_cos[i]==1 else 'companies'})"
        return "-"
    return {
        "INDUSTRY_NAME":     data["industry"],
        "TOTAL_CONSULTANTS": f"{data['total']:,}",
        "NUM_COMPANIES":     str(data["num_companies"]),
        "NUM_FUNCTIONS":     str(data["num_functions"]),
        "NUM_SKILLS":        str(data["num_skills"]),
        "TOP_SKILL_1":       _sk(0),
        "TOP_SKILL_2":       _sk(1),
        "TOP_SKILL_3":       _sk(2),
    }


def _mapping_skills(sk_list):
    lines = []
    for s in sk_list:
        top_cos = ", ".join(f"{co} ({cnt})" for co, cnt in s["companies"][:4])
        if len(s["companies"]) > 4:
            top_cos += f"  +{len(s['companies'])-4} more"
        lines.append(
            f"▸ {s['skill']}  -  {s['total']:,} consultants  ·  "
            f"{len(s['companies'])} {'company' if len(s['companies'])==1 else 'companies'}\n"
            f"   {top_cos}"
        )
    header = "  ·  ".join(s["skill"] for s in sk_list)
    return {
        "SKILLS_HEADER": header,
        "SKILL_SUMMARY": "\n\n".join(lines),
    }


def _mapping_company(data):
    fn_lines = "\n".join(
        f"▸ {fn}:  {cnt:,}" for fn, cnt in data["fn_breakdown"]
    )
    return {
        "COMPANY_NAME":      data["company"],
        "TOTAL_DEPLOYED":    f"{data['total']:,}",
        "NUM_FUNCTIONS_CO":  str(data["num_functions"]),
        "NUM_SKILLS_CO":     str(data["num_skills"]),
        "ENGAGEMENT_TYPE":   "Existing client",
        "FUNCTION_BREAKDOWN": fn_lines,
    }


# id-prefix -> profile dict key, for the 3-column target_skill_profile slide
_TSP_PREFIX = {"DOM": "domain_expertise", "TEC": "technical_stack",
              "ACA": "academic_professional"}


def _mapping_target_skill_profile(profile, n_slots=6):
    """6 card slots per column; unfilled slots are blanked (same convention as
    the case-study capability cards) rather than removed/reflowed."""
    mapping = {}
    for prefix, key in _TSP_PREFIX.items():
        items = (profile.get(key) or [])[:n_slots]
        for i in range(1, n_slots + 1):
            item = items[i - 1] if i <= len(items) else None
            mapping[f"{prefix}_{i}_T"] = item["name"] if item else ""
            mapping[f"{prefix}_{i}_D"] = item["description"] if item else ""
    return mapping


def _mapping_four_box(data):
    """title/heading text is forced to caps here (not at extraction) -- owner's
    spec: headings are always caps, kept as a display-layer transform so the
    AI extraction itself stays plain text, same separation used in reskin.py."""
    boxes = (data.get("boxes") or [])[:4]
    mapping = {"TITLE": (data.get("title") or "").upper(),
              "SUBHEAD": data.get("subhead") or ""}
    for i in range(1, 5):
        b = boxes[i - 1] if i <= len(boxes) else {}
        mapping[f"BOX{i}_HEAD"] = b.get("heading", "")
        mapping[f"BOX{i}_BODY"] = b.get("body", "")
    return mapping


def _mapping_roadmap_head(data):
    """Header markers only -- the columns/legend/footer are drawn programmatically
    by _draw_roadmap_columns since, unlike four_box, this shape has no fixed slot
    count for a static template to carry."""
    return {"TITLE": (data.get("title") or "").upper(),
            "SUBHEAD": data.get("subhead") or "",
            "INTRO": data.get("intro") or ""}


# ── Shared programmatic-drawing toolkit ────────────────────────────────────────
# Every "variable count" shape below (roadmap_board, box_grid, pillar_deepdive,
# scored_list, stat_overview, data_table) draws its OWN body directly with
# python-pptx instead of filling a fixed-slot template -- the slot count varies
# with the content, so there's no fixed template to fill. These are the shared
# low-level primitives (same pattern reskin.py's/create_skills_templates.py's
# bar/card/text helpers use), promoted to module level so every _draw_* below
# shares ONE implementation instead of five copies.
_CARD_BG = RGBColor(0xF5, 0xF5, 0xF5)
_LINE = RGBColor(0xDE, 0xDE, 0xDE)
_INK = RGBColor(0x11, 0x11, 0x10)
_WHITE = RGBColor(0xFF, 0xFF, 0xFF)
_BODY = RGBColor(0x3E, 0x3E, 0x3E)
_MUTE = RGBColor(0x6E, 0x6E, 0x69)
_RED = RGBColor(0xD6, 0x28, 0x39)
_TEAL = RGBColor(0x2A, 0x9D, 0x8F)
_NAVY = RGBColor(0x1C, 0x2B, 0x44)

# Owner's house rule (2026-07-09, restated from reskin.py's identical SZ_BODY/
# SZ_BODY_HEAD convention): plain content text is ALWAYS 11pt; a bold mini-
# heading nested inside a card/row is ALWAYS 13pt. Applied to every _draw_*
# function below with exactly ONE deliberate exception -- a stat tile's own
# big display number (_draw_stat_overview), which isn't prose or a heading,
# it's a hero figure, matching the pre-existing metric_tile() convention in
# create_skills_templates.py that this shape's design is modelled on.
_SZ_BODY = 11
_SZ_HEAD = 13


def _draw_bar(slide, l, t, w, h, fill):
    from pptx.util import Inches
    from pptx.enum.shapes import MSO_SHAPE
    s = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(l), Inches(t), Inches(w), Inches(h))
    s.fill.solid(); s.fill.fore_color.rgb = fill
    s.line.fill.background(); s.shadow.inherit = False
    return s


def _draw_card(slide, l, t, w, h, fill, line):
    from pptx.util import Inches, Pt
    from pptx.enum.shapes import MSO_SHAPE
    s = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(l), Inches(t), Inches(w), Inches(h))
    s.fill.solid(); s.fill.fore_color.rgb = fill
    s.line.color.rgb = line; s.line.width = Pt(0.75)
    s.shadow.inherit = False
    return s


def _draw_text(slide, l, t, w, h, text, size, color, bold=False, align=None, font="Raleway", anchor=None):
    from pptx.util import Inches, Pt
    from pptx.enum.text import PP_ALIGN
    align = align or PP_ALIGN.LEFT
    tb = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_top = tf.margin_bottom = Pt(2)
    if anchor:
        tf._txBody.bodyPr.set("anchor", anchor)   # 'ctr' for vertical centring
    p = tf.paragraphs[0]
    p.alignment = align
    r = p.add_run(); r.text = text
    r.font.name = font; r.font.size = Pt(size); r.font.color.rgb = color; r.font.bold = bold
    return tb


def _draw_bullets(slide, l, t, w, h, items, size, color, marker="• "):
    from pptx.util import Inches, Pt
    from pptx.enum.text import MSO_AUTO_SIZE
    tb = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_top = tf.margin_bottom = Pt(2)
    tf.auto_size = MSO_AUTO_SIZE.NONE
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        r = p.add_run(); r.text = marker + item
        r.font.name = "Raleway"; r.font.size = Pt(size); r.font.color.rgb = color
    return tb


def _fit_size(item_count, box_h_in, base=_SZ_BODY, floor=7, per_item_lines=2):
    """Pick a font size that keeps `item_count` short lines inside a box of
    `box_h_in` height, baked in rather than left to a renderer's autofit (same
    reasoning as reskin.py's _fit_body_text)."""
    if item_count <= 0:
        return base
    per_item_in = per_item_lines * base * 1.2 / 72.0
    needed = item_count * per_item_in
    if needed <= box_h_in:
        return base
    return max(floor, int(base * (box_h_in / needed)))


# First 2 distinct tag/category values get a filled colour badge (red, then
# dark navy); any further one falls back to a neutral outline style -- shared
# by roadmap_board's columns/legend and data_table's row pills.
_TAG_FILLS = [(_RED, _WHITE), (_NAVY, _WHITE)]
_TAG_OUTLINE = (_WHITE, _INK)


def _tag_style(tag, style_by_tag):
    if tag not in style_by_tag:
        i = len(style_by_tag)
        style_by_tag[tag] = _TAG_FILLS[i] if i < len(_TAG_FILLS) else _TAG_OUTLINE
    return style_by_tag[tag]


def _draw_footer_banner(slide, footer_title, footer_body, y, sw=13.33, margin=0.30):
    """The full-width closing summary bar shared by roadmap_board and
    stat_overview (e.g. 'THE FOUNDATION, ESTABLISHED IN PHASE 1'). Returns the
    banner's own height so the caller can advance past it."""
    from pptx.enum.text import PP_ALIGN
    if not (footer_title or footer_body):
        return 0.0
    h = 0.62
    _draw_bar(slide, margin, y, sw - 2 * margin, h, _INK)
    if footer_title:
        _draw_text(slide, margin + 0.2, y + 0.08, sw - 2 * margin - 0.4, 0.26,
                   footer_title.upper(), _SZ_HEAD, _WHITE, bold=True, align=PP_ALIGN.CENTER)
    if footer_body:
        _draw_text(slide, margin + 0.3, y + 0.34, sw - 2 * margin - 0.6, 0.26,
                   footer_body, _SZ_BODY, RGBColor(0xDD, 0xDD, 0xDD), align=PP_ALIGN.CENTER)
    return h


def _draw_roadmap_columns(slide, data):
    """Add the column board + legend + footer banner to a copied roadmap_board
    template slide (which already carries the filled TITLE/SUBHEAD/INTRO
    header via fill_markers). Pure python-pptx shape drawing -- no markers --
    since the column count is whatever the content has, not a fixed slot set."""
    from pptx.enum.text import PP_ALIGN

    SW = 13.33
    MARGIN = 0.30
    columns = data.get("columns") or []
    n = len(columns)
    if n == 0:
        return

    top = 1.55 if (data.get("intro") or "").strip() else 1.15
    footer_h = 0.62 if (data.get("footer_title") or data.get("footer_body")) else 0.0
    legend_h = 0.30 if data.get("legend") else 0.0
    bottom = 7.50 - 0.24 - footer_h - legend_h - (0.12 if footer_h or legend_h else 0)
    col_h = bottom - top

    # wrap into 2 rows once columns would otherwise get unreadably thin
    max_per_row = max(1, int((SW - 2 * MARGIN) / 1.35))
    rows = [columns[i:i + max_per_row] for i in range(0, n, max_per_row)] if n > max_per_row \
        else [columns]
    row_gap = 0.20
    row_h = (col_h - (len(rows) - 1) * row_gap) / len(rows) if len(rows) > 1 else col_h

    style_by_tag = {}
    gap = 0.16
    for ri, row_cols in enumerate(rows):
        m = len(row_cols)
        col_w = (SW - 2 * MARGIN - (m - 1) * gap) / m
        row_top = top + ri * (row_h + row_gap)
        head_h = 0.56
        for ci, col in enumerate(row_cols):
            l = MARGIN + ci * (col_w + gap)
            fill, text_color = _tag_style(col.get("tag", ""), style_by_tag)
            outline = (fill == _TAG_OUTLINE[0])
            _draw_card(slide, l, row_top, col_w, head_h, fill, _LINE) if outline else \
                _draw_bar(slide, l, row_top, col_w, head_h, fill)
            _draw_text(slide, l + 0.08, row_top + 0.05, col_w - 0.16, 0.28, col.get("name", ""),
                       _SZ_HEAD, text_color, bold=True)
            if col.get("tag"):
                _draw_text(slide, l + 0.08, row_top + 0.32, col_w - 0.16, 0.20, col["tag"].upper(),
                           8, text_color)
            body_top = row_top + head_h + 0.08
            body_h = row_h - head_h - 0.08
            _draw_card(slide, l, body_top, col_w, body_h, _CARD_BG, _LINE)
            items = col.get("items") or []
            size = _fit_size(len(items), body_h - 0.16)
            _draw_bullets(slide, l + 0.10, body_top + 0.08, col_w - 0.20, body_h - 0.16, items,
                         size, _BODY)

    y = top + col_h + 0.12
    y += _draw_footer_banner(slide, data.get("footer_title"), data.get("footer_body"), y, SW, MARGIN)
    if data.get("footer_title") or data.get("footer_body"):
        y += 0.10

    legend = data.get("legend") or []
    if legend:
        lw = (SW - 2 * MARGIN) / len(legend)
        for i, item in enumerate(legend):
            fill, _tc = _tag_style(item.get("tag", ""), style_by_tag)
            lx = MARGIN + i * lw
            _draw_bar(slide, lx, y + 0.03, 0.14, 0.14, fill)
            label = item.get("tag", "") + (", " + item["note"] if item.get("note") else "")
            _draw_text(slide, lx + 0.20, y, lw - 0.24, 0.24, label, 9, _BODY)


def _mapping_box_grid_head(data):
    return {"TITLE": (data.get("title") or "").upper(), "SUBHEAD": data.get("subhead") or ""}


def _draw_box_grid(slide, data):
    """N boxes (2-8), auto-arranged in a grid -- the same box_grid, evaluated
    on the header markers, does not need EXACTLY 4 (that's four_box; this is
    the generalisation used by 'Recreate with AI' for a source slide whose own
    grid has 3, 5, or 6 items)."""
    boxes = data.get("boxes") or []
    n = len(boxes)
    if n == 0:
        return
    cols = 2 if n <= 4 else (3 if n <= 6 else 4)
    cols = min(cols, n)
    rows = -(-n // cols)   # ceil

    SW, SH = 13.33, 7.50
    MARGIN = 0.30
    gap = 0.24
    top = 1.35
    bottom = SH - 0.24
    box_w = (SW - 2 * MARGIN - (cols - 1) * gap) / cols
    box_h = (bottom - top - (rows - 1) * gap) / rows

    for i, b in enumerate(boxes):
        r, c = divmod(i, cols)
        bl = MARGIN + c * (box_w + gap)
        bt = top + r * (box_h + gap)
        _draw_card(slide, bl, bt, box_w, box_h, _CARD_BG, _LINE)
        _draw_bar(slide, bl + 0.16, bt + 0.16, 0.06, 0.34, _RED)
        _draw_text(slide, bl + 0.16 + 0.16, bt + 0.16, box_w - 0.32 - 0.16, 0.34,
                   b.get("heading", ""), _SZ_HEAD, _INK, bold=True)
        _draw_text(slide, bl + 0.16, bt + 0.16 + 0.44, box_w - 0.32, box_h - 0.32 - 0.44,
                   b.get("body", ""), _SZ_BODY, _BODY)


def _mapping_pillar_head(data):
    return {"TITLE": (data.get("title") or "").upper(), "SUBHEAD": data.get("eyebrow") or ""}


_PILLAR_ACCENTS = [_TEAL, _RED, _NAVY]


def _draw_pillar_blocks(slide, data):
    """2-4 feature blocks stacked vertically, each an accent-bar card with a
    heading + one-line body on the left portion and its own bullet sub-points
    on the right -- matches the 'PILLAR 0N' deep-dive layout (owner's
    reference deck, 2026-07-09)."""
    blocks = data.get("blocks") or []
    n = len(blocks)
    if n == 0:
        return
    SW, SH = 13.33, 7.50
    MARGIN = 0.30
    top = 1.35
    bottom = SH - 0.24
    gap = 0.18
    block_h = (bottom - top - (n - 1) * gap) / n
    left_w = 4.20

    for i, b in enumerate(blocks):
        bt = top + i * (block_h + gap)
        _draw_card(slide, MARGIN, bt, SW - 2 * MARGIN, block_h, _WHITE, _LINE)
        _draw_bar(slide, MARGIN, bt, 0.06, block_h, _PILLAR_ACCENTS[i % len(_PILLAR_ACCENTS)])
        _draw_text(slide, MARGIN + 0.22, bt + 0.14, left_w - 0.3, 0.32, b.get("heading", ""),
                   _SZ_HEAD, _INK, bold=True)
        if b.get("body"):
            _draw_text(slide, MARGIN + 0.22, bt + 0.50, left_w - 0.3, block_h - 0.64,
                       b["body"], _SZ_BODY, _BODY)
        subs = b.get("subpoints") or []
        if subs:
            _draw_bullets(slide, MARGIN + left_w + 0.20, bt + 0.14,
                         SW - 2 * MARGIN - left_w - 0.40, block_h - 0.28, subs,
                         _fit_size(len(subs), block_h - 0.28, base=11, per_item_lines=1), _BODY)


def _mapping_scored_list_head(data):
    return {"TITLE": (data.get("title") or "").upper(), "SUBHEAD": data.get("subhead") or ""}


def _draw_scored_rows(slide, data):
    """2-8 stacked rows: name + description on the left, an optional stat chip
    right-aligned -- matches the 'Agent Architecture' style row list (owner's
    reference deck, 2026-07-09)."""
    from pptx.enum.text import PP_ALIGN
    rows = data.get("rows") or []
    n = len(rows)
    if n == 0:
        return
    SW, SH = 13.33, 7.50
    MARGIN = 0.30
    top = 1.35
    bottom = SH - 0.24
    gap = 0.14
    row_h = min(0.72, (bottom - top - (n - 1) * gap) / n)

    for i, row in enumerate(rows):
        rt = top + i * (row_h + gap)
        _draw_card(slide, MARGIN, rt, SW - 2 * MARGIN, row_h, _CARD_BG, _LINE)
        name_w = 3.0
        stat_w = 1.8 if row.get("stat") else 0.0
        desc_w = SW - 2 * MARGIN - name_w - stat_w - 0.4
        _draw_text(slide, MARGIN + 0.18, rt, name_w, row_h, row.get("name", ""),
                   _SZ_HEAD, _INK, bold=True, anchor="ctr")
        if row.get("description"):
            _draw_text(slide, MARGIN + 0.18 + name_w, rt, desc_w, row_h, row["description"],
                       _SZ_BODY, _BODY, anchor="ctr")
        if row.get("stat"):
            _draw_text(slide, SW - MARGIN - stat_w - 0.15, rt, stat_w, row_h, row["stat"],
                       _SZ_BODY, _TEAL, bold=True, align=PP_ALIGN.RIGHT, anchor="ctr")


def _mapping_stat_overview_head(data):
    return {"TITLE": (data.get("title") or "").upper(), "SUBHEAD": data.get("subhead") or "",
            "INTRO": data.get("intro") or ""}


def _draw_stat_overview(slide, data):
    """2-6 stat tiles + an optional named-items row + an optional closing
    banner -- matches the 'What is X?' overview layout (owner's reference
    deck, 2026-07-09)."""
    from pptx.enum.text import PP_ALIGN
    stats = data.get("stats") or []
    items = data.get("items") or []
    SW = 13.33
    MARGIN = 0.30
    top = 1.55 if (data.get("intro") or "").strip() else 1.15

    if stats:
        n = len(stats)
        gap = 0.20
        tile_w = (SW - 2 * MARGIN - (n - 1) * gap) / n
        tile_h = 1.15
        for i, s in enumerate(stats):
            l = MARGIN + i * (tile_w + gap)
            _draw_card(slide, l, top, tile_w, tile_h, _WHITE, _LINE)
            _draw_bar(slide, l, top, tile_w, 0.05, _TEAL)
            # the stat VALUE is a deliberate exception to the 11/13 rule -- a
            # hero display figure, not prose or a heading (matches the pre-
            # existing metric_tile() convention this shape is modelled on)
            _draw_text(slide, l + 0.06, top + 0.16, tile_w - 0.12, 0.5, s.get("value", ""),
                       24, _TEAL, bold=True, align=PP_ALIGN.CENTER)
            _draw_text(slide, l + 0.06, top + 0.72, tile_w - 0.12, 0.35, s.get("label", ""),
                       _SZ_BODY, _MUTE, align=PP_ALIGN.CENTER)
        y = top + tile_h + 0.25
    else:
        y = top

    if items:
        _draw_text(slide, MARGIN, y, SW - 2 * MARGIN, 0.24, "CORE CAPABILITIES", _SZ_BODY, _TEAL, bold=True)
        y += 0.32
        n = len(items)
        col_w = (SW - 2 * MARGIN) / n
        for i, label in enumerate(items):
            _draw_text(slide, MARGIN + i * col_w, y, col_w - 0.1, 0.3, label, _SZ_HEAD, _INK, bold=True)
        y += 0.45

    y += 0.15
    _draw_footer_banner(slide, data.get("footer_title"), data.get("footer_body"), y, SW, MARGIN)


def _mapping_data_table_head(data):
    return {"TITLE": (data.get("title") or "").upper(), "SUBHEAD": data.get("subhead") or ""}


def _draw_data_table(slide, data):
    """A left narrative panel (its own intro paragraph, drawn here -- NOT via
    a header marker, see build_slide11's docstring) + a right-side data table
    with a colour-coded pill per row's category -- matches the 'Market
    Heatmap' layout (owner's reference deck, 2026-07-09)."""
    from pptx.enum.text import PP_ALIGN
    rows = data.get("rows") or []
    if not rows:
        return
    SW, SH = 13.33, 7.50
    MARGIN = 0.30
    top = 1.35
    bottom = SH - 0.24
    panel_w = 5.60
    gap = 0.30
    table_l = MARGIN + panel_w + gap
    table_w = SW - table_l - MARGIN

    _draw_card(slide, MARGIN, top, panel_w, bottom - top, _CARD_BG, _LINE)
    _draw_bar(slide, MARGIN, top, panel_w, 0.05, _TEAL)
    if data.get("intro"):
        _draw_text(slide, MARGIN + 0.2, top + 0.2, panel_w - 0.4, bottom - top - 0.4,
                   data["intro"], _SZ_BODY, _BODY)

    cols = data.get("col_labels") or ["Item", "Category", "Value"]
    col_w = [table_w * 0.46, table_w * 0.34, table_w * 0.20]
    col_l = [table_l, table_l + col_w[0], table_l + col_w[0] + col_w[1]]
    head_h = 0.34
    _draw_text(slide, col_l[0], top, col_w[0], head_h, cols[0], _SZ_BODY, _MUTE, bold=True)
    _draw_text(slide, col_l[1], top, col_w[1], head_h, cols[1], _SZ_BODY, _MUTE, bold=True)
    _draw_text(slide, col_l[2], top, col_w[2], head_h, cols[2], _SZ_BODY, _MUTE, bold=True, align=PP_ALIGN.RIGHT)
    _draw_bar(slide, table_l, top + head_h, table_w, 0.015, _LINE)

    n = len(rows)
    row_h = min(0.42, (bottom - top - head_h - 0.1) / n)
    style_by_tag = {}
    for i, r in enumerate(rows):
        rt = top + head_h + 0.1 + i * row_h
        _draw_text(slide, col_l[0], rt, col_w[0], row_h, r.get("label", ""), _SZ_HEAD, _INK)
        if r.get("tag"):
            fill, text_color = _tag_style(r["tag"], style_by_tag)
            pill_w = min(col_w[1] - 0.1, 1.3)
            _draw_bar(slide, col_l[1], rt + 0.03, pill_w, row_h - 0.14, fill)
            _draw_text(slide, col_l[1], rt + 0.03, pill_w, row_h - 0.14, r["tag"], 8.5,
                       text_color, bold=True, align=PP_ALIGN.CENTER)
        _draw_text(slide, col_l[2], rt, col_w[2], row_h, r.get("value", ""), _SZ_BODY, _INK,
                   bold=True, align=PP_ALIGN.RIGHT)
        if i < n - 1:
            _draw_bar(slide, table_l, rt + row_h - 0.01, table_w, 0.012, _LINE)


# ------------------------------------------------------------------ #
# Charts
# ------------------------------------------------------------------ #
def _find_chart_placeholder(slide):
    for sh in slide.shapes:
        if sh.has_text_frame and "{{CHART}}" in sh.text_frame.text:
            return sh
    return None


def _add_pie_chart(slide, categories, values):
    """Pie chart for the industry_strength slide.
    Uses a distinct multi-color palette (no teal/green).
    Legend on the RIGHT shows category names; percentage labels sit INSIDE
    each slice — this avoids label overlap when there are many categories."""
    from pptx.enum.chart import XL_LABEL_POSITION

    holder = _find_chart_placeholder(slide)
    if holder is None:
        return False
    left, top, w, h = holder.left, holder.top, holder.width, holder.height
    holder._element.getparent().remove(holder._element)

    pairs = [(c, v) for c, v in zip(categories, values) if v]
    if not pairs:
        return False
    cats, vals = zip(*pairs)

    cd = CategoryChartData()
    cd.categories = list(cats)
    cd.add_series("Consultants", list(vals))

    gf    = slide.shapes.add_chart(XL_CHART_TYPE.PIE, left, top, w, h, cd)
    chart = gf.chart
    chart.has_title  = False
    # Legend on the RIGHT — shows function names with color squares, no overlap
    chart.has_legend = True
    chart.legend.position = XL_LEGEND_POSITION.RIGHT
    chart.legend.include_in_layout = False

    # Percentage labels inside each slice; at this position PowerPoint
    # never generates overlapping labels regardless of how many categories exist
    plot = chart.plots[0]
    plot.has_data_labels = True
    dl = plot.data_labels
    dl.show_category_name = False   # category shown in legend; no duplication
    dl.show_percentage    = True    # e.g. "27%"
    dl.show_value         = False
    dl.show_legend_key    = False
    dl.position           = XL_LABEL_POSITION.INSIDE_END

    # Apply distinct multi-color fill per slice
    for i, pt in enumerate(plot.series[0].points):
        pt.format.fill.solid()
        pt.format.fill.fore_color.rgb = CHART_COLORS_MULTI[i % len(CHART_COLORS_MULTI)]
    return True


def _add_bar(slide, categories, values):
    """Horizontal bar chart — categories on Y-axis, counts on X-axis."""
    holder = _find_chart_placeholder(slide)
    if holder is None:
        return False
    left, top, w, h = holder.left, holder.top, holder.width, holder.height
    holder._element.getparent().remove(holder._element)

    # Cap at 10 bars for readability
    cats = list(categories[:10])
    vals = list(values[:10])
    if not vals:
        return False

    cd = CategoryChartData()
    cd.categories = cats
    cd.add_series("Consultants", vals)
    gf = slide.shapes.add_chart(XL_CHART_TYPE.BAR_CLUSTERED, left, top, w, h, cd)
    chart = gf.chart
    chart.has_legend = False
    series = chart.plots[0].series[0]
    series.format.fill.solid()
    series.format.fill.fore_color.rgb = BRAND[0]
    return True


# ------------------------------------------------------------------ #
# Build slides into the assembled deck
# ------------------------------------------------------------------ #
def find_template(prs, name):
    tag = "J2W_TEMPLATE: " + name
    for s in prs.slides:
        if s.has_notes_slide:
            for line in (s.notes_slide.notes_text_frame.text or "").splitlines():
                if line.strip() == tag:
                    return s
    return None


def build_into(deck_path, order, cands, master_path=config.MASTER_DECK):
    """Copy filled skills slides into the assembled deck and reorder to match `order`."""
    from pptx import Presentation
    from deckengine.services.rendering import slide_generator
    from deckengine.services.rendering import assembler
    from deckengine.services.content.build_library import read_id

    cand_by_id = {c["id"]: c for c in cands if c["id"] in order}
    if not cand_by_id:
        return 0

    prs        = Presentation(deck_path)
    sld_id_lst = prs.slides._sldIdLst
    tfile      = Presentation(SKILLS_TEMPLATES)
    case_tpl   = {"slide": None, "loaded": False}   # lazy: only opened if needed

    skill_elem = {}
    for sid, c in cand_by_id.items():
        # ── content-store case study: the owner's ACTIVE learned template wins if
        #    one exists (see templatize.py); else the built-in case_study_v2 ──
        if c.get("template") == "case_study_v2":
            from deckengine.services.rendering import templatize as _templatize
            active = _templatize.active_template()
            if active:
                _templatize.fill_into(prs, active, c["record"])
                skill_elem[sid] = list(sld_id_lst)[-1]
                continue
            if not case_tpl["loaded"]:
                case_tpl["slide"] = find_template(Presentation(CASE_TEMPLATE), "case_study_v2")
                case_tpl["loaded"] = True
            if case_tpl["slide"] is None:
                print(f"  WARNING: case_study_v2 template not found in {CASE_TEMPLATE}")
                continue
            from deckengine.services.rendering import fill_case_study as _fcs
            new = slide_generator._copy_slide(prs, case_tpl["slide"])
            # per-run fill (NOT the run-collapsing fill_markers) so the red 'CASE STUDY:'
            # prefix and the black case title keep their distinct colours in the deck.
            _fcs.apply_markers(new, _fcs.build_mapping(c["record"]))
            skill_elem[sid] = list(sld_id_lst)[-1]
            continue

        t = find_template(tfile, c["template"])
        if t is None:
            print(f"  WARNING: template '{c['template']}' not found — run create_skills_templates.py")
            continue

        new  = slide_generator._copy_slide(prs, t)
        kind = c["kind"]
        data = c["data"]

        if kind == "industry_strength":
            fill_markers(new, _mapping_industry(data))
            fn_names = [fn for fn, _ in data["fn_chart"]]
            fn_vals  = [cnt for _, cnt in data["fn_chart"]]
            _add_pie_chart(new, fn_names, fn_vals)

        elif kind == "skill_deepdive":
            fill_markers(new, _mapping_skills(data))
            sk_names = [s["skill"] for s in data]
            sk_vals  = [s["total"] for s in data]
            _add_bar(new, sk_names, sk_vals)

        elif kind == "company_footprint":
            fill_markers(new, _mapping_company(data))
            fn_names = [fn for fn, _ in data["fn_breakdown"]]
            fn_vals  = [cnt for _, cnt in data["fn_breakdown"]]
            _add_bar(new, fn_names, fn_vals)

        elif kind == "target_skill_profile":
            fill_markers(new, _mapping_target_skill_profile(data))

        elif kind == "four_box":
            fill_markers(new, _mapping_four_box(data))

        elif kind == "roadmap_board":
            fill_markers(new, _mapping_roadmap_head(data))
            _draw_roadmap_columns(new, data)

        elif kind == "box_grid":
            fill_markers(new, _mapping_box_grid_head(data))
            _draw_box_grid(new, data)

        elif kind == "pillar_deepdive":
            fill_markers(new, _mapping_pillar_head(data))
            _draw_pillar_blocks(new, data)

        elif kind == "scored_list":
            fill_markers(new, _mapping_scored_list_head(data))
            _draw_scored_rows(new, data)

        elif kind == "stat_overview":
            fill_markers(new, _mapping_stat_overview_head(data))
            _draw_stat_overview(new, data)

        elif kind == "data_table":
            fill_markers(new, _mapping_data_table_head(data))
            _draw_data_table(new, data)

        skill_elem[sid] = list(sld_id_lst)[-1]

    # Reorder the whole deck to match `order`; use explicit None checks (lxml falsy-when-empty).
    cs_elem = {read_id(s): e for s, e in zip(prs.slides, list(sld_id_lst)) if read_id(s)}
    for sid in order:
        e = cs_elem.get(sid)
        if e is None:
            e = skill_elem.get(sid)
        if e is not None:
            sld_id_lst.append(e)

    assembler._atomic_save(prs, deck_path)
    return len(skill_elem)
