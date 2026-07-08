# -*- coding: utf-8 -*-
"""
templatize.py  --  "learn a template from a deck": upload any .pptx, pick one
slide, get an AI-proposed mapping of its text boxes onto the case-study content
schema (title / subtitle / challenge / solution / capabilities / results),
confirm or fix that mapping by hand, and save it as a reusable template. Once
ACTIVE, real case-study content (store picks, Create-with-AI) renders into it
instead of the built-in case_study_v2 -- so a deck follows the owner's own
house style, not a developer-authored one.

Deliberately NOT fully automatic: the AI only PROPOSES a mapping (one call,
cheap model); a human confirms it before it's ever used to render a real
slide. Getting a box's role wrong would silently misplace content on a
client-facing slide, so this mirrors the same human-review gate every other
AI-touched surface in this engine already has.

WHAT IT DOES NOT DO (scope, on purpose): it does not detect content TYPES
beyond the case-study schema (no icons, no charts, no arbitrary layouts) and
it does not auto-route between multiple learned templates by content shape --
only one template can be ACTIVE at a time (radio-button style), matching the
owner's stated goal ("every deck of mine follows the same box, same colours").
"""

import json
import os
import re
import time

from pptx import Presentation

from deckengine import config
from deckengine.services.rendering import slide_generator
from deckengine.services.rendering.fill_case_study import split_capability, split_result

INDEX_PATH = os.path.join(config.LEARNED_TEMPLATES_DIR, "_index.json")

# The case-study content schema, in the vocabulary the AI (and the confirm UI)
# assigns per text box. Repeatable roles may be assigned to more than one box
# (order preserved, top-to-bottom / left-to-right as read off the slide).
ROLE_VOCAB = [
    "title", "subtitle",
    "challenge_body", "solution_body",
    "capability_title", "capability_body",
    "result_stat", "result_caption",
    "skip",
]
ROLE_LABELS = {
    "title": "Title", "subtitle": "Subtitle (client/domain line)",
    "challenge_body": "Challenge text", "solution_body": "Solution text",
    "capability_title": "Capability heading", "capability_body": "Capability description",
    "result_stat": "Result stat/number", "result_caption": "Result caption",
    "skip": "Skip (leave as authored -- logo, page number, decoration, etc.)",
}


def _stream(data):
    import io
    return io.BytesIO(data) if isinstance(data, (bytes, bytearray)) else data


def _slug(name):
    s = re.sub(r"[^A-Za-z0-9]+", "-", (name or "").lower()).strip("-")
    return s or "template"


# --------------------------------------------------------------------------- #
# 1) inspect an uploaded deck (for the upload -> pick-a-slide step)
# --------------------------------------------------------------------------- #
def list_slides(deck_bytes):
    """[{index, shape_count, preview}] for every slide -- lets the owner pick
    which one to teach without opening the file."""
    prs = Presentation(_stream(deck_bytes))
    out = []
    for i, s in enumerate(prs.slides):
        texts = [sh.text_frame.text.strip() for sh in s.shapes
                if sh.has_text_frame and sh.text_frame.text.strip()]
        out.append({"index": i, "shape_count": len(texts),
                    "preview": [t[:60] for t in texts[:4]]})
    return out


def shapes_of(deck_bytes, slide_index):
    """[{idx, left, top, width, height, text}] for every text-bearing shape on
    ONE slide, in reading order (top-to-bottom, then left-to-right) -- the
    raw material for both the AI proposal and the confirm-screen table."""
    prs = Presentation(_stream(deck_bytes))
    s = prs.slides[slide_index]
    out = []
    for i, sh in enumerate(s.shapes):
        if sh.has_text_frame and sh.text_frame.text.strip():
            out.append({
                "idx": i,
                "left": round((sh.left or 0) / 914400, 2),
                "top": round((sh.top or 0) / 914400, 2),
                "width": round((sh.width or 0) / 914400, 2),
                "height": round((sh.height or 0) / 914400, 2),
                "text": sh.text_frame.text.strip()[:160],
            })
    out.sort(key=lambda d: (round(d["top"]), d["left"]))
    return out


# --------------------------------------------------------------------------- #
# 2) AI proposes a role per shape -- a PROPOSAL, never applied without review
# --------------------------------------------------------------------------- #
def propose_roles(shapes):
    """{str(idx): role} for every shape in `shapes` (from shapes_of()). Fails
    safe to all "skip" on any error -- never blocks the confirm screen, which
    just shows every box unassigned for the owner to set by hand."""
    if not shapes:
        return {}
    from deckengine.services.matching.ai_matcher import _client, MODEL

    lines = [
        "Below are the text boxes on ONE slide of a sales deck, in reading order "
        "(top to bottom, then left to right), each with its position/size in "
        "inches (left, top, width, height) and its current text.",
        "Assign each box ONE role from this list:",
        "  title              - the main slide heading",
        "  subtitle           - a client/domain line right under the heading "
        "(use AT MOST ONCE per slide -- a second box that looks like a section "
        "header belongs under 'skip', not a second subtitle)",
        "  challenge_body     - text describing a problem, pain point, or challenge",
        "  solution_body      - text describing the solution, approach, or advantage "
        "the vendor brings (an 'X vs Y' or 'before/after' slide's SECOND column is "
        "almost always this, not capability_title/body)",
        "  capability_title   - a short heading naming ONE SPECIFIC capability or "
        "feature the vendor delivers (there may be several) -- NOT a generic section "
        "label like 'KEY BENEFITS' or 'WHAT WE OFFER' that sits above a group of boxes",
        "  capability_body    - the one-line description under a capability_title",
        "  result_stat        - a big standalone number or metric (e.g. '45%', '2-4 Weeks')",
        "  result_caption     - the short caption under a result_stat",
        "  skip               - page numbers, logos/wordmarks, GENERIC SECTION LABELS "
        "that just organise a group of boxes below them (e.g. 'CURRENT PAIN POINTS', "
        "'J2W ADVANTAGE', 'KEY BENEFITS' -- these describe the GROUP, they are not "
        "content themselves), or anything that doesn't fit the case-study shape above",
        "",
        "A short ALL-CAPS heading sitting directly above several similar boxes is "
        "almost always a SECTION LABEL (skip), not itself a capability_title or "
        "result_stat -- look at what's grouped underneath it to decide.",
        "IMPORTANT: if the SAME role fits several boxes (e.g. six capability_title "
        "boxes, or four challenge_body bullets), assign it to all of them -- their "
        "given order will be preserved when real content is poured in later.",
        "",
    ]
    for s in shapes:
        lines.append(f"  idx={s['idx']}  pos=({s['left']},{s['top']},{s['width']},{s['height']})  "
                     f"text={s['text']!r}")
    lines.append('\nReturn ONLY JSON: {"roles": {"<idx>": "<role>", ...}} -- one entry per idx shown above.')

    valid_idx = {str(s["idx"]) for s in shapes}
    try:
        resp = _client().chat.completions.create(
            model=MODEL, temperature=0,
            response_format={"type": "json_object"},
            messages=[
                {"role": "system", "content": "You map slide text boxes onto a "
                 "case-study content schema. Reply with one JSON object only."},
                {"role": "user", "content": "\n".join(lines)},
            ],
        )
        data = json.loads(resp.choices[0].message.content)
        roles = data.get("roles") if isinstance(data, dict) else None
    except Exception:
        roles = None

    out = {}
    if isinstance(roles, dict):
        for k, v in roles.items():
            if k in valid_idx and v in ROLE_VOCAB:
                out[k] = v
    for s in shapes:
        out.setdefault(str(s["idx"]), "skip")   # every shape always has SOME role
    return out


# --------------------------------------------------------------------------- #
# 3) save / list / activate / delete learned templates
# --------------------------------------------------------------------------- #
def _load_index():
    try:
        with open(INDEX_PATH, encoding="utf-8") as f:
            return json.load(f)
    except (OSError, ValueError):
        return []


def _write_index(items):
    os.makedirs(config.LEARNED_TEMPLATES_DIR, exist_ok=True)
    with open(INDEX_PATH, "w", encoding="utf-8") as f:
        json.dump(items, f, ensure_ascii=False, indent=2)


def save_learned_template(deck_bytes, slide_index, role_map, name):
    """Extract ONE slide into its own single-slide .pptx (full fidelity --
    shapes, images, fonts, colours, positions all preserved exactly, same
    copy routine already used for skills/AI slides) and register its role map."""
    src_prs = Presentation(_stream(deck_bytes))
    src_slide = src_prs.slides[slide_index]

    new_prs = Presentation()
    new_prs.slide_width = src_prs.slide_width
    new_prs.slide_height = src_prs.slide_height
    slide_generator._copy_slide(new_prs, src_slide)

    tid = _slug(name) + "-" + str(int(time.time()))
    fname = tid + ".pptx"
    os.makedirs(config.LEARNED_TEMPLATES_DIR, exist_ok=True)
    new_prs.save(os.path.join(config.LEARNED_TEMPLATES_DIR, fname))

    row = {
        "id": tid, "name": (name or "").strip() or "Untitled template",
        "file": fname, "role_map": role_map, "active": False,
        "slide_w": new_prs.slide_width, "slide_h": new_prs.slide_height,
    }
    items = _load_index()
    items.insert(0, row)
    _write_index(items)
    return row


def all_templates():
    return _load_index()


def get(tid):
    return next((r for r in _load_index() if r["id"] == tid), None)


def active_template():
    return next((r for r in _load_index() if r.get("active")), None)


def set_active(tid):
    """Radio-button: activating one deactivates every other (the owner wants
    ONE consistent house style at a time, not per-content-type routing)."""
    items = _load_index()
    for r in items:
        r["active"] = (r["id"] == tid)
    _write_index(items)


def deactivate(tid):
    items = _load_index()
    for r in items:
        if r["id"] == tid:
            r["active"] = False
    _write_index(items)


def deactivate_all():
    items = _load_index()
    for r in items:
        r["active"] = False
    _write_index(items)


def delete(tid):
    items = _load_index()
    row = next((r for r in items if r["id"] == tid), None)
    if not row:
        return False
    try:
        os.remove(os.path.join(config.LEARNED_TEMPLATES_DIR, row["file"]))
    except OSError:
        pass
    _write_index([r for r in items if r["id"] != tid])
    return True


# --------------------------------------------------------------------------- #
# 4) fill a learned template with real case-study content
# --------------------------------------------------------------------------- #
def _pad(vals, n):
    vals = list(vals or [])[:n]
    return vals + [None] * (n - len(vals))


def _set_whole(shape, text):
    if not shape.has_text_frame:
        return
    tf = shape.text_frame
    if tf.paragraphs and tf.paragraphs[0].runs:
        tf.paragraphs[0].runs[0].text = text or ""
        for r in tf.paragraphs[0].runs[1:]:
            r.text = ""
        for p in tf.paragraphs[1:]:
            for r in p.runs:
                r.text = ""


def _apply_roles(slide, role_map, record):
    """Pour a case-study-shaped `record` (title/domain/client_descriptor/
    challenge/solution/capabilities/results) into `slide` by role. Shape index
    N in role_map corresponds to shape index N on this slide, since
    slide_generator._copy_slide() preserves shape order from the source."""
    by_role = {}
    for idx_str, role in (role_map or {}).items():
        if role == "skip":
            continue
        by_role.setdefault(role, []).append(int(idx_str))
    for role in by_role:
        by_role[role].sort()

    shapes = list(slide.shapes)

    def set_idx(idx, text):
        if 0 <= idx < len(shapes):
            _set_whole(shapes[idx], text)

    if by_role.get("title"):
        set_idx(by_role["title"][0], record.get("title", ""))

    if by_role.get("subtitle"):
        client = record.get("client_descriptor") or record.get("client", "")
        domain = record.get("domain", "")
        text = f"CLIENT: {client}  |  DOMAIN: {domain}" if (client or domain) else ""
        set_idx(by_role["subtitle"][0], text)

    # challenge / solution: the content may be a single string (the common
    # case, matches today's store schema) OR a list of bullet-worthy strings;
    # distributed across however many boxes the template actually has.
    for role, key in (("challenge_body", "challenge"), ("solution_body", "solution")):
        idxs = by_role.get(role, [])
        val = record.get(key, "")
        vals = val if isinstance(val, list) else ([val] if val else [])
        for i, idx in enumerate(idxs):
            set_idx(idx, vals[i] if i < len(vals) else "")

    # capabilities: repeatable {title, body} pairs
    cap_title_idxs = by_role.get("capability_title", [])
    cap_body_idxs = by_role.get("capability_body", [])
    n_caps = max(len(cap_title_idxs), len(cap_body_idxs))
    caps = _pad(record.get("capabilities") or [], n_caps)
    for i, idx in enumerate(cap_title_idxs):
        c = caps[i] if i < len(caps) else None
        t, _b = split_capability(c) if c else ("", "")
        set_idx(idx, t)
    for i, idx in enumerate(cap_body_idxs):
        c = caps[i] if i < len(caps) else None
        _t, b = split_capability(c) if c else ("", "")
        set_idx(idx, b)

    # results: repeatable {stat, caption} pairs (reuse the same highlight/
    # caption splitter every other J2W surface uses, so a plain results
    # sentence renders identically here too)
    stat_idxs = by_role.get("result_stat", [])
    cap_idxs = by_role.get("result_caption", [])
    n_res = max(len(stat_idxs), len(cap_idxs))
    results = _pad(record.get("results") or [], n_res)
    for i in range(n_res):
        top, bottom = split_result(results[i]) if results[i] else ("", "")
        if i < len(stat_idxs):
            set_idx(stat_idxs[i], top)
        if i < len(cap_idxs):
            set_idx(cap_idxs[i], bottom)


def fill_into(dest_prs, template_row, record):
    """Copy the learned slide into dest_prs (a Presentation already open for
    writing) and fill it by role. Returns the new slide."""
    tpl_path = os.path.join(config.LEARNED_TEMPLATES_DIR, template_row["file"])
    tpl_prs = Presentation(tpl_path)
    src_slide = tpl_prs.slides[0]
    new_slide = slide_generator._copy_slide(dest_prs, src_slide)
    _apply_roles(new_slide, template_row.get("role_map") or {}, record)
    return new_slide
