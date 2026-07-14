# -*- coding: utf-8 -*-
"""
image_normalize.py  --  make every picture in an uploaded deck a RASTER image.

An uploaded .pptx can carry pictures in formats the rest of the pipeline can't read:
  * SVG  -- modern icon sets. Often the PRIMARY image (the <a:blip> has no r:embed at
    all, only an <asvg:svgBlip> pointing at the .svg), so there is no PNG fallback to
    fall back to.
  * WDP  -- Windows HD Photo / JPEG-XR, used by Office "artistic effect" layers
    (<a14:imgLayer>).
  * EMF / WMF -- vector metafiles.
Pillow (and therefore python-pptx 1.x, which reads every image's size through Pillow)
cannot open any of these, so they were silently dropped -- icons vanished and, for a
primary SVG, the whole picture went blank (owner-reported, 2026-07-13,
'ServiceNow HTDD.pptx': slides whose icons are primary SVGs).

`normalize_deck(prs)` rewrites each picture so its displayed image is a PNG: it rasterises
the SVG / decodes the WDP / converts the metafile, registers the PNG in the slide, points
the <a:blip> at it, and strips the now-defunct svg/imgLayer extension. Everything
downstream (icon extraction, slide copy, restyle) then sees an ordinary raster and NOTHING
is missed, whatever the source format.

Conversion uses whatever tool is present, best first, and DEGRADES GRACEFULLY: with no
converter available a picture is simply left as-authored (the copy path still won't crash
-- see slide_generator._copy_slide). Nothing here is required for a deck that already
ships PNG/JPEG.
"""

import io
import os
import shutil
import subprocess
import tempfile

from pptx.enum.shapes import MSO_SHAPE_TYPE

# --- OOXML namespaces we touch (full URIs -- python-pptx's nsmap doesn't know asvg/a14) --
_A   = "http://schemas.openxmlformats.org/drawingml/2006/main"
_R   = "http://schemas.openxmlformats.org/officeDocument/2006/relationships"
_SVG = "http://schemas.microsoft.com/office/drawing/2016/SVG/main"
_A14 = "http://schemas.microsoft.com/office/drawing/2010/main"

_BLIP     = "{%s}blip" % _A
_EXTLST   = "{%s}extLst" % _A
_EMBED    = "{%s}embed" % _R
_SVGBLIP  = "{%s}svgBlip" % _SVG
_IMGLAYER = "{%s}imgLayer" % _A14

_RASTER = {"png", "jpeg", "gif", "bmp", "tiff"}


def _which(name):
    return shutil.which(name)


def _magic(blob):
    """Best-effort image format from the blob's leading bytes."""
    if not blob:
        return None
    b = blob[:16]
    if b[:8].startswith(b"\x89PNG"):
        return "png"
    if b[:3] == b"\xff\xd8\xff":
        return "jpeg"
    if b[:6] in (b"GIF87a", b"GIF89a"):
        return "gif"
    if b[:2] == b"BM":
        return "bmp"
    if b[:4] in (b"II*\x00", b"MM\x00*"):
        return "tiff"
    if b[:3] == b"II\xbc":                       # JPEG-XR / HD Photo (.wdp)
        return "wdp"
    if blob[40:44] == b" EMF":                   # EMF metafile
        return "emf"
    if b[:4] in (b"\x01\x00\x09\x00", b"\x02\x00\x09\x00") or b[:4] == b"\xd7\xcd\xc6\x9a":
        return "wmf"                             # WMF (standard / placeable)
    head = blob[:600].lstrip()
    if head[:4] == b"<svg" or b"<svg" in blob[:600]:
        return "svg"
    return None


def _pil_to_png(blob):
    try:
        from PIL import Image
        im = Image.open(io.BytesIO(blob)).convert("RGBA")
        out = io.BytesIO()
        im.save(out, "PNG")
        return out.getvalue()
    except Exception:
        return None


def _run(cmd, timeout=60):
    subprocess.run(cmd, check=True, capture_output=True, timeout=timeout)


def _svg_to_png(blob):
    # 1) cairosvg -- best fidelity, no external process (only if installed)
    try:
        import cairosvg
        return cairosvg.svg2png(bytestring=blob, output_width=512)
    except Exception:
        pass
    # 2) ImageMagick -- crisp for the flat icons these decks use
    if _which("convert"):
        try:
            with tempfile.TemporaryDirectory() as d:
                sp, pp = os.path.join(d, "i.svg"), os.path.join(d, "o.png")
                with open(sp, "wb") as f:
                    f.write(blob)
                _run(["convert", "-background", "none", "-density", "400", sp, pp], timeout=30)
                if os.path.getsize(pp) > 0:
                    with open(pp, "rb") as f:
                        return f.read()
        except Exception:
            pass
    # 3) LibreOffice -- always available where slides are rendered
    return _soffice_to_png(blob, "svg")


def _wdp_to_png(blob):
    # JxrDecApp (libjxr-tools) decodes JPEG-XR -> TIFF, which Pillow can read
    if _which("JxrDecApp"):
        try:
            with tempfile.TemporaryDirectory() as d:
                wp, tp = os.path.join(d, "i.wdp"), os.path.join(d, "o.tif")
                with open(wp, "wb") as f:
                    f.write(blob)
                _run(["JxrDecApp", "-i", wp, "-o", tp], timeout=30)
                with open(tp, "rb") as f:
                    return _pil_to_png(f.read())
        except Exception:
            pass
    return _pil_to_png(blob)             # Pillow with an imagecodecs plugin, if present


def _soffice_to_png(blob, ext):
    soffice = _which("libreoffice") or _which("soffice")
    if not soffice:
        return None
    try:
        with tempfile.TemporaryDirectory() as d:
            sp = os.path.join(d, "i." + ext)
            with open(sp, "wb") as f:
                f.write(blob)
            _run([soffice, "--headless", "--convert-to", "png", "--outdir", d, sp], timeout=90)
            pp = os.path.join(d, "i.png")
            if os.path.exists(pp):
                with open(pp, "rb") as f:
                    return f.read()
    except Exception:
        pass
    return None


def to_png(blob):
    """PNG bytes for any image blob, or None if no available tool can convert it.
    A blob that is already a raster is normalised through Pillow (PNG passthrough)."""
    fmt = _magic(blob)
    if fmt == "png":
        return blob
    if fmt in _RASTER:
        return _pil_to_png(blob)
    if fmt == "svg":
        return _svg_to_png(blob)
    if fmt == "wdp":
        return _wdp_to_png(blob)
    if fmt in ("emf", "wmf"):
        return _soffice_to_png(blob, fmt)
    return _pil_to_png(blob)             # unknown magic -- let Pillow try, else None


def _iter_image_shapes(shapes):
    """Every shape carrying an <a:blip> -- a <p:pic> picture OR an autoshape whose FILL is
    a picture (blipFill in spPr). The latter is how many decks store their icons, and was
    being missed entirely (owner-reported, 2026-07-13: slides whose icons are shape fills)."""
    for sh in shapes:
        try:
            if sh.shape_type == MSO_SHAPE_TYPE.GROUP:
                yield from _iter_image_shapes(sh.shapes)
                continue
            if next(iter(sh._element.iter(_BLIP)), None) is not None:
                yield sh
        except Exception:
            continue


def _strip_vector_exts(blip):
    """Remove the <a:ext> wrappers that carry an <asvg:svgBlip> or <a14:imgLayer> -- once
    the main blip points at a PNG they are defunct, and leaving a dangling one can make
    PowerPoint offer to 'repair' the file."""
    for extlst in list(blip):
        if extlst.tag != _EXTLST:
            continue
        for ext in list(extlst):
            if any(e.tag in (_SVGBLIP, _IMGLAYER) for e in ext.iter()):
                extlst.remove(ext)
        if len(extlst) == 0:
            blip.remove(extlst)


def _normalize_shape(sh):
    """Point every <a:blip> in a shape (picture or picture-fill) at a raster PNG. Returns
    True if it converted at least one."""
    part = sh.part
    converted = False
    for blip in list(sh._element.iter(_BLIP)):
        try:
            if _normalize_blip(blip, part):
                converted = True
        except Exception:
            continue
    return converted


def _normalize_blip(blip, part):
    """Point one <a:blip> at a raster PNG. Returns True if it converted one."""
    def blob_of(rid):
        if not rid:
            return None
        try:
            return part.related_part(rid).blob
        except Exception:
            return None

    main_rid = blip.get(_EMBED)
    main_blob = blob_of(main_rid)

    # Already a usable raster main image -> just drop the (defunct) svg/effect extension.
    if main_blob is not None and _magic(main_blob) in _RASTER:
        _strip_vector_exts(blip)
        return False

    # Otherwise source a raster: from a non-raster main blip, else from the svgBlip.
    svgblip = next(iter(blip.iter(_SVGBLIP)), None)
    svg_rid = svgblip.get(_EMBED) if svgblip is not None else None
    src_blob = main_blob if main_blob is not None else blob_of(svg_rid)
    if src_blob is None:
        _strip_vector_exts(blip)
        return False

    png = to_png(src_blob)
    if not png:
        return False                      # no converter available -- leave as-authored

    _img_part, new_rid = part.get_or_add_image_part(io.BytesIO(png))
    blip.set(_EMBED, new_rid)             # promote the raster to the picture's main image
    _strip_vector_exts(blip)
    return True


def normalize_deck(prs):
    """Convert every non-raster picture in the deck to PNG, in place. Returns the number
    converted. Fully fail-safe: any picture that can't be converted is left untouched."""
    converted = 0
    for slide in prs.slides:
        for sh in _iter_image_shapes(slide.shapes):
            try:
                if _normalize_shape(sh):
                    converted += 1
            except Exception:
                continue
    return converted
