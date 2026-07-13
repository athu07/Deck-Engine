# -*- coding: utf-8 -*-
"""templates.py  --  the /templates page: fill-on-demand templates, the "learn a
template" flow, and the deck re-skinner (upload -> J2W-branded -> download)."""

import io
import os
import re
import shutil
import time
import uuid

from flask import Blueprint, render_template, request, redirect
from pptx import Presentation

from deckengine import config
from deckengine.services.rendering import slide_generator, reskin, templatize, recreate
from .view_helpers import shell

bp = Blueprint("templates", __name__)

# temp holding area for an uploaded deck while the owner picks a slide + reviews
# the AI-proposed role mapping (the "learn a template" multi-step flow) -- keyed
# by a random token threaded through the form, cleaned up once saved/cancelled.
_PENDING_DIR = os.path.join(config.LEARNED_TEMPLATES_DIR, "_pending")


def _slide_shapes():
    """Every slide shape the engine can build, straight from the CONTENT_TEMPLATES
    registry -- the single source of truth the paste-a-document flow, the Custom Slide
    Builder, and Recreate-with-AI all classify against.

    This page used to list the tagged slides in the legacy `templates.pptx` instead: a
    stale file holding five tags, so it never showed four_box, box_grid, or any of the
    ten style-guide shapes. It was describing a file, not the app.
    """
    from deckengine.services.rendering import slide_schema

    frames = set()
    try:                                # which shapes have a template frame of their own
        frames = set(slide_generator.list_templates(config.SKILLS_TEMPLATES_PPTX))
    except Exception:
        pass

    out = []
    for t in slide_generator.CONTENT_TEMPLATES:
        key = t["key"]
        fields = slide_schema.fields_for(key)
        # the repeating part of the shape is what the salesperson actually fills in
        holds = [f["label"] for f in fields if f["type"] in ("objects", "strings")]
        out.append({
            "key": key,
            "label": t["label"],
            "what": t["classify_desc"],
            "holds": holds,
            "fields": len(fields),
            # case_study fills the branded case_study_v2 file, not a drawn frame
            "has_frame": key in frames or key == "case_study",
        })
    return out


@bp.route("/templates")
def templates_page():
    body = render_template("templates_page.html", items=_slide_shapes(),
                           learned=templatize.all_templates())
    return shell(body, active="templates", crumb="<b>Templates</b>")


# ── F2: re-skin an uploaded deck into the J2W design ──────────────────────────
def _slug(s):
    return re.sub(r"[^A-Za-z0-9_-]+", "_", s or "").strip("_") or "Deck"


# ── F2: "Re-skin a deck" -- DISABLED (owner's spec, 2026-07-13).
#    Re-skin restyles the uploaded deck's shapes in place. Recreate-with-AI now does
#    strictly more (it rebuilds each slide's LAYOUT to fit its content, and falls back to
#    exactly this restyle for a slide that doesn't fit) -- so the two overlapped, and
#    Re-skin was retired as a user feature. reskin.py STAYS: Recreate leans on its
#    restyle helpers (17 call sites) and its PNG renderer. To restore the feature,
#    uncomment this route and its card in templates_page.html.
#
# @bp.route("/template/reskin", methods=["POST"])
# def template_reskin():
#     f = request.files.get("deck_file")
#     if not f or not getattr(f, "filename", ""):
#         return redirect("/templates")
#     data = f.read()
#     try:                                          # reject non-pptx up front
#         Presentation(io.BytesIO(data))
#     except Exception:
#         body = render_template("templates_page.html", items=[],
#                                reskin_error="Couldn't read that file — please upload a valid .pptx.")
#         return shell(body, active="templates", crumb="<b>Templates</b>")
#     # rebrand the original deck in place (nothing dropped), then render true previews
#     os.makedirs(config.OUTPUT_DIR, exist_ok=True)
#     fname = "Reskinned_" + _slug(os.path.splitext(f.filename)[0]) + ".pptx"
#     out = os.path.join(config.OUTPUT_DIR, fname)
#     try:
#         reskin.restyle_deck(data, out)
#     except (PermissionError, OSError) as e:
#         from .view_helpers import file_busy_page
#         return file_busy_page(str(e))
#     rdir = os.path.join(config.RENDERS_DIR, "reskin")
#     shutil.rmtree(rdir, ignore_errors=True)
#     try:
#         pngs = reskin.render_pngs(out, rdir)
#     except Exception:
#         pngs = []
#     images = ["/static/renders/reskin/" + os.path.basename(p) for p in pngs]
#     nslides = len(Presentation(out).slides._sldIdLst)
#     body = render_template("reskin_preview.html", images=images, nslides=nslides,
#                            filename=fname, src_name=f.filename, mode="reskin")
#     return shell(body, active="templates", crumb="<b>Templates</b> / Re-skin preview")


# ── F5: "Recreate with AI" -- rebuild an uploaded deck in J2W's own layouts ────
@bp.route("/template/recreate", methods=["POST"])
def template_recreate():
    f = request.files.get("deck_file")
    if not f or not getattr(f, "filename", ""):
        return redirect("/templates")
    data = f.read()
    try:                                          # reject non-pptx up front
        Presentation(io.BytesIO(data))
    except Exception:
        body = render_template("templates_page.html", items=[],
                               reskin_error="Couldn't read that file — please upload a valid .pptx.")
        return shell(body, active="templates", crumb="<b>Templates</b>")
    industry = request.form.get("industry", "").strip()
    os.makedirs(config.OUTPUT_DIR, exist_ok=True)
    fname = "Recreated_" + _slug(os.path.splitext(f.filename)[0]) + ".pptx"
    out = os.path.join(config.OUTPUT_DIR, fname)
    try:
        _out, stats = recreate.recreate_deck(data, out, industry=industry)
    except (PermissionError, OSError) as e:
        from .view_helpers import file_busy_page
        return file_busy_page(str(e))
    rdir = os.path.join(config.RENDERS_DIR, "recreate")
    shutil.rmtree(rdir, ignore_errors=True)
    try:
        pngs = reskin.render_pngs(out, rdir)
    except Exception:
        pngs = []
    images = ["/static/renders/recreate/" + os.path.basename(p) for p in pngs]
    nslides = len(Presentation(out).slides._sldIdLst)
    stats_line = (f"{stats['recreated']} slide(s) redesigned, "
                 f"{stats['restyled']} restyled as-authored (didn't fit a "
                 f"content shape), {stats['skipped']} replaced with our own "
                 "title/closing.")
    body = render_template("reskin_preview.html", images=images, nslides=nslides,
                           filename=fname, src_name=f.filename, stats_line=stats_line,
                           report=stats.get("report", []), mode="recreate")
    return shell(body, active="templates", crumb="<b>Templates</b> / Recreate preview")


# ── F4: "learn a template" -- upload -> pick a slide -> AI-propose roles ──────
#       -> owner confirms/fixes -> save. See templatize.py for the full design
#       rationale (why this needs a human-confirm step before it's ever used).
def _pending_path(token):
    return os.path.join(_PENDING_DIR, token + ".pptx")


def _sweep_stale_pending(max_age_seconds=6 * 3600):
    """An abandoned upload (picked a slide, never saved) leaves a temp file
    behind forever otherwise -- swept opportunistically on the next upload
    rather than a scheduled job, since this is a low-traffic, low-stakes flow."""
    try:
        if not os.path.isdir(_PENDING_DIR):
            return
        now = time.time()
        for fn in os.listdir(_PENDING_DIR):
            p = os.path.join(_PENDING_DIR, fn)
            if now - os.path.getmtime(p) > max_age_seconds:
                os.remove(p)
    except OSError:
        pass


@bp.route("/templatize/upload", methods=["POST"])
def templatize_upload():
    f = request.files.get("deck_file")
    if not f or not getattr(f, "filename", ""):
        return redirect("/templates")
    data = f.read()
    try:
        Presentation(io.BytesIO(data))   # reject non-pptx up front
    except Exception:
        return redirect("/templates")
    _sweep_stale_pending()
    os.makedirs(_PENDING_DIR, exist_ok=True)
    token = uuid.uuid4().hex
    with open(_pending_path(token), "wb") as out:
        out.write(data)
    return redirect(f"/templatize/pick/{token}")


@bp.route("/templatize/pick/<token>")
def templatize_pick(token):
    p = _pending_path(token)
    if not os.path.exists(p):
        return redirect("/templates")
    with open(p, "rb") as f:
        data = f.read()
    slides = templatize.list_slides(data)
    body = render_template("templatize_pick.html", token=token, slides=slides)
    return shell(body, active="templates", crumb="<b>Templates</b> / Learn a template")


@bp.route("/templatize/review/<token>/<int:slide_index>")
def templatize_review(token, slide_index):
    p = _pending_path(token)
    if not os.path.exists(p):
        return redirect("/templates")
    with open(p, "rb") as f:
        data = f.read()
    shapes = templatize.shapes_of(data, slide_index)
    roles = templatize.propose_roles(shapes)
    rows = [dict(s, role=roles.get(str(s["idx"]), "skip")) for s in shapes]
    body = render_template("templatize_review.html", token=token, slide_index=slide_index,
                           rows=rows, role_vocab=templatize.ROLE_VOCAB,
                           role_labels=templatize.ROLE_LABELS)
    return shell(body, active="templates", crumb="<b>Templates</b> / Confirm mapping")


@bp.route("/templatize/save", methods=["POST"])
def templatize_save():
    token = request.form.get("token", "")
    slide_index = int(request.form.get("slide_index", "0"))
    name = request.form.get("name", "").strip()
    p = _pending_path(token)
    if not os.path.exists(p) or not name:
        return redirect("/templates")
    with open(p, "rb") as f:
        data = f.read()
    role_map = {}
    for key, val in request.form.items():
        if key.startswith("role__") and val in templatize.ROLE_VOCAB:
            role_map[key[len("role__"):]] = val
    templatize.save_learned_template(data, slide_index, role_map, name)
    try:
        os.remove(p)
    except OSError:
        pass
    return redirect("/templates")


@bp.route("/templatize/<tid>/activate", methods=["POST"])
def templatize_activate(tid):
    templatize.set_active(tid)
    return redirect("/templates")


@bp.route("/templatize/<tid>/deactivate", methods=["POST"])
def templatize_deactivate(tid):
    templatize.deactivate(tid)
    return redirect("/templates")


@bp.route("/templatize/<tid>/delete", methods=["POST"])
def templatize_delete(tid):
    templatize.delete(tid)
    return redirect("/templates")
