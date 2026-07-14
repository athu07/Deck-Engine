# -*- coding: utf-8 -*-
"""output.py  --  serve built decks (/output/<file>) and single-slide downloads."""

import io
import os
import tempfile

from flask import Blueprint, request, send_file, abort

from pptx import Presentation

from deckengine import config
from deckengine.constants import OUTPUT_DIR
from deckengine.services.content import case_library
from deckengine.services.rendering import assembler, fill_case_study, templatize
from .view_helpers import file_busy_page

bp = Blueprint("output", __name__)


def _render_slide(sid, path):
    """Render one slide id to `path` as its own .pptx. Content-store cases (AIP/WFS/MSS)
    render from the owner's ACTIVE learned template if one exists, else the built-in
    case_study_v2; master ids build from the master deck. Returns True on success,
    False if the id yields no slide."""
    rec = case_library.record(sid)
    if rec is not None:
        active = templatize.active_template()
        if active:
            prs = Presentation()
            prs.slide_width, prs.slide_height = active["slide_w"], active["slide_h"]
            templatize.fill_into(prs, active, rec)
            prs.save(path)
        else:
            fill_case_study.fill_row(rec, path)
        return True
    kept, _ = assembler.build_deck([sid], out=path)
    return bool(kept)


@bp.route("/favicon.ico")
def favicon():
    """Browsers ask for /favicon.ico at the root whatever the <link> says; without this
    every page load logged a 404."""
    return send_file(os.path.join(config.PROJECT_ROOT, "static", "favicon.ico"),
                     mimetype="image/vnd.microsoft.icon")


@bp.route("/output/<path:fname>")
def output_file(fname):
    fname = os.path.basename(fname)
    path = os.path.join(OUTPUT_DIR, fname)
    if not os.path.exists(path):
        abort(404)
    return send_file(path, as_attachment=bool(request.args.get("dl")), download_name=fname)


@bp.route("/slide/<sid>/preview.png")
def slide_preview(sid):
    """A rendered image of one master-deck slide, for the review page. Lazy: the browser
    asks for it, so the page paints immediately and the pictures fill in. The first ask
    renders the whole master deck once (~8.5s); every ask after that is a file read."""
    from deckengine.services.rendering import preview
    path = preview.master_slide_png(sid.upper())
    if not path:
        abort(404)                       # no preview available -> the page falls back to text
    mime = "image/webp" if path.endswith(".webp") else "image/png"
    return send_file(path, mimetype=mime, max_age=86400)


@bp.route("/slide/<sid>/download")
def slide_download(sid):
    sid = sid.upper()
    os.makedirs(OUTPUT_DIR, exist_ok=True)
    path = os.path.join(OUTPUT_DIR, f"Slide_{sid}.pptx")
    try:
        if not _render_slide(sid, path):
            abort(404)
    except (PermissionError, OSError) as e:
        return file_busy_page(str(e))
    return send_file(path, as_attachment=True, download_name=f"{sid}.pptx")


@bp.route("/slides/download")
def slides_download():
    """Bulk download: several library slides selected at once, combined into ONE .pptx in
    the selected order (owner's spec, 2026-07-14 -- previously each slide came as its own
    file inside a .zip). ids come in as a comma-separated ?ids= list. Served from memory so
    there's no shared file to lock between concurrent downloads."""
    from deckengine.services.rendering import slide_generator
    ids = [x.strip().upper() for x in request.args.get("ids", "").split(",") if x.strip()]
    if not ids:
        abort(400)
    dest = None
    added = 0
    with tempfile.TemporaryDirectory() as td:
        for sid in ids:
            tmp = os.path.join(td, f"{sid}.pptx")
            try:
                if not _render_slide(sid, tmp):
                    continue
            except (PermissionError, OSError) as e:
                return file_busy_page(str(e))
            src = Presentation(tmp)                 # each id rendered to its own 1-slide deck
            if dest is None:                        # combined deck takes the first slide's size
                dest = Presentation()
                dest.slide_width, dest.slide_height = src.slide_width, src.slide_height
            for s in src.slides:                    # copy every slide across, in order
                slide_generator._copy_slide(dest, s)
                added += 1
    if dest is None or not added:
        abort(404)
    mem = io.BytesIO()
    dest.save(mem)
    mem.seek(0)
    return send_file(mem, as_attachment=True, download_name="J2W_slides.pptx",
                     mimetype="application/vnd.openxmlformats-officedocument.presentationml.presentation")
