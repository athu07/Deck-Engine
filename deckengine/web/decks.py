# -*- coding: utf-8 -*-
"""
decks.py  --  the core deck-building flow.

/ (or /new) -> /build (pick slides) -> /review (edit) -> /finalize (assemble .pptx),
plus /deck to resume a deck-in-progress held in the browser's deck tray.
"""

import json
import logging
import re
import uuid
from zipfile import BadZipFile

logger = logging.getLogger(__name__)

from flask import Blueprint, request, render_template, abort, jsonify
from pptx import Presentation

from deckengine import config
from deckengine import constants
from deckengine.constants import (COVERAGE_THRESHOLD, CAPABILITY_COVER, _GENERIC_NEEDS,
                                  INDUSTRIES, FUNCTIONS, PHASES)
from deckengine.services.matching import matcher, relevance, ai_matcher, personas
from deckengine.services.matching.tagger import INDUSTRY as _BUILTIN_INDUSTRIES
from deckengine.services.rendering import (skills, staging, deck_build, fill_case_study,
                                           slide_generator, slide_schema, client_context)
from deckengine.services.content import case_library, editor
from deckengine.services.content.content_store import content_store
from deckengine.services.content.build_library import read_id
from deckengine.services.rendering import assembler
from deckengine.services import ingest as research
from deckengine.services import meeting_log
from deckengine.services import build_context
from .view_helpers import (shell, file_busy_page, current_salesperson,
                           legacy_case_ids as _legacy_case_ids,
                           resolve_industry, remember_custom_industry)

bp = Blueprint("decks", __name__)

# Two priority-picked cases whose case-to-case cosine is >= this are near-TWINS --
# the second is skipped so the deck never shows two look-alike proof points. Calibrated
# from the real store (2026-07-29): genuine duplicates land at 0.85-0.91 (e.g. identical
# DevSecOps-open-banking cases 0.91; WFS002/WFS025 workforce-scaling 0.87), while
# genuinely DISTINCT talent picks (RPO vs HTD vs accelerated-hiring) top out at ~0.74 --
# so 0.85 drops only true duplicates, never a real capability.
PRIORITY_DEDUP_SIM = 0.85

# Intro-deck per-work-type cap (owner-spec, 2026-07-31): on an Intro deck, cap each
# SELECTED work type (Workforce/AI Pods/MS) at this many case studies -- a CEILING with a
# quality gate, not a floor. The greedy pick loops already process candidates best-first
# (sorted by adj, which folds in cosine + industry + persona + function fit), so capping
# simply keeps the N best per work type and lets weaker ones fall away; nothing pads a
# work type back UP to this count if fewer genuinely clear the bar. Applies to literal
# priority-picks AND strategic bets COMBINED (one running count per work type, not a
# separate allowance for each pass). Only Intro -- First/Second Meeting and Proposal decks
# are unaffected. A need capped out this way is NOT shown as a "not in our library" gap
# (a real case exists for it; there just wasn't room) -- see `capped_needs` below.
INTRO_WORKTYPE_CAP = 3

# Gap/bet reconciliation (2026-07-31): the cheap embedding pre-filter before asking the
# LLM whether an already-picked case substantively covers a flagged gap (ai_matcher.
# resolve_gap_overlap). Deliberately LOW and inclusive -- this only decides whether a
# gap is worth ASKING about at all; the real yes/no comes from the grounded LLM read of
# the candidate's actual challenge/solution text, not this number. Calibrated so all
# three real Broadridge candidates (cosine 0.401 / 0.398 / 0.353) reach the LLM step.
GAP_PREFILTER_COVER = 0.30

# the pasted shapes that render through the shared inline-editor macro on /review
# (a case study has its own richer card; library slides have their own title/subtitle form)
_SHAPE_KINDS = tuple(k for k in slide_schema.SCHEMA if k != "case_study")


def _is_avoided(rec, avoid):
    """True if this content-store case is ABOUT a mismatch-flagged capability (so it must
    not be picked to answer a need)."""
    if not avoid or not rec:
        return False
    kw = rec.get("keywords") or []
    text = rec.get("title", "") + " " + (" ".join(kw) if isinstance(kw, list) else str(kw))
    head = relevance._tokens(text)
    for a in avoid:
        cap = a.get("capability", "") if isinstance(a, dict) else str(a)
        terms = relevance.specific_terms(cap)
        if terms and len(terms & head) >= max(1, len(terms) // 2):
            return True
    return False


@bp.route("/")
@bp.route("/new")
def home():
    try:
        lib_count = len(json.load(open(config.TAGGED_LIBRARY_JSON, encoding="utf-8")))
    except Exception:
        lib_count = 0
    body = render_template("new_form.html", industries=constants.all_industries(), functions=FUNCTIONS,
                                  phases=PHASES, library_count=lib_count, error="",
                                  content_templates=slide_generator.CONTENT_TEMPLATES)
    return shell(body, active="new", crumb="<b>New deck</b> / Context")


def _resume_page(reopen_seed=None):
    """Shared render for /deck (empty browser-tray resume) and /deck/reopen
    (server-persisted version reload) — both hydrate the SAME client-side
    resume mechanism (build.js reads localStorage j2w_deck); reopen_seed just
    pre-seeds that storage before build.js's init() runs (see build.html)."""
    titles = matcher._title_lookup()
    all_slides = sorted(((sid, t) for sid, t in titles.items()
                         if sid not in _legacy_case_ids()),
                        key=lambda kv: matcher._num(kv[0]))
    titles.update(case_library.title_map())   # so resumed store cases show their title
    case_lib = sorted(case_library.all_cases(), key=lambda c: (c["work_type"], c["title"]))
    empty_ctx = {"client_name": "", "industry": "", "transcript": "",
                 "phase": "", "recipient": "", "functions": [], "work_types": []}
    body = render_template("build.html", ctx=empty_ctx, picks=[], gaps=[],
                                  titles=titles, all_slides=all_slides, case_lib=case_lib,
                                  suggestions=[], suggested=[], ai_used=False,
                                  persona_labels=[], rationale=[], missing=[],
                                  research_read=False, research_failed=False, brief_weak=False,
                                  resume=True, build_id="", reopen_seed=reopen_seed)
    return shell(body, active="new", crumb="<b>New deck</b> / Your deck")


@bp.route("/deck")
def deck_resume():
    """Re-open the deck-in-progress (held in the browser's deck tray). The list and
    context are hydrated client-side from localStorage; the server just supplies
    the slide catalogue."""
    return _resume_page()


@bp.route("/deck/reopen")
def deck_reopen():
    """Re-open an already-finalized, already-downloaded deck for further editing
    (add/remove/reorder a slide, e.g. one more case study), continuing from its
    LATEST version. Finalizing again produces the next version — nothing is
    overwritten. If no version exists yet for this client+phase, falls back to
    a plain empty resume (same as /deck)."""
    client = request.args.get("client", "").strip()
    phase = request.args.get("phase", "").strip()
    rec, latest = meeting_log.latest_version(client, phase)
    if not latest:
        return _resume_page()
    seed = {
        "ctx": {
            "client_name": client,
            "industry": rec.get("industry", ""),
            "transcript": "",           # not persisted per-version; re-enter if needed
            "phase": phase,
            "recipient": latest.get("recipient", ""),
            "functions": latest.get("functions", []),
            "work_types": latest.get("work_types", []),
        },
        "order": latest.get("slide_ids", []),
        "active": True,
    }
    return _resume_page(reopen_seed=seed)


@bp.route("/build", methods=["POST"])
def build():
    ctx = {
        "client_name": request.form.get("client_name", "Client").strip(),
        # "Other…" -> the typed industry, never the raw sentinel (see resolve_industry)
        "industry": resolve_industry(request.form),
        "work_types": request.form.getlist("work_types"),
        "functions": request.form.getlist("functions"),
        "phase": request.form.get("phase", "").strip(),
        "recipient": request.form.get("recipient", "").strip(),
        "salesperson": current_salesperson(),
        "transcript": request.form.get("transcript", "").strip(),
    }
    # a free-typed "Other" industry (not in the built-in taxonomy) is remembered
    # so it shows up in the dropdown for every build after this one
    remember_custom_industry(ctx["industry"])
    # optional deep-research file (PDF/text) -> read alongside the notes for matching
    research_text = research.extract_text(request.files.get("research_file"))
    research_given = bool(request.files.get("research_file") and
                          getattr(request.files.get("research_file"), "filename", ""))
    research_failed = research_given and not research_text
    # optional STAKEHOLDER PROFILE (LinkedIn/bio) -> drives function/skill matching
    profile_text = research.extract_text(request.files.get("profile_file"))
    # optional AUTO-RESEARCH (the strategic brief from /research_account) -- the
    # salesperson reviewed/edited it client-side before submitting. Combined
    # ADDITIVELY with any uploaded research file text, never replacing it
    # (owner's spec, 2026-07-08: keep the manual upload option, add
    # auto-research alongside it).
    auto_brief = request.form.get("auto_company_text", "").strip()
    if auto_brief:
        research_text = (research_text + "\n\n" + auto_brief).strip() if research_text else auto_brief
    research_read = bool(research_text)                 # got text, from file and/or auto-research
    # optional CLIENT LOGO -> stamped onto the title slide next to the J2W wordmark
    # at finalize time (client_logo.stamp_into, called from deck_build.assemble).
    # Saved now (keyed by client name) so it survives all the way to /finalize
    # without needing new state threaded through /review's form. A manual
    # upload wins if both are present; otherwise, if "Find logo" was clicked
    # and confirmed a real image (see /logo_preview), its bytes ride along as
    # a data URI so /build uses EXACTLY what was previewed, not a re-rolled
    # search result.
    logo_file = request.files.get("client_logo_file")
    if logo_file and getattr(logo_file, "filename", ""):
        from deckengine.services.rendering import client_logo
        client_logo.save(ctx["client_name"], logo_file.read())
    else:
        data_uri = request.form.get("client_logo_data_uri", "").strip()
        if data_uri.startswith("data:") and "," in data_uri:
            import base64
            from deckengine.services.rendering import client_logo
            try:
                client_logo.save(ctx["client_name"], base64.b64decode(data_uri.split(",", 1)[1]))
            except Exception:
                pass
    match_notes = ctx["transcript"]
    if research_text:
        match_notes = (ctx["transcript"] + "\n\n[DEEP RESEARCH BRIEF]\n"
                       + research_text).strip()
    # Persist this build's full context (deep research + profile + FULL transcript),
    # keyed by build_id, so the AI case-study generator can synthesise from it later.
    build_id = uuid.uuid4().hex
    build_context.save(build_id, {
        "research": research_text, "profile": profile_text,
        "transcript": ctx["transcript"], "industry": ctx["industry"],
        "recipient": ctx["recipient"], "functions": ctx["functions"],
        "client_name": ctx["client_name"],
    })
    # Backstop for the browser's at-least-one-work-type check (e.g. JS disabled).
    if not ctx["work_types"]:
        try:
            lib_count = len(json.load(open(config.TAGGED_LIBRARY_JSON, encoding="utf-8")))
        except Exception:
            lib_count = 0
        body = render_template("new_form.html", industries=constants.all_industries(), functions=FUNCTIONS,
                                      phases=PHASES, library_count=lib_count,
                                      error="Please select at least one work type.",
                                      content_templates=slide_generator.CONTENT_TEMPLATES)
        return shell(body, active="new", crumb="<b>New deck</b> / Context")
    # The DEEP RESEARCH (or the notes) names the account's real interests. Extract
    # them, then split by coverage: each we HAVE a case for -> a PRIORITY pick
    # (pinned ahead of generic mail-thread matches); each we DON'T -> a build gap.
    # The account's needs come from BOTH the research brief AND the stakeholder's
    # profile (their function / current-role mandate) — balanced. For each need,
    # find OUR best case OF THE SELECTED WORK TYPES via a DIRECT skill->title match:
    # covered -> that case LEADS the deck; not covered -> a gap ("want to generate").
    priority_ids, missing = [], []
    profile_needs = []               # needs (name+description) — feeds lead_research
    avoid = []                       # the brief's mismatch flags
    priority_reasons = {}            # case_id -> {"why","signal"} for the rationale
    # Intro per-work-type cap state (INTRO_WORKTYPE_CAP) -- defined here, not inside the
    # literal-extraction branch below, so the (separate) strategic-bet pass can always see
    # it and keep counting against the SAME running total, even on a build where the
    # literal-needs branch never ran (e.g. extract_brief returned no usable needs).
    apply_worktype_cap = ctx.get("phase") == "Intro"
    wt_pick_counts = {}
    brief = {}                       # set below; kept defined even if extraction fails,
                                     # so the (separate) strategic-inference pass can use it
    # visible (not silent) degrade flag: True when we HAD research/profile/transcript
    # text to reason over but extract_brief() came back with nothing usable (a failed
    # OpenAI call, malformed reply, or a genuine no-needs read) -- the deck below still
    # gets built via the generic fallback ranker, but the salesperson should know the
    # "why this deck matches" picks are NOT research-driven this time.
    brief_weak = False
    has_source_text = bool(research_text.strip() or profile_text.strip() or ctx["transcript"].strip())
    # shared prep, hoisted above both the literal-extraction and strategic-inference
    # passes below (pure/no I/O, safe either way needs it)
    wanted = {w.upper() for w in ctx["work_types"]}
    wt_ids = {c["id"] for c in case_library.all_cases()
             if c.get("work_type", "").upper() in wanted}
    acct_fns = matcher._account_functions(set(ctx.get("functions", [])), match_notes)
    # role resonance (owner-spec, 2026-07-20): a detected persona (e.g. HR/Talent
    # Head) previously only nudged the FALLBACK rank_cases() path -- once profile/
    # research-driven priority picks fill the whole deck (the common case), it
    # never got a turn. Now threaded into shortlist_cases() below too. Detection
    # also scans the uploaded PROFILE text (not just recipient/transcript) --
    # a profile-only build (a name typed as recipient, no transcript) still
    # names the person's actual role/skills in the profile itself.
    # NOTE: deliberately ctx["transcript"], NOT match_notes (which also folds in the
    # full deep-research brief) -- confirmed 2026-07-23 (Vanderlande): a research doc
    # about the ACCOUNT broadly quotes several OTHER executives (its CTO, its CFO),
    # and scanning that whole document attributed CIO_CTO/FINANCE_HEAD/CEO_FOUNDER to
    # the actual recipient (an engineering-track stakeholder) purely because those
    # titles were mentioned describing someone else at the company. That persona
    # pollution then inflated unrelated cases in every capability's shortlist ranking,
    # crowding out the case that was actually the best content match. Only the
    # recipient's own name, the meeting transcript, and their own profile are reliable
    # signals of WHO WE'RE PITCHING TO; a company-wide research brief is not.
    persona_codes = personas.detect(ctx.get("recipient", ""), ctx["transcript"] + "\n" + profile_text)
    try:
        # ONE structured pass over research + profile + transcript: needs (with domain/
        # use-case), mismatch flags, expressed interest, account (honours the reference
        # priority: research primary, transcript = prior interest, profile-only otherwise).
        brief = ai_matcher.extract_brief(research_text, profile_text, ctx["transcript"])
        if brief and brief.get("needs"):
            avoid = brief.get("avoid") or []
            expressed = brief.get("expressed_interest") or []
            needs = brief["needs"]
            profile_needs = needs                       # names+descriptions for lead_research
            # expressed-interest topics -> lightweight needs so PRIOR INTEREST is weighted
            seen_names = {n["name"].lower() for n in needs}
            all_needs = needs + [{"name": x, "description": "", "domain": "", "use_case": ""}
                                 for x in expressed if x.lower() not in seen_names]
            # cheap shortlist per need — query is the bare CAPABILITY NAME (a crisp capability
            # embeds best; adding the use-case sentence or the client industry drags the match
            # toward the wrong cases). Industry is applied separately as a light boost.
            queries = [n["name"] for n in all_needs]
            shortlists = relevance.shortlist_cases(queries, industry=ctx.get("industry", ""),
                                                   functions=acct_fns, allowed_ids=wt_ids, top_n=8,
                                                   persona_codes=persona_codes)
            sl_by_name = {n["name"]: lst for n, lst in zip(all_needs, shortlists)}
            recs = {r["id"]: r for r in case_library._load()}
            expressed_lc = {x.lower() for x in expressed}
            # SELECTION is algorithmic (the semantic shortlist ranks capability reliably) --
            # covered -> priority pick; nothing clears the bar -> an honest gap.
            # GLOBAL greedy-by-strength assignment, not list order (owner-reported, 2026-07-23,
            # Vanderlande): needs used to be matched in EXTRACTION-LIST order, first-come-
            # first-claimed -- a case that was the OBVIOUS best fit for a later, more specific
            # need ("BOM Synchronization" for "BOM Harmonization", cosine 0.58) got sniped by
            # an earlier, broader need it was only a mediocre fit for ("PLM Automation", cosine
            # 0.49), which had a much better candidate of its own waiting. Collect every
            # (need, case) edge that clears the coverage bar across ALL needs, sort ALL of them
            # by match strength together, then assign greedily from the top -- a case goes to
            # whichever need it fits BEST, never just whichever need happened to be listed
            # first. Standard greedy approximation for maximum-weight bipartite matching:
            # simple, auditable, and provably no worse than list-order (which had no
            # relationship to fit quality at all).
            picked = []                       # [{need,id,title,blurb}] for the explainer
            claimed_needs = set()
            # capped_needs: needs whose only reason for going unfilled was the Intro
            # per-work-type cap (a real case existed; there just wasn't room for it) --
            # excluded from the "not in our library" gap list below, since that gap means
            # something different (no proof point exists at all).
            capped_needs = set()
            edges = []
            for n in all_needs:
                name = n["name"]
                if name.strip().lower() in _GENERIC_NEEDS:
                    continue
                for item in sl_by_name.get(name, []):
                    if _is_avoided(recs.get(item["id"], {}), avoid):
                        continue                        # skip a mismatch-flagged case
                    if not (item["cosine"] >= CAPABILITY_COVER or item["title_hits"] >= 1):
                        continue                        # never clears the coverage bar
                    edges.append((n, name, item))
            # Order eligible edges by ADJUSTED strength, not raw cosine. cosine is the
            # coverage GATE above (a case must clear CAPABILITY_COVER to be an edge at
            # all); among cases that clear it, `adj` is the real "match strength" -- it
            # folds in industry-vertical fit, the stakeholder persona, function, and
            # lexical overlap on top of meaning. Sorting by raw cosine here made those
            # signals inert at selection time: they only decided top-N shortlist
            # membership, never the pick (owner-reported, 2026-07-29, NTT DATA -- a
            # TECH_IT RPO case that was the obvious vertical + role fit sat in the
            # shortlist but lost the pick to an off-vertical BFSI case by 0.03 cosine,
            # because the +0.10 same-vertical boost lived only in adj). cosine stays as
            # the secondary tiebreak so meaning still breaks ties between equal-adj cases.
            edges.sort(key=lambda e: (-e[2]["adj"], -e[2]["cosine"], -e[2]["title_hits"]))
            for n, name, item in edges:
                if name in claimed_needs or item["id"] in priority_ids:
                    continue                            # need already filled, or case already used
                # MECHANISM-DEDUP: skip a case that is a near-TWIN (case-to-case cosine
                # >= PRIORITY_DEDUP_SIM) of one already picked, so the deck doesn't show
                # two look-alike proof points (e.g. WFS002 Zero-Disruption Scaling and
                # WFS025 Contract-to-Hire, cos 0.87; two identical DevSecOps-open-banking
                # cases, cos 0.91). A different, differentiated candidate can still fill
                # this need on a later edge; if none exists it stays an honest gap. The
                # 0.85 gate is calibrated ABOVE the ~0.74 max seen among genuinely
                # DISTINCT talent picks (RPO vs HTD vs accelerated-hiring), so it never
                # drops a real capability -- only true duplicates (2026-07-29).
                if relevance.max_similarity(item["id"], priority_ids) >= PRIORITY_DEDUP_SIM:
                    continue
                rc = recs.get(item["id"], {})
                if apply_worktype_cap:
                    wt = (rc.get("work_type") or "").upper()
                    if wt and wt_pick_counts.get(wt, 0) >= INTRO_WORKTYPE_CAP:
                        capped_needs.add(name)      # a real case existed -- just no room
                        continue
                claimed_needs.add(name)
                priority_ids.append(item["id"])
                if apply_worktype_cap:
                    wt = (rc.get("work_type") or "").upper()
                    if wt:
                        wt_pick_counts[wt] = wt_pick_counts.get(wt, 0) + 1
                picked.append({"need": name, "id": item["id"], "title": rc.get("title", ""),
                               "blurb": rc.get("challenge", ""),
                               "solution": rc.get("solution", ""),
                               "industry": rc.get("industry", ""),
                               "strength": item["cosine"]})
            for n in all_needs:
                name = n["name"]
                if (name.strip().lower() in _GENERIC_NEEDS or name in claimed_needs
                        or name in capped_needs):
                    continue
                if name.lower() not in expressed_lc and n.get("description"):
                    best_cos = max((item["cosine"] for item in sl_by_name.get(name, [])), default=0.0)
                    missing.append({"name": name, "description": n["description"],
                                    "domain": n.get("domain", ""), "use_case": n.get("use_case", ""),
                                    "sim": round(best_cos, 2)})
            missing = missing[:10]
            # rich, honest "why this was picked" line per pick (one LLM call; the model only
            # EXPLAINS the already-chosen case, so it cannot mis-pick). Templated fallback.
            # Now also a FACT-CHECK gate (owner-reported, 2026-07-23: the coverage bar
            # cleared on keyword/semantic overlap alone had let through cases whose real
            # content didn't actually match the need -- e.g. an electric-grid load-balancing
            # case explained as "capacity and demand matching" for a manufacturer, purely a
            # lexical collision on "demand"). explain_picks() now sees the case's real
            # SOLUTION text (not just a title) and can say fit=false; a rejected pick is
            # dropped back to an honest gap rather than kept with glossy reasoning.
            # the FORM industry is authoritative over the one extract_brief inferred: when an
            # account is described only through a finance stakeholder's profile + a logistics
            # mail (neither says what the COMPANY sells), the model guesses the account's
            # industry as "Finance" for what is really a retailer -- and the business-model fit
            # check below needs the account's REAL industry to reject investor / finance-
            # institution cases (ARKO, 2026-07-27). The salesperson's picked industry is ground
            # truth; use it.
            if ctx.get("industry"):
                brief.setdefault("account", {})["industry"] = ctx["industry"]
            ex = ai_matcher.explain_picks(brief, picked, profile_text, research_text)
            for it in picked:
                r = ex.get(it["id"])
                if r and r.get("fit") is False:
                    if it["id"] in priority_ids:
                        priority_ids.remove(it["id"])
                    missing.append({"name": it["need"], "description": it.get("blurb", ""),
                                    "domain": "", "use_case": "", "sim": round(it["strength"], 2)})
                    continue
                if r and r.get("reason"):
                    priority_reasons[it["id"]] = {"why": r["reason"], "signal": r.get("signal", ""),
                                                  "strength": it["strength"]}
                else:
                    priority_reasons[it["id"]] = {"why": "Proves " + it["need"].lower(),
                                                  "signal": "capability", "strength": it["strength"]}
            missing = missing[:10]
        else:
            # fail-safe: the old name-only extraction path (offline / brief parse failed).
            # This is a SILENT degrade unless flagged -- extract_brief() logs its own
            # failure server-side, but the salesperson only sees it via brief_weak below.
            if has_source_text:
                brief_weak = True
                logger.warning("extract_brief returned no usable needs for this build "
                               "(research=%d chars, profile=%d chars, transcript=%d chars) "
                               "-- falling back to name-only extraction / generic ranking",
                               len(research_text), len(profile_text), len(ctx["transcript"]))
            research_needs = ai_matcher.extract_accelerators(research_text) if research_text else []
            profile_needs = ai_matcher.extract_profile(profile_text) if profile_text else []
            if not research_needs and not profile_needs and ctx["transcript"].strip():
                research_needs = ai_matcher.extract_accelerators(ctx["transcript"])
            needs = profile_needs + research_needs
            if needs:
                best = relevance.best_cases([a["name"] for a in needs],
                                            industry=ctx.get("industry", ""), functions=acct_fns,
                                            allowed_ids=wt_ids)
                for a, (bid, _adj, cos, thits) in zip(needs, best):
                    if a["name"].strip().lower() in _GENERIC_NEEDS:
                        continue
                    covered = thits >= 1 or cos >= COVERAGE_THRESHOLD
                    if covered:
                        if bid and bid not in priority_ids:
                            priority_ids.append(bid)
                    else:
                        missing.append({**a, "sim": round(cos, 2)})
                missing = missing[:10]
    except Exception:
        logger.exception("research-driven priority-pick pipeline raised -- falling back "
                         "to the generic matcher.plan() ranker for this build")
        priority_ids, missing, avoid, priority_reasons = [], [], [], {}
        if has_source_text:
            brief_weak = True

    # Industry-level gap (owner's spec, 2026-07-14): a salesperson-typed "Other"
    # industry (constants.all_industries()'s custom entries) with genuinely NO
    # related case study anywhere in the selected work types gets its own honest
    # flag, folded into the same "not in our library" card -- a capability gap
    # says "we don't have a slide for X"; this says "we have no proof at all for
    # this INDUSTRY". Built-in industries (BFSI, HEALTHCARE...) are never
    # flagged this way -- they always have at least some tagged coverage.
    acct_industry = ctx.get("industry", "")
    if acct_industry and acct_industry.upper() not in _BUILTIN_INDUSTRIES:
        try:
            # matcher-SHAPED rows (keywords middot-joined, primary_industry set) --
            # the same shape relevance.py's other industry logic already expects,
            # not the raw store records (whose "keywords" is a list, not a string).
            wt_rows = [row for rows in case_library.candidate_rows(ctx["work_types"]).values()
                      for row in rows]
            if not relevance.industry_has_coverage(acct_industry, wt_rows):
                missing.append({
                    "name": "%s case studies" % acct_industry,
                    "description": "We don't yet have a proven case study specifically "
                                   "for %s -- this deck leans on function/capability "
                                   "matches instead of an industry-specific proof point."
                                   % acct_industry,
                    "domain": acct_industry, "use_case": "", "sim": 0.0,
                })
        except Exception:
            pass                          # never let this flag block a build

    # --- strategic inference: capability bets that are NEVER stated outright, but
    #     that a sharp account researcher would infer by connecting the STAKEHOLDER's
    #     own career/role to the COMPANY's business (owner's spec, 2026-07-08 — the
    #     Pankaj Kumar Pant / Waaree Group case: nothing in his profile says "we want
    #     predictive maintenance," but 2 years running a solar ingot/wafer plant + a
    #     career built on defect-catching makes it a reasoned bet). Kept as a fully
    #     SEPARATE pass from the literal extraction above — a failure here can never
    #     touch a literal pick, and a literal pick is never overridden by a bet.
    #     ai_matcher.infer_strategic_fit() only returns a bet when it can argue BOTH
    #     a personal angle (why HE'D want it) and a company angle (why THEY'D fund
    #     it) — either missing and it's dropped before it reaches here.
    try:
        bets = ai_matcher.infer_strategic_fit(research_text, profile_text, ctx["transcript"], brief)
        if bets:
            recs = {r["id"]: r for r in case_library._load()}
            bet_shortlists = relevance.shortlist_cases(
                [b["name"] for b in bets], industry=ctx.get("industry", ""),
                functions=acct_fns, allowed_ids=wt_ids, top_n=8,
                persona_codes=persona_codes)
            for b, shortlist in zip(bets, bet_shortlists):
                cand = None
                for item in shortlist:
                    if item["id"] in priority_ids:
                        continue                        # already used for a literal need
                    if _is_avoided(recs.get(item["id"], {}), avoid):
                        continue                        # skip a mismatch-flagged case
                    if apply_worktype_cap:
                        wt = (recs.get(item["id"], {}).get("work_type") or "").upper()
                        if wt and wt_pick_counts.get(wt, 0) >= INTRO_WORKTYPE_CAP:
                            continue                    # work type already at the Intro cap
                    cand = item
                    break
                if cand and (cand["cosine"] >= CAPABILITY_COVER or cand["title_hits"] >= 1):
                    # the narrative above was written BEFORE this real case was matched --
                    # re-ground it in what the case actually says (owner-reported, 2026-07-23:
                    # a "supplier collaboration" narrative had survived unmodified onto a
                    # banking finance-close case with no suppliers involved at all).
                    rc = recs.get(cand["id"], {})
                    v = ai_matcher.validate_bet_fit(b, rc.get("title", ""),
                                                    rc.get("challenge", ""), rc.get("solution", ""))
                    if not v["fit"]:
                        continue                # doesn't really hold up -- drop the bet, no gap
                    priority_ids.append(cand["id"])
                    if apply_worktype_cap:
                        wt = (rc.get("work_type") or "").upper()
                        if wt:
                            wt_pick_counts[wt] = wt_pick_counts.get(wt, 0) + 1
                    priority_reasons[cand["id"]] = {
                        "why": v["stakeholder_why"], "signal": "strategic bet",
                        "stakeholder_why": v["stakeholder_why"], "company_why": v["company_why"],
                        "strength": cand["cosine"]}
    except Exception:
        logger.exception("strategic-bet pass raised -- deck built with 0 strategic bets "
                         "this time (literal priority picks above are unaffected)")

    # GAP/BET RECONCILIATION (owner-reported, 2026-07-31, Broadridge India / Vivek Avvari):
    # `missing` is computed from LITERAL-need coverage above, BEFORE the strategic-bet pass
    # runs -- so a bet added afterward was never re-checked against gaps already flagged. A
    # "Wealth Management Technology" gap and an MSS052 strategic-bet pick (a wealth-management
    # conversational-analytics case) shipped in the SAME deck: the deck claimed MSS052 as a
    # proof point and, in the same breath, said we have no proof for that exact capability.
    # A bare cosine re-check can't safely resolve this (calibrated on this exact account:
    # "MSS052 covers Wealth Management" at cosine 0.401 and "MSS052 does NOT cover Capital
    # Markets Product Delivery" at 0.398 sit 0.003 apart -- no threshold gets both right), so
    # this uses the same grounded-LLM-judgment pattern as explain_picks/validate_bet_fit: a
    # cheap embedding pass only PRE-FILTERS to gaps with a plausible best-candidate (skip the
    # LLM ask entirely below GAP_PREFILTER_COVER -- obviously not covered, no need to ask),
    # then resolve_gap_overlap reads the candidate's REAL challenge/solution before deciding.
    if missing and priority_ids:
        try:
            recs = {r["id"]: r for r in case_library._load()}
            gap_texts = [f'{m["name"]}. {m.get("description", "")}' for m in missing]
            gap_shortlists = relevance.shortlist_cases(
                gap_texts, industry=ctx.get("industry", ""), functions=acct_fns,
                allowed_ids=set(priority_ids), top_n=1, persona_codes=persona_codes)
            candidates = []
            for m, sl in zip(missing, gap_shortlists):
                if sl and sl[0]["cosine"] >= GAP_PREFILTER_COVER:
                    cid = sl[0]["id"]
                    rc = recs.get(cid, {})
                    candidates.append({"name": m["name"], "description": m.get("description", ""),
                                       "candidate_id": cid, "candidate_title": rc.get("title", ""),
                                       "candidate_challenge": rc.get("challenge", ""),
                                       "candidate_solution": rc.get("solution", "")})
            resolved = ai_matcher.resolve_gap_overlap(candidates) if candidates else {}
            missing = [m for m in missing
                      if not resolved.get(m["name"], {}).get("covered")]
        except Exception:
            logger.exception("gap/bet reconciliation raised -- gaps kept unreconciled "
                             "(fail-safe: shown, not silently hidden)")

    # research + the profile's crisp focus areas LEAD the ranking; mail is secondary
    lead_research = (research_text + "\n"
                     + " ".join(f"{n['name']}. {n.get('description','')}" for n in profile_needs)).strip()
    result = matcher.plan({**ctx, "transcript": ctx["transcript"], "research": lead_research},
                          use_ai=True, priority_ids=priority_ids, avoid=avoid,
                          prefer_high_impact=bool(brief.get("prefer_high_impact")),
                          asked_flags={"asks_differentiation": bool(brief.get("asks_differentiation")),
                                      "asks_why_not_big_si": bool(brief.get("asks_why_not_big_si"))})
    # Gaps are FLAGS only now (no inline generation) — nothing to pre-fill here.
    titles = matcher._title_lookup()

    # --- skills slides (PURE Workforce only): auto-add after the standard block,
    #     before the case studies; labeled + removable in the panel ---
    sk = skills.candidates(ctx)
    if sk:
        picks = result["picks"]
        insert_at = next((i for i, p in enumerate(picks) if p["reason"].startswith("case")), None)
        if insert_at is None:
            insert_at = next((i for i, p in enumerate(picks) if p["slide_id"] in matcher.PIN_TO_END),
                             len(picks))
        sk_picks = []
        _chip = {"industry_strength": "IND", "skill_deepdive": "SKL",
                 "company_footprint": "FOOT", "target_skill_profile": "TSP"}
        for c in sk:
            titles[c["id"]] = c["label"]
            reason = {
                "industry_strength": "auto-added — RFI: industry strength slide",
                "skill_deepdive":    "auto-added — RFI: skills deployed slide",
                "company_footprint": "auto-added — RFI: existing client relationship",
                "target_skill_profile": "auto-added — skills mentioned: target skill profile slide",
            }.get(c["kind"], "auto-added — RFI data slide")
            sk_picks.append({"slide_id": c["id"], "reason": reason, "skill": True,
                             "tag": _chip.get(c["kind"], "SKL"),
                             "label": c["label"]})
        picks[insert_at:insert_at] = sk_picks

    # --- Client Context + Tailored Approach (Workforce, First/Second only):
    #     a real AI extraction from the notes, fails closed (see client_context.py) ---
    cc = client_context.candidates(ctx)
    if cc:
        insert_at = next((i for i, p in enumerate(result["picks"]) if p["reason"].startswith("case")), None)
        if insert_at is None:
            insert_at = next((i for i, p in enumerate(result["picks"]) if p["slide_id"] in matcher.PIN_TO_END),
                             len(result["picks"]))
        cc_picks = []
        for c in cc:
            titles[c["id"]] = c["label"]
            cc_picks.append({"slide_id": c["id"], "reason": "auto-added — client context drawn from your notes",
                             "skill": True, "tag": "CC", "label": c["label"]})
        result["picks"][insert_at:insert_at] = cc_picks

    # --- pre-built content slides (F1 "Already have the content?", queued
    #     across possibly several pastes): merge them in by shape, same
    #     boundary the skills/context auto-adds above use -- a case-study-
    #     shaped one joins the case-study zone; anything else (four_box,
    #     roadmap_board, box_grid, ...) sits just above it, with the standard
    #     slides (owner's spec, 2026-07-09) ---
    content_ids = [x for x in request.form.get("content_slide_ids", "").split(",") if x]
    if content_ids:
        picks = result["picks"]
        insert_at = next((i for i, p in enumerate(picks) if p["reason"].startswith("case")), None)
        if insert_at is None:
            insert_at = next((i for i, p in enumerate(picks) if p["slide_id"] in matcher.PIN_TO_END),
                             len(picks))
        standard_extra, case_extra = [], []
        for cid in content_ids:
            # A REUSED library case (the Custom Slide Builder's duplicate check offers
            # "use this one") rides as its own store id, not a staging id -- it needs no
            # NEW: prefix and no rebuild; it just joins the case-study zone. Skipping it
            # here would silently drop a slide the salesperson deliberately queued.
            store_rec = case_library.record(cid.upper())
            if store_rec:
                if cid.upper() in {p["slide_id"] for p in picks}:
                    continue                     # the matcher already picked it
                titles[cid.upper()] = store_rec.get("title", "")
                case_extra.append({"slide_id": cid.upper(), "reason": "you reused this slide",
                                   "skill": True, "tag": "YOU",
                                   "label": store_rec.get("title", "")})
                continue
            rec = staging.get(cid)
            if not rec:
                continue
            sid = "NEW:" + cid
            titles[sid] = rec.get("title", "Untitled")
            entry = {"slide_id": sid, "reason": "your pasted content", "skill": True,
                     "tag": "YOU", "label": rec.get("title", "Untitled")}
            if rec.get("content_type", "case_study") == "case_study":
                case_extra.append(entry)
            else:
                standard_extra.append(entry)
        picks[insert_at:insert_at] = standard_extra + case_extra

    all_slides = sorted(((sid, t) for sid, t in titles.items()
                         if sid not in _legacy_case_ids()),
                        key=lambda kv: matcher._num(kv[0]))
    # store-case titles power the picks/suggested display AND the JS title lookup
    # (so a manually-added or auto-picked AIP/WFS/MSS case shows its real title)
    titles.update(case_library.title_map())
    case_lib = sorted(case_library.all_cases(), key=lambda c: (c["work_type"], c["title"]))

    # --- "why this deck matches" rationale (deterministic, straight from picks) ---
    def _why(reason):
        r = re.sub(r"^case \[[^\]]*\] · (T\d · )?", "", reason or "")
        r = re.sub(r"\s*\(score [^)]*\)$", "", r)
        return r.strip() or "related to your notes"

    def _tier(reason):
        m = re.search(r"\bT(\d)\b", reason or "")
        return "T" + m.group(1) if m else ""

    def _raw_score(reason):
        # matcher.plan()'s own internal score (see matcher._case_reason) -- only used
        # as a fallback ranking signal for FILL picks that never went through the
        # priority pipeline below (so they have no captured cosine "strength"). Not
        # on the same 0..1 scale as strength; /10 keeps it roughly comparable for
        # ORDERING purposes only, never used as a selection threshold.
        m = re.search(r"\(score ([\d.]+)\)", reason or "")
        return float(m.group(1)) / 10.0 if m else 0.0

    rationale = [{"id": p["slide_id"], "title": titles.get(p["slide_id"], ""),
                  "why": _why(p["reason"]), "tier": _tier(p["reason"]),
                  "strength": _raw_score(p["reason"])}
                 for p in result["picks"] if p["slide_id"][:3] in ("AIP", "WFS", "MSS")]

    # "Why this resonates" per case: the LLM re-rank already justified each PRIORITY
    # pick (domain / prior-interest / role), so surface that as the fit line. Fill
    # cases (not tied to a specific need) keep their deterministic matcher reason.
    _bare = ("same industry", "related", "related to your notes")
    for r in rationale:
        pr = priority_reasons.get(r["id"])
        if pr and pr.get("why"):
            r["fit"] = pr["why"]
            if pr.get("signal"):
                r["signal"] = pr["signal"]
            if pr.get("company_why"):     # a strategic bet: show BOTH angles, not just one
                r["stakeholder_why"] = pr.get("stakeholder_why")
                r["company_why"] = pr["company_why"]
            if r.get("why", "").strip().lower() in _bare:
                r["why"] = ""          # the re-rank reason carries it; hide the bare fallback
            if "strength" in pr:       # real per-need cosine beats the coarser fallback score
                r["strength"] = pr["strength"]

    # owner-reported, 2026-07-23 (Kimberly-Clark): the panel showed picks in
    # whatever order needs were extracted, with no indication some matches are
    # stronger than others -- strongest match leads, so the best proof points
    # aren't buried under weaker ones.
    rationale.sort(key=lambda r: -r.get("strength", 0.0))

    # (priority picks + `missing` were computed before planning, above)
    body = render_template("build.html", ctx=ctx, picks=result["picks"],
                                  gaps=result["gaps"], titles=titles, all_slides=all_slides,
                                  case_lib=case_lib, rationale=rationale, missing=missing,
                                  research_read=research_read, research_failed=research_failed,
                                  brief_weak=brief_weak,
                                  suggestions=result.get("suggestions", []),
                                  suggested=result.get("suggested", []),
                                  ai_used=result.get("ai_used", False),
                                  persona_labels=result.get("persona_labels", []),
                                  resume=False, build_id=build_id, reopen_seed=None)
    return shell(body, active="new", crumb="<b>New deck</b> / Suggested slides")


def _case_slide(sid, rec):
    """Editable slide data for a case study (content-store or AI-drafted)."""
    domain = rec.get("domain") or rec.get("industry") or ""
    function = rec.get("function") or ""
    sub = rec.get("subhead") or ""
    if (not domain or not function) and sub:          # AI drafts carry a 'subhead' string
        dm = re.search(r"Domain:\s*([^|]+)", sub)
        fm = re.search(r"Function:\s*([^|]+)", sub)
        domain = domain or (dm.group(1).strip() if dm else "")
        function = function or (fm.group(1).strip() if fm else "")
    caps = fill_case_study._pad(rec.get("capabilities", []), 6)
    return {"kind": "case", "id": sid, "title": rec.get("title", ""),
            "domain": domain, "function": function,
            "challenge": rec.get("challenge", ""), "solution": rec.get("solution", ""),
            "caps": [fill_case_study.split_capability(c) for c in caps],
            "results": fill_case_study._pad(rec.get("results", []), 3)}


def _build_review_slides(ids, client):
    """One editable 'slide' per id, in deck order: library slides (title/subtitle),
    case studies (full body editable), and data-driven skills slides (read-only)."""
    prs = Presentation(assembler.SOURCE)
    by_id = {read_id(s): s for s in prs.slides if read_id(s)}
    store = content_store()          # {id: record} for AIP/WFS/MSS store cases
    slides = []
    for sid in ids:
        if sid in (client_context.CS_CONTEXT, client_context.CS_APPROACH):
            # Client Context / Tailored Approach: filled from an AI extraction at
            # FINALIZE time (see deck_build.assemble), not editable here — showing
            # the raw [BRACKET] placeholders as "editable fields" would be confusing.
            slides.append({"kind": "skills", "id": sid})
        elif sid in by_id:                                       # master library slide (CSxx)
            # A predefined slide is shown as a read-only PREVIEW -- the real rendered slide,
            # exactly as it appears in the download. It is not edited on /review; it goes
            # into the deck as-is (owner's spec, 2026-07-13).
            slides.append({"kind": "library", "id": sid})
        elif sid in store:                                     # content-store case
            slides.append(_case_slide(sid, store[sid]))
        elif sid.startswith("NEW:"):                           # pasted or AI-created slide
            rec = staging.get(sid[4:])
            if not rec:
                continue
            kind = rec.get("content_type", "case_study")
            if kind == "case_study":
                slides.append(_case_slide(sid, rec))
            else:
                # every other shape is drawn by the SHARED inline editor macro
                # (templates/_slide_editor.html), the same one the Custom Slide Builder
                # uses. Before this, only four_box and roadmap_board rendered here, and
                # both were read-only; the other five fell through to the "no text to
                # edit" placeholder below.
                slides.append({"kind": kind, "id": sid,
                               "vm": slide_schema.view_model(rec),
                               "schema": slide_schema.fields_for(kind)})
        elif sid.startswith("SK:") or sid.startswith("FP:") or sid.startswith("TSP:"):   # data-driven skills slide
            slides.append({"kind": "skills", "id": sid})
    return slides


def _render_review(client, ids, *, industry="", transcript="", phase="", recipient="",
                   functions=None, work_types=None):
    slides = _build_review_slides(ids, client)
    body = render_template("review.html", client=client, order=",".join(ids), slides=slides,
                           industry=industry, transcript=transcript, phase=phase, recipient=recipient,
                           functions=functions or [], work_types=work_types or [],
                           shape_kinds=_SHAPE_KINDS)
    return shell(body, active="new", crumb="<b>New deck</b> / Review &amp; edit")


@bp.route("/review", methods=["POST"])
def review():
    client = request.form.get("client_name", "Client").strip()
    ids = [x for x in request.form.get("order", "").split(",") if x]
    return _render_review(client, ids,
                          industry=request.form.get("industry", ""),
                          transcript=request.form.get("transcript", ""),
                          phase=request.form.get("phase", ""),
                          recipient=request.form.get("recipient", ""),
                          functions=request.form.getlist("functions"),
                          work_types=request.form.getlist("work_types"))


@bp.route("/from_content", methods=["POST"])
def from_content():
    """F1: the user already has the content for ONE slide — paste it or upload a
    document. Which template it becomes is either auto-detected (does this read as
    a case study, or a four-way breakdown?) or set explicitly via the template_hint
    dropdown (owner's spec, 2026-07-08: pasted content isn't always a case study).
    Builds ONE branded slide and stages it, returning JSON (not a page) — the
    caller is new-form.js's queue, which lets the salesperson repeat this for
    several slides (e.g. handed one-by-one by the MS team) before either taking
    just those straight to review, or continuing into the full context form
    below so they get merged into the generated deck in the right spot (owner's
    spec, 2026-07-09: content slides should slot in among the standard slides or
    the case studies, wherever their own shape belongs, not bolt on separately)."""
    client = request.form.get("client_name", "Client").strip()
    industry = request.form.get("industry", "")
    work_type = request.form.get("work_type", "").strip().upper()
    content = request.form.get("content", "").strip()
    ftext = research.extract_text(request.files.get("content_file"))
    if ftext:
        content = (content + "\n\n" + ftext).strip()
    # BOTH content and work_type are required -- work_type used to be silently
    # optional here, which meant staging.add() got "" and case_library.
    # promote_ai_case() (called unconditionally at /finalize) would look up an
    # empty-string prefix, find nothing, and return None -- the slide built and
    # downloaded fine, but NEVER reached the content store, so it was invisible
    # in Slide Library / search / the Excel registry with no error anywhere
    # (owner-reported, 2026-07-08: a real case study went missing this way).
    if not content or work_type not in ("WORKFORCE", "AI_POD", "MS"):
        err = ("Paste the case-study content or attach a document first." if not content
              else "Choose a work type — it's needed to save this into the shared library.")
        return jsonify({"ok": False, "error": err})
    rec, tdef = slide_generator.build_content_slide(
        content, industry, request.form.get("template_hint", "auto").strip())
    staged = staging.add(rec, work_type, industry, client)
    return jsonify({"ok": True, "id": staged["id"], "title": staged.get("title", "Untitled"),
                    "content_type": staged.get("content_type", "case_study"),
                    "template_label": tdef["label"]})


@bp.route("/finalize", methods=["POST"])
def finalize():
    client = request.form.get("client_name", "Client").strip()
    phase = request.form.get("phase", "")
    ids = [x for x in request.form.get("order", "").split(",") if x]
    if not ids:
        abort(400)
    import os
    os.makedirs(config.OUTPUT_DIR, exist_ok=True)
    # every version gets its own file, never overwritten (owner: keep every
    # version forever) -- e.g. "Joulestowatts_Acme Bank FV1.pptx", then FV2...
    # reserve_version() claims the number+filename ATOMICALLY, before the build
    # runs, so two near-simultaneous finalizes for the same client+phase can't
    # be handed the same version and silently overwrite each other's file.
    version, filename = meeting_log.reserve_version(client, phase)
    path = os.path.join(config.OUTPUT_DIR, filename)

    # ---- AI slides: apply edits, then ACCEPT (promote -> library) or REJECT ----
    final_ids = list(ids)
    accepted = []
    for gid in [x for x in request.form.get("ai_ids", "").split(",") if x]:
        decision = request.form.get("ai_decision__" + gid, "accept")
        if decision == "reject":
            staging.discard(gid)
            continue
        staging.update_content(
            gid,
            title=request.form.get("ai_title__" + gid),
            keywords=request.form.get("ai_keywords__" + gid),
            bullets=[b.strip() for b in request.form.get("ai_bullets__" + gid, "").splitlines() if b.strip()],
        )
        new_cs = staging.promote(gid)          # full sign-off -> real, client-ready slide
        if new_cs:
            accepted.append(new_cs)
    # Slot the AI slides in BEFORE the closing slides (Next Steps / Let's win together),
    # not at the very end of the deck.
    if accepted:
        insert_at = next((i for i, s in enumerate(final_ids) if s in matcher.PIN_TO_END),
                         len(final_ids))
        final_ids[insert_at:insert_at] = accepted

    # Pasted non-case shapes (four-box, roadmap, box grid, deep-dive, scored list, stat
    # overview, data table) are edited through the shared inline-editor macro, which posts
    # ONE json field per slide. Apply them to the staged record straight away -- that
    # record is what deck_build reads, so the edit needs no further threading. Each goes
    # through its own normalizer (slide_schema.apply_edits), which also means a hand-edit
    # can't hand the renderer a shape it can't fill.
    for key, val in request.form.items():
        if not key.startswith("shape__NEW:"):
            continue
        stg_id = key[len("shape__NEW:"):]
        rec = staging.get(stg_id)
        if not rec:
            continue
        try:
            fields = slide_schema.apply_edits(rec, json.loads(val),
                                              request.form.get("industry", ""))
        except (ValueError, TypeError):
            continue                # a malformed payload leaves the slide as it was built
        staging.update_fields(stg_id, fields)

    # collect inline edits from the review slide-view: edit__ = library title/subtitle,
    # cs__ = case-study body fields (title/domain/function/challenge/solution/caps/results).
    edits = {}
    raw_cases = {}
    for key, val in request.form.items():
        if key.startswith("edit__"):
            _, sid, idx = key.split("__")
            edits.setdefault(sid, {})[int(idx)] = val
        elif key.startswith("cs__"):
            parts = key.split("__")
            if len(parts) != 3:
                continue
            _, sid, field = parts
            d = raw_cases.setdefault(sid, {"caps": {}, "results": {}})
            if field in ("title", "domain", "function", "challenge", "solution"):
                d[field] = val
            elif field.startswith("cap") and field.endswith("_t"):
                d["caps"].setdefault(int(field[3:-2]), {})["t"] = val
            elif field.startswith("cap") and field.endswith("_b"):
                d["caps"].setdefault(int(field[3:-2]), {})["b"] = val
            elif field.startswith("res"):
                d["results"][int(field[3:])] = val
    case_edits = {}
    for sid, d in raw_cases.items():
        ov = {k: d[k] for k in ("title", "domain", "function", "challenge", "solution") if k in d}
        caps = []
        for i in sorted(d["caps"]):
            t = (d["caps"][i].get("t") or "").strip()
            b = (d["caps"][i].get("b") or "").strip()
            if t or b:
                caps.append({"title": t, "body": b})
        if caps:
            ov["capabilities"] = caps
        results = [(d["results"][i] or "").strip() for i in sorted(d["results"])]
        results = [r for r in results if r]
        if results:
            ov["results"] = results
        if ov:
            case_edits[sid] = ov

    try:
        deck_build.assemble(final_ids, path, client=client,
                            industry=request.form.get("industry", ""),
                            work_types=request.form.getlist("work_types"),
                            transcript=request.form.get("transcript", ""),
                            edits=edits, case_edits=case_edits,
                            phase=request.form.get("phase", ""))
    except (PermissionError, BadZipFile) as e:
        # the build never produced a real deck -- free the claimed version
        # number/filename so the next attempt doesn't skip one for nothing.
        meeting_log.cancel_reservation(client, phase, version)
        return file_busy_page(e)
    meeting_log.finalize_version(               # auto-log this version (no extra step)
        client=client, phase=phase, version=version,
        industry=request.form.get("industry", ""),
        functions=request.form.getlist("functions"),
        work_types=request.form.getlist("work_types"),
        recipient=request.form.get("recipient", ""),
        salesperson=current_salesperson(),
        slide_ids=final_ids,
        edits=edits,
        case_edits=case_edits,
    )
    count = len(list(Presentation(path).slides))
    body = render_template("preview.html", client=client, filename=filename, count=count)
    return shell(body, active="new", crumb="<b>New deck</b> / Done")
