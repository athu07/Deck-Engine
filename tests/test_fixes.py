# -*- coding: utf-8 -*-
"""
test_fixes.py  --  hermetic guards for the two core fixes (no network).

  1. Issue 1: a mismatch-flagged capability in the WRONG domain is demoted by
     rank_cases (so it can't win the fill).
  2. Issue 2: build_context saves/loads the rich context by build_id.
  3. Issue 2: draft_case_study actually feeds research + profile + notes into the
     generation prompt (they used to be discarded).

Plus the Custom Slide Builder guards (2026-07-10):
  4. dedupe.similar_cases fails safe to [] with no API key.
  5. staging.add round-trips EVERY shape's fields (it used to drop five of them).
  6. deck_build.staged_item routes all 8 shapes (five used to become case studies).
  7. promote_case is idempotent -- download AND deck-build saves the case once.
"""

import os
import pathlib

REPO = pathlib.Path(__file__).resolve().parent.parent


def _chdir():
    os.chdir(REPO)


def _isolate_staging(tmp_path, monkeypatch):
    """Point staging at a throwaway file -- these tests must never touch staging.json."""
    from deckengine.services.rendering import staging
    monkeypatch.setattr(staging, "STAGE_DIR", str(tmp_path))
    monkeypatch.setattr(staging, "STAGE_JSON", str(tmp_path / "staging.json"))
    return staging


# ── Issue 1: mismatch-flag demotion ───────────────────────────────────────────
def test_avoid_flag_demotes_cross_domain_case():
    _chdir()
    from deckengine.services.matching import relevance
    rows = [
        {"slide_id": "SAME", "title": "Computer Vision Inspection", "keywords": "computer vision",
         "primary_industry": "MANUFACTURING", "primary_function": "", "work_types": "MS",
         "search_text": "computer vision inspection manufacturing"},
        {"slide_id": "CROSS", "title": "Computer Vision for Web QA", "keywords": "computer vision",
         "primary_industry": "TECH_IT", "primary_function": "", "work_types": "MS",
         "search_text": "computer vision web qa"},
    ]
    ranked = relevance.rank_cases(
        "we need computer vision", rows, industry="MANUFACTURING", use_semantic=False,
        avoid=[{"capability": "computer vision", "reason": "meant for web QA, not this domain"}])
    score = {r["row"]["slide_id"]: r["score"] for r in ranked}
    # the flagged, wrong-domain case is pushed well below the same-domain one
    assert score["CROSS"] < score["SAME"]
    assert score["CROSS"] < 0            # penalised below the keep floor

    # without the avoid flag the cross-domain case is NOT penalised
    plain = relevance.rank_cases("we need computer vision", rows,
                                 industry="MANUFACTURING", use_semantic=False)
    plain_score = {r["row"]["slide_id"]: r["score"] for r in plain}
    assert plain_score["CROSS"] > score["CROSS"]


# ── Issue 2: build-context persistence ────────────────────────────────────────
def test_build_context_roundtrip(tmp_path, monkeypatch):
    from deckengine import config
    from deckengine.services import build_context
    monkeypatch.setattr(config, "BUILD_CONTEXT_DIR", str(tmp_path))
    build_context.save("abc123DEF", {"research": "R", "profile": "P", "transcript": "T",
                                     "industry": "MANUFACTURING"})
    got = build_context.load("abc123DEF")
    assert got["research"] == "R" and got["profile"] == "P" and got["industry"] == "MANUFACTURING"
    assert build_context.load("does-not-exist") == {}   # missing -> {} (never raises)
    assert build_context.load("") == {}                 # empty id -> {}


# ── Issue 2: generation prompt actually receives the rich context ─────────────
def test_draft_case_study_feeds_research_profile_notes(monkeypatch):
    _chdir()
    import openai

    captured = {}

    class _FakeClient:
        class chat:
            class completions:
                @staticmethod
                def create(**kw):
                    captured["prompt"] = kw["messages"][-1]["content"]
                    raise RuntimeError("stop after capture")   # -> draft takes fail-safe path

    monkeypatch.setattr(openai, "OpenAI", lambda *a, **k: _FakeClient())
    from deckengine.services.rendering import slide_generator
    out = slide_generator.draft_case_study(
        "computer vision inspection",
        {"industry": "MANUFACTURING", "recipient": "CTO", "function": "Quality",
         "research": "RESEARCH_MARKER_XYZ", "profile": "PROFILE_MARKER_XYZ",
         "notes": "TRANSCRIPT_MARKER_XYZ"})
    p = captured.get("prompt", "")
    assert "RESEARCH_MARKER_XYZ" in p, "deep research not in generation prompt"
    assert "PROFILE_MARKER_XYZ" in p, "profile not in generation prompt"
    assert "TRANSCRIPT_MARKER_XYZ" in p, "transcript not in generation prompt"
    # fail-safe still returns a well-formed placeholder (6 caps / 3 results)
    assert len(out["capabilities"]) == 6 and len(out["results"]) == 3


# ── F2: re-skin restyles in place and preserves ALL content ───────────────────
def test_reskin_preserves_content(tmp_path):
    _chdir()
    import io
    from pptx import Presentation
    from pptx.util import Inches
    from deckengine.services.rendering import reskin
    src = Presentation()
    s = src.slides.add_slide(src.slide_layouts[5])
    s.shapes.title.text = "Hello Deck"
    tbl = s.shapes.add_table(2, 2, Inches(1), Inches(2), Inches(4), Inches(1)).table
    tbl.cell(0, 0).text = "A"; tbl.cell(0, 1).text = "B"
    tbl.cell(1, 0).text = "1"; tbl.cell(1, 1).text = "ZZZVAL"
    buf = io.BytesIO(); src.save(buf)
    out = str(tmp_path / "reskinned.pptx")
    reskin.restyle_deck(buf.getvalue(), out)
    r = Presentation(out)
    text = " ".join(sh.text_frame.text for sl in r.slides for sh in sl.shapes if sh.has_text_frame)
    cells = " ".join(c.text for sl in r.slides for sh in sl.shapes if sh.has_table
                     for row in sh.table.rows for c in row.cells)
    assert "HELLO DECK" in text          # title preserved (headings are uppercased by design)
    assert "ZZZVAL" in cells             # table data preserved untouched
    # J2W branding applied: the title run is now Oswald
    title_fonts = {run.font.name for sl in r.slides for sh in sl.shapes
                   if sh.has_text_frame and "HELLO DECK" in sh.text_frame.text
                   for p in sh.text_frame.paragraphs for run in p.runs}
    assert "Oswald" in title_fonts


# ── Library bulk download: several slides -> ONE combined .pptx (2026-07-14) ───
def test_bulk_slides_combined_pptx():
    _chdir()
    import io
    from pptx import Presentation
    import app as appmod
    c = appmod.app.test_client()
    r = c.get("/slides/download?ids=CS01,MSS001")            # master + content-store case
    assert r.status_code == 200
    assert "presentationml.presentation" in r.headers["Content-Type"]   # a .pptx, not a .zip
    prs = Presentation(io.BytesIO(r.data))
    assert len(prs.slides._sldIdLst) == 2                   # both slides in one file
    assert c.get("/slides/download?ids=").status_code == 400  # nothing selected


# ── Custom Slide Builder (2026-07-10) ─────────────────────────────────────────
def test_dedupe_fails_safe_without_a_key(monkeypatch):
    """A duplicate check must never be the reason a slide can't be built."""
    _chdir()
    from deckengine.services.matching import dedupe, relevance
    monkeypatch.setattr(relevance, "embed_texts", lambda texts: None)   # offline
    assert dedupe.similar_cases("any pasted content at all") == []
    assert dedupe.similar_cases("") == []                               # nothing to check


def test_dedupe_ranks_and_gates_by_work_type(monkeypatch):
    """Above the bar -> a match, best first; below it -> silence. Work type gates."""
    _chdir()
    from deckengine.services.matching import dedupe, relevance
    from deckengine.services.content import case_library

    monkeypatch.setattr(relevance, "embed_texts", lambda texts: [[1.0, 0.0]])
    monkeypatch.setattr(relevance, "_load_case_embeddings",
                        lambda: {"MSS900": [1.0, 0.0],        # identical -> 1.00
                                 "MSS901": [0.9, 0.436],      # ~0.90 -> above the bar
                                 "MSS902": [0.0, 1.0],        # orthogonal -> 0.00
                                 "WFS900": [1.0, 0.0]})       # identical, wrong work type
    monkeypatch.setattr(case_library, "_load", lambda: [
        {"id": "MSS900", "title": "Exact", "work_type": "MS", "domain": ""},
        {"id": "MSS901", "title": "Close", "work_type": "MS", "domain": ""},
        {"id": "MSS902", "title": "Unrelated", "work_type": "MS", "domain": ""},
        {"id": "WFS900", "title": "Other work type", "work_type": "WORKFORCE", "domain": ""},
    ])

    hits = dedupe.similar_cases("text", allowed_work_types=["MS"])
    assert [h["id"] for h in hits] == ["MSS900", "MSS901"]      # ranked, orthogonal dropped
    assert hits[0]["score"] > hits[1]["score"] >= dedupe.THRESHOLD
    assert "WFS900" not in [h["id"] for h in hits]              # gated out by work type
    # an MS case is not a duplicate of a Workforce one, even word-for-word
    assert [h["id"] for h in dedupe.similar_cases("text", allowed_work_types=["WORKFORCE"])] == ["WFS900"]


def test_staging_round_trips_every_shape_field(tmp_path, monkeypatch):
    """staging.add() used to drop eyebrow/blocks/rows/stats/items/col_labels, so those
    five shapes rendered empty (or fell through to the case-study path)."""
    _chdir()
    staging = _isolate_staging(tmp_path, monkeypatch)

    grid = staging.add({"content_type": "box_grid", "template": "box_grid", "title": "G",
                        "boxes": [{"heading": "h", "body": "b"}]}, "MS", "BFSI")
    assert staging.get(grid["id"])["boxes"] == [{"heading": "h", "body": "b"}]

    stats = staging.add({"content_type": "stat_overview", "template": "stat_overview",
                         "title": "S", "intro": "i", "stats": [{"value": "42%", "label": "x"}],
                         "items": ["a", "b"]}, "MS", "BFSI")
    back = staging.get(stats["id"])
    assert back["stats"] == [{"value": "42%", "label": "x"}] and back["items"] == ["a", "b"]

    table = staging.add({"content_type": "data_table", "template": "data_table", "title": "T",
                         "col_labels": ["A", "B", "C"], "rows": [{"label": "r"}]}, "MS", "BFSI")
    assert staging.get(table["id"])["col_labels"] == ["A", "B", "C"]

    pillar = staging.add({"content_type": "pillar_deepdive", "template": "pillar_deepdive",
                          "title": "P", "eyebrow": "e", "blocks": [{"heading": "h"}]}, "MS", "BFSI")
    back = staging.get(pillar["id"])
    assert back["eyebrow"] == "e" and back["blocks"] == [{"heading": "h"}]


def test_staged_item_routes_all_shapes():
    """Only case_study carries a content-store record; every other shape goes down the
    generic draw path. Before this, five shapes were silently built AS case studies."""
    _chdir()
    from deckengine.services.rendering import deck_build

    for shape in ("four_box", "roadmap_board", "box_grid", "pillar_deepdive",
                  "scored_list", "stat_overview", "data_table"):
        item = deck_build.staged_item("NEW:G1", {"content_type": shape, "template": shape,
                                                 "title": "X"}, "BFSI")
        assert item["kind"] == shape and item["template"] == shape
        assert "record" not in item and item["data"]["title"] == "X"

    case = deck_build.staged_item("NEW:G2", {"content_type": "case_study", "title": "C",
                                             "challenge": "ch", "solution": "so"}, "BFSI")
    assert case["kind"] == "case_study" and case["template"] == "case_study_v2"
    assert case["record"]["challenge"] == "ch" and case["record"]["industry"] == "BFSI"

    # a record with no content_type is a case study (the historical default)
    assert deck_build.staged_item("NEW:G3", {"title": "C"}, "")["kind"] == "case_study"

    # /review edits win over the drafted text
    edited = deck_build.staged_item("NEW:G4", {"content_type": "case_study", "title": "old"},
                                    "", {"title": "new"})
    assert edited["record"]["title"] == "new"


def test_promote_case_saves_once_across_both_commit_points(tmp_path, monkeypatch):
    """A slide downloaded from the builder AND folded into a deck must reach the
    library exactly once -- promote_ai_case mints a fresh id on every call."""
    _chdir()
    staging = _isolate_staging(tmp_path, monkeypatch)
    from deckengine.services.rendering import deck_build
    from deckengine.services.content import case_library

    calls = []
    monkeypatch.setattr(case_library, "promote_ai_case",
                        lambda rec, wt, ind: (calls.append((wt, ind)), "MSS999")[1])

    stg = staging.add({"content_type": "case_study", "title": "A Case", "challenge": "c",
                       "solution": "s", "capabilities": ["One: a"], "results": ["r"]},
                      "MS", "BFSI")
    record = deck_build.staged_item("NEW:" + stg["id"], stg, "BFSI")["record"]

    first = deck_build.promote_case(stg["id"], staging.get(stg["id"]), record, "", "BFSI")
    second = deck_build.promote_case(stg["id"], staging.get(stg["id"]), record, "WORKFORCE", "BFSI")

    assert first == second == "MSS999"
    assert calls == [("MS", "BFSI")]                    # saved once; the record's OWN work type won
    assert staging.get(stg["id"])["promoted_id"] == "MSS999"


def test_promote_case_never_raises(tmp_path, monkeypatch):
    """A library-save failure must never block the deck."""
    _chdir()
    staging = _isolate_staging(tmp_path, monkeypatch)
    from deckengine.services.rendering import deck_build
    from deckengine.services.content import case_library

    def boom(*a, **k):
        raise OSError("Excel is open in another program")
    monkeypatch.setattr(case_library, "promote_ai_case", boom)

    stg = staging.add({"content_type": "case_study", "title": "A"}, "MS", "BFSI")
    assert deck_build.promote_case(stg["id"], stg, {"title": "A"}, "MS", "BFSI") is None
    assert not staging.get(stg["id"]).get("promoted_id")   # not marked -> retried next time


def test_paste_parser_splits_a_document_into_slides():
    """A whole deck's content, marked up by hand, cut into slides. A label that names a
    category outright picks the template; anything else is the slide's HEADING and the
    shape is left to the content (AUTO), never guessed from the title."""
    _chdir()
    from deckengine.services.content import paste_parser as pp

    slides = pp.parse(
        "slide 1 - casestudy\nA bank could not close procurement cycles.\n\n"
        "slide 2 - for box section\nFour pillars: intake, triage, delivery, assurance.\n\n"
        "Slide 3: road map\nPhase 1 discovery. Phase 2 pilot.\n\n"
        "slide 4 - How we think before we build\nThree principles guide us.\n\n"
        "slide 5\nNo label at all.\n")
    assert [s["template"] for s in slides] == [
        "case_study", "four_box", "roadmap_board", pp.AUTO, pp.AUTO]
    assert [s["number"] for s in slides] == [1, 2, 3, 4, 5]
    assert [s["matched"] for s in slides] == [True, True, True, False, False]
    assert slides[0]["content"].startswith("A bank")
    assert slides[1]["content"].endswith("assurance.")   # cut at the NEXT header, not beyond

    # a heading is kept AND prepended to the content, so the slide gets titled what the
    # author called it instead of the generator inventing a title
    assert slides[3]["heading"] == "How we think before we build"
    assert slides[3]["content"].startswith("How we think before we build\n")
    assert slides[4]["heading"] == ""                     # no label -> no heading


def test_paste_parser_has_no_length_limits():
    """A real header is not a short one. "Slide 4: Case Study 1, Reconciliation automation
    with agentic AI" is 64 characters, and a 60-char cap silently swallowed two slides out
    of nine in the owner's first real document. Length is never the signal."""
    _chdir()
    from deckengine.services.content import paste_parser as pp

    long_title = ("An extremely long slide title that goes on at length describing exactly "
                  "what the author intends to convey on this particular slide")
    slides = pp.parse("slide 1 - case study\nbody one\n\nslide 2: %s\nbody two\n" % long_title)
    assert len(slides) == 2
    assert slides[1]["heading"] == long_title

    real = pp.parse("Slide 4: Case Study 1, Reconciliation automation with agentic AI\nbody\n")
    assert len(real) == 1 and real[0]["number"] == 4


def test_paste_parser_uses_slide_numbers_to_reject_stray_lines():
    """Slide numbers validate each other: real headers form an increasing chain. A stray
    line that merely reads like a header can't extend it, so it stays as content."""
    _chdir()
    from deckengine.services.content import paste_parser as pp

    slides = pp.parse("slide 1 - case study\nreal body one\n\n"
                      "slide 2 - road map\nreal body two\n\n"
                      "Slide 1: as discussed above.\nthis is prose, not a slide\n\n"
                      "slide 3 - data table\nreal body three\n")
    assert [s["number"] for s in slides] == [1, 2, 3]         # the out-of-sequence one is gone
    assert "as discussed above" in slides[1]["content"]        # absorbed as content

    # gaps are fine -- increasing, not necessarily consecutive
    assert [s["number"] for s in pp.parse(
        "slide 2 - case study\na\n\nslide 5 - road map\nb\n\nslide 9 - data table\nc\n")] == [2, 5, 9]

    # prose with no separator after the number is never even a candidate
    one = pp.parse("Just one slide.\nSlide 2 of the printed deck covers the roadmap in detail.")
    assert len(one) == 1 and "covers the roadmap" in one[0]["content"]

    # a header with nothing under it is not a slide
    assert len(pp.parse("slide 1 - case study\n\nslide 2 - road map\nreal content")) == 1
    assert pp.parse("") == []


def test_paste_parser_only_matches_a_category_exactly():
    """Fuzzy label matching is gone. The heading "How we would engage with Voya" once
    matched the alias "named list with stats" on the single word "with" -- a preposition
    picked the template. A label either names a category, or it doesn't."""
    _chdir()
    from deckengine.services.content import paste_parser as pp
    assert pp.match_template("Road Map") == "roadmap_board"
    assert pp.match_template("FOUR BOX") == "four_box"
    assert pp.match_template("headline stats") == "stat_overview"
    assert pp.match_template("deep dive") == "pillar_deepdive"
    assert pp.match_template("for box section") == "four_box"
    # headings, not categories -> the content decides, not the title
    assert pp.match_template("How we would engage with Voya") == pp.AUTO
    assert pp.match_template("Case Study 1, Reconciliation automation") == pp.AUTO
    assert pp.match_template("What a first step looks like") == pp.AUTO
    assert pp.match_template("nonsense words") == pp.AUTO
    assert pp.match_template("") == pp.AUTO


def test_builder_parse_route_reports_the_split_without_building(tmp_path, monkeypatch):
    """/builder/parse must be cheap: no slide generation, no staging, no rendering.
    It reports the split and any library duplicate so a bad split is caught first.

    A slide whose category the author NAMED must not cost an AI call at all."""
    _chdir()
    staging = _isolate_staging(tmp_path, monkeypatch)
    import app as appmod
    from deckengine.services.matching import relevance
    from deckengine.services.rendering import slide_generator, preview

    monkeypatch.setattr(relevance, "embed_texts", lambda texts: None)      # offline
    def never(*a, **k):
        raise AssertionError("parse must not build, render, or classify a named category")
    monkeypatch.setattr(slide_generator, "build_content_slide", never)
    monkeypatch.setattr(preview, "case_png", never)
    monkeypatch.setattr(slide_generator, "classify_content_many", never)

    c = appmod.app.test_client()
    d = c.post("/builder/parse", data={
        "work_type": "MS",
        "content": "slide 1 - case study\nA procurement story.\n\nslide 2 - road map\nPhase 1.\n",
    }).get_json()

    assert d["ok"] is True
    assert [s["template"] for s in d["slides"]] == ["case_study", "roadmap_board"]
    assert all(s["matched"] for s in d["slides"])           # the author's own labels won
    assert all(s["matches"] == [] for s in d["slides"])     # offline -> no dup check, no block
    assert staging.all_items() == []                        # nothing staged
    assert "auto" not in {t["key"] for t in d["templates"]}  # every slide has a real shape

    # the same guards as every other builder route
    assert c.post("/builder/parse", data={"content": "x"}).get_json()["ok"] is False   # no work type
    assert c.post("/builder/parse", data={"work_type": "MS", "content": ""}).get_json()["ok"] is False


def test_parse_route_reads_the_shape_of_unlabelled_slides(tmp_path, monkeypatch):
    """The intelligence layer: where the author gave a category we honour it untouched;
    where they gave a heading or nothing, we read the CONTENT and map it to a category
    we actually have. One call for the whole document, and the heading rides along as a
    hint. No slide ever reaches the review screen as 'auto'."""
    _chdir()
    _isolate_staging(tmp_path, monkeypatch)
    import app as appmod
    from deckengine.services.matching import relevance
    from deckengine.services.rendering import slide_generator

    monkeypatch.setattr(relevance, "embed_texts", lambda texts: None)      # offline
    seen = {}
    def fake_classify(slides):
        seen["slides"] = slides
        return ["pillar_deepdive", "box_grid"]
    monkeypatch.setattr(slide_generator, "classify_content_many", fake_classify)

    c = appmod.app.test_client()
    d = c.post("/builder/parse", data={
        "work_type": "MS",
        "content": "slide 1 - case study\nA story.\n\n"
                   "slide 2: How we think before we build\nThree principles.\n\n"
                   "slide 3\nFour reassurances, unlabelled.\n",
    }).get_json()

    assert [s["template"] for s in d["slides"]] == ["case_study", "pillar_deepdive", "box_grid"]
    assert [s["matched"] for s in d["slides"]] == [True, False, False]
    assert [s.get("inferred", False) for s in d["slides"]] == [False, True, True]
    # ONE call, carrying only the slides that needed reading, heading passed as a hint
    assert len(seen["slides"]) == 2
    assert seen["slides"][0]["heading"] == "How we think before we build"
    assert seen["slides"][1]["heading"] == ""


def test_classify_content_many_fails_safe(monkeypatch):
    """Offline, or on a reply that lost count, every slide falls back to the registry's
    first shape -- never a partial mapping silently applied to the wrong slides."""
    _chdir()
    from deckengine.services.rendering import slide_generator
    import openai

    slides = [{"heading": "", "content": "a"}, {"heading": "", "content": "b"}]
    monkeypatch.setattr(openai, "OpenAI", lambda *a, **k: (_ for _ in ()).throw(RuntimeError("offline")))
    assert slide_generator.classify_content_many(slides) == ["case_study", "case_study"]
    assert slide_generator.classify_content_many([]) == []


def test_parse_route_only_flags_duplicates_for_case_study_shapes(tmp_path, monkeypatch):
    """The library holds only case studies, so a roadmap can never be a duplicate of
    one -- offering to swap it in would replace the slide with a different slide."""
    _chdir()
    _isolate_staging(tmp_path, monkeypatch)
    import app as appmod
    from deckengine.services.rendering import slide_generator
    from deckengine.services.matching import dedupe

    monkeypatch.setattr(slide_generator, "classify_content_many",
                        lambda slides: ["case_study"] * len(slides))
    monkeypatch.setattr(dedupe, "similar_cases_many",
                        lambda texts, **k: [[{"id": "MSS001", "title": "T", "work_type": "MS",
                                              "domain": "", "score": 0.95}] for _ in texts])
    c = appmod.app.test_client()
    d = c.post("/builder/parse", data={
        "work_type": "MS",
        "content": "slide 1 - case study\nA story.\n\nslide 2 - road map\nPhase 1.\n"
                   "\nslide 3\nUnlabelled, reads as a case study.\n",
    }).get_json()

    by_num = {s["number"]: s for s in d["slides"]}
    assert by_num[1]["matches"][0]["id"] == "MSS001"        # a real possible duplicate
    assert by_num[1]["matches"][0]["percent"] == 95
    assert by_num[2]["matches"] == []                       # roadmap cannot duplicate a case study
    assert by_num[3]["matches"][0]["id"] == "MSS001"        # inferred case study -> still checked


def test_similar_cases_many_is_one_embedding_call(monkeypatch):
    """N slides must cost ONE embedding call, not N. Blanks keep their own empty slot."""
    _chdir()
    from deckengine.services.matching import dedupe, relevance
    from deckengine.services.content import case_library

    calls = []
    def fake_embed(texts):
        calls.append(list(texts))
        return [[1.0, 0.0] for _ in texts]
    monkeypatch.setattr(relevance, "embed_texts", fake_embed)
    monkeypatch.setattr(relevance, "_load_case_embeddings", lambda: {"MSS1": [1.0, 0.0]})
    monkeypatch.setattr(case_library, "_load",
                        lambda: [{"id": "MSS1", "title": "T", "work_type": "MS", "domain": ""}])

    out = dedupe.similar_cases_many(["alpha", "   ", "beta"], allowed_work_types=["MS"])
    assert len(calls) == 1 and calls[0] == ["alpha", "beta"]   # one call, blanks not sent
    assert out[1] == []                                        # the blank keeps its slot
    assert out[0][0]["id"] == "MSS1" and out[2][0]["id"] == "MSS1"


def test_staging_ids_never_collide(tmp_path, monkeypatch):
    """Ids used to be "G%03d" % (len(items)+1), so a discard freed a number that the
    next add re-used -- two records with one id, and staging.get() returned the wrong
    one. The builder also adds concurrently, which the old scheme couldn't survive."""
    _chdir()
    staging = _isolate_staging(tmp_path, monkeypatch)
    a = staging.add({"title": "A"}, "MS", "")
    b = staging.add({"title": "B"}, "MS", "")
    staging.discard(a["id"])                       # frees G001
    c = staging.add({"title": "C"}, "MS", "")
    assert c["id"] != b["id"]                      # must not collide with the live record
    assert staging.get(b["id"])["title"] == "B"
    assert staging.get(c["id"])["title"] == "C"


def test_slide_schema_covers_every_shape():
    """Every shape the builder can produce must be editable. A shape with no schema
    builds and renders but silently offers no fields."""
    _chdir()
    from deckengine.services.rendering import slide_generator, slide_schema
    for t in slide_generator.CONTENT_TEMPLATES:
        assert slide_schema.fields_for(t["key"]), "no editable fields for %s" % t["key"]


def test_every_shape_is_wired_end_to_end():
    """A slide shape is only usable if ALL of these exist. Half-wiring one means it
    builds and then renders empty, or renders and then can't be edited, or can't be
    named in a pasted document -- each a silent failure the user only sees in the .pptx.

    (case_study is the exception: it fills the branded case_study_v2 template through
    fill_case_study, not a drawn body, so it has no J2W_TEMPLATE frame of its own.)"""
    _chdir()
    from pptx import Presentation
    from deckengine import config
    from deckengine.services.content import paste_parser
    from deckengine.services.rendering import draw_templates, skills, slide_generator, slide_schema

    tags = set()
    for slide in Presentation(config.SKILLS_TEMPLATES_PPTX).slides:
        if slide.has_notes_slide:
            for line in (slide.notes_slide.notes_text_frame.text or "").splitlines():
                if line.strip().startswith("J2W_TEMPLATE:"):
                    tags.add(line.split(":", 1)[1].strip())

    drawn_by_skills = {"four_box", "roadmap_board", "box_grid", "pillar_deepdive",
                       "scored_list", "stat_overview", "data_table"}
    problems = []
    for t in slide_generator.CONTENT_TEMPLATES:
        key = t["key"]
        if key == "case_study":
            continue
        if key not in tags:
            problems.append("%s: no template slide in skills_templates.pptx" % key)
        if not (key in drawn_by_skills or key in draw_templates.DRAWERS):
            problems.append("%s: nothing draws its body" % key)
        if not slide_schema.fields_for(key):
            problems.append("%s: not editable (no schema)" % key)
        try:
            slide_schema.normalize(key, {})
        except Exception as e:
            problems.append("%s: normalizer blows up on an empty record (%s)" % (key, e))
        if not callable(t.get("builder")):
            problems.append("%s: no builder to turn pasted content into it" % key)
    assert not problems, "\n".join(problems)

    # every new shape can be named in a pasted document ("slide 3 - governance layer")
    for key in draw_templates.DRAWERS:
        assert paste_parser.match_template(key.replace("_", " ")) == key, \
            "no paste_parser alias resolves to %s" % key


def test_every_shape_is_reachable_from_every_entry_point():
    """A shape is only 'wired up all over the application' if it appears in ALL of these.
    Miss one and the shape works in the builder but renders as a header on /template/
    recreate, or never shows on /templates, or can't be edited on /review."""
    _chdir()
    from deckengine.services.rendering import draw_templates, recreate, slide_generator, slide_schema
    from deckengine.web import decks as decks_mod
    from deckengine.web.templates import _slide_shapes

    on_templates_page = {s["key"] for s in _slide_shapes()}
    # four_box is marker-filled only -- it has no drawn body in skills.py either
    recreate_known = set(recreate._MAPPING_FNS) | set(recreate._DRAW_FNS)

    problems = []
    for t in slide_generator.CONTENT_TEMPLATES:
        key = t["key"]
        if key not in on_templates_page:
            problems.append("%s: absent from the /templates page" % key)
        if not slide_schema.fields_for(key):
            problems.append("%s: not editable" % key)
        if key != "case_study":
            if key not in recreate_known:
                problems.append("%s: Recreate-with-AI would render only its header" % key)
            if key not in decks_mod._SHAPE_KINDS:
                problems.append("%s: /review can't edit it" % key)
    assert not problems, "\n".join(problems)


def test_recreate_classifier_has_a_letter_per_template():
    """The classifier indexed a fixed 11-letter string; the registry grew to 18 and the
    whole Recreate feature died with IndexError before the AI was even called. The
    alphabet (26) plus the NONE slot must always outrun the registry."""
    _chdir()
    import string
    from deckengine.services.rendering import slide_generator
    assert len(slide_generator.CONTENT_TEMPLATES) + 1 <= len(string.ascii_uppercase)

    # offline, classify_slides fails safe to _NONE_KEY for every slide, never raises
    from deckengine.services.rendering import recreate
    keys = recreate.classify_slides(["four pillars", "a data table", ""])
    valid = {t["key"] for t in slide_generator.CONTENT_TEMPLATES} | {recreate._NONE_KEY}
    assert len(keys) == 3 and all(k in valid for k in keys)


def test_recreate_white_icon_gets_a_dark_chip():
    """A white icon vanished on the pale chip (owner-reported, 2026-07-13). A light icon
    now makes the chip the full accent colour; a coloured icon keeps the pale tint."""
    _chdir()
    import io
    from PIL import Image
    from deckengine.services.rendering import draw_templates as dt

    def blob(rgb):
        b = io.BytesIO(); Image.new("RGBA", (16, 16), rgb + (255,)).save(b, "PNG")
        return b.getvalue()

    assert dt._is_light_icon(blob((255, 255, 255))) is True     # white -> dark chip
    assert dt._is_light_icon(blob((40, 60, 80))) is False       # dark -> pale chip
    assert dt._is_light_icon(b"not an image") is False          # fail-safe


def test_recreate_detects_a_headline_score():
    """A source 'N/100' shown as a filled circle was flattened to a bullet; it's now
    detected and redrawn as a ring. Don't catch 'IA/IB' or a '0-100' range."""
    _chdir()
    from deckengine.services.rendering import recreate
    assert recreate._detect_score("Sample Validation 88 /100 Proceed") == (88, 100)
    assert recreate._detect_score("score 4/5 stars") == (4, 5)
    assert recreate._detect_score("Type IA/IB classification") is None
    assert recreate._detect_score("Scored output (0-100)") is None
    assert recreate._detect_score("148+ checks") is None


def test_recreate_reads_tables_charts_and_groups():
    """The core Recreate bug: _slide_text saw only plain text boxes, so a table slide or
    a chart slide reached the classifier as a bare heading, was misjudged, and fell to
    restyle-in-place -- which is why Recreate produced Reskin's output on a data deck
    (owner-reported, 2026-07-13)."""
    _chdir()
    from pptx import Presentation
    from pptx.util import Inches
    from pptx.chart.data import CategoryChartData
    from pptx.enum.chart import XL_CHART_TYPE
    from deckengine.services.rendering import recreate

    prs = Presentation()

    s = prs.slides.add_slide(prs.slide_layouts[6])
    s.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.8)).text_frame.text = "COMPARISON"
    t = s.shapes.add_table(2, 2, Inches(0.5), Inches(1.5), Inches(9), Inches(3)).table
    t.cell(0, 0).text = "Cost"; t.cell(0, 1).text = "$40"
    t.cell(1, 0).text = "Speed"; t.cell(1, 1).text = "2 days"
    table_text = recreate._slide_text(s)
    assert "Cost | $40" in table_text and "Speed | 2 days" in table_text, table_text

    s2 = prs.slides.add_slide(prs.slide_layouts[6])
    s2.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.8)).text_frame.text = "GROWTH"
    cd = CategoryChartData(); cd.categories = ["Q1", "Q2"]; cd.add_series("Rev", (10, 20))
    s2.shapes.add_chart(XL_CHART_TYPE.COLUMN_CLUSTERED, Inches(1), Inches(1.5),
                        Inches(8), Inches(4), cd)
    chart_text = recreate._slide_text(s2)
    assert "Q1 10" in chart_text and "Q2 20" in chart_text, chart_text


def test_new_shapes_normalize_and_draw_from_an_empty_record(tmp_path):
    """Every drawer must survive a record with nothing in it -- the AI builder fails safe
    to a placeholder, and a blank slide is far better than a stack trace at download."""
    _chdir()
    from pptx import Presentation
    from deckengine import config
    from deckengine.services.rendering import draw_templates, skills, slide_schema

    tpl = Presentation(config.SKILLS_TEMPLATES_PPTX)
    for kind in draw_templates.DRAWERS:
        rec = slide_schema.normalize(kind, {})
        assert rec["content_type"] == kind
        blank = Presentation()
        blank.slide_width, blank.slide_height = tpl.slide_width, tpl.slide_height
        path = str(tmp_path / (kind + ".pptx"))
        blank.save(path)
        built = skills.build_into(path, [kind],
                                  [{"id": kind, "template": kind, "kind": kind, "data": rec}])
        assert built == 1, "%s drew nothing from an empty record" % kind


def test_slide_schema_edits_round_trip_and_normalize():
    """An edit is re-run through the shape's own normalizer, so a hand-edit can't break
    an invariant the renderer depends on."""
    _chdir()
    from deckengine.services.rendering import slide_schema as ss

    case = {"content_type": "case_study", "industry": "BFSI", "title": "T", "subhead": "s",
            "challenge": "c", "solution": "so", "capabilities": ["A: x"], "results": ["r"]}
    out = ss.apply_edits(case, dict(ss.extract("case_study", case), title="New"))
    assert out["title"] == "New"
    # emptying them must NOT produce a slide the renderer can't fill
    bad = ss.apply_edits(case, {"capabilities": [], "results": [], "title": ""})
    assert len(bad["capabilities"]) == 6 and len(bad["results"]) == 3
    assert bad["title"]                                   # never blank

    roadmap = {"content_type": "roadmap_board", "title": "T", "subhead": "", "intro": "",
               "columns": [{"name": "n", "tag": "P1", "items": ["a"]}], "legend": [],
               "footer_title": "", "footer_body": ""}
    out = ss.apply_edits(roadmap, dict(ss.extract("roadmap_board", roadmap),
                                       columns=[{"name": "Phase 1", "tag": "P1", "items": ["x", "y"]}]))
    assert out["columns"][0]["items"] == ["x", "y"]       # nested list survives the trip


def test_slide_schema_ignores_fields_the_browser_invents():
    """The editor may only touch declared content fields -- never a record's bookkeeping,
    and never its shape."""
    _chdir()
    from deckengine.services.rendering import slide_schema as ss
    rec = {"content_type": "four_box", "id": "G001", "work_type": "MS", "promoted_id": "MSS9",
           "title": "T", "subhead": "s", "boxes": [{"heading": "h", "body": "b"}]}
    out = ss.apply_edits(rec, {"title": "ok", "content_type": "case_study",
                               "id": "HACK", "promoted_id": "STOLEN", "work_type": "AI_POD"})
    assert set(out) == {"title", "subhead", "boxes"}
    assert out["title"] == "ok"


def test_builder_returns_an_editable_slide_card_not_an_image(tmp_path, monkeypatch):
    """A built slide comes back as the SAME editable slide card the review page shows --
    click straight onto the text. Not a picture with an Edit button beside it."""
    _chdir()
    _isolate_staging(tmp_path, monkeypatch)
    import app as appmod
    from deckengine.services.rendering import slide_generator

    monkeypatch.setattr(slide_generator, "build_content_slide", lambda content, industry="", hint="auto": (
        {"content_type": "roadmap_board", "template": "roadmap_board", "title": "Plan",
         "subhead": "", "intro": "",
         "columns": [{"name": "Discovery", "tag": "P1", "items": ["interviews"]}],
         "legend": [], "footer_title": "", "footer_body": ""},
        {"key": "roadmap_board", "label": "Phased roadmap / board"}))

    d = appmod.app.test_client().post("/builder/slide", data={
        "content": "Phase 1 discovery.", "work_type": "MS", "industry": "BFSI"}).get_json()

    assert d["ok"] and d["content_type"] == "roadmap_board"
    assert "png" not in d                                    # no image; the card IS the slide
    html = d["html"]
    assert 'class="se ' in html                              # the review page's editable input
    assert 'data-path="title"' in html
    assert 'data-path="columns.0.name"' in html              # nested record position
    assert 'data-path="columns.0.items.0"' in html           # a list inside a record
    assert "Discovery" in html and "interviews" in html      # current values, not blanks


def test_builder_edit_saves_and_persists(tmp_path, monkeypatch):
    """Saving edits typed onto the card updates the staged record -- the single source of
    truth that the download, the deck, and the library all read."""
    _chdir()
    staging = _isolate_staging(tmp_path, monkeypatch)
    import app as appmod

    stg = staging.add({"content_type": "stat_overview", "template": "stat_overview",
                       "title": "Old", "subhead": "", "intro": "",
                       "stats": [{"value": "42%", "label": "l"}], "items": []}, "MS", "BFSI")
    c = appmod.app.test_client()

    r = c.post("/builder/slide/%s" % stg["id"],
               json={"title": "New", "stats": [{"value": "61%", "label": "faster"}]}).get_json()
    assert r["ok"] and r["title"] == "New"

    back = staging.get(stg["id"])                 # persisted, and bookkeeping untouched
    assert back["title"] == "New" and back["stats"] == [{"value": "61%", "label": "faster"}]
    assert back["work_type"] == "MS" and back["id"] == stg["id"]

    assert c.post("/builder/slide/NOPE", json={}).status_code == 404


def test_review_page_can_now_edit_every_pasted_shape(tmp_path, monkeypatch):
    """The five shapes that used to fall through to "no text to edit", plus the two that
    rendered read-only ("re-paste to change the text"), are all inline-editable now --
    through the same macro, and /finalize applies their edits to the staged record."""
    _chdir()
    staging = _isolate_staging(tmp_path, monkeypatch)
    import app as appmod
    c = appmod.app.test_client()

    stg = staging.add({"content_type": "box_grid", "template": "box_grid", "title": "Old title",
                       "subhead": "s", "boxes": [{"heading": "h1", "body": "b1"},
                                                 {"heading": "h2", "body": "b2"}]}, "MS", "BFSI")
    sid = "NEW:" + stg["id"]

    r = c.post("/review", data={"client_name": "Acme", "order": sid})
    html = r.data.decode()
    assert r.status_code == 200
    assert 'data-shape-id="%s"' % sid in html                 # rendered by the shared macro
    assert 'data-path="boxes.0.heading"' in html
    assert "not inline-editable yet" not in html              # the old read-only footer is gone

    # /finalize applies the edit to the staged record, which is what deck_build reads
    import json as _json
    c.post("/finalize", data={
        "client_name": "Acme", "phase": "Intro", "order": sid, "work_types": "MS",
        "shape__" + sid: _json.dumps({"title": "Edited title",
                                      "boxes": [{"heading": "H1", "body": "B1"}]}),
    })
    back = staging.get(stg["id"])
    assert back["title"] == "Edited title"
    assert back["boxes"][0] == {"heading": "H1", "body": "B1"}
    assert len(back["boxes"]) >= 2            # normalizer re-padded: box_grid needs 2+


def test_master_previews_ship_and_serve_without_libreoffice(monkeypatch):
    """The review page shows a picture of each master-deck slide. Those previews are
    COMMITTED (static/previews/), so a server with no LibreOffice -- a plain Render web
    service -- still serves them. Rendering at runtime meant every preview 404'd there
    (owner-reported, 2026-07-10), and paid a 40-slide conversion on every cold start
    where the tools DID exist."""
    _chdir()
    from deckengine.services.rendering import preview, reskin

    # every CSxx slide in the master must have a shipped preview keyed by content hash
    from pptx import Presentation
    from deckengine import config
    from deckengine.services.content.build_library import read_id
    ids = [i for i in (read_id(s) for s in Presentation(config.MASTER_DECK).slides) if i]
    shipped = preview._shipped_dir()
    import os
    missing = [i for i in ids if not os.path.exists(os.path.join(shipped, i + ".webp"))]
    assert not missing, ("no shipped preview for %s -- run scripts/prerender_master.py "
                         "and commit static/previews/" % missing)

    # and they resolve with the render tools entirely absent
    monkeypatch.setattr(reskin, "_on_path", lambda name: False)
    assert preview.master_slide_png("CS01")
    assert preview.master_slide_png("NOPE99") is None    # unknown id, still no crash


def test_master_preview_key_is_content_not_mtime(tmp_path, monkeypatch):
    """The key must be the deck's CONTENT, not its mtime: git checkout rewrites mtimes,
    so an mtime key could never match a build-time render."""
    _chdir()
    from deckengine.services.rendering import preview
    k1 = preview.master_key()
    # touching the file (new mtime, same bytes) must NOT change the key
    import os, time
    os.utime(__import__("deckengine").config.MASTER_DECK, None)
    preview._master_key_cache.clear()
    assert preview.master_key() == k1


def test_draft_with_ai_button_survives_quotes_in_the_gap_text():
    """The button used to be onclick="draftAI({{ m.name|tojson }}, ...)". |tojson emits
    real double quotes, so the HTML attribute ended at the first one and the browser only
    ever saw `onclick="draftAI("` -- every click threw "Unexpected end of input" and did
    nothing. The gap text is arbitrary prose, so it rides on data-* attributes now."""
    _chdir()
    import html as htmlmod
    import re
    from flask import render_template
    import app as appmod

    nasty = 'Catch "micro-cracks" & O\'Brien\'s <defects> before scrap'
    with appmod.app.test_request_context():
        page = render_template(
            "build.html",
            ctx={"client_name": "V", "industry": "BFSI", "transcript": "", "phase": "",
                 "recipient": "", "functions": [], "work_types": ["MS"]},
            picks=[], gaps=[], titles={}, all_slides=[], case_lib=[], suggestions=[],
            suggested=[], ai_used=True, persona_labels=[], rationale=[],
            missing=[{"name": 'Predictive "maintenance"', "description": nasty,
                      "domain": "Manufacturing", "use_case": "ingot line"}],
            research_read=True, research_failed=False, resume=False, build_id="abc",
            reopen_seed=None)

    assert "onclick=\"draftAI" not in page          # no JS source inside an HTML attribute
    assert 'class="btn draft-btn"' in page
    desc = re.search(r'data-desc="([^"]*)"', page)
    assert desc, "the data-desc attribute is not well-formed"
    assert htmlmod.unescape(desc.group(1)) == nasty  # every character survives the trip

    # the card only exists when there IS a gap to draft, and the manual form stays disabled
    assert 'id="ca-preview"' in page and "ca-topic" not in page


def test_static_assets_are_versioned():
    """Without an mtime stamp the browser keeps the old app.css / build.js after a change,
    and the user sees the OLD page against the NEW server. That is exactly how a working
    Draft-with-AI looked broken."""
    _chdir()
    import re
    import app as appmod
    c = appmod.app.test_client()
    page = c.get("/builder").data.decode()
    assert re.search(r'/static/js/builder\.js\?v=\d+', page)
    assert re.search(r'/static/app\.css\?v=\d+', page)
    # and the versioned URL still serves the file
    url = re.search(r'(/static/js/builder\.js\?v=\d+)', page).group(1)
    assert c.get(url).status_code == 200


def test_staging_reads_never_see_a_half_written_file(tmp_path, monkeypatch):
    """`json.dump(x, open(path,"w"))` truncates first, then streams. A reader landing in
    that window got a syntax error, _load() swallowed it and returned [], and the caller
    concluded the record did not exist -- so six of nine concurrent saves 404'd on slides
    that were in the file the whole time. Measured before the fix: 207/300 reads blind."""
    _chdir()
    import threading
    staging = _isolate_staging(tmp_path, monkeypatch)

    for i in range(60):                       # enough records that a write isn't instant
        staging.add({"title": "S%d" % i, "content_type": "case_study",
                     "challenge": "x" * 400, "solution": "y" * 400}, "MS", "BFSI")
    ids = [it["id"] for it in staging.all_items()]
    target = ids[30]

    stop = threading.Event()
    def churn():
        while not stop.is_set():
            staging.update_fields(ids[0], {"title": "churn"})
    writer = threading.Thread(target=churn, daemon=True)
    writer.start()
    try:
        misses = sum(1 for _ in range(400) if staging.get(target) is None)
    finally:
        stop.set()
        writer.join(timeout=2)

    assert misses == 0, "%d reads saw a torn staging file" % misses
    assert len(staging.all_items()) == 60      # and no record was lost to the churn


def test_concurrent_slide_saves_all_succeed(tmp_path, monkeypatch):
    """The builder saves every card at once when you download. Each POST reads the record
    then writes it back; with a non-atomic write most of them 404'd."""
    _chdir()
    staging = _isolate_staging(tmp_path, monkeypatch)
    import app as appmod
    from concurrent.futures import ThreadPoolExecutor

    ids = [staging.add({"content_type": "four_box", "template": "four_box",
                        "title": "T%d" % i, "subhead": "",
                        "boxes": [{"heading": "h", "body": "b" * 300}]}, "MS", "BFSI")["id"]
           for i in range(9)]
    c = appmod.app.test_client()

    def save(sid):
        return c.post("/builder/slide/%s" % sid,
                      json={"title": "Edited " + sid}).status_code

    with ThreadPoolExecutor(max_workers=9) as pool:
        codes = list(pool.map(save, ids))

    assert codes == [200] * 9, "some saves 404'd: %r" % codes
    assert all(staging.get(sid)["title"] == "Edited " + sid for sid in ids)


def test_builder_routes(monkeypatch):
    """The builder page renders, and every endpoint validates its input."""
    _chdir()
    import app as appmod
    from deckengine.services.matching import relevance
    monkeypatch.setattr(relevance, "embed_texts", lambda texts: None)   # offline
    c = appmod.app.test_client()

    assert c.get("/builder").status_code == 200
    # a work type is REQUIRED: without it the case study can't be saved to the library
    assert c.post("/builder/slide", data={"content": "x", "work_type": ""}).get_json()["ok"] is False
    assert c.get("/builder/preview/NOPE001").status_code == 404
    assert c.get("/builder/download?ids=").status_code == 400


def test_concurrent_slide_builds_get_distinct_ids(tmp_path, monkeypatch):
    """The browser fires four /builder/slide requests at once. staging.add() is a
    read-modify-write on one JSON file, so without its lock two of them would mint the
    same id and one slide would silently vanish from the queue."""
    _chdir()
    staging = _isolate_staging(tmp_path, monkeypatch)
    from concurrent.futures import ThreadPoolExecutor

    with ThreadPoolExecutor(max_workers=8) as pool:
        ids = list(pool.map(lambda i: staging.add({"title": "S%d" % i}, "MS", "")["id"],
                            range(24)))

    assert len(set(ids)) == 24                       # no collisions
    assert len(staging.all_items()) == 24            # and nothing was overwritten



# ── The "Other…" industry (owner-reported, 2026-07-17) ────────────────────────
# Picking "Other…" makes the select's value the literal "__OTHER__" sentinel and the
# real industry is typed into the companion `industry_other` field. Two symptoms, one
# root cause -- the typed text was not a real form field, so anything reading the form
# BEFORE submit saw the sentinel: "Research this account" researched the account as
# though its industry were "(__OTHER__)" (and that junk brief then drove matching and
# every AI-written slide), and the typed industry was only ever remembered by a
# completed /build. These pin the resolution, both endpoints, and the store guard.

def _isolate_custom_industries(tmp_path, monkeypatch):
    """Point the custom-industry store at a throwaway file -- never touch the real one."""
    from deckengine.services.content import industries
    monkeypatch.setattr(industries, "_PATH", str(tmp_path / "custom_industries.json"))
    return industries


def test_resolve_industry_unwraps_the_other_sentinel():
    _chdir()
    from deckengine.web.view_helpers import resolve_industry
    # "Other…" picked -> the TYPED industry, never the sentinel
    assert resolve_industry({"industry": "__OTHER__",
                             "industry_other": " Renewable Diesel Production "}) \
        == "Renewable Diesel Production"
    # a built-in pick is untouched, even with stale text left in the box
    assert resolve_industry({"industry": "BFSI", "industry_other": "leftover"}) == "BFSI"
    # "Other…" with nothing typed is simply no industry -- never the sentinel
    assert resolve_industry({"industry": "__OTHER__", "industry_other": ""}) == ""
    assert resolve_industry({}) == ""


def test_research_account_researches_the_typed_industry_not_the_sentinel(tmp_path, monkeypatch):
    """The regression itself: strategic_brief() must receive the typed industry."""
    _chdir()
    _isolate_custom_industries(tmp_path, monkeypatch)
    from deckengine.web import api

    seen = {}

    def _spy(company_name, stakeholder_name="", industry="", **kw):
        seen["industry"] = industry
        return "BRIEF"

    monkeypatch.setattr(api.deep_research, "strategic_brief", _spy)

    import app as appmod
    c = appmod.app.test_client()
    r = c.post("/research_account", data={"client_name": "Neste", "recipient": "Head of Eng",
                                          "industry": "__OTHER__",
                                          "industry_other": "Renewable Diesel Production"})
    assert r.status_code == 200
    assert seen["industry"] == "Renewable Diesel Production"


def test_a_typed_other_industry_is_remembered_for_next_time(tmp_path, monkeypatch):
    """Typed once -> pickable from the dropdown the second time, from EITHER the
    research button or a completed build (it used to need a completed build)."""
    _chdir()
    industries = _isolate_custom_industries(tmp_path, monkeypatch)
    from deckengine.web.view_helpers import remember_custom_industry
    from deckengine import constants

    remember_custom_industry("Renewable Diesel Production")
    assert "Renewable Diesel Production" in industries.load()
    assert "Renewable Diesel Production" in constants.all_industries()   # in the dropdown

    remember_custom_industry("renewable diesel production")   # same industry, retyped
    assert len(industries.load()) == 1                        # not duplicated
    remember_custom_industry("BFSI")                          # a built-in code
    assert len(industries.load()) == 1                        # never shadows the taxonomy


def test_the_other_sentinel_can_never_be_stored_as_an_industry(tmp_path, monkeypatch):
    """Last line before disk: a stale cached new-form.js or a no-JS submit must not
    put a literal "__OTHER__" row in the dropdown for every salesperson, forever."""
    _chdir()
    industries = _isolate_custom_industries(tmp_path, monkeypatch)
    assert industries.add("__OTHER__") is False
    assert industries.load() == []


# ══════════════════════════════════════════════════════════════════════════════
# /build pick-count + typed-context regressions (owner-reported, 2026-08-04)
#
#   1. A deck built from a stakeholder profile came back with 3 case studies
#      whatever the input -- two separate ceilings were fighting each other and
#      one of them silently threw away picks the pipeline had already vetted.
#   2. Whatever the salesperson typed into the Context box left no trace in the
#      built deck: topics the client actually raised were neither picked nor
#      flagged as a gap.
# ══════════════════════════════════════════════════════════════════════════════
def _cand(cid, adj=0.90, cosine=0.80, title_hits=1):
    """One shortlist_cases() candidate row."""
    return {"id": cid, "adj": adj, "cosine": cosine, "title_hits": title_hits}


def _hermetic_build(monkeypatch, tmp_path, *, brief, shortlists):
    """Pin every AI seam /build touches to a fixed answer, so a POST is offline and
    deterministic. `shortlists` maps a need NAME -> the candidates shortlist_cases
    would have returned for it. Returns a dict that captures what reached
    matcher.plan()."""
    from deckengine import config
    from deckengine.services.matching import ai_matcher, matcher, relevance

    monkeypatch.setattr(config, "BUILD_CONTEXT_DIR", str(tmp_path))
    monkeypatch.setattr(ai_matcher, "extract_brief",
                        lambda research="", profile="", transcript="": dict(brief))
    monkeypatch.setattr(ai_matcher, "infer_strategic_fit",
                        lambda research="", profile="", transcript="", brief=None: [])
    monkeypatch.setattr(ai_matcher, "resolve_gap_overlap", lambda gaps: {})
    monkeypatch.setattr(ai_matcher, "explain_picks",
                        lambda brief, items, profile="", research="":
                        {i["id"]: {"fit": True, "reason": "proves " + i["need"],
                                   "signal": "capability"} for i in items})
    monkeypatch.setattr(relevance, "embed_texts", lambda texts: None)   # no network
    monkeypatch.setattr(relevance, "max_similarity", lambda sid, others: 0.0)
    monkeypatch.setattr(relevance, "shortlist_cases",
                        lambda texts, industry="", functions=None, allowed_ids=None,
                        top_n=8, persona_codes=(): [list(shortlists.get(t, [])) for t in texts])

    seen = {}
    real_plan = matcher.plan

    def spy_plan(context, *a, **kw):
        seen["priority_ids"] = list(kw.get("priority_ids") or [])
        seen["research"] = context.get("research") or ""
        seen["transcript"] = context.get("transcript") or ""
        return real_plan(context, *a, **kw)

    monkeypatch.setattr(matcher, "plan", spy_plan)
    return seen


def _post_build(client, **over):
    import re
    form = {"client_name": "Acme Corp", "industry": "RETAIL", "phase": "Intro",
            "recipient": "FP&A Manager", "transcript": "notes", "work_types": ["MS"],
            "functions": []}
    form.update(over)
    r = client.post("/build", data=form, content_type="multipart/form-data")
    assert r.status_code == 200, r.status_code
    html = r.get_data(as_text=True)
    ids = re.findall(r'data-id="([^"]+)"', html)
    return {
        "html": html,
        "cases": [i for i in ids if i[:3] in ("AIP", "WFS", "MSS")],
        "gaps": re.findall(r'data-name="([^"]*)"', html),
    }


def _brief(needs, expressed=()):
    return {"needs": [{"name": n, "description": n + " description", "domain": "Retail",
                       "use_case": n} for n in needs],
            "avoid": [], "expressed_interest": list(expressed),
            "account": {"industry": "RETAIL", "role": "FP&A Manager",
                        "company_context": "A convenience retailer."},
            "prefer_high_impact": False, "asks_differentiation": False,
            "asks_why_not_big_si": False}


def test_plan_never_silently_drops_a_vetted_priority_pick():
    """matcher.plan applied its OWN Intro ceiling (a flat 4) on top of the per-work-type
    cap the pick pipeline had already enforced (3 EACH), so a mixed Intro deck lost every
    pick past the 4th -- no gap, no log, no trace (owner-reported, 2026-08-04)."""
    _chdir()
    from deckengine.services.matching import matcher
    ctx = {"client_name": "Acme", "industry": "RETAIL", "phase": "Intro",
           "work_types": ["MS", "WORKFORCE"], "functions": [], "recipient": "CFO",
           "transcript": "notes"}
    vetted = ["MSS050", "MSS051", "MSS048", "WFS019", "WFS032", "WFS043"]
    ids = [p["slide_id"] for p in matcher.plan(ctx, use_ai=False, priority_ids=vetted)["picks"]]
    dropped = [v for v in vetted if v not in ids]
    assert not dropped, "plan() threw away vetted picks: %s" % dropped


def test_intro_deck_carries_three_case_studies_per_selected_work_type(tmp_path, monkeypatch):
    """The owner's Intro rule is 3 case studies per SELECTED work type (a ceiling with a
    quality gate). A Workforce+MS Intro deck must therefore be able to carry 6, not 4."""
    _chdir()
    import app as appmod
    needs = ["Month-End Close", "Financial Reconciliation", "Close Orchestration",
             "Managed Recruitment", "Hire Train Deploy", "Accelerated Hiring"]
    picks = ["MSS050", "MSS051", "MSS048", "WFS019", "WFS032", "WFS043"]
    _hermetic_build(monkeypatch, tmp_path, brief=_brief(needs),
                    shortlists={n: [_cand(c)] for n, c in zip(needs, picks)})
    out = _post_build(appmod.app.test_client(), work_types=["MS", "WORKFORCE"], phase="Intro")
    assert sorted(out["cases"]) == sorted(picks), out["cases"]


def test_a_topic_the_client_raised_is_flagged_when_we_have_no_case(tmp_path, monkeypatch):
    """A capability the salesperson typed into the Context box that our library can't
    prove was dropped on the floor -- never picked, never flagged. It is the single most
    important gap there is: the client asked for it out loud."""
    _chdir()
    import app as appmod
    _hermetic_build(monkeypatch, tmp_path,
                    brief=_brief(["Financial Reconciliation"],
                                 expressed=["SAP S/4HANA migration"]),
                    shortlists={"Financial Reconciliation": [_cand("MSS051")],
                                "SAP S/4HANA migration": []})     # nothing covers it
    out = _post_build(appmod.app.test_client(),
                      transcript="Their CIO asked twice about an SAP S/4HANA migration.")
    assert "MSS051" in out["cases"]
    assert any("SAP S/4HANA migration" in g for g in out["gaps"]), out["gaps"]


def test_a_topic_the_client_raised_gets_first_claim_on_a_case(tmp_path, monkeypatch):
    """When slots are scarce (an Intro deck caps at 3 per work type), a topic the client
    actually raised must claim its proof point ahead of one merely inferred from the
    stakeholder's profile -- even though the profile match scores higher. It still has to
    clear the same coverage bar; this only decides who claims a case first."""
    _chdir()
    import app as appmod
    profile_needs = ["Month-End Close", "Financial Reconciliation", "Management Reporting"]
    _hermetic_build(
        monkeypatch, tmp_path,
        brief=_brief(profile_needs, expressed=["Predictive Maintenance"]),
        shortlists={"Month-End Close": [_cand("MSS048", adj=0.95, cosine=0.90)],
                    "Financial Reconciliation": [_cand("MSS051", adj=0.94, cosine=0.89)],
                    "Management Reporting": [_cand("MSS050", adj=0.93, cosine=0.88)],
                    # genuinely covered, but the weakest of the four
                    "Predictive Maintenance": [_cand("MSS084", adj=0.55, cosine=0.52)]})
    out = _post_build(appmod.app.test_client(), phase="Intro", work_types=["MS"],
                      transcript="They want predictive maintenance for the fuel fleet.")
    assert "MSS084" in out["cases"], out["cases"]


def test_topics_the_client_raised_reach_the_fill_ranker(tmp_path, monkeypatch):
    """`lead_research` -- the full-weight signal the fill ranker sorts on -- was built
    from the profile's needs only, so the salesperson's own notes carried no weight
    there at all."""
    _chdir()
    import app as appmod
    seen = _hermetic_build(monkeypatch, tmp_path,
                           brief=_brief(["Financial Reconciliation"],
                                        expressed=["Warehouse Automation"]),
                           shortlists={"Financial Reconciliation": [_cand("MSS051")]})
    _post_build(appmod.app.test_client(),
                transcript="They asked about warehouse automation for the fleet.")
    assert "Warehouse Automation" in seen["research"], seen["research"]


def test_an_unreadable_profile_file_warns_instead_of_matching_on_nothing(tmp_path, monkeypatch):
    """A scanned/image profile PDF yields no text. The research file has warned about
    this since day one; the profile file silently matched on nothing instead."""
    _chdir()
    import io
    import app as appmod
    _hermetic_build(monkeypatch, tmp_path, brief={}, shortlists={})
    out = _post_build(appmod.app.test_client(), transcript="",
                      profile_file=(io.BytesIO(b"%PDF-1.4 not really a pdf"), "profile.pdf"))
    assert "couldn't read any text from the stakeholder profile" in out["html"].lower(), \
        "no warning shown for an unreadable profile file"


# ── 2026-08-11: the "only 3 case studies for a 4-domain ask" report ────────────
# One salesperson build (Managed Services / AUTOMOTIVE / Second Meeting, notes = "AI
# accelerators mapped to Finance, Supply Chain, Production Planning, and Manufacturing")
# came back with 3 case studies while 46 MS cases sat above the relevance floor. Five
# separate mechanisms each threw away correct matches; one test per mechanism.

def test_fill_loop_is_not_capped_by_the_priority_pick_count():
    """`pick_cap = min(ceiling, max(len(priority), 3))` read as a floor of 3 but ALSO
    acted as a ceiling equal to the priority count: with 3 priority picks the fill loop
    ran zero times and MAX_CASE_PICKS=12 was unreachable on every research-driven deck."""
    _chdir()
    from deckengine.services.matching import matcher
    ctx = {"client_name": "Acme", "industry": "MANUFACTURING", "phase": "Second Meeting",
           "work_types": ["MS"], "functions": [], "recipient": "CTO",
           "transcript": "We need test automation, predictive maintenance and digital "
                         "twin across our plants."}
    picks = matcher.plan(ctx, use_ai=False, priority_ids=["MSS012", "MSS005", "MSS034"])["picks"]
    cases = [p["slide_id"] for p in picks if p["slide_id"][:3] in ("AIP", "WFS", "MSS")]
    assert len(cases) > 3, f"fill loop still capped at the priority count: {cases}"
    assert len(cases) <= matcher.MAX_CASE_PICKS, cases


def test_a_narrow_ask_still_produces_a_tight_deck():
    """The counterpart to the test above: lifting the cap must NOT pad every deck to 12.
    Breadth is controlled by the RELATIVE relevance floor, so a single-topic ask stays
    tight on its own."""
    _chdir()
    from deckengine.services.matching import matcher
    ctx = {"client_name": "Acme", "industry": "MANUFACTURING", "phase": "Second Meeting",
           "work_types": ["MS"], "functions": [], "recipient": "CFO",
           "transcript": "We are only interested in touchless invoice processing."}
    picks = matcher.plan(ctx, use_ai=False)["picks"]
    cases = [p["slide_id"] for p in picks if p["slide_id"][:3] in ("AIP", "WFS", "MSS")]
    assert len(cases) <= 6, f"a single-topic ask was padded out to {len(cases)} cases: {cases}"


def test_fill_loop_never_ships_two_near_twin_cases():
    """The soft MMR demotion was enough only while the fill loop stopped at 3 picks;
    filling to the real ceiling let a 0.887-similar case win a late slot on score alone
    (WFS002 Zero-Disruption Scaling + WFS025 Contract-to-Hire in one deck)."""
    _chdir()
    import itertools
    from deckengine.services.matching import matcher, relevance
    ctx = {"client_name": "HPE", "industry": "HEALTHCARE", "phase": "Second Meeting",
           "work_types": ["WORKFORCE"], "functions": [], "recipient": "Head of Talent",
           "transcript": "Looking for RPO and cloud migration staffing for healthcare IT."}
    picks = matcher.plan(ctx, use_ai=False)["picks"]
    cases = [p["slide_id"] for p in picks if p["slide_id"][:3] in ("AIP", "WFS", "MSS")]
    twins = [(a, b) for a, b in itertools.combinations(cases, 2)
             if relevance.max_similarity(a, [b]) >= matcher.NEAR_TWIN_SIM]
    assert not twins, f"deck ships look-alike proof points: {twins}"


def test_one_need_can_claim_more_than_one_case(tmp_path, monkeypatch):
    """A need claimed exactly one case, so every other case that cleared the coverage
    bar for it was discarded -- six cleared it on the reported build, five were dropped."""
    _chdir()
    import app as appmod
    from deckengine.constants import MAX_CASES_PER_NEED
    seen = _hermetic_build(monkeypatch, tmp_path,
                           brief=_brief(["Supply Chain Optimization"]),
                           shortlists={"Supply Chain Optimization": [
                               _cand("MSS125"), _cand("MSS032"), _cand("MSS002"),
                               _cand("MSS025")]})
    _post_build(appmod.app.test_client(), phase="Second Meeting",
                transcript="They asked about supply chain optimization.")
    # assert on the PRIORITY picks -- the need-claiming path this fix changed. (The deck
    # itself can legitimately carry more: matcher.plan's generic fill loop adds
    # top-ranked cases on top, and that is a separate, independently-capped mechanism.)
    claimed = seen["priority_ids"]
    assert len(claimed) > 1, f"one need still claims only one case: {claimed}"
    assert len(claimed) <= MAX_CASES_PER_NEED, claimed


def test_every_area_the_salesperson_enumerated_becomes_its_own_need():
    """'AI accelerators mapped to Finance, Supply Chain, Production Planning and
    Manufacturing' arrived as ONE expressed topic, so it could claim ONE case for the
    whole list -- the deck proved Finance and silently dropped the other three areas."""
    _chdir()
    from deckengine.services.matching.ai_matcher import split_enumerated_interest as split
    got = split("AI accelerators mapped to Finance, Supply Chain, Production Planning, "
                "and Manufacturing.")
    assert got == ["AI accelerators for Finance", "AI accelerators for Supply Chain",
                   "AI accelerators for Production Planning",
                   "AI accelerators for Manufacturing"], got
    # a bare list with no usable shared lead-in -> the areas on their own
    assert split("Interested in RPO, hire-train-deploy, and high-volume hiring") == [
        "RPO", "hire-train-deploy", "high-volume hiring"]
    # ordinary prose must NEVER be chopped into needs
    assert split("They want predictive maintenance and better uptime") == []
    assert split("We had a long conversation about their plans, and they mentioned "
                 "that the current ERP is failing them badly.") == []
    assert split("They want predictive maintenance") == []
    assert split("") == []


def test_an_adjacent_vertical_is_not_demoted_as_a_wrong_industry():
    """No Managed-Services case is tagged AUTOMOTIVE, so an automotive account could
    earn the industry boost from nothing while every correctly-tagged MANUFACTURING case
    took the wrong-vertical demotion -- leaving UNTAGGED records the only un-demoted
    ones, i.e. missing data became an advantage."""
    _chdir()
    from deckengine.services.matching import relevance
    assert relevance._same_family("AUTOMOTIVE", "MANUFACTURING")
    assert relevance._same_family("AUTOMOTIVE", "PROCESS_MFG")
    assert not relevance._same_family("AUTOMOTIVE", "BFSI")
    rows = [
        {"slide_id": "MFG", "title": "Production Scheduling", "keywords": "production planning",
         "primary_industry": "MANUFACTURING", "primary_function": "", "work_types": "MS",
         "search_text": "ai production scheduling across plants"},
        {"slide_id": "UNTAGGED", "title": "Production Scheduling", "keywords": "production planning",
         "primary_industry": "", "primary_function": "", "work_types": "MS",
         "search_text": "ai production scheduling across plants"},
    ]
    ranked = relevance.rank_cases("ai production planning", rows, industry="AUTOMOTIVE",
                                  wanted={"MS"}, use_semantic=False)
    by_id = {it["row"]["slide_id"]: it for it in ranked}
    assert by_id["MFG"]["score"] > by_id["UNTAGGED"]["score"], \
        "an untagged case still outranks the same case correctly tagged as an adjacent vertical"
    assert by_id["MFG"]["eligible"], "adjacent-vertical case is not even eligible"


def test_persona_boost_cannot_outweigh_meaning_in_the_shortlist():
    """shortlist_cases scores on COSINE (0..1) but reused rank_cases' persona weight
    (0.3 x up to 4 = 1.2), so persona was 2-3x the entire semantic signal it was meant
    to be a tiebreak on -- a 0.224-cosine case outranked a 0.521-cosine exact match."""
    _chdir()
    from deckengine.services.matching import relevance
    from deckengine.services.matching.personas import score_boost
    max_persona = relevance.W_SHORT_PERSONA * 4          # score_boost caps at 4
    other_tiebreaks = 0.12 * 3 + 0.10 * 2 + 0.10 + relevance.W_SHORT_IND_MATCH
    assert max_persona <= other_tiebreaks / 2, (
        "persona still dominates the shortlist tiebreaks", max_persona, other_tiebreaks)
    assert max_persona < 0.30, "persona can still outweigh a typical cosine gap"
    assert score_boost(("PROCUREMENT_OPS",), {"slide_id": "X"})[0] <= 4
